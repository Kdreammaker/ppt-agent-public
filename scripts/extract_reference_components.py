from __future__ import annotations

import argparse
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE_DIR = Path(__file__).resolve().parents[1]
if str(BASE_DIR) not in sys.path:
    sys.path.insert(0, str(BASE_DIR))

from scripts.reference_pipeline import LIBRARY_ROOT, load_json, write_json
from system.pptx_system import shape_bounds_inches

EXTRACTABLE_TYPES = {
    MSO_SHAPE_TYPE.AUTO_SHAPE,
    MSO_SHAPE_TYPE.LINE,
    MSO_SHAPE_TYPE.FREEFORM,
}


def extract_shape_specs(pptx_path: Path) -> list[dict[str, Any]]:
    prs = Presentation(str(pptx_path))
    slide = prs.slides[0]
    specs: list[dict[str, Any]] = []
    for index, shape in enumerate(slide.shapes, start=1):
        if shape.shape_type not in EXTRACTABLE_TYPES:
            continue
        text = " ".join(getattr(shape, "text", "").split())
        if text:
            continue
        specs.append(
            {
                "component_type": "shape_spec",
                "shape_index": index,
                "shape_name": getattr(shape, "name", f"shape_{index}"),
                "shape_type": str(shape.shape_type),
                "bounds": shape_bounds_inches(shape),
                "notes": "Simple vector-like PPT shape captured as metadata. This is not a replacement for PPTX templates or blueprints.",
            }
        )
    return specs


def extract_from_metadata(path: Path, *, dry_run: bool) -> dict[str, Any]:
    metadata = load_json(path)
    written: list[str] = []
    for record in metadata.get("slides", []):
        identity = record.get("identity", {})
        slide_id = identity.get("slide_id")
        pptx_path = BASE_DIR / str(identity.get("one_slide_pptx", ""))
        if not slide_id or not pptx_path.exists():
            continue
        specs = extract_shape_specs(pptx_path)
        if not specs:
            continue
        quality = record.get("curation", {}).get("quality_label", "unknown")
        output = LIBRARY_ROOT / str(quality) / "components" / f"{slide_id}_components.json"
        payload = {
            "schema_version": "1.0",
            "slide_id": slide_id,
            "source_metadata": path.resolve().relative_to(BASE_DIR).as_posix(),
            "source_pptx": identity.get("one_slide_pptx"),
            "components": specs,
        }
        if not dry_run:
            write_json(output, payload)
        written.append(output.resolve().relative_to(BASE_DIR).as_posix())
    return {"metadata_path": path.resolve().relative_to(BASE_DIR).as_posix(), "written": written}


def collect_metadata(paths: list[str]) -> list[Path]:
    if paths:
        return [(BASE_DIR / path).resolve() if not Path(path).is_absolute() else Path(path) for path in paths]
    return list(LIBRARY_ROOT.glob("*/metadata/*.json"))


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Optionally extract simple reference components as shape specs.")
    parser.add_argument("paths", nargs="*")
    parser.add_argument("--dry-run", action="store_true")
    args = parser.parse_args(argv)
    results = [extract_from_metadata(path, dry_run=args.dry_run) for path in collect_metadata(args.paths) if path.exists()]
    print(f"metadata_files={len(results)} component_files={sum(len(item['written']) for item in results)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

from __future__ import annotations

import argparse
import copy
import json
import os
import sys
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

BASE_DIR = Path(__file__).resolve().parents[1]
if str(BASE_DIR) not in sys.path:
    sys.path.insert(0, str(BASE_DIR))

from system.deck_models import validate_deck_spec
from system.blueprint_loader import load_blueprints
from system.pptx_system import (
    ThemeConfig,
    add_image,
    add_simple_bar_chart,
    add_textbox,
    find_blank_layout,
    import_slide,
    replace_placeholder_text,
    resolve_color,
    set_all_fonts,
)
from system.template_engine import render_template_slide, slot_text_budget
from system.template_text_dna import load_template_text_dna_cleanup, slide_spec_with_text_dna_cleanup

WORKDIR = BASE_DIR
REPORTS_DIR = BASE_DIR / "outputs" / "reports"
OUTPUTS_ROOT = BASE_DIR / "outputs"
MODE_POLICY_ROOT = BASE_DIR / "config" / "mode_policies"
REFERENCE_KNOWLEDGE_GRAPH_PATH = BASE_DIR / "config" / "reference_knowledge_graph.json"
TEMPLATE_PATTERN_CATALOG_PATH = BASE_DIR / "config" / "template_pattern_catalog.json"
THEME_XML_PATH = "ppt/theme/theme1.xml"
DRAWINGML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def resolve_path(base_dir: Path, value: str | None) -> Path | None:
    if value is None:
        return None
    path = Path(value)
    return path if path.is_absolute() else (base_dir / path).resolve()


def require_repo_output_path(path: Path | None, *, label: str) -> Path:
    if path is None:
        raise ValueError(f"{label} is required")
    resolved = path.resolve()
    allowed_roots = [OUTPUTS_ROOT.resolve()]
    workspace = os.environ.get("PPT_AGENT_WORKSPACE")
    if workspace:
        allowed_roots.append((Path(workspace) / "outputs").resolve())
    if not any(resolved == root or resolved.is_relative_to(root) for root in allowed_roots):
        roots = ", ".join(root.as_posix() for root in allowed_roots)
        raise ValueError(f"{label} must stay under one of [{roots}]: {resolved}")
    return resolved


def load_json(path: Path | None) -> dict[str, Any] | None:
    if path is None or not path.exists():
        return None
    return json.loads(path.read_text(encoding="utf-8"))


def normalize_hex_color(value: str) -> str:
    color = value.strip().lstrip("#")
    if len(color) != 6 or any(char not in "0123456789abcdefABCDEF" for char in color):
        raise ValueError(f"Expected 6-digit hex color, got: {value}")
    return color.upper()


def normalize_accent_slot(slot: str) -> str:
    normalized = slot.replace("_", "")
    valid = {f"accent{index}" for index in range(1, 7)}
    if normalized not in valid:
        raise ValueError(f"Unsupported theme accent slot: {slot}")
    return normalized


def patched_theme_xml(theme_xml: bytes, accent_overrides: dict[str, str]) -> bytes:
    ET.register_namespace("a", DRAWINGML_NS)
    root = ET.fromstring(theme_xml)
    clr_scheme = root.find(f".//{{{DRAWINGML_NS}}}clrScheme")
    if clr_scheme is None:
        raise ValueError("theme1.xml does not contain a clrScheme")

    by_name = {child.tag.rsplit("}", 1)[-1]: child for child in list(clr_scheme)}
    for raw_slot, raw_color in accent_overrides.items():
        slot = normalize_accent_slot(raw_slot)
        if slot not in by_name:
            raise ValueError(f"theme1.xml does not define {slot}")
        element = by_name[slot]
        for child in list(element):
            element.remove(child)
        ET.SubElement(element, f"{{{DRAWINGML_NS}}}srgbClr", {"val": normalize_hex_color(raw_color)})
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)


def apply_theme_accent_overrides(output_path: Path, accent_overrides: dict[str, str] | None) -> None:
    if not accent_overrides:
        return

    temp_path = output_path.with_name(f"{output_path.name}.tmp")
    with zipfile.ZipFile(output_path, "r") as source, zipfile.ZipFile(temp_path, "w") as target:
        found_theme = False
        for item in source.infolist():
            data = source.read(item.filename)
            if item.filename == THEME_XML_PATH:
                data = patched_theme_xml(data, accent_overrides)
                found_theme = True
            target.writestr(item, data)
        if not found_theme:
            raise ValueError(f"PPTX package does not contain {THEME_XML_PATH}")
    os.replace(temp_path, output_path)


def apply_template_design_dna(reference_catalog: dict[str, Any] | None) -> None:
    if not reference_catalog:
        return
    design_dna_path = BASE_DIR / "config" / "template_design_dna.json"
    design_dna = load_json(design_dna_path)
    if not design_dna:
        return
    dna_by_slide_id = design_dna.get("slides", {})
    dna_fields = {
        "density",
        "tone",
        "structure",
        "visual_weight",
        "content_capacity",
        "footer_supported",
        "best_for",
        "avoid_for",
        "review_notes",
        "override_applied",
    }
    for slide in reference_catalog.get("slides", []):
        dna = dna_by_slide_id.get(slide.get("slide_id"))
        if not dna:
            continue
        for field in dna_fields:
            if field in dna:
                slide[field] = dna[field]


def apply_extras(slide, slide_spec: dict[str, Any], theme: ThemeConfig) -> None:
    geometry_map = {
        "rect": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        "round_rect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        "ellipse": MSO_AUTO_SHAPE_TYPE.OVAL,
        "triangle": MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        "arc": MSO_AUTO_SHAPE_TYPE.ARC,
        "wave": MSO_AUTO_SHAPE_TYPE.WAVE,
    }
    for shape_spec in slide_spec.get("shapes", []):
        geometry = shape_spec.get("geometry") or ("round_rect" if shape_spec.get("radius") else "rect")
        shape = slide.shapes.add_shape(
            geometry_map.get(geometry, MSO_AUTO_SHAPE_TYPE.RECTANGLE),
            Inches(shape_spec["left"]),
            Inches(shape_spec["top"]),
            Inches(shape_spec["width"]),
            Inches(shape_spec["height"]),
        )
        if shape_spec.get("rotation") is not None:
            shape.rotation = float(shape_spec["rotation"])
        fill = shape_spec.get("fill")
        if fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = resolve_color(theme, fill)
        else:
            shape.fill.background()
        line = shape_spec.get("line")
        if line:
            shape.line.color.rgb = resolve_color(theme, line)
            if shape_spec.get("line_width") is not None:
                shape.line.width = Inches(float(shape_spec["line_width"]))
        else:
            shape.line.fill.background()

    for image in slide_spec.get("images", []):
        add_image(
            slide,
            image["path"],
            left=image["left"],
            top=image["top"],
            width=image["width"],
            height=image["height"],
        )

    for chart in slide_spec.get("bar_charts", []):
        add_simple_bar_chart(
            slide,
            left=chart["left"],
            top=chart["top"],
            width=chart["width"],
            height=chart["height"],
            categories=chart["categories"],
            values=chart["values"],
            theme=theme,
        )

    alignments = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    for box in slide_spec.get("text_boxes", []):
        font_size = box.get("font_size")
        if font_size is None:
            role = box.get("font_role") or "body"
            font_size = theme.sizes.get(role, theme.sizes["body"])
        add_textbox(
            slide,
            box["left"],
            box["top"],
            box["width"],
            box["height"],
            box["text"],
            font_name=theme.font_family,
            font_size=font_size,
            font_color=resolve_color(theme, box.get("color"), "primary"),
            bold=box.get("bold", False),
            align=alignments.get(box.get("align")),
            max_chars_per_line=box.get("max_chars_per_line"),
        )


def resolve_slide_asset_paths(slide_spec: dict[str, Any], spec_dir: Path) -> dict[str, Any]:
    resolved = copy.deepcopy(slide_spec)
    for image in resolved.get("images", []):
        path = resolve_path(spec_dir, image.get("path"))
        if path is not None:
            image["path"] = str(path)
    if "image_slots" in resolved:
        resolved["image_slots"] = {
            slot: str(resolve_path(spec_dir, value) or value)
            for slot, value in resolved.get("image_slots", {}).items()
        }
    return resolved


def workspace_asset_root() -> Path | None:
    value = os.environ.get("PPT_AGENT_WORKSPACE")
    if not value:
        return None
    path = Path(value)
    if not path.is_absolute():
        path = (Path.cwd() / path).resolve()
    return path


def apply_workspace_asset_intents(spec: dict[str, Any], workspace: Path | None) -> list[dict[str, Any]]:
    usage: list[dict[str, Any]] = []
    if workspace is None:
        return usage
    for intent in spec.get("asset_intents", []):
        if not isinstance(intent, dict) or intent.get("source_type") != "user_upload":
            continue
        asset_id = str(intent.get("asset_id") or "")
        relative_path = str(intent.get("workspace_relative_path") or "")
        if not asset_id or not relative_path or Path(relative_path).is_absolute() or ".." in Path(relative_path).parts:
            usage.append(
                {
                    "asset_id": asset_id,
                    "source_type": "user_upload",
                    "status": "invalid_workspace_relative_path",
                    "workspace_relative_path": relative_path,
                }
            )
            continue
        asset_path = (workspace / relative_path).resolve()
        try:
            asset_path.relative_to(workspace.resolve())
        except ValueError:
            usage.append(
                {
                    "asset_id": asset_id,
                    "source_type": "user_upload",
                    "status": "path_outside_workspace",
                    "workspace_relative_path": relative_path,
                }
            )
            continue
        status = "resolved" if asset_path.exists() else "missing"
        slide_number = intent.get("slide_number")
        slot = intent.get("slot")
        if status == "resolved" and intent.get("asset_class") == "image" and slide_number and slot:
            slides = spec.get("slides", [])
            index = int(slide_number) - 1
            if 0 <= index < len(slides) and isinstance(slides[index], dict):
                slides[index].setdefault("image_slots", {})[str(slot)] = asset_path.as_posix()
                status = "resolved_to_image_slot"
        usage.append(
            {
                "asset_id": asset_id,
                "source_type": "user_upload",
                "asset_class": intent.get("asset_class"),
                "slide_number": slide_number,
                "slot": slot,
                "workspace_relative_path": relative_path,
                "status": status,
                "usage_rationale": intent.get("usage_rationale") or intent.get("notes"),
                "private_upload_allowed": bool(intent.get("private_upload_allowed", False)),
            }
        )
    return usage


def write_asset_usage_report(output_path: Path, spec: dict[str, Any], usage: list[dict[str, Any]]) -> None:
    reports_dir = BASE_DIR / "outputs" / "reports"
    reports_dir.mkdir(parents=True, exist_ok=True)
    all_intents = []
    for intent in spec.get("asset_intents", []):
        if not isinstance(intent, dict):
            continue
        all_intents.append(
            {
                "asset_id": intent.get("asset_id"),
                "source_type": intent.get("source_type") or ("user_upload" if intent.get("source_policy") == "workspace_user_asset" else "catalog_or_connector"),
                "asset_class": intent.get("asset_class"),
                "role": intent.get("role"),
                "slide_number": intent.get("slide_number"),
                "slot": intent.get("slot"),
                "workspace_relative_path": intent.get("workspace_relative_path"),
                "materialization": intent.get("materialization"),
                "license_action": intent.get("license_action"),
                "risk_level": intent.get("risk_level"),
                "semantic_context": intent.get("semantic_context"),
                "template_media_policy": intent.get("template_media_policy"),
                "usage_rationale": intent.get("usage_rationale") or intent.get("notes"),
            }
        )
    payload = {
        "schema_version": "1.0",
        "deck": output_path.name,
        "deck_name": spec.get("name"),
        "summary": {
            "asset_intents": len(all_intents),
            "workspace_assets": len([item for item in all_intents if item.get("source_type") == "user_upload"]),
            "resolved_workspace_assets": len([item for item in usage if str(item.get("status", "")).startswith("resolved")]),
        },
        "assets": all_intents,
        "workspace_resolution": usage,
        "public_boundary": {
            "uses_asset_ids_not_absolute_paths": True,
            "private_upload_allowed_default": False,
            "raw_private_payloads_included": False,
        },
    }
    report_json = reports_dir / f"{output_path.stem}_asset_usage_summary.json"
    report_md = reports_dir / f"{output_path.stem}_asset_usage_summary.md"
    latest_json = reports_dir / "asset_usage_summary.json"
    for path in (report_json, latest_json):
        path.write_text(json.dumps(payload, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    lines = [
        "# Asset Usage Summary",
        "",
        f"- Deck: `{output_path.name}`",
        f"- Asset intents: {payload['summary']['asset_intents']}",
        f"- Workspace assets: {payload['summary']['workspace_assets']}",
        "",
        "| Asset ID | Source | Class | Slide | Slot | Status |",
        "| --- | --- | --- | ---: | --- | --- |",
    ]
    status_by_id = {item.get("asset_id"): item.get("status") for item in usage}
    for asset in all_intents:
        lines.append(
            f"| {asset.get('asset_id') or ''} | {asset.get('source_type') or ''} | {asset.get('asset_class') or ''} | "
            f"{asset.get('slide_number') or ''} | {asset.get('slot') or ''} | {status_by_id.get(asset.get('asset_id'), '') or ''} |"
        )
    report_md.write_text("\n".join(lines) + "\n", encoding="utf-8")


def load_spec(spec_path: str | Path) -> tuple[dict[str, Any], Path]:
    spec_path = Path(spec_path).resolve()
    spec = json.loads(spec_path.read_text(encoding="utf-8"))
    validate_deck_spec(spec)
    return spec, spec_path.parent


def validate_supported_aspect_ratio(spec: dict[str, Any]) -> None:
    aspect_ratio = spec.get("aspect_ratio", "16:9")
    if aspect_ratio == "16:9":
        return
    if aspect_ratio == "4:3":
        raise ValueError(
            "4:3 output requires dedicated 4:3 template libraries and aspect-specific blueprints. "
            "Automatic stretching or cropping of 16:9 templates is not supported."
        )
    raise ValueError(f"Unsupported aspect_ratio: {aspect_ratio}")


def configure_presentation_size(prs: Presentation, spec: dict[str, Any]) -> None:
    validate_supported_aspect_ratio(spec)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)


def relative_to_base(path: Path) -> str:
    resolved = path.resolve()
    try:
        return resolved.relative_to(BASE_DIR).as_posix()
    except ValueError:
        return Path(os.path.relpath(resolved, BASE_DIR)).as_posix()


def old_catalog_slide_map(old_catalog: dict[str, Any] | None) -> dict[str, int]:
    if not old_catalog:
        return {}
    return {item["key"]: item["slide_no"] for item in old_catalog.get("validated_slide_library", [])}


def indexed_reference_catalog(reference_catalog: dict[str, Any] | None) -> tuple[dict[str, dict[str, Any]], list[dict[str, Any]]]:
    if not reference_catalog:
        return {}, []
    slides = reference_catalog.get("slides", [])
    by_key = {item["template_key"]: item for item in slides}
    return by_key, slides


def reference_slide_summary(item: dict[str, Any]) -> dict[str, Any]:
    return {
        "template_key": item.get("template_key"),
        "slide_id": item.get("slide_id"),
        "library_id": item.get("library_id"),
        "purpose": item.get("purpose"),
        "scope": item.get("scope"),
        "variant": item.get("variant"),
        "quality_score": item.get("quality_score"),
        "design_tier": item.get("design_tier"),
        "usage_policy": item.get("usage_policy"),
        "default_rank": item.get("default_rank"),
        "style_tags": item.get("style_tags", []),
        "density": item.get("density"),
        "tone": item.get("tone", []),
        "structure": item.get("structure"),
        "visual_weight": item.get("visual_weight"),
        "content_capacity": item.get("content_capacity"),
        "footer_supported": item.get("footer_supported"),
        "best_for": item.get("best_for", []),
        "avoid_for": item.get("avoid_for", []),
        "review_notes": item.get("review_notes"),
        "override_applied": item.get("override_applied", False),
    }


def load_mode_policy_context(spec: dict[str, Any]) -> dict[str, Any]:
    mode = spec.get("mode_policy", "auto")
    policy_path = MODE_POLICY_ROOT / f"{mode}_mode_policy.json"
    policy = load_json(policy_path) or {
        "schema_version": "1.0",
        "mode": mode,
        "uses": {},
        "default_assumptions": [],
        "fallback_rules": [],
    }
    graph = load_json(REFERENCE_KNOWLEDGE_GRAPH_PATH) or {}
    pattern_catalog = load_json(TEMPLATE_PATTERN_CATALOG_PATH) or {}
    return {
        "mode": mode,
        "policy_path": relative_to_base(policy_path),
        "policy": policy,
        "graph": graph,
        "pattern_catalog": pattern_catalog,
    }


def normalized_list(value: Any) -> list[str]:
    if value is None:
        return []
    if isinstance(value, str):
        return [value]
    if isinstance(value, list):
        return [item for item in value if isinstance(item, str)]
    return []


def template_graph_edge_count(graph: dict[str, Any], slide_id: str | None) -> int:
    if not slide_id:
        return 0
    node_id = f"production_template:{slide_id}"
    return sum(
        1
        for edge in graph.get("edges", [])
        if edge.get("source") == node_id or edge.get("target") == node_id
    )


def pattern_fit_score(pattern_catalog: dict[str, Any], item: dict[str, Any], selector: dict[str, Any]) -> float:
    patterns = pattern_catalog.get("patterns", [])
    if not isinstance(patterns, list) or not patterns:
        return 0.0
    purpose = selector.get("purpose")
    scope = selector.get("scope")
    matches = 0
    for pattern in patterns:
        if not isinstance(pattern, dict):
            continue
        if pattern.get("purpose") == purpose or pattern.get("category") == purpose:
            matches += 1
        if scope and pattern.get("scope") == scope:
            matches += 1
        if pattern.get("template_key") == item.get("template_key"):
            matches += 2
    return min(1.0, matches / 4)


def variation_score_components(
    item: dict[str, Any],
    selection_context: dict[str, Any] | None = None,
) -> dict[str, float]:
    previous = (selection_context or {}).get("previous_reference")
    if not isinstance(previous, dict):
        return {}

    penalty = 0.0
    bonus = 0.0
    previous_density = previous.get("density")
    previous_visual_weight = previous.get("visual_weight")
    previous_structure = previous.get("structure")
    current_density = item.get("density")
    current_visual_weight = item.get("visual_weight")
    current_structure = item.get("structure")

    if previous_density == "dense" and current_density == "dense":
        penalty -= 0.35
    if previous_visual_weight == "text_heavy" and current_visual_weight == "text_heavy":
        penalty -= 0.30
    if previous_structure and previous_structure == current_structure:
        penalty -= 0.20
    if previous_density == "dense" and current_density in {"light", "medium"}:
        bonus += 0.15
    if previous_visual_weight == "text_heavy" and current_visual_weight in {"balanced", "visual_first"}:
        bonus += 0.15

    components: dict[str, float] = {}
    if penalty:
        components["variation_penalty"] = round(penalty, 4)
    if bonus:
        components["variation_bonus"] = round(min(0.25, bonus), 4)
    return components


def mode_policy_score_components(
    item: dict[str, Any],
    selector: dict[str, Any],
    mode_policy_context: dict[str, Any] | None,
    selection_context: dict[str, Any] | None = None,
) -> dict[str, float]:
    if not mode_policy_context:
        return {}
    policy = mode_policy_context.get("policy", {})
    weights = policy.get("template_retrieval_weights", {}) if isinstance(policy, dict) else {}
    if not weights:
        weights = {
            "purpose_match": 1.0,
            "deck_type_fit": 0.8,
            "tone_fit": 0.6,
            "pattern_fit": 0.5,
            "quality_score": 0.7,
        }

    purpose = selector.get("purpose")
    scope = selector.get("scope")
    tone_targets = set(normalized_list(selector.get("tone") or selector.get("tones")))
    item_tones = set(normalized_list(item.get("tone"))) | set(normalized_list(item.get("style_tags")))
    preferred_density = selector.get("preferred_density")
    preferred_structure = selector.get("preferred_structure")
    preferred_visual_weight = selector.get("preferred_visual_weight")

    purpose_match = 1.0 if item.get("purpose") == purpose else 0.0
    if scope is None:
        deck_type_fit = 0.75
    elif item.get("scope") == scope:
        deck_type_fit = 1.0
    elif item.get("scope") == "generic":
        deck_type_fit = 0.65
    else:
        deck_type_fit = 0.0
    tone_fit = 0.0 if tone_targets else 0.5
    if tone_targets:
        tone_fit = len(tone_targets & item_tones) / max(1, len(tone_targets))
    density_fit = 1.0 if preferred_density and item.get("density") == preferred_density else 0.0
    structure_fit = 1.0 if preferred_structure and item.get("structure") == preferred_structure else 0.0
    visual_weight_fit = 1.0 if preferred_visual_weight and item.get("visual_weight") == preferred_visual_weight else 0.0
    quality_score = max(0.0, min(1.0, float(item.get("quality_score", 0) or 0)))
    pattern_fit = pattern_fit_score(mode_policy_context.get("pattern_catalog", {}), item, selector)

    components = {
        "purpose_match": purpose_match * float(weights.get("purpose_match", 0)),
        "deck_type_fit": deck_type_fit * float(weights.get("deck_type_fit", 0)),
        "tone_fit": tone_fit * float(weights.get("tone_fit", 0)),
        "pattern_fit": pattern_fit * float(weights.get("pattern_fit", 0)),
        "quality_score": quality_score * float(weights.get("quality_score", 0)),
    }
    if preferred_density:
        components["density_fit"] = density_fit * 0.4
    if preferred_structure:
        components["structure_fit"] = structure_fit * 0.4
    if preferred_visual_weight:
        components["visual_weight_fit"] = visual_weight_fit * 0.3
    if selector.get("use_variation_penalty", True):
        components.update(variation_score_components(item, selection_context))
    return components


def mode_policy_score(
    item: dict[str, Any],
    selector: dict[str, Any],
    mode_policy_context: dict[str, Any] | None,
    selection_context: dict[str, Any] | None = None,
) -> float:
    return round(sum(mode_policy_score_components(item, selector, mode_policy_context, selection_context).values()), 4)


def mode_policy_selection_details(
    selected: dict[str, Any],
    selector: dict[str, Any] | None,
    mode_policy_context: dict[str, Any] | None,
    selection_context: dict[str, Any] | None = None,
) -> dict[str, Any] | None:
    if not mode_policy_context:
        return None
    policy = mode_policy_context.get("policy", {})
    graph = mode_policy_context.get("graph", {})
    pattern_catalog = mode_policy_context.get("pattern_catalog", {})
    selector = selector or {}
    components = mode_policy_score_components(selected, selector, mode_policy_context, selection_context)
    previous = (selection_context or {}).get("previous_reference")
    return {
        "mode": mode_policy_context.get("mode"),
        "policy_path": mode_policy_context.get("policy_path"),
        "default_assumptions": policy.get("default_assumptions", []),
        "fallback_rules": policy.get("fallback_rules", []),
        "retrieval_weights": policy.get("template_retrieval_weights", {}),
        "graph_query_strategy": policy.get("graph_query_strategy"),
        "artifact_status": {
            "reference_knowledge_graph_nodes": len(graph.get("nodes", [])),
            "reference_knowledge_graph_edges": len(graph.get("edges", [])),
            "template_graph_edges_for_selected": template_graph_edge_count(graph, selected.get("slide_id")),
            "template_pattern_count": len(pattern_catalog.get("patterns", [])),
            "design_dna_present_for_selected": bool(selected.get("density") or selected.get("tone") or selected.get("structure")),
        },
        "selected_score": mode_policy_score(selected, selector, mode_policy_context, selection_context),
        "selected_score_components": components,
        "variation_context": (
            None
            if not isinstance(previous, dict)
            else {
                "previous_template_key": previous.get("template_key"),
                "previous_density": previous.get("density"),
                "previous_structure": previous.get("structure"),
                "previous_visual_weight": previous.get("visual_weight"),
            }
        ),
        "notes": "Mode policy weights are applied after hard selector constraints so raw/candidate references remain unavailable to runtime selection.",
    }


def ranked_reference_slides(
    reference_slides: list[dict[str, Any]],
    selector: dict[str, Any],
    mode_policy_context: dict[str, Any] | None = None,
    selection_context: dict[str, Any] | None = None,
) -> list[dict[str, Any]]:
    purpose = selector["purpose"]
    required_tags = set(selector.get("required_tags", []))
    source_library = selector.get("source_library")
    scope = selector.get("scope")
    preferred_variant = selector.get("preferred_variant")
    fallback_variants = selector.get("fallback_variants", [])
    min_quality_score = selector.get("min_quality_score")
    preferred_design_tier = selector.get("preferred_design_tier")
    usage_policies = set(selector.get("usage_policies", []))
    prefer_high_quality = selector.get("prefer_high_quality", True)
    selector_tones = set(normalized_list(selector.get("tone") or selector.get("tones")))
    preferred_density = selector.get("preferred_density")
    preferred_structure = selector.get("preferred_structure")
    preferred_visual_weight = selector.get("preferred_visual_weight")

    candidates = [slide for slide in reference_slides if slide["purpose"] == purpose]
    if source_library:
        candidates = [slide for slide in candidates if slide["library_id"] == source_library]
    if required_tags:
        candidates = [slide for slide in candidates if required_tags.issubset(set(slide.get("style_tags", [])))]
    if min_quality_score is not None:
        candidates = [slide for slide in candidates if float(slide.get("quality_score", 0)) >= float(min_quality_score)]
    if usage_policies:
        candidates = [slide for slide in candidates if slide.get("usage_policy") in usage_policies]
    if selector_tones:
        candidates = [
            slide
            for slide in candidates
            if selector_tones & (set(normalized_list(slide.get("tone"))) | set(normalized_list(slide.get("style_tags"))))
        ]
    if preferred_density:
        candidates = [slide for slide in candidates if slide.get("density") == preferred_density]
    if preferred_structure:
        candidates = [slide for slide in candidates if slide.get("structure") == preferred_structure]
    if preferred_visual_weight:
        candidates = [slide for slide in candidates if slide.get("visual_weight") == preferred_visual_weight]
    if not candidates:
        raise LookupError(f"No template candidates found for selector: {selector}")

    def scope_bucket(item: dict[str, Any]) -> int:
        if scope is None:
            return 0
        if item["scope"] == scope:
            return 0
        if item["scope"] == "generic":
            return 1
        return 2

    variant_order: dict[str, int] = {}
    if preferred_variant:
        variant_order[preferred_variant] = 0
    for index, variant in enumerate(fallback_variants, start=1):
        variant_order.setdefault(variant, index)

    def variant_bucket(item: dict[str, Any]) -> int:
        if item["variant"] in variant_order:
            return variant_order[item["variant"]]
        if preferred_variant:
            return len(variant_order) + 1
        return 0

    def design_tier_bucket(item: dict[str, Any]) -> int:
        if preferred_design_tier is None:
            return 0
        return 0 if item.get("design_tier") == preferred_design_tier else 1

    def quality_bucket(item: dict[str, Any]) -> float:
        return -float(item.get("quality_score", 0)) if prefer_high_quality else 0.0

    return sorted(
        candidates,
        key=lambda item: (
            scope_bucket(item),
            variant_bucket(item),
            design_tier_bucket(item),
            quality_bucket(item),
            -mode_policy_score(item, selector, mode_policy_context, selection_context),
            item.get("default_rank", 999),
            item["template_key"],
        ),
    )


def select_reference_slide(
    reference_slides: list[dict[str, Any]],
    selector: dict[str, Any],
    mode_policy_context: dict[str, Any] | None = None,
) -> dict[str, Any]:
    return ranked_reference_slides(reference_slides, selector, mode_policy_context)[0]


def selector_selection_factors(selected: dict[str, Any], selector: dict[str, Any]) -> dict[str, Any]:
    scope = selector.get("scope")
    preferred_variant = selector.get("preferred_variant")
    fallback_variants = selector.get("fallback_variants", [])
    preferred_design_tier = selector.get("preferred_design_tier")

    if scope is None:
        scope_match = "not_requested"
    elif selected.get("scope") == scope:
        scope_match = "exact"
    elif selected.get("scope") == "generic":
        scope_match = "generic_fallback"
    else:
        scope_match = "other_scope"

    if not preferred_variant:
        variant_match = "not_requested"
    elif selected.get("variant") == preferred_variant:
        variant_match = "preferred"
    elif selected.get("variant") in fallback_variants:
        variant_match = "fallback"
    else:
        variant_match = "other_variant"

    return {
        "selector": selector,
        "priority_order": [
            "scope",
            "variant",
            "design_tier",
            "quality_score",
            "mode_policy_score",
            "variation_penalty",
            "default_rank",
            "template_key",
        ],
        "selected_matches": {
            "scope_match": scope_match,
            "variant_match": variant_match,
            "design_tier_match": (
                None if preferred_design_tier is None else selected.get("design_tier") == preferred_design_tier
            ),
            "quality_score": selected.get("quality_score"),
            "usage_policy": selected.get("usage_policy"),
            "tone": selected.get("tone", []),
            "density": selected.get("density"),
            "structure": selected.get("structure"),
            "visual_weight": selected.get("visual_weight"),
        },
    }


def selector_confidence(selected: dict[str, Any], selector: dict[str, Any]) -> float:
    factors = selector_selection_factors(selected, selector)["selected_matches"]
    scope_match = factors["scope_match"]
    variant_match = factors["variant_match"]
    if scope_match == "exact" and variant_match in {"preferred", "not_requested"}:
        return 0.95
    if scope_match == "exact" and variant_match == "fallback":
        return 0.85
    if scope_match == "generic_fallback" and variant_match in {"preferred", "fallback", "not_requested"}:
        return 0.75
    return 0.55


def selector_rationale(
    reference_slides: list[dict[str, Any]],
    selector: dict[str, Any],
    mode_policy_context: dict[str, Any] | None = None,
    selection_context: dict[str, Any] | None = None,
) -> tuple[dict[str, Any], dict[str, Any]]:
    ranked = ranked_reference_slides(reference_slides, selector, mode_policy_context, selection_context)
    selected = ranked[0]
    rationale = {
        "source_type": "slide_selector",
        "purpose": selector.get("purpose"),
        "scope": selector.get("scope"),
        "selected_template_key": selected.get("template_key"),
        "selected_slide_id": selected.get("slide_id"),
        "selected_library": selected.get("library_id"),
        "selected_variant": selected.get("variant"),
        "candidate_templates": [
            {"rank": index, **reference_slide_summary(candidate)}
            for index, candidate in enumerate(ranked[:8], start=1)
        ],
        "selection_factors": selector_selection_factors(selected, selector),
        "mode_policy": mode_policy_selection_details(selected, selector, mode_policy_context, selection_context),
        "rejected_candidates": [
            {
                **reference_slide_summary(candidate),
                "mode_policy_score": mode_policy_score(candidate, selector, mode_policy_context, selection_context),
                "rejection_reason": "Ranked lower by selector priority order.",
            }
            for candidate in ranked[1:6]
        ],
        "confidence": selector_confidence(selected, selector),
        "notes": "Selected by slide_selector against the production reference catalog with mode policy scoring as an ambiguity resolver.",
    }
    return selected, rationale


def direct_template_rationale(
    entry: dict[str, Any],
    mode_policy_context: dict[str, Any] | None = None,
) -> dict[str, Any]:
    return {
        "source_type": "direct_template_key",
        "purpose": entry.get("purpose"),
        "scope": entry.get("scope"),
        "selected_template_key": entry.get("template_key"),
        "selected_slide_id": entry.get("slide_id"),
        "selected_library": entry.get("library_id"),
        "selected_variant": entry.get("variant"),
        "candidate_templates": [reference_slide_summary(entry)],
        "selection_factors": {
            "selector": None,
            "priority_order": [],
            "selected_matches": {
                "template_key": "fixed_by_spec",
            },
        },
        "mode_policy": mode_policy_selection_details(entry, None, mode_policy_context),
        "rejected_candidates": [],
        "confidence": 1.0,
        "notes": "Template key was explicitly provided by the deck spec.",
    }


def legacy_source_rationale(
    source_type: str,
    slide_spec: dict[str, Any],
    mode_policy_context: dict[str, Any] | None = None,
) -> dict[str, Any]:
    mode_policy = None
    if mode_policy_context:
        mode_policy = {
            "mode": mode_policy_context.get("mode"),
            "policy_path": mode_policy_context.get("policy_path"),
            "notes": "Compatibility source path; mode policy recorded for audit but did not rank a production template.",
        }
    return {
        "source_type": source_type,
        "purpose": None,
        "scope": None,
        "selected_template_key": slide_spec.get("template_key") or slide_spec.get("base_slide_key"),
        "selected_slide_id": None,
        "selected_library": None,
        "selected_variant": None,
        "candidate_templates": [],
        "selection_factors": {
            "selector": None,
            "priority_order": [],
            "selected_matches": {},
        },
        "mode_policy": mode_policy,
        "rejected_candidates": [],
        "confidence": 1.0 if source_type == "blank_source" else 0.6,
        "notes": "Resolved through a compatibility source path rather than reference-catalog selection.",
    }


def resolve_slide_source(
    slide_spec: dict[str, Any],
    *,
    spec: dict[str, Any],
    spec_dir: Path,
    old_catalog: dict[str, Any] | None,
    old_slide_map: dict[str, int],
    reference_catalog: dict[str, Any] | None,
    reference_by_key: dict[str, dict[str, Any]],
    mode_policy_context: dict[str, Any] | None = None,
    selection_context: dict[str, Any] | None = None,
) -> dict[str, Any]:
    if slide_spec.get("source_mode") == "blank":
        return {
            "presentation_path": None,
            "slide_no": None,
            "slide_id": None,
            "template_key": None,
            "rationale": legacy_source_rationale("blank_source", slide_spec, mode_policy_context),
        }

    if "base_slide_no" in slide_spec:
        source_template = resolve_path(spec_dir, spec.get("source_template")) or resolve_path(spec_dir, old_catalog["template_path"])  # type: ignore[index]
        return {
            "presentation_path": source_template,
            "slide_no": slide_spec["base_slide_no"],
            "slide_id": None,
            "template_key": None,
            "rationale": legacy_source_rationale("base_slide_no", slide_spec, mode_policy_context),
        }

    template_key = slide_spec.get("template_key") or slide_spec.get("base_slide_key")
    if template_key and template_key in reference_by_key:
        entry = reference_by_key[template_key]
        return {
            "presentation_path": resolve_path(BASE_DIR, entry["library_path"]),
            "slide_no": entry["library_slide_no"],
            "slide_id": entry["slide_id"],
            "template_key": entry["template_key"],
            "rationale": direct_template_rationale(entry, mode_policy_context),
        }

    if slide_spec.get("slide_selector"):
        if reference_catalog is None:
            raise ValueError("slide_selector requires reference_catalog_path")
        entry, rationale = selector_rationale(
            reference_catalog["slides"],
            slide_spec["slide_selector"],
            mode_policy_context,
            selection_context,
        )
        return {
            "presentation_path": resolve_path(BASE_DIR, entry["library_path"]),
            "slide_no": entry["library_slide_no"],
            "slide_id": entry["slide_id"],
            "template_key": entry["template_key"],
            "rationale": rationale,
        }

    if template_key and template_key in old_slide_map:
        source_template = resolve_path(spec_dir, spec.get("source_template")) or resolve_path(spec_dir, old_catalog["template_path"])  # type: ignore[index]
        return {
            "presentation_path": source_template,
            "slide_no": old_slide_map[template_key],
            "slide_id": None,
            "template_key": template_key,
            "rationale": legacy_source_rationale("legacy_catalog_key", slide_spec, mode_policy_context),
        }

    raise KeyError(f"Unable to resolve base slide for spec: {slide_spec}")


def resolved_slide_sources(
    spec: dict[str, Any],
    spec_dir: Path,
    old_catalog: dict[str, Any] | None,
    old_slide_map: dict[str, int],
    reference_catalog: dict[str, Any] | None,
    reference_by_key: dict[str, dict[str, Any]],
    mode_policy_context: dict[str, Any] | None = None,
) -> list[dict[str, Any]]:
    sources: list[dict[str, Any]] = []
    previous_reference: dict[str, Any] | None = None
    for slide_spec in spec["slides"]:
        resolved = resolve_slide_source(
            slide_spec,
            spec=spec,
            spec_dir=spec_dir,
            old_catalog=old_catalog,
            old_slide_map=old_slide_map,
            reference_catalog=reference_catalog,
            reference_by_key=reference_by_key,
            mode_policy_context=mode_policy_context,
            selection_context={"previous_reference": previous_reference} if previous_reference else None,
        )
        sources.append(resolved)
        selected_reference = reference_by_key.get(str(resolved.get("template_key")))
        if selected_reference:
            previous_reference = selected_reference
    return sources


def validate_recipe_membership(
    spec: dict[str, Any],
    slide_specs: list[dict[str, Any]],
    old_catalog: dict[str, Any] | None,
    old_slide_map: dict[str, int],
) -> None:
    if not spec.get("recipe") or not old_catalog:
        return

    recipe_name = spec["recipe"]
    recipe = set(old_catalog.get("recipes", {}).get(recipe_name, []))
    if not recipe:
        return

    requested_slide_nos = [
        slide_spec["base_slide_no"]
        if "base_slide_no" in slide_spec
        else old_slide_map.get(slide_spec.get("base_slide_key", ""), -1)
        for slide_spec in slide_specs
        if slide_spec.get("base_slide_key") in old_slide_map or "base_slide_no" in slide_spec
    ]
    missing = [slide_no for slide_no in requested_slide_nos if slide_no not in recipe]
    if missing:
        raise ValueError(f"Slides {missing} are not part of recipe '{recipe_name}'")


def add_source_slides(
    prs: Presentation,
    slide_specs: list[dict[str, Any]],
    sources: list[dict[str, Any]],
) -> None:
    source_cache: dict[Path, Presentation] = {}
    for slide_spec, resolved in zip(slide_specs, sources, strict=False):
        if slide_spec.get("source_mode") == "blank":
            prs.slides.add_slide(find_blank_layout(prs))
            continue

        path = resolved["presentation_path"]
        if path is None or not Path(path).exists():
            prs.slides.add_slide(find_blank_layout(prs))
            continue
        if path not in source_cache:
            source_cache[path] = Presentation(str(path))
        src_slide = source_cache[path].slides[resolved["slide_no"] - 1]
        import_slide(prs, src_slide)


def render_slide(
    *,
    slide,
    slide_spec: dict[str, Any],
    resolved: dict[str, Any],
    index: int,
    total_pages: int,
    theme: ThemeConfig,
    blueprint_index: dict[str, dict[str, Any]],
    overflow_events: list[dict[str, Any]],
) -> None:
    blueprint = blueprint_index.get(resolved["slide_id"]) if resolved["slide_id"] else None
    layout = slide_spec["layout"]

    if slide_spec.get("source_mode") == "blank":
        apply_extras(slide, slide_spec, theme)
        return

    if blueprint is None and resolved.get("slide_id") is None:
        apply_extras(slide, slide_spec, theme)
        return

    if layout in {"template_slide", "blueprint_overlay"}:
        if blueprint is None:
            raise ValueError(f"{layout} requires blueprint metadata for {resolved['template_key']}")
        render_template_slide(slide, slide_spec, theme, blueprint, overflow_events=overflow_events)
    else:
        raise ValueError(f"Unsupported non-template layout: {layout}")

    apply_extras(slide, slide_spec, theme)


def write_text_overflow_report(output_path: Path, events: list[dict[str, Any]]) -> None:
    reports_dir = BASE_DIR / "outputs" / "reports"
    reports_dir.mkdir(parents=True, exist_ok=True)
    report_json = reports_dir / f"{output_path.stem}_text_overflow.json"
    report_md = reports_dir / f"{output_path.stem}_text_overflow.md"
    payload = {
        "deck": output_path.as_posix(),
        "summary": {
            "cutoff_events": len(events),
            "affected_slots": len({(event["slide_id"], event["slot"]) for event in events}),
        },
        "events": events,
    }
    report_json.write_text(json.dumps(payload, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")

    lines = [
        "# Text Overflow Policy Report",
        "",
        f"- Deck: `{output_path.name}`",
        f"- Cut-off events: {len(events)}",
        "",
    ]
    if events:
        lines.extend(
            [
                "| Slide ID | Slot | Resolution | LLM | Fallback | Original chars | Budget | Final text |",
                "| --- | --- | --- | --- | --- | ---: | ---: | --- |",
            ]
        )
        for event in events:
            truncated = str(event["truncated_text"]).replace("\n", " / ").replace("|", "\\|")
            llm_state = "used" if event.get("llm_used") else "attempted" if event.get("llm_attempted") else "off"
            lines.append(
                f"| {event['slide_id']} | {event['slot']} | {event.get('resolution', '')} | "
                f"{llm_state} | {event.get('fallback_reason') or ''} | {event['original_chars']} | "
                f"{event['budget']} | {truncated[:120]} |"
            )
    else:
        lines.append("_No cut-off events._")
    report_md.write_text("\n".join(lines) + "\n", encoding="utf-8")


def slide_selection_rationale_rows(
    slide_specs: list[dict[str, Any]],
    sources: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    rows: list[dict[str, Any]] = []
    for slide_no, (slide_spec, resolved) in enumerate(zip(slide_specs, sources, strict=False), start=1):
        rationale = dict(resolved.get("rationale", {}))
        rows.append(
            {
                "slide_number": slide_no,
                "layout": slide_spec.get("layout"),
                "slide_id": resolved.get("slide_id"),
                "template_key": resolved.get("template_key"),
                "purpose": rationale.get("purpose"),
                "scope": rationale.get("scope"),
                "selected_template_key": rationale.get("selected_template_key"),
                "selected_library": rationale.get("selected_library"),
                "selected_variant": rationale.get("selected_variant"),
                "source_type": rationale.get("source_type"),
                "candidate_templates": rationale.get("candidate_templates", []),
                "selection_factors": rationale.get("selection_factors", {}),
                "mode_policy": rationale.get("mode_policy"),
                "rejected_candidates": rationale.get("rejected_candidates", []),
                "confidence": rationale.get("confidence"),
                "notes": rationale.get("notes", ""),
            }
        )
    return rows


def write_slide_selection_rationale_report(
    output_path: Path,
    spec: dict[str, Any],
    slide_specs: list[dict[str, Any]],
    sources: list[dict[str, Any]],
) -> None:
    reports_dir = BASE_DIR / "outputs" / "reports"
    reports_dir.mkdir(parents=True, exist_ok=True)
    rows = slide_selection_rationale_rows(slide_specs, sources)
    by_source_type: dict[str, int] = {}
    by_mode_policy: dict[str, int] = {}
    for row in rows:
        source_type = str(row.get("source_type") or "unknown")
        by_source_type[source_type] = by_source_type.get(source_type, 0) + 1
        mode_policy = row.get("mode_policy") or {}
        mode = str(mode_policy.get("mode") or "unrecorded")
        by_mode_policy[mode] = by_mode_policy.get(mode, 0) + 1

    payload = {
        "deck": output_path.as_posix(),
        "deck_name": spec.get("name"),
        "summary": {
            "slides": len(rows),
            "by_source_type": by_source_type,
            "by_mode_policy": by_mode_policy,
        },
        "slides": rows,
    }

    report_json = reports_dir / f"{output_path.stem}_slide_selection_rationale.json"
    report_md = reports_dir / f"{output_path.stem}_slide_selection_rationale.md"
    latest_json = reports_dir / "slide_selection_rationale.json"
    latest_md = reports_dir / "slide_selection_rationale.md"
    for path in (report_json, latest_json):
        path.write_text(json.dumps(payload, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")

    lines = [
        "# Slide Selection Rationale",
        "",
        f"- Deck: `{output_path.name}`",
        f"- Slides: {len(rows)}",
        f"- Mode policy: `{spec.get('mode_policy', 'auto')}`",
        "",
        "| Slide | Source | Purpose | Scope | Selected template | Library | Variant | Confidence |",
        "| ---: | --- | --- | --- | --- | --- | --- | ---: |",
    ]
    for row in rows:
        confidence = row.get("confidence")
        confidence_text = "" if confidence is None else f"{float(confidence):.2f}"
        lines.append(
            f"| {row['slide_number']} | {row.get('source_type') or ''} | {row.get('purpose') or ''} | "
            f"{row.get('scope') or ''} | {row.get('selected_template_key') or ''} | "
            f"{row.get('selected_library') or ''} | {row.get('selected_variant') or ''} | {confidence_text} |"
        )

    lines.extend(["", "## Notes", ""])
    for row in rows:
        candidate_count = len(row.get("candidate_templates", []))
        rejected_count = len(row.get("rejected_candidates", []))
        mode_policy = row.get("mode_policy") or {}
        mode_suffix = f" Mode policy: {mode_policy.get('mode')}." if mode_policy.get("mode") else ""
        lines.append(
            f"- Slide {row['slide_number']}: {row.get('notes') or 'No notes.'} "
            f"Candidates recorded: {candidate_count}; rejected candidates recorded: {rejected_count}.{mode_suffix}"
        )

    markdown = "\n".join(lines) + "\n"
    for path in (report_md, latest_md):
        path.write_text(markdown, encoding="utf-8")


def slot_bounds(slot_def: dict[str, Any]) -> dict[str, Any] | None:
    bounds = slot_def.get("bounds")
    if not bounds:
        return None
    return {
        "left": bounds.get("left"),
        "top": bounds.get("top"),
        "width": bounds.get("width"),
        "height": bounds.get("height"),
    }


def text_values_for_slide(slide_spec: dict[str, Any], blueprint: dict[str, Any]) -> tuple[dict[str, Any], dict[str, str]]:
    text_values = dict(slide_spec.get("text_slots", {}))
    source_kind_by_slot = {slot_name: "text_slots" for slot_name in text_values}
    sequential_values = slide_spec.get("text_values", [])
    sequential_slots = [
        item["slot"]
        for item in blueprint.get("editable_text_slots", [])
        if item["slot"] not in text_values
    ]
    for slot_name, value in zip(sequential_slots, sequential_values, strict=False):
        text_values[slot_name] = value
        source_kind_by_slot[slot_name] = "text_values"
    return text_values, source_kind_by_slot


def deck_slot_map_rows(
    slide_specs: list[dict[str, Any]],
    sources: list[dict[str, Any]],
    blueprint_index: dict[str, dict[str, Any]],
) -> list[dict[str, Any]]:
    rows: list[dict[str, Any]] = []
    for slide_no, (slide_spec, resolved) in enumerate(zip(slide_specs, sources, strict=False), start=1):
        slide_id = resolved.get("slide_id")
        if not slide_id:
            continue
        blueprint = blueprint_index.get(slide_id)
        if not blueprint:
            continue

        text_values, source_kind_by_slot = text_values_for_slide(slide_spec, blueprint)
        clear_unfilled_slots = slide_spec.get("clear_unfilled_slots", True)
        for slot_def in blueprint.get("editable_text_slots", []):
            slot_name = slot_def["slot"]
            override = slide_spec.get("slot_overrides", {}).get(slot_name, {})
            current_value = text_values.get(slot_name, "" if clear_unfilled_slots else None)
            rows.append(
                {
                    "slide_number": slide_no,
                    "slide_id": slide_id,
                    "template_key": resolved.get("template_key"),
                    "slot_name": slot_name,
                    "slot_kind": "text",
                    "current_value": current_value,
                    "value_source": source_kind_by_slot.get(slot_name, "unfilled"),
                    "bounds": slot_bounds(slot_def),
                    "font_role": override.get("font_role") or slot_def.get("font_role"),
                    "fit_strategy": override.get("fit_strategy") or slot_def.get("fit_strategy") or "preserve_template",
                    "budget": slot_text_budget(slot_def, override),
                    "shape_name": slot_def.get("shape_name"),
                    "shape_index": slot_def.get("shape_index"),
                }
            )

        image_values = slide_spec.get("image_slots", {})
        for slot_def in blueprint.get("editable_image_slots", []):
            slot_name = slot_def["slot"]
            rows.append(
                {
                    "slide_number": slide_no,
                    "slide_id": slide_id,
                    "template_key": resolved.get("template_key"),
                    "slot_name": slot_name,
                    "slot_kind": "image",
                    "current_value": image_values.get(slot_name),
                    "value_source": "image_slots" if slot_name in image_values else "unfilled",
                    "bounds": slot_bounds(slot_def),
                    "font_role": None,
                    "fit_strategy": None,
                    "budget": None,
                    "shape_name": slot_def.get("shape_name"),
                    "shape_index": slot_def.get("shape_index"),
                }
            )

        chart_values = slide_spec.get("chart_slots", {})
        for slot_def in blueprint.get("editable_chart_slots", []):
            slot_name = slot_def["slot"]
            rows.append(
                {
                    "slide_number": slide_no,
                    "slide_id": slide_id,
                    "template_key": resolved.get("template_key"),
                    "slot_name": slot_name,
                    "slot_kind": "chart",
                    "current_value": chart_values.get(slot_name),
                    "value_source": "chart_slots" if slot_name in chart_values else "unfilled",
                    "bounds": slot_bounds(slot_def),
                    "font_role": None,
                    "fit_strategy": None,
                    "budget": None,
                    "shape_name": slot_def.get("shape_name"),
                    "shape_index": slot_def.get("shape_index"),
                }
            )
        table_values = slide_spec.get("table_slots", {})
        for slot_def in blueprint.get("editable_table_slots", []):
            slot_name = slot_def["slot"]
            rows.append(
                {
                    "slide_number": slide_no,
                    "slide_id": slide_id,
                    "template_key": resolved.get("template_key"),
                    "slot_name": slot_name,
                    "slot_kind": "table",
                    "current_value": table_values.get(slot_name),
                    "value_source": "table_slots" if slot_name in table_values else "unfilled",
                    "bounds": slot_bounds(slot_def),
                    "font_role": None,
                    "fit_strategy": None,
                    "budget": None,
                    "shape_name": slot_def.get("shape_name"),
                    "shape_index": slot_def.get("shape_index"),
                }
            )
    return rows


def value_preview(value: Any, limit: int = 80) -> str:
    if value is None:
        return ""
    if isinstance(value, (dict, list)):
        text = json.dumps(value, ensure_ascii=False)
    else:
        text = str(value)
    text = " ".join(text.split())
    text = text.replace("|", "\\|")
    return text if len(text) <= limit else text[: limit - 3].rstrip() + "..."


def write_deck_slot_map_report(
    output_path: Path,
    spec: dict[str, Any],
    slide_specs: list[dict[str, Any]],
    sources: list[dict[str, Any]],
    blueprint_index: dict[str, dict[str, Any]],
) -> None:
    reports_dir = BASE_DIR / "outputs" / "reports"
    reports_dir.mkdir(parents=True, exist_ok=True)
    rows = deck_slot_map_rows(slide_specs, sources, blueprint_index)
    by_kind: dict[str, int] = {}
    filled_slots = 0
    for row in rows:
        by_kind[row["slot_kind"]] = by_kind.get(row["slot_kind"], 0) + 1
        if row.get("current_value") not in (None, ""):
            filled_slots += 1

    payload = {
        "deck": output_path.as_posix(),
        "deck_name": spec.get("name"),
        "summary": {
            "slides": len(slide_specs),
            "mapped_slots": len(rows),
            "filled_slots": filled_slots,
            "by_kind": by_kind,
        },
        "slots": rows,
    }

    report_json = reports_dir / f"{output_path.stem}_deck_slot_map.json"
    report_md = reports_dir / f"{output_path.stem}_deck_slot_map.md"
    latest_json = reports_dir / "deck_slot_map.json"
    latest_md = reports_dir / "deck_slot_map.md"
    for path in (report_json, latest_json):
        path.write_text(json.dumps(payload, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")

    lines = [
        "# Deck Slot Map",
        "",
        f"- Deck: `{output_path.name}`",
        f"- Slides: {len(slide_specs)}",
        f"- Mapped slots: {len(rows)}",
        f"- Filled slots: {filled_slots}",
        "",
        "| Slide | Template | Kind | Slot | Value source | Budget | Current value |",
        "| ---: | --- | --- | --- | --- | ---: | --- |",
    ]
    for row in rows:
        budget = "" if row.get("budget") is None else str(row["budget"])
        lines.append(
            f"| {row['slide_number']} | {row.get('template_key') or ''} | {row['slot_kind']} | "
            f"{row['slot_name']} | {row.get('value_source') or ''} | {budget} | "
            f"{value_preview(row.get('current_value'))} |"
        )
    markdown = "\n".join(lines) + "\n"
    for path in (report_md, latest_md):
        path.write_text(markdown, encoding="utf-8")


def build_deck_from_spec(
    spec_path: str | Path,
    run_id: str | None = None,
    project_id: str | None = None,
) -> Path:
    """Build the deck and deck-specific reports.

    run_id and project_id are accepted for backward compatibility, but output
    mirroring is owned by CLI entrypoints so manifests are written once.
    """
    spec_path = Path(spec_path).resolve()
    spec, spec_dir = load_spec(spec_path)
    validate_supported_aspect_ratio(spec)
    workspace_usage = apply_workspace_asset_intents(spec, workspace_asset_root())

    old_catalog_path = resolve_path(spec_dir, spec.get("catalog_path"))
    reference_catalog_path = resolve_path(spec_dir, spec.get("reference_catalog_path", "../../config/reference_catalog.json"))
    blueprint_path = resolve_path(spec_dir, spec.get("blueprint_path", "../../config/template_blueprints.json"))

    old_catalog = load_json(old_catalog_path)
    reference_catalog = load_json(reference_catalog_path)
    apply_template_design_dna(reference_catalog)
    template_text_dna_cleanup = load_template_text_dna_cleanup()
    mode_policy_context = load_mode_policy_context(spec)
    blueprints = load_blueprints(blueprint_path) if blueprint_path else {"slides": {}}
    old_slide_map = old_catalog_slide_map(old_catalog)
    reference_by_key, _ = indexed_reference_catalog(reference_catalog)

    theme = ThemeConfig.from_json(resolve_path(spec_dir, spec["theme_path"]))
    output_path = require_repo_output_path(resolve_path(spec_dir, spec["output_path"]), label="deck output_path")
    output_path.parent.mkdir(parents=True, exist_ok=True)

    slide_specs = spec["slides"]
    validate_recipe_membership(spec, slide_specs, old_catalog, old_slide_map)
    sources = resolved_slide_sources(
        spec,
        spec_dir,
        old_catalog,
        old_slide_map,
        reference_catalog,
        reference_by_key,
        mode_policy_context,
    )

    prs = Presentation()
    configure_presentation_size(prs, spec)
    add_source_slides(prs, slide_specs, sources)

    if old_catalog:
        replace_placeholder_text(prs, old_catalog.get("placeholder_translation", {}))
    set_all_fonts(prs, theme.font_family)

    blueprint_index = blueprints.get("slides", {})
    total_pages = len(prs.slides)
    overflow_events: list[dict[str, Any]] = []
    render_slide_specs = [resolve_slide_asset_paths(slide_spec, spec_dir) for slide_spec in slide_specs]
    for index, (slide_spec, resolved) in enumerate(zip(render_slide_specs, sources, strict=False)):
        slide_spec = slide_spec_with_text_dna_cleanup(slide_spec, resolved, template_text_dna_cleanup)
        render_slide(
            slide=prs.slides[index],
            slide_spec=slide_spec,
            resolved=resolved,
            index=index,
            total_pages=total_pages,
            theme=theme,
            blueprint_index=blueprint_index,
            overflow_events=overflow_events,
        )

    prs.save(str(output_path))
    apply_theme_accent_overrides(output_path, spec.get("theme_accent_overrides"))
    write_text_overflow_report(output_path, overflow_events)
    write_slide_selection_rationale_report(output_path, spec, slide_specs, sources)
    write_deck_slot_map_report(output_path, spec, slide_specs, sources, blueprint_index)
    write_asset_usage_report(output_path, spec, workspace_usage)
    return output_path


def main(argv: list[str] | None = None) -> int:
    argv = argv or sys.argv[1:]
    parser = argparse.ArgumentParser()
    parser.add_argument("spec_path", nargs="?", default=str(WORKDIR / "data" / "specs" / "jb_meeting_deck_spec.json"))
    args = parser.parse_args(argv)
    spec_path = Path(args.spec_path).resolve()
    output = build_deck_from_spec(spec_path)
    print(output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

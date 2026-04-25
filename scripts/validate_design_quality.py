from __future__ import annotations

import json
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation

BASE_DIR = Path(__file__).resolve().parents[1]
if str(BASE_DIR) not in sys.path:
    sys.path.insert(0, str(BASE_DIR))

from scripts.build_deck import (
    indexed_reference_catalog,
    load_json,
    load_spec,
    old_catalog_slide_map,
    resolve_path,
    resolved_slide_sources,
)
from system.blueprint_loader import load_blueprints
from system.pptx_system import EMU_PER_INCH

PLACEHOLDER_PATTERNS = [
    "please fill",
    "please write",
    "please enter",
    "input the",
    "provide a concise",
    "provide a brief",
    "write a detailed",
    "copyrights",
    "office suite",
    "lorem",
    "click to",
]


def shape_bounds(shape) -> dict[str, float]:
    return {
        "left": shape.left / EMU_PER_INCH,
        "top": shape.top / EMU_PER_INCH,
        "width": shape.width / EMU_PER_INCH,
        "height": shape.height / EMU_PER_INCH,
    }


def font_size(shape) -> float:
    sizes: list[float] = []
    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size is not None:
                    sizes.append(run.font.size.pt)
    return max(sizes) if sizes else 12.0


def estimated_capacity(bounds: dict[str, float], size: float) -> tuple[int, int]:
    # Conservative Korean capacity estimate for 16:9 slides.
    chars_per_line = max(5, int((bounds["width"] * 96) / max(size * 0.88, 1)))
    lines = max(1, int((bounds["height"] * 72) / max(size * 1.25, 1)))
    return chars_per_line, lines


def normalize_text(value: str) -> str:
    return " ".join((value or "").split())


def inspect_shape(slide_no: int, shape_no: int, shape) -> list[dict[str, Any]]:
    issues: list[dict[str, Any]] = []
    text = getattr(shape, "text", "")
    normalized = " ".join(text.split())
    if not normalized:
        return issues

    lowered = normalized.lower()
    for pattern in PLACEHOLDER_PATTERNS:
        if pattern in lowered:
            issues.append(
                {
                    "severity": "error",
                    "slide": slide_no,
                    "shape": shape_no,
                    "type": "placeholder",
                    "message": f"Residual placeholder pattern found: {pattern}",
                    "text": normalized[:120],
                }
            )

    bounds = shape_bounds(shape)
    size = font_size(shape)
    chars_per_line, max_lines = estimated_capacity(bounds, size)
    actual_lines = normalized.count("\n") + 1
    estimated_lines = max(actual_lines, (len(normalized) + chars_per_line - 1) // chars_per_line)
    # Strict template slots are governed by blueprint budgets, overflow reports,
    # readback diagnostics, and visual smoke. The generic geometry heuristic is
    # intentionally reserved for non-slot static text because some imported PPTX
    # slot shapes expose misleading bounds through python-pptx.
    is_strict_slot_shape = str(getattr(shape, "name", "")).startswith("slot:")
    if not is_strict_slot_shape and estimated_lines > max_lines + 1:
        issues.append(
            {
                "severity": "warning",
                "slide": slide_no,
                "shape": shape_no,
                "type": "text_overflow_risk",
                "message": f"Estimated {estimated_lines} lines for a box that safely fits about {max_lines}.",
                "text": normalized[:120],
            }
        )

    if bounds["top"] < 0.8 and len(normalized) > 42 and size >= 12:
        issues.append(
            {
                "severity": "warning",
                "slide": slide_no,
                "shape": shape_no,
                "type": "header_collision_risk",
                "message": "Long text appears inside the header band.",
                "text": normalized[:120],
            }
        )

    if bounds["top"] > 6.35 and len(normalized) > 95:
        issues.append(
            {
                "severity": "warning",
                "slide": slide_no,
                "shape": shape_no,
                "type": "footer_overload",
                "message": "Footer text is long and may collide with page number or footer bar.",
                "text": normalized[:120],
            }
        )

    return issues


def template_expected_texts(spec_path: Path) -> list[tuple[int, str, str]]:
    spec, spec_dir = load_spec(spec_path)
    old_catalog_path = resolve_path(spec_dir, spec.get("catalog_path"))
    reference_catalog_path = resolve_path(spec_dir, spec.get("reference_catalog_path", "../../config/reference_catalog.json"))
    blueprint_path = resolve_path(spec_dir, spec.get("blueprint_path", "../../config/template_blueprints.json"))

    old_catalog = load_json(old_catalog_path)
    reference_catalog = load_json(reference_catalog_path)
    blueprints = load_blueprints(blueprint_path) if blueprint_path else {"slides": {}}
    old_slide_map = old_catalog_slide_map(old_catalog)
    reference_by_key, _ = indexed_reference_catalog(reference_catalog)
    sources = resolved_slide_sources(
        spec,
        spec_dir,
        old_catalog,
        old_slide_map,
        reference_catalog,
        reference_by_key,
    )

    expected: list[tuple[int, str, str]] = []
    for slide_no, (slide_spec, resolved) in enumerate(zip(spec["slides"], sources, strict=False), start=1):
        if slide_spec["layout"] != "template_slide" or not resolved.get("slide_id"):
            continue
        blueprint = blueprints.get("slides", {}).get(resolved["slide_id"], {})
        text_slots = {item["slot"]: item for item in blueprint.get("editable_text_slots", [])}
        text_values = dict(slide_spec.get("text_slots", {}))
        sequential_values = slide_spec.get("text_values", [])
        sequential_slots = [item["slot"] for item in blueprint.get("editable_text_slots", []) if item["slot"] not in text_values]
        for slot_name, value in zip(sequential_slots, sequential_values, strict=False):
            text_values[slot_name] = value

        for slot_name, value in text_values.items():
            if not value or "\n" in value or slot_name not in text_slots:
                continue
            override = slide_spec.get("slot_overrides", {}).get(slot_name, {})
            fit_strategy = override.get("fit_strategy") or text_slots[slot_name].get("fit_strategy") or "preserve_template"
            if fit_strategy != "manual_wrap":
                expected.append((slide_no, slot_name, value))
    return expected


def inspect_template_wrap(prs: Presentation, spec_path: Path) -> list[dict[str, Any]]:
    issues: list[dict[str, Any]] = []
    for slide_no, slot_name, expected in template_expected_texts(spec_path):
        if slide_no > len(prs.slides):
            issues.append(
                {
                    "severity": "error",
                    "slide": slide_no,
                    "shape": None,
                    "type": "template_slot_missing_slide",
                    "message": f"Spec references template slide {slide_no}, but output has fewer slides.",
                    "text": slot_name,
                }
            )
            continue
        expected_normalized = normalize_text(expected)
        matching_texts = [
            getattr(shape, "text", "")
            for shape in prs.slides[slide_no - 1].shapes
            if normalize_text(getattr(shape, "text", "")) == expected_normalized
        ]
        if any("\n" in text for text in matching_texts):
            issues.append(
                {
                    "severity": "error",
                    "slide": slide_no,
                    "shape": None,
                    "type": "template_forced_wrap",
                    "message": f"Template slot '{slot_name}' appears to contain code-inserted line breaks.",
                    "text": expected[:120],
                }
            )
        elif not matching_texts:
            issues.append(
                {
                    "severity": "warning",
                    "slide": slide_no,
                    "shape": None,
                    "type": "template_slot_readback_missing",
                    "message": f"Could not read back exact template slot value for '{slot_name}'.",
                    "text": expected[:120],
                }
            )
    return issues


def validate(path: Path, template_spec_path: Path | None = None) -> dict[str, Any]:
    prs = Presentation(str(path))
    issues: list[dict[str, Any]] = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        for shape_no, shape in enumerate(slide.shapes, start=1):
            issues.extend(inspect_shape(slide_no, shape_no, shape))
    if template_spec_path is not None:
        issues.extend(inspect_template_wrap(prs, template_spec_path))

    return {
        "file": str(path),
        "slide_count": len(prs.slides),
        "issue_count": len(issues),
        "errors": sum(1 for issue in issues if issue["severity"] == "error"),
        "warnings": sum(1 for issue in issues if issue["severity"] == "warning"),
        "issues": issues,
    }


def main(argv: list[str] | None = None) -> int:
    argv = argv or sys.argv[1:]
    if not argv:
        print(
            "Usage: python scripts/validate_design_quality.py <pptx> [output_json] [--template-spec <spec_json>]",
            file=sys.stderr,
        )
        return 2
    template_spec_path = None
    if "--template-spec" in argv:
        flag_index = argv.index("--template-spec")
        try:
            template_spec_path = Path(argv[flag_index + 1]).resolve()
        except IndexError:
            print("--template-spec requires a path", file=sys.stderr)
            return 2
        argv = argv[:flag_index] + argv[flag_index + 2 :]

    pptx_path = Path(argv[0]).resolve()
    result = validate(pptx_path, template_spec_path)
    if len(argv) >= 2:
        output_path = Path(argv[1]).resolve()
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(json.dumps(result, indent=2, ensure_ascii=False), encoding="utf-8")
        print(output_path)
    else:
        print(json.dumps(result, indent=2, ensure_ascii=False))
    return 1 if result["errors"] else 0


if __name__ == "__main__":
    raise SystemExit(main())

from __future__ import annotations

import argparse
import json
import subprocess
import sys
from pathlib import Path
from typing import Any

import fitz
from PIL import Image, ImageDraw
from pptx import Presentation

BASE_DIR = Path(__file__).resolve().parents[1]
REFERENCE_CATALOG_PATH = BASE_DIR / "config" / "reference_catalog.json"
OUTPUT_DIR = BASE_DIR / "outputs" / "previews" / "template_index"
REPORT_PATH = BASE_DIR / "outputs" / "reports" / "template_thumbnail_index.json"


def soffice_path() -> Path:
    candidates = [
        Path("C:/Program Files/LibreOffice/program/soffice.exe"),
        Path("C:/Program Files (x86)/LibreOffice/program/soffice.exe"),
    ]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    raise FileNotFoundError("LibreOffice soffice.exe not found")


def render_pdf_to_pngs(pdf_path: Path, output_dir: Path) -> list[Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(pdf_path)
    created: list[Path] = []
    for index in range(len(doc)):
        page = doc.load_page(index)
        pix = page.get_pixmap(matrix=fitz.Matrix(1.6, 1.6), alpha=False)
        output_path = output_dir / f"slide_{index + 1:02d}.png"
        pix.save(output_path)
        created.append(output_path)
    return created


def extract_slide_preview_lines(slide) -> list[str]:
    lines: list[str] = []
    for shape in slide.shapes:
        text = getattr(shape, "text", "")
        normalized = " ".join(text.split())
        if normalized:
            lines.append(normalized)
    return lines[:6]


def render_text_previews(library_path: Path, output_dir: Path, slides: list[dict[str, str]]) -> list[Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(library_path))
    created: list[Path] = []
    for index, slide_meta in enumerate(slides, start=1):
        slide = prs.slides[index - 1]
        image = Image.new("RGB", (1280, 720), color=(16, 33, 59))
        draw = ImageDraw.Draw(image)
        draw.rectangle((40, 40, 1240, 680), fill=(255, 255, 255))
        draw.text((80, 80), slide_meta["template_key"], fill=(16, 33, 59))
        draw.text((80, 120), f"{slide_meta['purpose']} | {slide_meta['variant']}", fill=(29, 73, 243))
        draw.text((80, 160), f"Fallback preview from PPT text | {library_path.name}", fill=(122, 132, 148))

        y = 230
        for line in extract_slide_preview_lines(slide):
            draw.text((100, y), f"- {line[:120]}", fill=(39, 39, 39))
            y += 52

        output_path = output_dir / f"slide_{index:02d}.png"
        image.save(output_path)
        created.append(output_path)
    return created


def build_contact_sheet(image_paths: list[Path], output_path: Path, labels: list[str]) -> None:
    if not image_paths:
        return
    images = [Image.open(path).convert("RGB") for path in image_paths]
    thumb_width = 360
    thumb_height = int(images[0].height * (thumb_width / images[0].width))
    columns = 3
    padding = 24
    label_height = 60
    rows = (len(images) + columns - 1) // columns
    canvas = Image.new(
        "RGB",
        (
            columns * thumb_width + (columns + 1) * padding,
            rows * (thumb_height + label_height) + (rows + 1) * padding,
        ),
        color=(246, 248, 252),
    )
    draw = ImageDraw.Draw(canvas)

    for index, image in enumerate(images):
        row = index // columns
        col = index % columns
        x = padding + col * (thumb_width + padding)
        y = padding + row * (thumb_height + label_height + padding)
        resized = image.resize((thumb_width, thumb_height))
        canvas.paste(resized, (x, y))
        draw.rectangle((x, y + thumb_height, x + thumb_width, y + thumb_height + label_height), fill=(255, 255, 255))
        draw.text((x + 12, y + thumb_height + 12), labels[index], fill=(16, 33, 59))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    canvas.save(output_path)


def convert_pptx_to_pdf(input_path: Path, output_dir: Path) -> Path:
    output_dir.mkdir(parents=True, exist_ok=True)
    process = subprocess.run(
        [
            str(soffice_path()),
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(input_path),
        ],
        check=True,
        capture_output=True,
        text=True,
    )
    pdf_path = output_dir / f"{input_path.stem}.pdf"
    if not pdf_path.exists():
        raise FileNotFoundError(f"Failed to render {input_path.name} to PDF: {process.stdout} {process.stderr}")
    return pdf_path


def build_thumbnail_index(
    *,
    reference_catalog_path: Path = REFERENCE_CATALOG_PATH,
    output_dir: Path = OUTPUT_DIR,
    report_path: Path = REPORT_PATH,
    check: bool = False,
) -> dict[str, Any]:
    reference_catalog = json.loads(reference_catalog_path.read_text(encoding="utf-8"))
    output_dir.mkdir(parents=True, exist_ok=True)
    markdown_lines = ["# Template Index", ""]
    report: dict[str, Any] = {
        "schema_version": "1.0",
        "reference_catalog_path": str(reference_catalog_path.relative_to(BASE_DIR)),
        "output_dir": str(output_dir.relative_to(BASE_DIR)),
        "libraries": [],
        "slides": [],
        "summary": {
            "libraries": 0,
            "slides": 0,
            "image_paths": 0,
            "contact_sheets": 0,
            "fallback_libraries": 0,
        },
    }

    slides_by_library: dict[str, list[dict[str, str]]] = {}
    library_lookup = {library["library_id"]: library for library in reference_catalog["libraries"]}
    for slide in reference_catalog["slides"]:
        slides_by_library.setdefault(slide["library_id"], []).append(slide)

    for library_id, slides in slides_by_library.items():
        library = library_lookup[library_id]
        library_path = BASE_DIR / Path(library["library_path"])
        pdf_dir = output_dir / "pdf"
        image_dir = output_dir / library_id
        fallback_reason = ""
        try:
            pdf_path = convert_pptx_to_pdf(library_path, pdf_dir)
            image_paths = render_pdf_to_pngs(pdf_path, image_dir)
            render_mode = "libreoffice_pdf"
        except Exception as exc:
            image_paths = render_text_previews(library_path, image_dir, slides)
            render_mode = "text_fallback"
            fallback_reason = str(exc)
        labels = [f"{slide['template_key']} | {slide['purpose']} | {slide['variant']}" for slide in slides]
        build_contact_sheet(image_paths, image_dir / "contact_sheet.png", labels)
        contact_sheet = image_dir / "contact_sheet.png"

        library_report = {
            "library_id": library_id,
            "library_path": str(Path(library["library_path"])),
            "scope": library.get("scope"),
            "render_mode": render_mode,
            "fallback_reason": fallback_reason,
            "slide_count": len(slides),
            "image_count": len(image_paths),
            "contact_sheet": str(contact_sheet.relative_to(BASE_DIR)),
        }
        report["libraries"].append(library_report)
        if render_mode == "text_fallback":
            report["summary"]["fallback_libraries"] += 1
        report["summary"]["contact_sheets"] += 1 if contact_sheet.exists() else 0
        report["summary"]["image_paths"] += len(image_paths)

        markdown_lines.extend(
            [
                f"## {library_id}",
                "",
                f"- Scope: `{library['scope']}`",
                f"- Source: `{library['source_path']}`",
                f"- Contact sheet: `{(image_dir / 'contact_sheet.png').as_posix()}`",
                "",
            ]
        )
        for index, slide in enumerate(slides, start=1):
            thumbnail_path = image_dir / f"slide_{index:02d}.png"
            report["slides"].append(
                {
                    "slide_id": slide["slide_id"],
                    "template_key": slide["template_key"],
                    "library_id": library_id,
                    "library_slide_no": slide.get("library_slide_no", index),
                    "purpose": slide.get("purpose"),
                    "variant": slide.get("variant"),
                    "scope": slide.get("scope"),
                    "density": slide.get("density"),
                    "usage_policy": slide.get("usage_policy"),
                    "thumbnail_path": str(thumbnail_path.relative_to(BASE_DIR)),
                    "thumbnail_exists": thumbnail_path.exists(),
                }
            )
            markdown_lines.append(
                f"- `{index:02d}` `{slide['template_key']}` | purpose=`{slide['purpose']}` | variant=`{slide['variant']}` | density=`{slide['density']}`"
            )
        markdown_lines.append("")

    index_path = output_dir / "INDEX.md"
    index_path.write_text("\n".join(markdown_lines), encoding="utf-8")
    report["index_path"] = str(index_path.relative_to(BASE_DIR))
    report["summary"]["libraries"] = len(report["libraries"])
    report["summary"]["slides"] = len(report["slides"])
    report_path.parent.mkdir(parents=True, exist_ok=True)
    report_path.write_text(json.dumps(report, indent=2, ensure_ascii=False), encoding="utf-8")

    if check:
        expected_slides = len(reference_catalog.get("slides", []))
        missing_thumbnails = [slide["slide_id"] for slide in report["slides"] if not slide["thumbnail_exists"]]
        if not index_path.exists():
            raise FileNotFoundError(f"Template thumbnail index missing: {index_path}")
        if report["summary"]["slides"] != expected_slides:
            raise AssertionError(
                f"Thumbnail slide count mismatch: {report['summary']['slides']} != {expected_slides}"
            )
        if missing_thumbnails:
            raise AssertionError(f"Missing template thumbnails: {missing_thumbnails[:5]}")
        if report["summary"]["contact_sheets"] != len(slides_by_library):
            raise AssertionError(
                f"Contact sheet count mismatch: {report['summary']['contact_sheets']} != {len(slides_by_library)}"
            )
    return report


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Render or validate the production template thumbnail index.")
    parser.add_argument("--reference-catalog", default=str(REFERENCE_CATALOG_PATH.relative_to(BASE_DIR)))
    parser.add_argument("--output-dir", default=str(OUTPUT_DIR.relative_to(BASE_DIR)))
    parser.add_argument("--report-json", default=str(REPORT_PATH.relative_to(BASE_DIR)))
    parser.add_argument("--check", action="store_true", help="Render and assert thumbnail/contact-sheet completeness.")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    report = build_thumbnail_index(
        reference_catalog_path=(BASE_DIR / args.reference_catalog).resolve(),
        output_dir=(BASE_DIR / args.output_dir).resolve(),
        report_path=(BASE_DIR / args.report_json).resolve(),
        check=args.check,
    )
    print(BASE_DIR / report["index_path"])
    print(f"template_thumbnail_slides={report['summary']['slides']}")
    print(f"template_thumbnail_fallback_libraries={report['summary']['fallback_libraries']}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

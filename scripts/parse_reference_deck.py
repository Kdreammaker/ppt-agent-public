from __future__ import annotations

import argparse
import json
import shutil
import sys
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE_DIR = Path(__file__).resolve().parents[1]
if str(BASE_DIR) not in sys.path:
    sys.path.insert(0, str(BASE_DIR))

from scripts.reference_pipeline import (
    canonical_slide_id,
    file_sha256,
    find_reference_record_for_path,
    slide_no_from_asset_stem,
    update_reference_record_status,
    utc_now,
    write_json,
)
from scripts.validate_visual_smoke import render_pptx


REFERENCE_ROOT = BASE_DIR / "assets" / "slides" / "references"
DEFAULT_CURATION_ROOT = BASE_DIR / "assets" / "slides" / "categorizing"
BUCKETS = [
    "unlabeled",
    "good_candidate",
    "normal_candidate",
    "weak_candidate",
]


def slugify(value: str) -> str:
    slug = "".join(char.lower() if char.isalnum() else "_" for char in value).strip("_")
    while "__" in slug:
        slug = slug.replace("__", "_")
    return slug or "reference_deck"


def base_relative(path: Path) -> str:
    return path.resolve().relative_to(BASE_DIR).as_posix()


def ensure_safe_curation_path(path: Path, root: Path) -> None:
    resolved = path.resolve()
    root_resolved = root.resolve()
    if resolved != root_resolved and root_resolved not in resolved.parents:
        raise ValueError(f"Refusing to write outside curation root: {resolved}")


def make_bucket_dirs(curation_root: Path) -> None:
    for media_kind in ("png", "pptx"):
        for bucket in BUCKETS:
            (curation_root / media_kind / bucket).mkdir(parents=True, exist_ok=True)
    (curation_root / "records").mkdir(parents=True, exist_ok=True)


def asset_stem(deck_id: str, slide_no: int) -> str:
    return canonical_slide_id(deck_id, slide_no)


def matching_deck_assets(curation_root: Path, deck_id: str) -> list[Path]:
    paths: list[Path] = []
    for media_kind, suffix in (("png", ".png"), ("pptx", ".pptx")):
        for bucket in BUCKETS:
            bucket_dir = curation_root / media_kind / bucket
            paths.extend(sorted(bucket_dir.glob(f"{deck_id}_s*{suffix}")))
            paths.extend(sorted(bucket_dir.glob(f"{deck_id}__slide_*{suffix}")))
    records_dir = curation_root / "records"
    paths.extend(sorted(records_dir.glob(f"{deck_id}_*")))
    return paths


def prepare_deck_outputs(curation_root: Path, deck_id: str, *, force: bool) -> None:
    existing = [path for path in matching_deck_assets(curation_root, deck_id) if path.exists()]
    if not existing:
        return
    if not force:
        sample = existing[0].relative_to(BASE_DIR).as_posix()
        raise FileExistsError(f"Existing parsed assets found for {deck_id}: {sample}. Use --force to rebuild them.")
    for path in existing:
        if path.is_file():
            path.unlink()


def slide_text_preview(slide: Any) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        text = " ".join(getattr(shape, "text", "").split())
        if text:
            texts.append(text)
    return texts[:8]


def slide_text_joined(slide: Any) -> str:
    return " ".join(" ".join(getattr(shape, "text", "").split()) for shape in slide.shapes).strip()


def infer_story_position(slide_no: int, slide_count: int) -> str:
    if slide_no == 1:
        return "opening"
    if slide_no == slide_count:
        return "closing"
    if slide_no <= max(2, int(slide_count * 0.25)):
        return "early"
    if slide_no >= max(1, int(slide_count * 0.75)):
        return "late"
    return "middle"


def guess_story_role(slide_no: int, slide_count: int, texts: list[str]) -> str:
    joined = " ".join(texts).casefold()
    if slide_no == 1:
        return "cover"
    if slide_no == slide_count:
        return "closing"
    if any(term in joined for term in ("agenda", "contents", "index", "목차")):
        return "toc"
    if any(term in joined for term in ("case", "사례", "result", "성과")):
        return "proof"
    if any(term in joined for term in ("process", "timeline", "roadmap", "프로세스")):
        return "process"
    return "content"


def guess_category(texts: list[str]) -> str:
    joined = " ".join(texts).casefold()
    if any(term in joined for term in ("chart", "market", "metric", "data", "시장")):
        return "data_story"
    if any(term in joined for term in ("team", "profile", "member", "팀")):
        return "team"
    if any(term in joined for term in ("case", "project", "portfolio", "프로젝트")):
        return "case_study"
    if any(term in joined for term in ("proposal", "solution", "제안", "솔루션")):
        return "proposal"
    return "general"


def slide_counts(slide: Any) -> dict[str, Any]:
    text_shape_count = 0
    picture_count = 0
    chart_count = 0
    table_count = 0
    for shape in slide.shapes:
        if " ".join(getattr(shape, "text", "").split()):
            text_shape_count += 1
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture_count += 1
        if getattr(shape, "has_chart", False):
            chart_count += 1
        if getattr(shape, "has_table", False):
            table_count += 1
    return {
        "shape_count": len(slide.shapes),
        "text_shape_count": text_shape_count,
        "picture_count": picture_count,
        "chart_count": chart_count,
        "table_count": table_count,
    }


def save_single_slide_decks(source: Path, curation_root: Path, deck_id: str, *, force: bool) -> list[Path]:
    source_prs = Presentation(str(source))
    created: list[Path] = []
    for keep_index in range(len(source_prs.slides)):
        slide_no = keep_index + 1
        single_path = curation_root / "pptx" / "unlabeled" / f"{asset_stem(deck_id, slide_no)}.pptx"
        if single_path.exists() and not force:
            raise FileExistsError(f"{single_path} already exists. Use --force to rebuild it.")
        work = Presentation(str(source))
        slide_id_list = work.slides._sldIdLst
        for remove_index in reversed(range(len(work.slides))):
            if remove_index == keep_index:
                continue
            rel_id = slide_id_list[remove_index].rId
            work.part.drop_rel(rel_id)
            slide_id_list.remove(slide_id_list[remove_index])
        work.save(str(single_path))
        created.append(single_path)
    return created


def render_previews(source: Path, curation_root: Path, deck_id: str, *, force: bool) -> list[Path]:
    render_root = curation_root / "_render" / deck_id
    if render_root.exists():
        shutil.rmtree(render_root)
    rendered = render_pptx(source, render_root)
    created: list[Path] = []
    for path in rendered:
        slide_no = int(path.stem.split("_")[-1])
        target = curation_root / "png" / "unlabeled" / f"{asset_stem(deck_id, slide_no)}.png"
        if target.exists() and not force:
            raise FileExistsError(f"{target} already exists. Use --force to rebuild it.")
        if target.exists():
            target.unlink()
        path.replace(target)
        created.append(target)
    if render_root.exists():
        shutil.rmtree(render_root)
    return created


def write_contact_sheet(previews: list[Path], curation_root: Path, deck_id: str) -> Path:
    thumbs: list[tuple[int, Image.Image]] = []
    for preview in sorted(previews):
        slide_no = slide_no_from_asset_stem(preview.stem)
        image = Image.open(preview).convert("RGB")
        image.thumbnail((360, 205))
        thumbs.append((slide_no, image.copy()))

    cols = 3
    rows = max(1, (len(thumbs) + cols - 1) // cols)
    cell_w, cell_h = 400, 260
    sheet = Image.new("RGB", (cols * cell_w, rows * cell_h), "#F3F4F6")
    draw = ImageDraw.Draw(sheet)
    try:
        font = ImageFont.truetype("arial.ttf", 18)
    except OSError:
        font = ImageFont.load_default()

    for index, (slide_no, image) in enumerate(thumbs):
        x = (index % cols) * cell_w
        y = (index // cols) * cell_h
        draw.rectangle([x + 12, y + 12, x + cell_w - 12, y + cell_h - 12], fill="white", outline="#D1D5DB")
        draw.text((x + 24, y + 22), f"Slide {slide_no:02d}", fill="#111827", font=font)
        sheet.paste(image, (x + 20, y + 52))

    output_path = curation_root / "records" / f"{deck_id}_contact_sheet.png"
    sheet.save(output_path)
    return output_path


def build_inventory(source: Path, curation_root: Path, deck_id: str, previews: list[Path], one_slide_decks: list[Path]) -> dict[str, Any]:
    prs = Presentation(str(source))
    preview_by_slide = {slide_no_from_asset_stem(path.stem): path for path in previews}
    pptx_by_slide = {slide_no_from_asset_stem(path.stem): path for path in one_slide_decks}
    slides = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        slide_id = asset_stem(deck_id, slide_no)
        slides.append(
            {
                "slide": slide_no,
                "slide_id": slide_id,
                "reference_id": deck_id,
                "quality": "unlabeled",
                "status": "unlabeled",
                "asset_stem": slide_id,
                "preview_png": base_relative(preview_by_slide[slide_no]),
                "one_slide_pptx": base_relative(pptx_by_slide[slide_no]),
                "intent": None,
                "notes": None,
                "text_preview": slide_text_preview(slide),
            }
        )
    return {
        "schema_version": "1.0",
        "deck_id": deck_id,
        "reference_id": deck_id,
        "source_pptx": base_relative(source),
        "curation_root": base_relative(curation_root),
        "slide_count": len(slides),
        "buckets": BUCKETS,
        "slides": slides,
    }


def build_deck_metadata(source: Path, deck_id: str, reference_record: dict[str, Any] | None, prs: Presentation) -> dict[str, Any]:
    slide_texts = [slide_text_joined(slide) for slide in prs.slides]
    non_empty = [text for text in slide_texts if text]
    summary = " ".join(non_empty)[:1200]
    return {
        "schema_version": "1.0",
        "reference_id": deck_id,
        "source_pptx": base_relative(source),
        "source_checksum": file_sha256(source),
        "slide_count": len(prs.slides),
        "deck_title": non_empty[0][:160] if non_empty else None,
        "deck_text_summary": summary,
        "deck_type_primary": (reference_record or {}).get("deck_type_primary", "unknown"),
        "inferred_audience": {
            "value": None,
            "confidence": 0,
            "method": "not_inferred_at_parse_time",
            "requires_human_review": True,
        },
        "reference_record": reference_record,
        "parser": {
            "name": "parse_reference_deck.py",
            "version": "2.0",
            "rendering": "validate_visual_smoke.render_pptx",
        },
        "created_at": utc_now(),
        "updated_at": utc_now(),
    }


def build_slide_metadata(
    source: Path,
    curation_root: Path,
    deck_id: str,
    reference_record: dict[str, Any] | None,
    prs: Presentation,
) -> dict[str, Any]:
    slide_texts = [slide_text_preview(slide) for slide in prs.slides]
    source_checksum = file_sha256(source)
    records: list[dict[str, Any]] = []
    slide_count = len(prs.slides)
    for slide_no, slide in enumerate(prs.slides, start=1):
        slide_id = asset_stem(deck_id, slide_no)
        counts = slide_counts(slide)
        records.append(
            {
                "schema_version": "1.0",
                "identity": {
                    "slide_id": slide_id,
                    "reference_id": deck_id,
                    "source_file": base_relative(source),
                    "source_checksum": source_checksum,
                    "source_slide_no": slide_no,
                    "original_sequence_index": slide_no - 1,
                    "asset_stem": slide_id,
                    "preview_png": base_relative(curation_root / "png" / "unlabeled" / f"{slide_id}.png"),
                    "one_slide_pptx": base_relative(curation_root / "pptx" / "unlabeled" / f"{slide_id}.pptx"),
                },
                "deck_context": {
                    "deck_id": deck_id,
                    "reference_id": deck_id,
                    "deck_type_primary": (reference_record or {}).get("deck_type_primary", "unknown"),
                    "slide_count": slide_count,
                },
                "curation": {
                    "quality_label": "unlabeled",
                    "candidate_bucket": None,
                    "curation_stage": "parsed",
                    "promotion_status": "not_promoted",
                    "human_review_status": "unreviewed",
                    "review_notes": None,
                    "updated_at": utc_now(),
                },
                "semantic_seed": {
                    "source_section": None,
                    "previous_slide_text_preview": slide_texts[slide_no - 2] if slide_no > 1 else [],
                    "next_slide_text_preview": slide_texts[slide_no] if slide_no < slide_count else [],
                    "original_story_position": infer_story_position(slide_no, slide_count),
                    "initial_category_guess": {
                        "value": guess_category(slide_texts[slide_no - 1]),
                        "confidence": 0.35,
                        "method": "keyword_seed",
                        "requires_human_review": True,
                    },
                    "initial_story_role_guess": {
                        "value": guess_story_role(slide_no, slide_count, slide_texts[slide_no - 1]),
                        "confidence": 0.4,
                        "method": "position_and_keyword_seed",
                        "requires_human_review": True,
                    },
                },
                "deterministic_counts": counts,
                "semantic_classification": {},
                "visual_structure": {},
                "style_system": {},
                "elements": {},
                "reuse_and_graph_hints": {},
            }
        )
    return {
        "schema_version": "1.0",
        "reference_id": deck_id,
        "source_pptx": base_relative(source),
        "slides": records,
        "created_at": utc_now(),
        "updated_at": utc_now(),
    }


def parse_reference_deck(source: Path, curation_root: Path, *, force: bool = False) -> dict[str, Any]:
    source = source.resolve()
    if not source.exists():
        raise FileNotFoundError(source)
    if source.suffix.lower() != ".pptx":
        raise ValueError(f"Expected a .pptx file: {source}")

    reference_record = find_reference_record_for_path(source)
    deck_id = str(reference_record.get("reference_id")) if reference_record else slugify(source.stem)
    curation_root = curation_root.resolve()
    ensure_safe_curation_path(curation_root, curation_root)
    make_bucket_dirs(curation_root)
    prepare_deck_outputs(curation_root, deck_id, force=force)

    previews = render_previews(source, curation_root, deck_id, force=force)
    one_slide_decks = save_single_slide_decks(source, curation_root, deck_id, force=force)
    contact_sheet = write_contact_sheet(previews, curation_root, deck_id)
    inventory = build_inventory(source, curation_root, deck_id, previews, one_slide_decks)
    inventory["contact_sheet"] = base_relative(contact_sheet)
    prs = Presentation(str(source))
    deck_metadata = build_deck_metadata(source, deck_id, reference_record, prs)
    slide_metadata = build_slide_metadata(source, curation_root, deck_id, reference_record, prs)

    records_dir = curation_root / "records"
    write_json(records_dir / f"{deck_id}_slide_inventory.json", inventory)
    write_json(records_dir / f"{deck_id}_slide_quality_labels.json", inventory)
    write_json(records_dir / f"{deck_id}_deck_metadata.json", deck_metadata)
    write_json(records_dir / f"{deck_id}_slide_metadata.json", slide_metadata)
    update_reference_record_status(deck_id, "parsed")
    return inventory


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Split a raw reference PPTX into PNG previews and one-slide PPTX candidates.")
    parser.add_argument("pptx_path")
    parser.add_argument("--curation-root", default=str(DEFAULT_CURATION_ROOT))
    parser.add_argument("--force", action="store_true", help="Rebuild an existing curation deck directory.")
    args = parser.parse_args(argv)

    source = Path(args.pptx_path)
    if not source.is_absolute():
        source = (BASE_DIR / source).resolve()
    curation_root = Path(args.curation_root)
    if not curation_root.is_absolute():
        curation_root = (BASE_DIR / curation_root).resolve()
    curation_root.mkdir(parents=True, exist_ok=True)

    inventory = parse_reference_deck(source, curation_root, force=args.force)
    print(
        json.dumps(
            {
                "deck_id": inventory["deck_id"],
                "curation_root": inventory["curation_root"],
                "slide_count": inventory["slide_count"],
                "contact_sheet": inventory["contact_sheet"],
            },
            indent=2,
            ensure_ascii=False,
        )
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

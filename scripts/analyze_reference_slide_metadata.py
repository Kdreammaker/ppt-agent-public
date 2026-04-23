from __future__ import annotations

import argparse
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE_DIR = Path(__file__).resolve().parents[1]
if str(BASE_DIR) not in sys.path:
    sys.path.insert(0, str(BASE_DIR))

from scripts.reference_pipeline import DEFAULT_CURATION_ROOT, LIBRARY_ROOT, load_json, utc_now, write_json
from system.pptx_system import shape_bounds_inches


def density_from_counts(text_count: int, picture_count: int, shape_count: int) -> str:
    if shape_count >= 24 or text_count >= 12:
        return "high"
    if shape_count >= 12 or text_count >= 5 or picture_count >= 2:
        return "medium"
    return "low"


def analyze_one_slide_pptx(path: Path) -> dict[str, Any]:
    prs = Presentation(str(path))
    slide = prs.slides[0]
    text_items: list[dict[str, Any]] = []
    pictures: list[dict[str, Any]] = []
    shape_items: list[dict[str, Any]] = []
    chart_count = 0
    table_count = 0
    font_sizes: list[float] = []
    font_families: set[str] = set()

    for index, shape in enumerate(slide.shapes, start=1):
        bounds = shape_bounds_inches(shape)
        text = " ".join(getattr(shape, "text", "").split())
        if text:
            text_items.append({"shape_index": index, "text_preview": text[:120], "bounds": bounds})
            if getattr(shape, "text_frame", None):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            font_families.add(run.font.name)
                        if run.font.size:
                            font_sizes.append(float(run.font.size.pt))
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pictures.append({"shape_index": index, "bounds": bounds})
        else:
            shape_items.append({"shape_index": index, "shape_type": str(shape.shape_type), "bounds": bounds})
        if getattr(shape, "has_chart", False):
            chart_count += 1
        if getattr(shape, "has_table", False):
            table_count += 1

    shape_count = len(slide.shapes)
    density = density_from_counts(len(text_items), len(pictures), shape_count)
    return {
        "deterministic_counts": {
            "shape_count": shape_count,
            "text_shape_count": len(text_items),
            "picture_count": len(pictures),
            "chart_count": chart_count,
            "table_count": table_count,
        },
        "visual_structure": {
            "visual_density": density,
            "layout_type": "image_led" if pictures and len(text_items) <= 3 else "text_grid" if len(text_items) >= 5 else "mixed",
            "background_type": "unknown",
        },
        "style_system": {
            "font_families": sorted(font_families),
            "font_size_profile": {
                "min": min(font_sizes) if font_sizes else None,
                "max": max(font_sizes) if font_sizes else None,
                "count": len(font_sizes),
            },
            "dominant_colors": [],
            "accent_colors": [],
        },
        "elements": {
            "texts": text_items,
            "images": pictures,
            "shapes": shape_items[:50],
            "charts": chart_count,
            "tables": table_count,
        },
        "reuse_and_graph_hints": {
            "reuse_potential": {
                "value": "review_required",
                "confidence": 0.25,
                "method": "deterministic_seed",
                "requires_human_review": True,
            },
            "production_readiness": {
                "value": "not_ready",
                "confidence": 0.9,
                "method": "requires_explicit_promotion",
                "requires_human_review": False,
            },
            "best_for": [],
            "avoid_for": [],
        },
    }


def enrich_metadata_file(path: Path) -> dict[str, Any]:
    metadata = load_json(path)
    enriched = 0
    for record in metadata.get("slides", []):
        identity = record.get("identity", {})
        pptx_path = BASE_DIR / str(identity.get("one_slide_pptx", ""))
        if not pptx_path.exists():
            continue
        analysis = analyze_one_slide_pptx(pptx_path)
        for key, value in analysis.items():
            record[key] = value
        record["analysis"] = {
            "analyzed_at": utc_now(),
            "analyzer": "analyze_reference_slide_metadata.py",
            "version": "1.0",
        }
        enriched += 1
    metadata["updated_at"] = utc_now()
    write_json(path, metadata)
    return {"path": path.resolve().relative_to(BASE_DIR).as_posix(), "enriched": enriched}


def collect_metadata(paths: list[str]) -> list[Path]:
    if paths:
        return [(BASE_DIR / path).resolve() if not Path(path).is_absolute() else Path(path) for path in paths]
    return [
        *DEFAULT_CURATION_ROOT.glob("records/*_slide_metadata.json"),
        *LIBRARY_ROOT.glob("*/metadata/*.json"),
    ]


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Enrich reference slide metadata with deterministic PPTX analysis.")
    parser.add_argument("paths", nargs="*")
    args = parser.parse_args(argv)
    results = [enrich_metadata_file(path) for path in collect_metadata(args.paths) if path.exists()]
    print(f"metadata_files={len(results)} enriched_slides={sum(item['enriched'] for item in results)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

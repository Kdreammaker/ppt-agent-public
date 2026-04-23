from __future__ import annotations

import json
import sys
from collections import Counter
from pathlib import Path
from statistics import mean
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE_DIR = Path(__file__).resolve().parents[1]
if str(BASE_DIR) not in sys.path:
    sys.path.insert(0, str(BASE_DIR))

from system.pptx_system import EMU_PER_INCH, shape_bounds_inches

CATALOG_PATH = BASE_DIR / "config" / "reference_catalog.json"
STYLE_PROFILE_PATH = BASE_DIR / "config" / "reference_style_profiles.json"
AUDIT_PATH = BASE_DIR / "docs" / "REFERENCE_DESIGN_AUDIT.md"


def resolve_path(value: str) -> Path:
    path = Path(value)
    return path if path.is_absolute() else (BASE_DIR / path).resolve()


def safe_rgb(color_format) -> str | None:
    try:
        rgb = color_format.rgb
    except Exception:
        return None
    if rgb is None:
        return None
    return f"#{rgb}"


def collect_fill(shape) -> str | None:
    try:
        if not shape.fill or shape.fill.type is None:
            return None
        return safe_rgb(shape.fill.fore_color)
    except Exception:
        return None


def collect_line(shape) -> str | None:
    try:
        if not shape.line:
            return None
        return safe_rgb(shape.line.color)
    except Exception:
        return None


def collect_text_runs(shape) -> list[dict[str, Any]]:
    runs: list[dict[str, Any]] = []
    if not getattr(shape, "has_text_frame", False):
        return runs
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text.strip():
                continue
            font = run.font
            size = font.size.pt if font.size is not None else None
            color = None
            try:
                color = safe_rgb(font.color)
            except Exception:
                color = None
            runs.append(
                {
                    "text": run.text.strip(),
                    "font_name": font.name,
                    "font_size": size,
                    "font_color": color,
                    "bold": font.bold,
                }
            )
    return runs


def classify_slide(slide) -> dict[str, Any]:
    shape_count = len(slide.shapes)
    text_shape_count = 0
    picture_count = 0
    group_count = 0
    fill_colors: Counter[str] = Counter()
    line_colors: Counter[str] = Counter()
    font_names: Counter[str] = Counter()
    font_colors: Counter[str] = Counter()
    font_sizes: list[float] = []
    text_chars = 0
    left_edges: list[float] = []
    top_edges: list[float] = []

    for shape in slide.shapes:
        bounds = shape_bounds_inches(shape)
        left_edges.append(bounds["left"])
        top_edges.append(bounds["top"])

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture_count += 1
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_count += 1

        fill = collect_fill(shape)
        if fill:
            fill_colors[fill] += 1
        line = collect_line(shape)
        if line:
            line_colors[line] += 1

        text = getattr(shape, "text", "")
        if text.strip():
            text_shape_count += 1
            text_chars += len(text.strip())
        for run in collect_text_runs(shape):
            if run["font_name"]:
                font_names[run["font_name"]] += 1
            if run["font_color"]:
                font_colors[run["font_color"]] += 1
            if run["font_size"]:
                font_sizes.append(float(run["font_size"]))

    layout_signature = "open-canvas"
    if picture_count >= 2:
        layout_signature = "image-led"
    elif shape_count >= 28:
        layout_signature = "dense-card-system"
    elif text_shape_count >= 8:
        layout_signature = "text-card-system"
    elif shape_count <= 8 and text_shape_count <= 3:
        layout_signature = "minimal-cover-or-divider"

    return {
        "shape_count": shape_count,
        "text_shape_count": text_shape_count,
        "picture_count": picture_count,
        "group_count": group_count,
        "text_chars": text_chars,
        "layout_signature": layout_signature,
        "dominant_fills": fill_colors.most_common(8),
        "dominant_lines": line_colors.most_common(8),
        "dominant_font_names": font_names.most_common(8),
        "dominant_font_colors": font_colors.most_common(8),
        "font_size_min": round(min(font_sizes), 1) if font_sizes else None,
        "font_size_max": round(max(font_sizes), 1) if font_sizes else None,
        "font_size_avg": round(mean(font_sizes), 1) if font_sizes else None,
        "left_edge_avg": round(mean(left_edges), 2) if left_edges else None,
        "top_edge_avg": round(mean(top_edges), 2) if top_edges else None,
    }


def analyze_library(library: dict[str, Any]) -> dict[str, Any]:
    path = resolve_path(library["library_path"])
    prs = Presentation(str(path))
    slide_profiles = []
    aggregate_fill: Counter[str] = Counter()
    aggregate_line: Counter[str] = Counter()
    aggregate_font_names: Counter[str] = Counter()
    aggregate_font_colors: Counter[str] = Counter()
    layout_signatures: Counter[str] = Counter()
    shape_counts = []
    text_chars = []
    image_counts = []

    for index, slide in enumerate(prs.slides, start=1):
        profile = classify_slide(slide)
        slide_profiles.append({"slide_no": index, **profile})
        aggregate_fill.update(dict(profile["dominant_fills"]))
        aggregate_line.update(dict(profile["dominant_lines"]))
        aggregate_font_names.update(dict(profile["dominant_font_names"]))
        aggregate_font_colors.update(dict(profile["dominant_font_colors"]))
        layout_signatures[profile["layout_signature"]] += 1
        shape_counts.append(profile["shape_count"])
        text_chars.append(profile["text_chars"])
        image_counts.append(profile["picture_count"])

    return {
        "library_id": library["library_id"],
        "scope": library["scope"],
        "library_path": library["library_path"],
        "source_path": library["source_path"],
        "slide_count": len(prs.slides),
        "avg_shape_count": round(mean(shape_counts), 1) if shape_counts else 0,
        "avg_text_chars": round(mean(text_chars), 1) if text_chars else 0,
        "avg_image_count": round(mean(image_counts), 1) if image_counts else 0,
        "layout_signatures": layout_signatures.most_common(),
        "dominant_fills": aggregate_fill.most_common(10),
        "dominant_lines": aggregate_line.most_common(10),
        "dominant_font_names": aggregate_font_names.most_common(10),
        "dominant_font_colors": aggregate_font_colors.most_common(10),
        "slides": slide_profiles,
    }


def design_reading(profile: dict[str, Any]) -> list[str]:
    library_id = profile["library_id"]
    if library_id == "template_library_04_v1":
        return [
            "구조와 목적 분류에는 유용하지만, 시각적 완성도는 낮은 편입니다.",
            "원본 placeholder의 빈 여백과 작은 텍스트가 많아 그대로 쓰기보다 layout builder의 골격으로 쓰는 것이 안전합니다.",
            "커버, 목차, 2x2, 4-step처럼 목적형 구조를 빌려오고 실제 디자인은 별도 규칙으로 덮는 전략이 적합합니다.",
        ]
    if library_id == "template_library_01_v1":
        return [
            "흰 배경, 네이비 룰, 연한 카드, 작은 아이콘으로 구성된 정돈된 B2B 세일즈 문법이 강합니다.",
            "상단 얇은 라인과 우상단 페이지, 하단의 작은 브랜드 문구가 반복되어 deck 전체 rhythm을 만듭니다.",
            "내부 공유 자료에서는 과도한 장식보다 요약 카드/절차/근거 블록의 기준 디자인으로 쓰기 좋습니다.",
        ]
    if library_id == "template_library_02_v1":
        return [
            "네이비와 오렌지 accent 조합이 가장 명확하며, 차트와 proof 중심 페이지의 완성도가 높습니다.",
            "제목은 크고 짧게, 본문은 좌측 narrative와 우측 chart/diagram으로 나누는 구조가 반복됩니다.",
            "시장/유가/리스크처럼 데이터가 있는 보고서에는 이 스타일을 우선 적용하는 것이 좋습니다.",
        ]
    if library_id == "template_library_03_v1":
        return [
            "스크린샷 또는 실적 이미지를 크게 두고 오른쪽에 결과 bullet을 배치하는 case-study 문법이 강합니다.",
            "파란 CTA/pill, 회색 메타 태그, 넓은 여백이 포트폴리오형 신뢰감을 만듭니다.",
            "프로젝트 성과, 적용 사례, before/after 정리 슬라이드의 reference로 좋습니다.",
        ]
    return ["추가 분석이 필요한 reference입니다."]


def write_audit(profiles: list[dict[str, Any]]) -> None:
    lines = [
        "# Reference Design Audit",
        "",
        "이 문서는 `assets/slides/templates/decks/*.pptx`를 기준으로 원본 reference의 시각 언어를 추출한 결과입니다.",
        "자동 생성 deck의 디자인 품질을 높이기 위해, 단순 slide purpose뿐 아니라 `design_tier`, `quality_score`, `usage_policy`를 함께 사용합니다.",
        "",
        "## 핵심 결론",
        "- `template_library_04_v1`는 구조 참조용입니다. 그대로 디자인 기준으로 쓰면 완성도가 낮아지기 쉽습니다.",
        "- `template_library_01_v1`는 내부/외부 공유용 B2B 카드 시스템의 기준으로 적합합니다.",
        "- `template_library_02_v1`는 데이터, 차트, 시장/리스크 서사의 기준으로 가장 강합니다.",
        "- `template_library_03_v1`는 스크린샷, 사례, 성과 증명 중심 슬라이드의 기준으로 좋습니다.",
        "- 앞으로 자동 선택은 purpose만 보지 말고 `quality_score`, `usage_policy`, `style_tags`를 같이 봐야 합니다.",
        "",
    ]

    for profile in profiles:
        lines.extend(
            [
                f"## {profile['library_id']}",
                "",
                f"- Scope: `{profile['scope']}`",
                f"- Slides: `{profile['slide_count']}`",
                f"- Avg shapes: `{profile['avg_shape_count']}`",
                f"- Avg text chars: `{profile['avg_text_chars']}`",
                f"- Avg images: `{profile['avg_image_count']}`",
                f"- Layout signatures: `{profile['layout_signatures']}`",
                f"- Dominant fills: `{profile['dominant_fills'][:5]}`",
                f"- Dominant font colors: `{profile['dominant_font_colors'][:5]}`",
                "",
                "### 디자인 해석",
            ]
        )
        lines.extend(f"- {item}" for item in design_reading(profile))
        lines.append("")

    lines.extend(
        [
            "## 시스템 반영 규칙",
            "",
            "- 목적형 구조만 필요한 경우: `template_library_04_v1`를 사용하되 `layout builder`가 실제 시각 구성을 다시 그립니다.",
            "- 디자인 완성도가 중요한 경우: `usage_policy=production_ready`, `min_quality_score>=4.0` 조건을 우선 적용합니다.",
            "- 데이터/차트 슬라이드: `template_library_02_v1`를 우선 후보로 둡니다.",
            "- 사례/성과 슬라이드: `template_library_03_v1`를 우선 후보로 둡니다.",
            "- 프로세스/서비스 소개: `template_library_01_v1`를 우선 후보로 둡니다.",
            "",
            "## 다음 큐레이션 기준",
            "",
            "- slide별 실제 사용 가능 상태를 `production_ready`, `curate_before_use`, `structure_only`로 분류합니다.",
            "- `structure_only` 슬라이드는 원본 요소를 보존하지 않고, custom layout builder 또는 overlay safe zone으로만 활용합니다.",
            "- `production_ready` 슬라이드는 원본의 footer, rule, card, icon, chart rhythm을 보존하는 방식으로 치환합니다.",
            "- 같은 purpose라도 최소 2개 이상의 visual variant를 유지하되, default는 품질 점수가 높은 슬라이드로 둡니다.",
        ]
    )
    AUDIT_PATH.write_text("\n".join(lines) + "\n", encoding="utf-8")


def main() -> int:
    catalog = json.loads(CATALOG_PATH.read_text(encoding="utf-8"))
    profiles = [analyze_library(library) for library in catalog["libraries"]]
    output = {
        "version": "0.1",
        "workspace": BASE_DIR.as_posix(),
        "profiles": profiles,
    }
    STYLE_PROFILE_PATH.write_text(json.dumps(output, indent=2, ensure_ascii=False), encoding="utf-8")
    write_audit(profiles)
    print(STYLE_PROFILE_PATH)
    print(AUDIT_PATH)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

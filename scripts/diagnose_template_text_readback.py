from __future__ import annotations

import argparse
import json
import sys
from collections import Counter, defaultdict
from pathlib import Path
from typing import Any

from pptx import Presentation

BASE_DIR = Path(__file__).resolve().parents[1]
if str(BASE_DIR) not in sys.path:
    sys.path.insert(0, str(BASE_DIR))

from scripts.build_deck import load_spec, resolve_path
from scripts.validate_visual_smoke import normalize_text, slide_texts, template_expected_text_rows

DEFAULT_REPORT_JSON = BASE_DIR / "outputs" / "reports" / "template_text_readback_diagnostics.json"
DEFAULT_REPORT_MD = BASE_DIR / "outputs" / "reports" / "template_text_readback_diagnostics.md"
DEFAULT_SPEC_PATHS = [
    BASE_DIR / "data" / "specs" / "jb_meeting_blueprint_spec.json",
    BASE_DIR / "data" / "specs" / "portfolio_auto_selection_sample.json",
    BASE_DIR / "data" / "specs" / "report_auto_selection_sample.json",
    BASE_DIR / "data" / "specs" / "sales_auto_selection_sample.json",
    BASE_DIR / "data" / "specs" / "template_slide_sample_spec.json",
]


def deck_path_for_spec(spec_path: Path) -> Path:
    spec, spec_dir = load_spec(spec_path)
    output_path = resolve_path(spec_dir, spec.get("output_path"))
    if output_path is None:
        raise ValueError(f"Spec has no output_path: {spec_path}")
    return output_path


def default_pairs() -> list[tuple[Path, Path]]:
    return [(spec_path, deck_path_for_spec(spec_path)) for spec_path in DEFAULT_SPEC_PATHS]


def parse_pair(value: str) -> tuple[Path, Path]:
    separator = ".json:"
    if separator in value:
        split_at = value.index(separator) + len(".json")
        spec_value = value[:split_at]
        deck_value = value[split_at + 1 :]
    elif ":" in value:
        spec_value, deck_value = value.split(":", 1)
    else:
        raise argparse.ArgumentTypeError("Pairs must use <spec_json>:<pptx>")
    return Path(spec_value).resolve(), Path(deck_value).resolve()


def classify_row(row: dict[str, Any], actual_texts: list[str], matched: bool) -> str:
    if matched:
        return "present"
    if row.get("source_kind") == "text_values":
        return "sequential_text_values_risk"
    if row.get("slot") in {"eyebrow", "subtitle"} or str(row.get("slot", "")).endswith("_number"):
        return "likely_unmapped_slot"
    return "missing"


def diagnose_pair(spec_path: Path, deck_path: Path) -> list[dict[str, Any]]:
    prs = Presentation(str(deck_path))
    rows: list[dict[str, Any]] = []
    for expected in template_expected_text_rows(spec_path):
        slide_no = expected["slide"]
        if slide_no > len(prs.slides):
            actual_texts: list[str] = []
            matched = False
        else:
            actual_texts = slide_texts(prs.slides[slide_no - 1])
            matched = normalize_text(expected["expected_text"]) in actual_texts
        classification = classify_row(expected, actual_texts, matched)
        rows.append(
            {
                "spec_path": str(spec_path),
                "deck_path": str(deck_path),
                "slide": slide_no,
                "slot": expected["slot"],
                "expected_text": expected["expected_text"],
                "source_kind": expected["source_kind"],
                "classification": classification,
                "resolved_slide_id": expected.get("resolved_slide_id"),
                "template_key": expected.get("template_key"),
                "shape_index": expected.get("shape_index"),
                "fit_strategy": expected.get("fit_strategy"),
                "actual_slide_texts": actual_texts[:12],
            }
        )
    return rows


def summarize(rows: list[dict[str, Any]]) -> dict[str, Any]:
    by_classification = Counter(row["classification"] for row in rows)
    by_deck = defaultdict(Counter)
    by_template = defaultdict(Counter)
    for row in rows:
        by_deck[Path(row["deck_path"]).name][row["classification"]] += 1
        by_template[row["resolved_slide_id"] or "<none>"][row["classification"]] += 1
    return {
        "total_rows": len(rows),
        "by_classification": dict(sorted(by_classification.items())),
        "by_deck": {deck: dict(counter) for deck, counter in sorted(by_deck.items())},
        "by_template": {template: dict(counter) for template, counter in sorted(by_template.items())},
    }


def markdown_table(headers: list[str], rows: list[list[Any]]) -> list[str]:
    if not rows:
        return ["_None._"]
    lines = [
        "| " + " | ".join(headers) + " |",
        "| " + " | ".join("---" for _ in headers) + " |",
    ]
    for row in rows:
        lines.append("| " + " | ".join(str(item).replace("\n", " ") for item in row) + " |")
    return lines


def write_json_report(path: Path, rows: list[dict[str, Any]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    payload = {"summary": summarize(rows), "rows": rows}
    path.write_text(json.dumps(payload, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def write_markdown_report(path: Path, rows: list[dict[str, Any]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    summary = summarize(rows)
    problem_rows = [row for row in rows if row["classification"] != "present"]
    lines = [
        "# Template Text Readback Diagnostics",
        "",
        "This report compares expected template slot text with the generated PPTX readback text.",
        "",
        "## Summary",
        "",
        f"- Total expected rows: {summary['total_rows']}",
        f"- Present: {summary['by_classification'].get('present', 0)}",
        f"- Missing: {sum(count for key, count in summary['by_classification'].items() if key != 'present')}",
        "",
        "## By Classification",
        "",
        *markdown_table(
            ["Classification", "Count"],
            [[key, value] for key, value in summary["by_classification"].items()],
        ),
        "",
        "## By Deck",
        "",
        *markdown_table(
            ["Deck", "present", "likely_unmapped_slot", "sequential_text_values_risk", "missing"],
            [
                [
                    deck,
                    counter.get("present", 0),
                    counter.get("likely_unmapped_slot", 0),
                    counter.get("sequential_text_values_risk", 0),
                    counter.get("missing", 0),
                ]
                for deck, counter in summary["by_deck"].items()
            ],
        ),
        "",
        "## Missing / Risk Rows",
        "",
        *markdown_table(
            ["Deck", "Slide", "Slot", "Class", "Template", "Expected", "Actual text sample"],
            [
                [
                    Path(row["deck_path"]).name,
                    row["slide"],
                    row["slot"],
                    row["classification"],
                    row["resolved_slide_id"],
                    row["expected_text"][:80],
                    " / ".join(row["actual_slide_texts"][:4])[:140],
                ]
                for row in problem_rows
            ],
        ),
    ]
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Diagnose template slot text readback coverage.")
    parser.add_argument(
        "--pair",
        action="append",
        type=parse_pair,
        help="Spec/deck pair in the form <spec_json>:<pptx>. Defaults to five template/sample decks.",
    )
    parser.add_argument("--output-json", default=str(DEFAULT_REPORT_JSON))
    parser.add_argument("--output-md", default=str(DEFAULT_REPORT_MD))
    args = parser.parse_args(argv)

    pairs = args.pair or default_pairs()
    rows: list[dict[str, Any]] = []
    for spec_path, deck_path in pairs:
        rows.extend(diagnose_pair(spec_path, deck_path))

    write_json_report(Path(args.output_json).resolve(), rows)
    write_markdown_report(Path(args.output_md).resolve(), rows)
    summary = summarize(rows)
    print(Path(args.output_json).resolve())
    print(Path(args.output_md).resolve())
    print(json.dumps(summary["by_classification"], ensure_ascii=False, sort_keys=True))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

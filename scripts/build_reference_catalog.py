from __future__ import annotations

import json
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE_DIR = Path(__file__).resolve().parents[1]
if str(BASE_DIR) not in sys.path:
    sys.path.insert(0, str(BASE_DIR))

from system.blueprint_loader import save_blueprints
from system.pptx_system import DEFAULT_PROTECTED_TOKENS, normalize_fingerprint, rel_path, shape_bounds_inches


SOURCE_CATALOG_PATH = BASE_DIR / "config" / "template_library_catalog.json"
REFERENCE_CATALOG_PATH = BASE_DIR / "config" / "reference_catalog.json"
BLUEPRINT_PATH = BASE_DIR / "config" / "template_blueprints.json"
OVERRIDES_PATH = BASE_DIR / "config" / "template_blueprint_overrides.json"
QUALITY_OVERRIDES_PATH = BASE_DIR / "config" / "template_quality_overrides.json"

REMOVE_TEXT_PATTERNS = [
    "office suite",
    "all rights reserved",
    "by ",
    "copyright",
]
PLACEHOLDER_TEXT_PATTERNS = [
    "please fill",
    "please write",
    "please enter",
    "input the",
    "provide a concise",
    "provide a brief",
    "write a detailed",
    "text here",
    "main text",
    "sub text",
    "description",
    "title",
    "keyword",
    "index",
    "contents",
]
HEADER_ZONE = {"left": 0.58, "top": 0.42, "width": 3.38, "height": 0.28}
FOOTER_ZONE = {"left": 0.0, "top": 6.66, "width": 13.33, "height": 0.84}
PAGE_ZONE = {"left": 11.7, "top": 6.84, "width": 0.7, "height": 0.34}


def resolve_path(value: str) -> Path:
    path = Path(value)
    return path if path.is_absolute() else (BASE_DIR / path).resolve()


def normalized_slug(value: str) -> str:
    return value.replace("/", "_").replace("-", "_")


def density_bucket(shape_count: int, text_chars: int, image_count: int) -> str:
    if shape_count >= 24 or text_chars >= 800:
        return "dense"
    if image_count >= 4 or text_chars >= 350 or shape_count >= 12:
        return "medium"
    return "light"


def is_removable_text(text: str) -> bool:
    lowered = text.lower()
    return any(pattern in lowered for pattern in REMOVE_TEXT_PATTERNS)


def is_editable_text(text: str, purpose: str) -> bool:
    lowered = text.lower()
    if purpose == "toc" and lowered in {"index", "contents"}:
        return True
    if any(pattern in lowered for pattern in PLACEHOLDER_TEXT_PATTERNS):
        return True
    if len(text) <= 24:
        return True
    return False


def font_role_for_slot(slot_name: str, purpose: str, bounds: dict[str, float]) -> str:
    if slot_name in {"title", "hero_title"}:
        return "cover_title" if purpose == "cover" else "hero_title"
    if slot_name in {"subtitle", "summary_kicker"}:
        return "body_strong"
    if slot_name.endswith("_title") or slot_name.startswith("card_") and slot_name.endswith("_title"):
        return "card_title"
    if bounds["top"] < 1.1:
        return "section_header"
    if bounds["height"] <= 0.3:
        return "body_strong"
    return "body"


def max_chars_for_bounds(bounds: dict[str, float], font_role: str) -> int:
    base = int(bounds["width"] * 7.5)
    if font_role in {"cover_title", "hero_title"}:
        base = int(bounds["width"] * 5.2)
    return max(12, min(base, 60))


def slot_name_series(purpose: str, count: int) -> list[str]:
    if purpose == "cover":
        seeds = ["title", "subtitle", "summary_kicker", "footer_note", "badge"]
    elif purpose == "toc":
        seeds = ["title", "subtitle", "toc_item_1", "toc_item_2", "toc_item_3", "toc_item_4", "toc_item_5", "toc_item_6"]
    elif purpose == "summary":
        seeds = ["title", "subtitle", "card_1_title", "card_1_body", "card_2_title", "card_2_body", "card_3_title", "card_3_body", "card_4_title", "card_4_body"]
    elif purpose == "issue":
        seeds = ["title", "subtitle", "card_1_title", "card_1_body", "card_2_title", "card_2_body", "card_3_title", "card_3_body", "card_4_title", "card_4_body"]
    elif purpose in {"strategy", "process", "timeline"}:
        seeds = ["title", "subtitle", "step_1_title", "step_1_body", "step_2_title", "step_2_body", "step_3_title", "step_3_body", "step_4_title", "step_4_body", "footer_note"]
    elif purpose in {"closing"}:
        seeds = ["title", "subtitle", "cta_title", "cta_body", "footer_note"]
    elif purpose in {"team"}:
        seeds = ["title", "subtitle", "profile_1_name", "profile_1_body", "profile_2_name", "profile_2_body", "profile_3_name", "profile_3_body", "profile_4_name", "profile_4_body"]
    else:
        seeds = ["title", "subtitle", "body_1", "body_2", "body_3", "body_4", "body_5", "body_6"]

    if count <= len(seeds):
        return seeds[:count]
    extras = [f"text_{index}" for index in range(1, count - len(seeds) + 1)]
    return seeds + extras


def default_overlay_zones(purpose: str) -> list[dict[str, Any]]:
    if purpose == "cover":
        return [
            {"slot": "hero_panel", "left": 0.55, "top": 0.78, "width": 7.55, "height": 5.95},
            {"slot": "footer_meta", "left": 0.95, "top": 5.7, "width": 5.5, "height": 0.4},
        ]
    if purpose == "toc":
        return [
            {"slot": "toc_content", "left": 0.8, "top": 1.2, "width": 11.5, "height": 4.9},
        ]
    return [
        {"slot": "content_primary", "left": 0.8, "top": 1.05, "width": 11.7, "height": 5.3},
    ]


def analyse_slide(
    slide,
    *,
    library: dict[str, Any],
    slide_entry: dict[str, Any],
    library_slide_no: int,
) -> tuple[dict[str, Any], dict[str, Any]]:
    slide_id = f"{library['library_id']}.{slide_entry['key']}"
    all_indices: list[int] = []
    remove_shapes: list[int] = []
    editable_text_slots: list[dict[str, Any]] = []
    editable_image_slots: list[dict[str, Any]] = []
    editable_chart_slots: list[dict[str, Any]] = []
    text_candidates: list[tuple[float, float, int, Any, str, dict[str, float]]] = []
    picture_count = 0
    group_count = 0
    text_chars = 0

    for index, shape in enumerate(slide.shapes, start=1):
        all_indices.append(index)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_count += 1

        text = normalize_fingerprint(getattr(shape, "text", ""))
        bounds = shape_bounds_inches(shape)
        if text:
            text_chars += len(text)
            if is_removable_text(text):
                remove_shapes.append(index)
                continue
            if is_editable_text(text, slide_entry["purpose"]):
                text_candidates.append((bounds["top"], bounds["left"], index, shape, text, bounds))
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            picture_count += 1
            editable_image_slots.append(
                {
                    "slot": f"image_{picture_count}",
                    "shape_index": index,
                    "shape_name": getattr(shape, "name", f"shape_{index}"),
                    "bounds": bounds,
                }
            )

    text_candidates.sort(key=lambda item: (item[0], item[1], item[2]))
    slot_names = slot_name_series(slide_entry["purpose"], len(text_candidates))
    for slot_name, (_, _, index, shape, text, bounds) in zip(slot_names, text_candidates, strict=False):
        font_role = font_role_for_slot(slot_name, slide_entry["purpose"], bounds)
        editable_text_slots.append(
            {
                "slot": slot_name,
                "shape_index": index,
                "shape_name": getattr(shape, "name", f"shape_{index}"),
                "fingerprint": text,
                "bounds": bounds,
                "font_role": font_role,
                "color_token": "primary",
                "bold": slot_name.endswith("_title") or slot_name in {"title", "subtitle", "summary_kicker", "cta_title"},
                "align": "center" if slide_entry["purpose"] in {"cover", "toc"} else "left",
                "max_lines": 3 if slot_name in {"title", "hero_title"} else 5,
                "allow_shrink": True,
                "fit_strategy": "preserve_template",
                "max_chars_per_line": max_chars_for_bounds(bounds, font_role),
                "protect_tokens": DEFAULT_PROTECTED_TOKENS,
            }
        )

    preserve_shapes = [index for index in all_indices if index not in remove_shapes]
    shape_count = len(all_indices)
    density = density_bucket(shape_count, text_chars, picture_count)
    style_tags = sorted(
        {
            *library.get("style_tags", []),
            *slide_entry.get("style_tags", []),
            slide_entry["purpose"],
            density,
        }
    )
    mode = "overlay" if group_count > 0 or shape_count >= 18 else "patch"

    reference_slide = {
        "slide_id": slide_id,
        "template_key": slide_entry["key"],
        "library_id": library["library_id"],
        "library_path": rel_path(BASE_DIR, resolve_path(library["output_path"])),
        "source_path": rel_path(BASE_DIR, resolve_path(library["source_path"])),
        "library_slide_no": library_slide_no,
        "source_slide_no": slide_entry["slide_no"],
        "purpose": slide_entry["purpose"],
        "variant": slide_entry["variant"],
        "scope": slide_entry.get("scope_override", library["scope"]),
        "style_tags": style_tags,
        "density": density,
        "orientation": "landscape",
        "display_name": f"{slide_entry['purpose']} / {slide_entry['variant'].split('/')[-1]}",
        "shape_count": shape_count,
        "text_shape_count": len(editable_text_slots),
        "image_count": picture_count,
        "mode": mode,
    }

    blueprint = {
        "slide_id": slide_id,
        "template_key": slide_entry["key"],
        "library_id": library["library_id"],
        "library_path": rel_path(BASE_DIR, resolve_path(library["output_path"])),
        "library_slide_no": library_slide_no,
        "purpose": slide_entry["purpose"],
        "variant": slide_entry["variant"],
        "scope": slide_entry.get("scope_override", library["scope"]),
        "mode": mode,
        "preserve_shapes": preserve_shapes,
        "remove_shapes": remove_shapes,
        "editable_text_slots": editable_text_slots,
        "editable_image_slots": editable_image_slots,
        "editable_chart_slots": editable_chart_slots,
        "overlay_safe_zones": default_overlay_zones(slide_entry["purpose"]),
        "header_zone": HEADER_ZONE,
        "footer_zone": FOOTER_ZONE,
        "page_zone": PAGE_ZONE,
        "text_rules": {
            "wrap_language": "ko",
            "protect_tokens": DEFAULT_PROTECTED_TOKENS,
            "section_number_mode": "header",
            "page_number_mode": "page_number",
            "page_number_position": "bottom_right",
        },
    }
    return reference_slide, blueprint


def assign_default_ranks(slides: list[dict[str, Any]]) -> None:
    grouped: dict[tuple[str, str], list[dict[str, Any]]] = {}
    for slide in slides:
        grouped.setdefault((slide["purpose"], slide["scope"]), []).append(slide)
    for items in grouped.values():
        items.sort(key=lambda item: (item["library_id"], item["library_slide_no"]))
        for index, item in enumerate(items, start=1):
            item["default_rank"] = index


def apply_reference_quality_overrides(libraries: list[dict[str, Any]], slides: list[dict[str, Any]]) -> None:
    if not QUALITY_OVERRIDES_PATH.exists():
        for slide in slides:
            slide.setdefault("design_tier", "candidate")
            slide.setdefault("quality_score", 3.0)
            slide.setdefault("usage_policy", "candidate")
        return

    overrides = json.loads(QUALITY_OVERRIDES_PATH.read_text(encoding="utf-8"))
    library_defaults = overrides.get("library_defaults", {})
    slide_overrides = overrides.get("slides", {})

    for library in libraries:
        default = library_defaults.get(library["library_id"], {})
        if default:
            library.update(
                {
                    key: value
                    for key, value in default.items()
                    if key in {"design_tier", "quality_score", "usage_policy", "design_notes"}
                }
            )

    for slide in slides:
        default = library_defaults.get(slide["library_id"], {})
        specific = slide_overrides.get(slide["slide_id"], {})
        merged = {
            "design_tier": "candidate",
            "quality_score": 3.0,
            "usage_policy": "candidate",
            "design_notes": "",
            **default,
            **specific,
        }
        slide["design_tier"] = merged["design_tier"]
        slide["quality_score"] = merged["quality_score"]
        slide["usage_policy"] = merged["usage_policy"]
        if merged.get("design_notes"):
            slide["design_notes"] = merged["design_notes"]


def merge_slot(existing_slots: list[dict[str, Any]], override_slot: dict[str, Any]) -> dict[str, Any]:
    shape_index = override_slot.get("shape_index")
    slot_name = override_slot.get("slot")
    match = None
    for slot in existing_slots:
        if shape_index is not None and slot.get("shape_index") == shape_index:
            match = slot
            break
        if slot_name is not None and slot.get("slot") == slot_name:
            match = slot
            break
    merged = dict(match or {})
    if (
        match is not None
        and shape_index is not None
        and shape_index != match.get("shape_index")
    ):
        for identity_key in ("shape_name", "fingerprint", "bounds"):
            if identity_key not in override_slot:
                merged.pop(identity_key, None)
    merged.update(override_slot)
    return merged


def apply_blueprint_overrides(blueprints: dict[str, Any]) -> None:
    if not OVERRIDES_PATH.exists():
        return

    overrides = json.loads(OVERRIDES_PATH.read_text(encoding="utf-8"))
    for slide_id, override in overrides.get("slides", {}).items():
        if slide_id not in blueprints["slides"]:
            raise KeyError(f"Override references missing slide_id: {slide_id}")
        blueprint = blueprints["slides"][slide_id]

        for key in ("mode", "header_zone", "footer_zone", "page_zone", "text_rules"):
            if key in override:
                blueprint[key] = override[key]
        for key in ("preserve_shapes", "remove_shapes", "overlay_safe_zones", "editable_image_slots", "editable_chart_slots"):
            if key in override:
                blueprint[key] = override[key]

        if "editable_text_slots" in override:
            existing_slots = blueprint.get("editable_text_slots", [])
            replacement = [
                merge_slot(existing_slots, override_slot)
                for override_slot in override["editable_text_slots"]
            ]
            if override.get("append_unlisted_text_slots", False):
                used_indices = {slot.get("shape_index") for slot in replacement}
                used_names = {slot.get("slot") for slot in replacement}
                for slot in existing_slots:
                    if slot.get("shape_index") not in used_indices and slot.get("slot") not in used_names:
                        replacement.append(slot)
            blueprint["editable_text_slots"] = replacement


def main() -> int:
    source_catalog = json.loads(SOURCE_CATALOG_PATH.read_text(encoding="utf-8"))
    reference_slides: list[dict[str, Any]] = []
    blueprints: dict[str, Any] = {"version": "0.1", "slides": {}}

    library_summaries = []
    for library in source_catalog["libraries"]:
        library_path = resolve_path(library["output_path"])
        prs = Presentation(str(library_path))
        if len(prs.slides) != len(library["slides"]):
            raise ValueError(
                f"Library slide count mismatch for {library['library_id']}: "
                f"{len(prs.slides)} != {len(library['slides'])}"
            )

        library_summaries.append(
            {
                "library_id": library["library_id"],
                "scope": library["scope"],
                "library_path": rel_path(BASE_DIR, library_path),
                "source_path": rel_path(BASE_DIR, resolve_path(library["source_path"])),
                "description": library["description"],
                "style_tags": library.get("style_tags", []),
                "slide_count": len(library["slides"]),
            }
        )

        for library_slide_no, slide_entry in enumerate(library["slides"], start=1):
            slide = prs.slides[library_slide_no - 1]
            ref_slide, blueprint = analyse_slide(
                slide,
                library=library,
                slide_entry=slide_entry,
                library_slide_no=library_slide_no,
            )
            reference_slides.append(ref_slide)
            blueprints["slides"][ref_slide["slide_id"]] = blueprint

    apply_reference_quality_overrides(library_summaries, reference_slides)
    assign_default_ranks(reference_slides)
    apply_blueprint_overrides(blueprints)

    reference_catalog = {
        "version": "0.1",
        "workspace": BASE_DIR.as_posix(),
        "libraries": library_summaries,
        "slides": sorted(reference_slides, key=lambda item: (item["library_id"], item["library_slide_no"])),
    }

    REFERENCE_CATALOG_PATH.write_text(json.dumps(reference_catalog, indent=2, ensure_ascii=False), encoding="utf-8")
    save_blueprints(BLUEPRINT_PATH, blueprints)
    print(REFERENCE_CATALOG_PATH)
    print(BLUEPRINT_PATH)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

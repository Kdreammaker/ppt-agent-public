from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation

BASE_DIR = Path(__file__).resolve().parents[1]

WATCH_PURPOSES = {"chart", "market", "analysis"}
ARTIFACT_PATTERNS = [
    "TEXT",
    "Financial",
    "Funding History",
    "Office Suite",
    "Market Index",
    "Alpha",
    "Client",
    "Sentiment",
    "Back-testing",
    "Source:",
]


def load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def normalize(value: str) -> str:
    return " ".join((value or "").split())


def shape_texts(pptx_path: Path, slide_no: int) -> list[dict[str, Any]]:
    prs = Presentation(str(pptx_path))
    slide = prs.slides[slide_no - 1]
    rows: list[dict[str, Any]] = []
    for index, shape in enumerate(slide.shapes, start=1):
        text = normalize(getattr(shape, "text", ""))
        if not text:
            continue
        rows.append(
            {
                "shape_index": index,
                "shape_name": getattr(shape, "name", ""),
                "text": text,
            }
        )
    return rows


def inspect_templates(reference_catalog: dict[str, Any]) -> list[dict[str, Any]]:
    issues: list[dict[str, Any]] = []
    for slide in reference_catalog.get("slides", []):
        if slide.get("purpose") not in WATCH_PURPOSES:
            continue
        library_path = BASE_DIR / slide["library_path"]
        for row in shape_texts(library_path, int(slide["library_slide_no"])):
            matched = [pattern for pattern in ARTIFACT_PATTERNS if pattern.lower() in row["text"].lower()]
            if not matched:
                continue
            issues.append(
                {
                    "slide_id": slide["slide_id"],
                    "template_key": slide["template_key"],
                    "purpose": slide.get("purpose"),
                    "variant": slide.get("variant"),
                    "library_path": slide["library_path"],
                    "library_slide_no": slide["library_slide_no"],
                    "shape_index": row["shape_index"],
                    "shape_name": row["shape_name"],
                    "matched_patterns": matched,
                    "text_preview": row["text"][:180],
                    "recommended_action": "Expose as a slot, remove from the curated template, or mark the slide as structure-only until cleaned.",
                }
            )
    return issues


def write_markdown(path: Path, payload: dict[str, Any]) -> None:
    lines = [
        "# Template Static Artifact Audit",
        "",
        f"- Issues: {payload['summary']['issues']}",
        f"- Slides affected: {payload['summary']['slides_affected']}",
        "",
    ]
    if payload["issues"]:
        lines.extend(
            [
                "| Slide ID | Shape | Patterns | Preview | Action |",
                "| --- | ---: | --- | --- | --- |",
            ]
        )
        for issue in payload["issues"]:
            preview = issue["text_preview"].replace("|", "\\|")
            lines.append(
                f"| `{issue['slide_id']}` | {issue['shape_index']} | "
                f"{', '.join(issue['matched_patterns'])} | {preview} | {issue['recommended_action']} |"
            )
    else:
        lines.append("_No static artifact patterns found._")
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--reference-catalog", default="config/reference_catalog.json")
    parser.add_argument("--output-json", default="outputs/reports/template_static_artifact_audit.json")
    parser.add_argument("--output-md", default="outputs/reports/template_static_artifact_audit.md")
    args = parser.parse_args(argv)

    reference_catalog = load_json(BASE_DIR / args.reference_catalog)
    issues = inspect_templates(reference_catalog)
    payload = {
        "summary": {
            "issues": len(issues),
            "slides_affected": len({issue["slide_id"] for issue in issues}),
            "purposes_checked": sorted(WATCH_PURPOSES),
            "patterns": ARTIFACT_PATTERNS,
        },
        "issues": issues,
    }
    output_json = BASE_DIR / args.output_json
    output_md = BASE_DIR / args.output_md
    output_json.parent.mkdir(parents=True, exist_ok=True)
    output_json.write_text(json.dumps(payload, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    write_markdown(output_md, payload)
    print(output_json)
    print(output_md)
    print(f"static_artifact_issues={len(issues)}")
    print(f"static_artifact_slides={payload['summary']['slides_affected']}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

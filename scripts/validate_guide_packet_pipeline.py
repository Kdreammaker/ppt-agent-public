from __future__ import annotations

import argparse
import json
import sys
import tempfile
from pathlib import Path

from PIL import Image, ImageChops, ImageStat
from pptx import Presentation

BASE_DIR = Path(__file__).resolve().parents[1]
if str(BASE_DIR) not in sys.path:
    sys.path.insert(0, str(BASE_DIR))

from system.guide_packet import build_auto_variants, build_from_guide_packet, load_guide_packet, safe_rel


REQUIRED_REPORTS = [
    "generated.pptx",
    "deck-plan.json",
    "renderer-contract.json",
    "guide-compliance-report.json",
    "final-qa.json",
    "used-assets-report.json",
    "html-guide-render-report.json",
]

REQUIRED_FINAL_QA_FIELDS = [
    "approved_assets_used",
    "asset_slots_fallbacks",
    "unresolved_blockers",
    "overflow_overlap_scan",
    "contrast_scan",
    "checksum_validation_result",
    "layout_archetypes_used",
]


def assert_true(condition: bool, message: str) -> None:
    if not condition:
        raise AssertionError(message)


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def validate_project(project_dir: Path, expected_slides: int) -> dict:
    missing = [name for name in REQUIRED_REPORTS if not (project_dir / name).exists()]
    assert_true(not missing, f"missing required outputs: {missing}")
    prs = Presentation(str(project_dir / "generated.pptx"))
    assert_true(len(prs.slides) == expected_slides, f"slide count mismatch: {len(prs.slides)} != {expected_slides}")
    qa = load_json(project_dir / "final-qa.json")
    missing_qa_fields = [field for field in REQUIRED_FINAL_QA_FIELDS if field not in qa]
    assert_true(not missing_qa_fields, f"final QA missing fields: {missing_qa_fields}")
    assert_true(not qa["unresolved_blockers"], f"final QA has unresolved blockers: {qa['unresolved_blockers']}")
    assert_true(qa["html_screenshot_used_in_pptx"] is False, "HTML screenshot flag must be false")
    assert_true(qa["native_powerpoint_rendering"] is True, "native PPT rendering flag must be true")
    assert_true(bool(qa["palette_roles_used"]), "palette roles were not reported")
    assert_true(bool(qa["typography_roles_used"]), "typography roles were not reported")
    deck_plan = load_json(project_dir / "deck-plan.json")
    private_markers = ["C:\\", "/Users/", "/home/", "drive_id", "source_attachment"]
    raw = json.dumps(deck_plan, ensure_ascii=False)
    assert_true(not any(marker in raw for marker in private_markers), "deck-plan contains private marker")
    final_raw = json.dumps(qa, ensure_ascii=False)
    assert_true("C:\\" not in final_raw and "/Users/" not in final_raw and "/home/" not in final_raw, "final QA contains local path")
    previews = list((project_dir / "previews").glob("*.png"))
    assert_true(len(previews) == expected_slides, f"preview count mismatch: {len(previews)} != {expected_slides}")
    used_assets = load_json(project_dir / "used-assets-report.json")
    expected_slots = sum(len(slide.get("asset_slots", [])) for slide in load_json(project_dir / "guide-data.public.json")["slide_plan"]["slides"])
    assert_true(used_assets["recorded_asset_slots"] == expected_slots, "not all asset slots were reported")
    return {"project_dir": safe_rel(project_dir), "slides": expected_slides, "previews": len(previews)}


def validate_invalid_packet_blocks(valid_guide: Path) -> None:
    with tempfile.TemporaryDirectory(prefix="guide_packet_invalid_") as tmp:
        tmp_path = Path(tmp)
        data = load_json(valid_guide)
        data["guide_identity"]["slide_count"] = data["guide_identity"]["slide_count"] + 1
        invalid = tmp_path / "guide-data.public.json"
        invalid.write_text(json.dumps(data, indent=2), encoding="utf-8")
        try:
            load_guide_packet(invalid)
        except Exception:
            return
        raise AssertionError("invalid packet did not block before PPTX generation")


def validate_auto_variant_difference(auto_root: Path) -> dict:
    plan_a = load_json(auto_root / "variant-a" / "deck-plan.json")
    plan_b = load_json(auto_root / "variant-b" / "deck-plan.json")
    contract_a = load_json(auto_root / "variant-a" / "renderer-contract.json")
    contract_b = load_json(auto_root / "variant-b" / "renderer-contract.json")
    assert_true(
        plan_a.get("variant_strategy", {}).get("id") == "investor_open",
        "variant A deck-plan is missing investor_open strategy",
    )
    assert_true(
        plan_b.get("variant_strategy", {}).get("id") == "operator_dense",
        "variant B deck-plan is missing operator_dense strategy",
    )
    assert_true(
        contract_a.get("variant_strategy", {}).get("id") != contract_b.get("variant_strategy", {}).get("id"),
        "renderer contracts do not contain distinct variant strategies",
    )
    assert_true(
        contract_a.get("strategy_contract") != contract_b.get("strategy_contract"),
        "renderer contracts do not contain distinct strategy_contract values",
    )
    required_plan_fields = [
        "strategy_id",
        "layout_recipe",
        "content_emphasis",
        "evidence_treatment",
        "visual_asset_role",
        "palette_emphasis",
        "typography_role_bias",
        "chart_or_table_style",
        "density_budget",
    ]
    for plan_name, plan in {"variant-a": plan_a, "variant-b": plan_b}.items():
        for slide in plan.get("slides", []):
            missing = [field for field in required_plan_fields if field not in slide]
            assert_true(not missing, f"{plan_name} slide {slide.get('slide_no')} missing strategy plan fields: {missing}")
    recipes_a = [
        slide.get("renderer_metadata", {}).get("layout_strategy", {}).get("layout_recipe")
        for slide in plan_a.get("slides", [])
    ]
    recipes_b = [
        slide.get("renderer_metadata", {}).get("layout_strategy", {}).get("layout_recipe")
        for slide in plan_b.get("slides", [])
    ]
    changed_recipes = sum(1 for a, b in zip(recipes_a, recipes_b, strict=False) if a != b)
    assert_true(changed_recipes >= 8, f"machine-facing layout recipes are too similar: changed_recipes={changed_recipes}")

    a_paths = sorted((auto_root / "variant-a" / "previews").glob("slide_*.png"))
    b_paths = sorted((auto_root / "variant-b" / "previews").glob("slide_*.png"))
    assert_true(a_paths and len(a_paths) == len(b_paths), "auto variant preview sets are incomplete")
    mean_diffs: list[float] = []
    changed_slides = 0
    for a_path, b_path in zip(a_paths, b_paths, strict=True):
        a_img = Image.open(a_path).convert("RGB").resize((320, 180))
        b_img = Image.open(b_path).convert("RGB").resize((320, 180))
        diff = ImageChops.difference(a_img, b_img)
        mean = sum(ImageStat.Stat(diff).mean) / 3.0
        mean_diffs.append(mean)
        if mean >= 5.0:
            changed_slides += 1
    average_mean_diff = sum(mean_diffs) / len(mean_diffs)
    assert_true(
        average_mean_diff >= 8.0 and changed_slides >= max(8, len(mean_diffs) // 2),
        f"auto variants are not visually distinct enough: average_mean_diff={average_mean_diff:.2f}, changed_slides={changed_slides}",
    )
    comparison = load_json(auto_root / "variant-comparison-report.json")
    required_claim_terms = ["deck-plan.json", "renderer-contract.json", "strategy", "evidence treatment"]
    raw = json.dumps(comparison, ensure_ascii=False)
    missing_claims = [term for term in required_claim_terms if term not in raw]
    assert_true(not missing_claims, f"variant comparison report missing concrete claims: {missing_claims}")
    return {
        "average_mean_diff": round(average_mean_diff, 3),
        "slides_over_threshold": changed_slides,
        "slide_count": len(mean_diffs),
        "machine_plan_changed_recipes": changed_recipes,
        "variant_a_strategy": plan_a.get("variant_strategy", {}).get("id"),
        "variant_b_strategy": plan_b.get("variant_strategy", {}).get("id"),
    }


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Validate B167-B178 guide packet pipeline.")
    parser.add_argument("--guide", default="data/fixtures/lumaloop_guide/guide-data.public.json")
    parser.add_argument("--output-root", default="outputs/projects/guide-pipeline-validation")
    args = parser.parse_args(argv)

    guide = Path(args.guide).resolve()
    packet, _, _ = load_guide_packet(guide)
    validate_invalid_packet_blocks(guide)

    output_root = Path(args.output_root).resolve()
    assistant = build_from_guide_packet(guide, mode="assistant", output_root=output_root, project_id="assistant", html_guide_requested=True)
    assistant_project = output_root / "assistant"
    assistant_summary = validate_project(assistant_project, packet.guide_identity.slide_count)

    auto = build_auto_variants(guide, output_root=output_root, project_id="auto")
    auto_root = output_root / "auto"
    assert_true((auto_root / "variant-comparison-report.json").exists(), "missing auto comparison report")
    assert_true((auto_root / "auto-mode-recommendation.md").exists(), "missing auto recommendation")
    variant_a = validate_project(auto_root / "variant-a", packet.guide_identity.slide_count)
    variant_b = validate_project(auto_root / "variant-b", packet.guide_identity.slide_count)
    variant_difference = validate_auto_variant_difference(auto_root)

    result = {
        "status": "pass",
        "assistant": assistant_summary,
        "auto": auto,
        "variant_a": variant_a,
        "variant_b": variant_b,
        "variant_difference": variant_difference,
    }
    print(json.dumps(result, indent=2, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

from __future__ import annotations

import argparse
import json
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Any

import fitz
from PIL import Image, ImageStat
from pptx import Presentation

BASE_DIR = Path(__file__).resolve().parents[1]
if str(BASE_DIR) not in sys.path:
    sys.path.insert(0, str(BASE_DIR))

from scripts.build_deck import (
    indexed_reference_catalog,
    load_json,
    load_spec,
    old_catalog_slide_map,
    resolve_path,
    resolved_slide_sources,
)
from system.blueprint_loader import load_blueprints
from system.template_engine import coerce_text_to_slot_budget

DEFAULT_RENDER_ROOT = BASE_DIR / "outputs" / "previews" / "visual_smoke"
BLANK_STDDEV_THRESHOLD = 3.0
FOREGROUND_RATIO_THRESHOLD = 0.0025
MIN_RENDER_SIZE = 100

try:
    fitz.TOOLS.mupdf_display_errors(False)
except AttributeError:
    pass


def soffice_path() -> Path:
    env_path = os.environ.get("LIBREOFFICE_SOFFICE")
    if env_path:
        candidate = Path(env_path)
        if candidate.exists():
            return candidate
    path_candidate = shutil.which("soffice") or shutil.which("libreoffice")
    if path_candidate:
        return Path(path_candidate)
    candidates = [
        Path("C:/Program Files/LibreOffice/program/soffice.exe"),
        Path("C:/Program Files (x86)/LibreOffice/program/soffice.exe"),
        Path("/usr/bin/soffice"),
        Path("/usr/bin/libreoffice"),
        Path("/snap/bin/libreoffice"),
    ]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    raise FileNotFoundError("LibreOffice soffice.exe not found")


def convert_pptx_to_pdf(input_path: Path, output_dir: Path) -> Path:
    output_dir.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="lo-profile-") as profile_dir:
        profile_uri = Path(profile_dir).resolve().as_uri()
        process = subprocess.run(
            [
                str(soffice_path()),
                "--headless",
                f"-env:UserInstallation={profile_uri}",
                "--convert-to",
                "pdf",
                "--outdir",
                str(output_dir),
                str(input_path),
            ],
            check=True,
            capture_output=True,
            text=True,
        )
    pdf_path = output_dir / f"{input_path.stem}.pdf"
    if not pdf_path.exists():
        raise FileNotFoundError(
            f"Failed to render {input_path.name} to PDF: {process.stdout} {process.stderr}"
        )
    return pdf_path


def render_pdf_to_pngs(pdf_path: Path, output_dir: Path) -> list[Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(pdf_path)
    created: list[Path] = []
    for index in range(len(doc)):
        page = doc.load_page(index)
        pix = page.get_pixmap(matrix=fitz.Matrix(1.6, 1.6), alpha=False)
        output_path = output_dir / f"slide_{index + 1:02d}.png"
        pix.save(output_path)
        created.append(output_path)
    return created


def render_pptx(input_path: Path, output_dir: Path) -> list[Path]:
    pdf_dir = output_dir / "_pdf"
    pdf_path = convert_pptx_to_pdf(input_path, pdf_dir)
    return render_pdf_to_pngs(pdf_path, output_dir)


def normalize_text(value: str) -> str:
    return " ".join((value or "").split())


def slide_texts(slide) -> list[str]:
    return [
        normalize_text(getattr(shape, "text", ""))
        for shape in slide.shapes
        if normalize_text(getattr(shape, "text", ""))
    ]


def template_expected_text_rows(spec_path: Path) -> list[dict[str, Any]]:
    spec, spec_dir = load_spec(spec_path)
    old_catalog_path = resolve_path(spec_dir, spec.get("catalog_path"))
    reference_catalog_path = resolve_path(spec_dir, spec.get("reference_catalog_path", "../../config/reference_catalog.json"))
    blueprint_path = resolve_path(spec_dir, spec.get("blueprint_path", "../../config/template_blueprints.json"))

    old_catalog = load_json(old_catalog_path)
    reference_catalog = load_json(reference_catalog_path)
    blueprints = load_blueprints(blueprint_path) if blueprint_path else {"slides": {}}
    old_slide_map = old_catalog_slide_map(old_catalog)
    reference_by_key, _ = indexed_reference_catalog(reference_catalog)
    sources = resolved_slide_sources(
        spec,
        spec_dir,
        old_catalog,
        old_slide_map,
        reference_catalog,
        reference_by_key,
    )

    rows: list[dict[str, Any]] = []
    for slide_no, (slide_spec, resolved) in enumerate(zip(spec["slides"], sources, strict=False), start=1):
        if slide_spec["layout"] != "template_slide" or not resolved.get("slide_id"):
            continue
        blueprint = blueprints.get("slides", {}).get(resolved["slide_id"], {})
        text_slots = {item["slot"]: item for item in blueprint.get("editable_text_slots", [])}
        text_values = dict(slide_spec.get("text_slots", {}))
        source_kind_by_slot = {slot_name: "text_slots" for slot_name in text_values}
        sequential_values = slide_spec.get("text_values", [])
        sequential_slots = [item["slot"] for item in blueprint.get("editable_text_slots", []) if item["slot"] not in text_values]
        for slot_name, value in zip(sequential_slots, sequential_values, strict=False):
            text_values[slot_name] = value
            source_kind_by_slot[slot_name] = "text_values"

        for slot_name, value in text_values.items():
            if not value or "\n" in value or slot_name not in text_slots:
                continue
            override = slide_spec.get("slot_overrides", {}).get(slot_name, {})
            fit_strategy = override.get("fit_strategy") or text_slots[slot_name].get("fit_strategy") or "preserve_template"
            if fit_strategy == "manual_wrap":
                continue
            expected_value, _ = coerce_text_to_slot_budget(value, text_slots[slot_name], override)
            rows.append(
                {
                    "slide": slide_no,
                    "slot": slot_name,
                    "expected_text": expected_value,
                    "source_kind": source_kind_by_slot.get(slot_name, "unknown"),
                    "resolved_slide_id": resolved.get("slide_id"),
                    "template_key": resolved.get("template_key"),
                    "shape_index": text_slots[slot_name].get("shape_index"),
                    "fit_strategy": fit_strategy,
                }
            )
    return rows


def average_color(colors: list[tuple[int, int, int]]) -> tuple[int, int, int]:
    return tuple(int(sum(color[index] for color in colors) / len(colors)) for index in range(3))


def foreground_ratio(image: Image.Image) -> float:
    rgb = image.convert("RGB")
    width, height = rgb.size
    sample = rgb.resize((min(width, 320), max(1, int(height * min(width, 320) / width))))
    sample_width, sample_height = sample.size
    corner_colors = [
        sample.getpixel((0, 0)),
        sample.getpixel((sample_width - 1, 0)),
        sample.getpixel((0, sample_height - 1)),
        sample.getpixel((sample_width - 1, sample_height - 1)),
    ]
    background = average_color(corner_colors)
    if hasattr(sample, "get_flattened_data"):
        pixels = list(sample.get_flattened_data())
    else:
        pixels = list(sample.getdata())
    foreground_pixels = sum(
        1 for pixel in pixels if sum(abs(pixel[index] - background[index]) for index in range(3)) > 30
    )
    return foreground_pixels / max(len(pixels), 1)


def collect_image_metrics(image_path: Path, slide_no: int) -> dict[str, Any]:
    image = Image.open(image_path).convert("RGB")
    gray = image.convert("L")
    stat = ImageStat.Stat(gray)
    return {
        "slide": slide_no,
        "path": str(image_path),
        "width": image.width,
        "height": image.height,
        "stddev": round(float(stat.stddev[0]), 4),
        "foreground_ratio": round(foreground_ratio(image), 6),
    }


def inspect_image_metrics(metrics: list[dict[str, Any]]) -> list[dict[str, Any]]:
    issues: list[dict[str, Any]] = []
    for item in metrics:
        if item["width"] < MIN_RENDER_SIZE or item["height"] < MIN_RENDER_SIZE:
            issues.append(
                {
                    "severity": "error",
                    "slide": item["slide"],
                    "type": "tiny_render",
                    "message": f"Rendered image is too small: {item['width']}x{item['height']}.",
                }
            )
        if item["stddev"] < BLANK_STDDEV_THRESHOLD or item["foreground_ratio"] < FOREGROUND_RATIO_THRESHOLD:
            issues.append(
                {
                    "severity": "error",
                    "slide": item["slide"],
                    "type": "blank_slide",
                    "message": (
                        "Rendered slide appears blank "
                        f"(stddev={item['stddev']}, foreground_ratio={item['foreground_ratio']})."
                    ),
                }
            )
    return issues


def inspect_major_text(prs: Presentation, spec_path: Path | None) -> list[dict[str, Any]]:
    if spec_path is None:
        return []

    issues: list[dict[str, Any]] = []
    for row in template_expected_text_rows(spec_path):
        slide_no = row["slide"]
        slot_name = row["slot"]
        expected = row["expected_text"]
        if slide_no > len(prs.slides):
            issues.append(
                {
                    "severity": "warning",
                    "slide": slide_no,
                    "type": "major_text_missing",
                    "message": f"Spec references slide {slide_no}, but output has fewer slides.",
                    "slot": slot_name,
                    "expected_text": expected,
                    "text": expected[:120],
                    "actual_slide_texts": [],
                    "resolved_slide_id": row.get("resolved_slide_id"),
                    "template_key": row.get("template_key"),
                }
            )
            continue
        expected_normalized = normalize_text(expected)
        actual_texts = slide_texts(prs.slides[slide_no - 1])
        matched = any(text == expected_normalized for text in actual_texts)
        if not matched:
            issues.append(
                {
                    "severity": "warning",
                    "slide": slide_no,
                    "type": "major_text_missing",
                    "message": f"Could not read back expected template text for slot '{slot_name}'.",
                    "slot": slot_name,
                    "expected_text": expected,
                    "text": expected[:120],
                    "actual_slide_texts": actual_texts[:12],
                    "resolved_slide_id": row.get("resolved_slide_id"),
                    "template_key": row.get("template_key"),
                }
            )
    return issues


def validate(pptx_path: Path, spec_path: Path | None = None, render_dir: Path | None = None) -> dict[str, Any]:
    prs = Presentation(str(pptx_path))
    issues: list[dict[str, Any]] = []
    image_metrics: list[dict[str, Any]] = []
    rendered_paths: list[Path] = []
    temp_dir: tempfile.TemporaryDirectory[str] | None = None
    active_render_dir = render_dir
    if active_render_dir is None:
        temp_dir = tempfile.TemporaryDirectory(prefix="ppt_visual_smoke_")
        active_render_dir = Path(temp_dir.name)

    try:
        rendered_paths = render_pptx(pptx_path, active_render_dir)
        image_metrics = [
            collect_image_metrics(path, slide_no)
            for slide_no, path in enumerate(rendered_paths, start=1)
        ]
        issues.extend(inspect_image_metrics(image_metrics))
    except Exception as exc:
        issues.append(
            {
                "severity": "error",
                "slide": None,
                "type": "render_failed",
                "message": str(exc),
            }
        )
    finally:
        if temp_dir is not None:
            temp_dir.cleanup()

    if rendered_paths and len(rendered_paths) != len(prs.slides):
        issues.append(
            {
                "severity": "error",
                "slide": None,
                "type": "slide_count_mismatch",
                "message": f"PPTX has {len(prs.slides)} slides but renderer produced {len(rendered_paths)} images.",
            }
        )

    issues.extend(inspect_major_text(prs, spec_path))

    return {
        "file": str(pptx_path),
        "slide_count": len(prs.slides),
        "rendered_slide_count": len(rendered_paths),
        "errors": sum(1 for issue in issues if issue["severity"] == "error"),
        "warnings": sum(1 for issue in issues if issue["severity"] == "warning"),
        "issue_count": len(issues),
        "issues": issues,
        "image_metrics": image_metrics,
    }


def default_render_dir(pptx_path: Path) -> Path:
    return DEFAULT_RENDER_ROOT / pptx_path.stem


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Run low-noise rendered visual smoke checks on a PPTX.")
    parser.add_argument("pptx")
    parser.add_argument("output_json", nargs="?")
    parser.add_argument("--spec", help="Optional deck spec for template major-text readback checks.")
    parser.add_argument("--render-dir", help="Directory for rendered image output.")
    parser.add_argument("--keep-images", action="store_true", help="Persist rendered images under outputs/previews/visual_smoke.")
    args = parser.parse_args(argv)

    pptx_path = Path(args.pptx).resolve()
    spec_path = Path(args.spec).resolve() if args.spec else None
    render_dir = None
    temp_dir: tempfile.TemporaryDirectory[str] | None = None
    if args.render_dir:
        render_dir = Path(args.render_dir).resolve()
        if render_dir.exists():
            shutil.rmtree(render_dir)
    elif args.keep_images:
        render_dir = default_render_dir(pptx_path)
        if render_dir.exists():
            shutil.rmtree(render_dir)
    else:
        temp_dir = tempfile.TemporaryDirectory(prefix="ppt_visual_smoke_")
        render_dir = Path(temp_dir.name)

    try:
        result = validate(pptx_path, spec_path, render_dir)
    finally:
        if temp_dir is not None:
            temp_dir.cleanup()

    if args.output_json:
        output_path = Path(args.output_json).resolve()
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(json.dumps(result, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
        print(output_path)
    else:
        print(json.dumps(result, indent=2, ensure_ascii=False))

    return 1 if result["errors"] else 0


if __name__ == "__main__":
    raise SystemExit(main())

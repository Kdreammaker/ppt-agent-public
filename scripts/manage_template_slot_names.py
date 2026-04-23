from __future__ import annotations

import argparse
import json
import sys
from collections import Counter, defaultdict
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation

BASE_DIR = Path(__file__).resolve().parents[1]
if str(BASE_DIR) not in sys.path:
    sys.path.insert(0, str(BASE_DIR))

from system.blueprint_loader import load_blueprints, save_blueprints
from system.pptx_system import get_shape_by_index

BLUEPRINT_PATH = BASE_DIR / "config" / "template_blueprints.json"
DEFAULT_MANIFEST_PATH = BASE_DIR / "config" / "template_slot_name_manifest.json"
DEFAULT_REPORT_JSON = BASE_DIR / "outputs" / "reports" / "template_slot_name_audit.json"
DEFAULT_REPORT_MD = BASE_DIR / "outputs" / "reports" / "template_slot_name_audit.md"

LOW_RISK_PURPOSES = {"cover", "summary", "closing"}
MEDIUM_RISK_PURPOSES = {"issue", "process", "strategy", "team"}
HIGH_RISK_PURPOSES = {"analysis", "chart", "market", "timeline"}


@dataclass
class SlotAudit:
    library_id: str
    purpose: str
    variant: str
    slide_id: str
    slot: str
    slot_kind: str
    shape_index: int | None
    current_shape_name: str | None
    blueprint_shape_name: str | None
    desired_shape_name: str
    status: str
    risk: str
    recommended_action: str

    def as_dict(self) -> dict[str, Any]:
        return {
            "library_id": self.library_id,
            "purpose": self.purpose,
            "variant": self.variant,
            "slide_id": self.slide_id,
            "slot": self.slot,
            "slot_kind": self.slot_kind,
            "shape_index": self.shape_index,
            "current_shape_name": self.current_shape_name,
            "blueprint_shape_name": self.blueprint_shape_name,
            "desired_shape_name": self.desired_shape_name,
            "status": self.status,
            "risk": self.risk,
            "recommended_action": self.recommended_action,
        }


def load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def save_json(path: Path, data: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def desired_shape_name(slot_name: str) -> str:
    return f"slot:{slot_name}"


def iter_slots(blueprint: dict[str, Any]):
    for slot in blueprint.get("editable_text_slots", []):
        yield "text", slot
    for slot in blueprint.get("editable_image_slots", []):
        yield "image", slot


def slot_list_name(slot_kind: str) -> str:
    if slot_kind == "text":
        return "editable_text_slots"
    if slot_kind == "image":
        return "editable_image_slots"
    raise ValueError(f"Unsupported slot_kind for manifest: {slot_kind}")


def find_blueprint_slot(
    blueprint: dict[str, Any],
    slot_name: str,
    slot_kind: str,
    shape_index: int | None,
) -> dict[str, Any] | None:
    for slot in blueprint.get(slot_list_name(slot_kind), []):
        if slot.get("slot") == slot_name and slot.get("shape_index") == shape_index:
            return slot
    return None


def selected_slide_ids(blueprints: dict[str, Any], slide_id: str | None) -> list[str]:
    if slide_id:
        if slide_id not in blueprints["slides"]:
            raise KeyError(f"Unknown slide_id: {slide_id}")
        return [slide_id]
    return sorted(blueprints["slides"])


def selected_slots(blueprint: dict[str, Any], slot_names: set[str] | None):
    for slot_kind, slot in iter_slots(blueprint):
        if slot_names is None or slot["slot"] in slot_names:
            yield slot_kind, slot


def classify_risk(blueprint: dict[str, Any], status: str) -> str:
    if status in {"missing_shape", "manual_review"}:
        return "high"
    purpose = blueprint.get("purpose", "")
    if purpose in HIGH_RISK_PURPOSES or blueprint.get("mode") == "overlay":
        return "high"
    if purpose in MEDIUM_RISK_PURPOSES:
        return "medium"
    if purpose in LOW_RISK_PURPOSES:
        return "low"
    return "medium"


def recommended_action(status: str, risk: str) -> str:
    if status == "ok":
        return "no_action"
    if status == "blueprint_stale":
        return "update_blueprint_shape_name"
    if status == "missing_shape":
        return "repair_blueprint_shape_index"
    if status == "manual_review":
        return "inspect_slot_identity"
    if status == "rename_available":
        if risk == "low":
            return "apply_in_low_risk_batch"
        if risk == "medium":
            return "queue_after_low_risk_batch"
        return "manual_review_before_apply"
    return "inspect_slot_identity"


def audit_slot(
    slide,
    blueprint: dict[str, Any],
    slide_id: str,
    slot_kind: str,
    slot: dict[str, Any],
    desired_name_counts: Counter[str],
) -> SlotAudit:
    shape_index = slot.get("shape_index")
    shape = get_shape_by_index(slide, shape_index) if shape_index else None
    current_name = getattr(shape, "name", None) if shape is not None else None
    desired = desired_shape_name(slot["slot"])
    desired_conflict = desired_name_counts[desired] > 1
    desired_used_by_other_shape = any(
        candidate_index != shape_index and getattr(candidate, "name", None) == desired
        for candidate_index, candidate in enumerate(slide.shapes, start=1)
    )

    if shape is None:
        status = "missing_shape"
    elif desired_conflict or desired_used_by_other_shape:
        status = "manual_review"
    elif current_name and current_name.startswith("slot:") and current_name != desired:
        status = "manual_review"
    elif current_name == desired and slot.get("shape_name") == desired:
        status = "ok"
    elif current_name == desired:
        status = "blueprint_stale"
    else:
        status = "rename_available"

    risk = classify_risk(blueprint, status)
    return SlotAudit(
        library_id=blueprint["library_id"],
        purpose=blueprint["purpose"],
        variant=blueprint["variant"],
        slide_id=slide_id,
        slot=slot["slot"],
        slot_kind=slot_kind,
        shape_index=shape_index,
        current_shape_name=current_name,
        blueprint_shape_name=slot.get("shape_name"),
        desired_shape_name=desired,
        status=status,
        risk=risk,
        recommended_action=recommended_action(status, risk),
    )


def summarize(rows: list[SlotAudit], field: str) -> Counter[str]:
    return Counter(str(getattr(row, field)) for row in rows)


def summarize_pair(rows: list[SlotAudit], left: str, right: str) -> dict[str, Counter[str]]:
    grouped: dict[str, Counter[str]] = defaultdict(Counter)
    for row in rows:
        grouped[str(getattr(row, left))][str(getattr(row, right))] += 1
    return dict(sorted(grouped.items()))


def markdown_table(headers: list[str], rows: list[list[Any]]) -> list[str]:
    if not rows:
        return ["_None._"]
    lines = [
        "| " + " | ".join(headers) + " |",
        "| " + " | ".join("---" for _ in headers) + " |",
    ]
    for row in rows:
        lines.append("| " + " | ".join(str(item) for item in row) + " |")
    return lines


def counter_rows(counter: Counter[str]) -> list[list[Any]]:
    return [[key, counter[key]] for key in sorted(counter)]


def grouped_counter_rows(grouped: dict[str, Counter[str]]) -> tuple[list[str], list[list[Any]]]:
    keys = sorted({key for counter in grouped.values() for key in counter})
    rows = []
    for group, counter in grouped.items():
        rows.append([group, *[counter[key] for key in keys], sum(counter.values())])
    return ["Group", *keys, "total"], rows


def write_json_report(path: Path, rows: list[SlotAudit]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    payload = {
        "summary": {
            "total_slots": len(rows),
            "by_status": dict(summarize(rows, "status")),
            "by_risk": dict(summarize(rows, "risk")),
            "by_library": dict(summarize(rows, "library_id")),
            "by_purpose": dict(summarize(rows, "purpose")),
        },
        "rows": [row.as_dict() for row in rows],
    }
    path.write_text(json.dumps(payload, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def manifest_entries(rows: list[SlotAudit], blueprints: dict[str, Any]) -> list[dict[str, Any]]:
    entries: list[dict[str, Any]] = []
    for row in sorted(
        (candidate for candidate in rows if candidate.status == "ok"),
        key=lambda item: (item.library_id, item.slide_id, item.slot_kind, item.slot),
    ):
        blueprint = blueprints["slides"][row.slide_id]
        entries.append(
            {
                "slide_id": row.slide_id,
                "library_id": row.library_id,
                "template_key": blueprint["template_key"],
                "library_path": blueprint["library_path"],
                "library_slide_no": blueprint["library_slide_no"],
                "slot": row.slot,
                "slot_kind": row.slot_kind,
                "shape_index": row.shape_index,
                "shape_name": row.desired_shape_name,
            }
        )
    return entries


def write_manifest(path: Path, rows: list[SlotAudit], blueprints: dict[str, Any]) -> None:
    entries = manifest_entries(rows, blueprints)
    payload = {
        "version": "0.1",
        "description": "Stable template slot shape names reapplied after curated library rebuilds.",
        "entries": entries,
    }
    save_json(path, payload)
    print(f"saved {path} ({len(entries)} entries)")


def apply_manifest(manifest_path: Path, blueprint_path: Path = BLUEPRINT_PATH) -> None:
    manifest = load_json(manifest_path)
    blueprints = load_blueprints(blueprint_path)
    presentation_cache: dict[Path, Presentation] = {}
    changed_libraries: set[Path] = set()
    changed_blueprint = False
    errors: list[str] = []
    seen_keys: set[tuple[str, str, str, int | None]] = set()

    for entry in manifest.get("entries", []):
        slide_id = entry.get("slide_id")
        slot_name = entry.get("slot")
        slot_kind = entry.get("slot_kind")
        shape_index = entry.get("shape_index")
        shape_name = entry.get("shape_name")
        key = (slide_id, slot_kind, slot_name, shape_index)

        if key in seen_keys:
            errors.append(f"{slide_id}:{slot_name} duplicate manifest entry")
            continue
        seen_keys.add(key)

        if not shape_name or not str(shape_name).startswith("slot:"):
            errors.append(f"{slide_id}:{slot_name} manifest shape_name must be slot:<name>")
            continue
        if slide_id not in blueprints.get("slides", {}):
            errors.append(f"{slide_id}:{slot_name} slide_id is not present in {blueprint_path}")
            continue

        blueprint = blueprints["slides"][slide_id]
        blueprint_slot = find_blueprint_slot(blueprint, slot_name, slot_kind, shape_index)
        if blueprint_slot is None:
            errors.append(
                f"{slide_id}:{slot_name} no blueprint {slot_kind} slot at shape_index={shape_index}"
            )
            continue

        library_path = BASE_DIR / Path(entry["library_path"])
        library_slide_no = int(entry["library_slide_no"])
        if not library_path.exists():
            errors.append(f"{slide_id}:{slot_name} library does not exist: {library_path}")
            continue
        if library_path not in presentation_cache:
            presentation_cache[library_path] = Presentation(str(library_path))
        prs = presentation_cache[library_path]
        if library_slide_no < 1 or library_slide_no > len(prs.slides):
            errors.append(
                f"{slide_id}:{slot_name} library_slide_no={library_slide_no} is out of range"
            )
            continue

        slide = prs.slides[library_slide_no - 1]
        shape = get_shape_by_index(slide, shape_index)
        if shape is None:
            errors.append(f"{slide_id}:{slot_name} missing PPTX shape at index={shape_index}")
            continue

        current_shape_name = getattr(shape, "name", None)
        if (
            current_shape_name
            and current_shape_name.startswith("slot:")
            and current_shape_name != shape_name
        ):
            errors.append(
                f"{slide_id}:{slot_name} shape index={shape_index} already has "
                f"{current_shape_name!r}, expected {shape_name!r}"
            )
            continue

        if current_shape_name != shape_name:
            shape.name = shape_name
            changed_libraries.add(library_path)
        if blueprint_slot.get("shape_name") != shape_name:
            blueprint_slot["shape_name"] = shape_name
            changed_blueprint = True

    if errors:
        joined = "\n".join(f"- {error}" for error in errors)
        raise RuntimeError(f"Manifest apply failed:\n{joined}")

    for library_path in sorted(changed_libraries):
        presentation_cache[library_path].save(str(library_path))
        print(f"saved {library_path}")
    if changed_blueprint:
        save_blueprints(blueprint_path, blueprints)
        print(f"saved {blueprint_path}")
    print(
        f"applied {len(manifest.get('entries', []))} manifest entries "
        f"from {manifest_path}"
    )


def write_markdown_report(path: Path, rows: list[SlotAudit]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    status_counter = summarize(rows, "status")
    risk_counter = summarize(rows, "risk")
    low_risk_candidates = [
        row for row in rows if row.status == "rename_available" and row.risk == "low"
    ]
    manual_review = [
        row
        for row in rows
        if row.status in {"manual_review", "missing_shape"} or row.recommended_action.startswith("manual")
    ]

    library_headers, library_rows = grouped_counter_rows(summarize_pair(rows, "library_id", "status"))
    purpose_headers, purpose_rows = grouped_counter_rows(summarize_pair(rows, "purpose", "status"))

    lines = [
        "# Template Slot Name Audit",
        "",
        "This report classifies editable text/image slots by current shape-name stability.",
        "",
        "## Summary",
        "",
        f"- Total slots: {len(rows)}",
        f"- Low-risk rename candidates: {len(low_risk_candidates)}",
        f"- Manual-review rows: {len(manual_review)}",
        "",
        "### By Status",
        "",
        *markdown_table(["Status", "Count"], counter_rows(status_counter)),
        "",
        "### By Risk",
        "",
        *markdown_table(["Risk", "Count"], counter_rows(risk_counter)),
        "",
        "## Library Summary",
        "",
        *markdown_table(library_headers, library_rows),
        "",
        "## Purpose Summary",
        "",
        *markdown_table(purpose_headers, purpose_rows),
        "",
        "## Low-Risk Rename Candidates",
        "",
        *markdown_table(
            ["Slide", "Slot", "Kind", "Current", "Desired", "Action"],
            [
                [
                    row.slide_id,
                    row.slot,
                    row.slot_kind,
                    row.current_shape_name,
                    row.desired_shape_name,
                    row.recommended_action,
                ]
                for row in low_risk_candidates[:120]
            ],
        ),
    ]
    if len(low_risk_candidates) > 120:
        lines.append(f"\n_Showing first 120 of {len(low_risk_candidates)} low-risk candidates._")

    lines.extend(
        [
            "",
            "## Manual Review",
            "",
            *markdown_table(
                ["Slide", "Slot", "Kind", "Status", "Risk", "Current", "Desired", "Action"],
                [
                    [
                        row.slide_id,
                        row.slot,
                        row.slot_kind,
                        row.status,
                        row.risk,
                        row.current_shape_name,
                        row.desired_shape_name,
                        row.recommended_action,
                    ]
                    for row in manual_review[:120]
                ],
            ),
        ]
    )
    if len(manual_review) > 120:
        lines.append(f"\n_Showing first 120 of {len(manual_review)} manual-review rows._")

    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def print_summary(rows: list[SlotAudit]) -> None:
    print(f"total_slots={len(rows)}")
    for label, counter in (
        ("status", summarize(rows, "status")),
        ("risk", summarize(rows, "risk")),
        ("library", summarize(rows, "library_id")),
        ("purpose", summarize(rows, "purpose")),
    ):
        print(f"{label}:")
        for key in sorted(counter):
            print(f"  {key}: {counter[key]}")


def collect_audit_rows(blueprints: dict[str, Any], slide_id: str | None, slot_filter: set[str] | None) -> list[SlotAudit]:
    presentation_cache: dict[Path, Presentation] = {}
    audit_rows: list[SlotAudit] = []

    for current_slide_id in selected_slide_ids(blueprints, slide_id):
        blueprint = blueprints["slides"][current_slide_id]
        library_path = BASE_DIR / Path(blueprint["library_path"])
        if library_path not in presentation_cache:
            presentation_cache[library_path] = Presentation(str(library_path))
        prs = presentation_cache[library_path]
        slide = prs.slides[blueprint["library_slide_no"] - 1]
        filtered_slots = list(selected_slots(blueprint, slot_filter))
        desired_name_counts = Counter(desired_shape_name(slot["slot"]) for _, slot in filtered_slots)

        for slot_kind, slot in filtered_slots:
            audit_rows.append(
                audit_slot(slide, blueprint, current_slide_id, slot_kind, slot, desired_name_counts)
            )
    return audit_rows


def apply_slot_names(blueprint_path: Path, blueprints: dict[str, Any], rows: list[SlotAudit]) -> None:
    presentation_cache: dict[Path, Presentation] = {}
    changed_libraries: set[Path] = set()
    changed_blueprint = False
    rows_by_slot = {(row.slide_id, row.slot, row.slot_kind): row for row in rows}

    for slide_id, blueprint in blueprints["slides"].items():
        library_path = BASE_DIR / Path(blueprint["library_path"])
        for slot_kind, slot in selected_slots(blueprint, None):
            row = rows_by_slot.get((slide_id, slot["slot"], slot_kind))
            if row is None or row.status in {"missing_shape", "manual_review"}:
                continue
            if library_path not in presentation_cache:
                presentation_cache[library_path] = Presentation(str(library_path))
            prs = presentation_cache[library_path]
            slide = prs.slides[blueprint["library_slide_no"] - 1]
            shape = get_shape_by_index(slide, slot.get("shape_index"))
            if shape is None:
                continue
            if shape.name != row.desired_shape_name:
                shape.name = row.desired_shape_name
                changed_libraries.add(library_path)
            if slot.get("shape_name") != row.desired_shape_name:
                slot["shape_name"] = row.desired_shape_name
                changed_blueprint = True

    for library_path in sorted(changed_libraries):
        presentation_cache[library_path].save(str(library_path))
        print(f"saved {library_path}")
    if changed_blueprint:
        save_blueprints(blueprint_path, blueprints)
        print(f"saved {blueprint_path}")


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Audit or apply stable template slot shape names.")
    parser.add_argument("--blueprint-path", default=str(BLUEPRINT_PATH))
    parser.add_argument("--slide-id")
    parser.add_argument("--slots", nargs="*", help="Slot names to audit/apply. Defaults to all selected slide slots.")
    parser.add_argument("--apply", action="store_true", help="Persist shape_name changes to the library PPTX and blueprint JSON.")
    parser.add_argument("--json", action="store_true", help="Print audit rows as JSON.")
    parser.add_argument("--output-json", help=f"Write full audit report JSON. Default path: {DEFAULT_REPORT_JSON}")
    parser.add_argument("--output-md", help=f"Write full audit report Markdown. Default path: {DEFAULT_REPORT_MD}")
    parser.add_argument("--summary", action="store_true", help="Print status/risk/library/purpose summary.")
    parser.add_argument("--write-manifest", help=f"Write ok slot identities to a manifest. Default path: {DEFAULT_MANIFEST_PATH}")
    parser.add_argument("--manifest", default=str(DEFAULT_MANIFEST_PATH), help="Manifest path for --apply-manifest.")
    parser.add_argument("--apply-manifest", action="store_true", help="Reapply slot names from a manifest to library PPTX and blueprint JSON.")
    args = parser.parse_args(argv)

    blueprint_path = Path(args.blueprint_path).resolve()
    blueprints = load_blueprints(blueprint_path)
    slot_filter = set(args.slots) if args.slots else None
    audit_rows = collect_audit_rows(blueprints, args.slide_id, slot_filter)

    if args.apply:
        apply_slot_names(blueprint_path, blueprints, audit_rows)
    if args.write_manifest:
        write_manifest(Path(args.write_manifest).resolve(), audit_rows, blueprints)
    if args.apply_manifest:
        apply_manifest(Path(args.manifest).resolve(), blueprint_path)

    if args.output_json:
        write_json_report(Path(args.output_json).resolve(), audit_rows)
    if args.output_md:
        write_markdown_report(Path(args.output_md).resolve(), audit_rows)
    if args.summary:
        print_summary(audit_rows)

    rows = [row.as_dict() for row in audit_rows]
    if args.json:
        print(json.dumps(rows, indent=2, ensure_ascii=False))
    elif (
        not args.summary
        and not args.output_json
        and not args.output_md
        and not args.write_manifest
        and not args.apply_manifest
    ):
        for row in rows:
            print(
                f"{row['slide_id']}:{row['slot']} "
                f"{row['slot_kind']} index={row['shape_index']} "
                f"current={row['current_shape_name']!r} desired={row['desired_shape_name']!r} "
                f"status={row['status']} risk={row['risk']} action={row['recommended_action']}"
            )
    return 1 if any(row.status == "missing_shape" for row in audit_rows) else 0


if __name__ == "__main__":
    raise SystemExit(main())

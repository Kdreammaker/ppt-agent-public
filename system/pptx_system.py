from __future__ import annotations

import json
import shutil
from io import BytesIO
from copy import deepcopy
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Inches, Pt


EMU_PER_INCH = 914400
DEFAULT_PROTECTED_TOKENS = [
    "API",
    "WebSocket",
    "UI/UX",
    "TAM",
    "SAM",
    "SOM",
    "PoC",
    "ROI",
    "MAU",
    "KPI",
]


def hex_to_rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    if len(value) != 6:
        raise ValueError(f"Expected 6-digit hex color, got: {value}")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def resolve_color(theme: "ThemeConfig", token: str | None, fallback: str = "primary") -> RGBColor:
    if token is None:
        return theme.colors[fallback]
    if token in theme.colors:
        return theme.colors[token]
    if token.startswith("#"):
        return hex_to_rgb(token)
    raise KeyError(f"Unknown color token: {token}")


def rel_path(base_dir: Path, path: str | Path) -> str:
    path = Path(path).resolve()
    try:
        return path.relative_to(base_dir.resolve()).as_posix()
    except ValueError:
        return path.as_posix()


def shape_bounds_inches(shape) -> dict[str, float]:
    return {
        "left": round(shape.left / EMU_PER_INCH, 3),
        "top": round(shape.top / EMU_PER_INCH, 3),
        "width": round(shape.width / EMU_PER_INCH, 3),
        "height": round(shape.height / EMU_PER_INCH, 3),
    }


def normalize_fingerprint(text: str) -> str:
    return " ".join((text or "").split())[:140]


def remove_shape(shape) -> None:
    element = shape._element
    parent = element.getparent()
    if parent is not None:
        parent.remove(element)


def find_blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if layout.name.lower() == "blank":
            return layout
    return prs.slide_layouts[len(prs.slide_layouts) - 1]


def remove_all_slides(prs: Presentation) -> None:
    sld_id_list = prs.slides._sldIdLst  # pyright: ignore[reportAttributeAccessIssue]
    original_items = list(sld_id_list)
    for item in original_items:
        prs.part.drop_rel(item.rId)
    for item in list(sld_id_list):
        sld_id_list.remove(item)


def reorder_and_trim_slides(prs: Presentation, keep_order: list[int]) -> None:
    sld_id_list = prs.slides._sldIdLst  # pyright: ignore[reportAttributeAccessIssue]
    original_items = list(sld_id_list)
    kept_items = [original_items[index - 1] for index in keep_order]

    for item in original_items:
        if item not in kept_items:
            prs.part.drop_rel(item.rId)

    for item in list(sld_id_list):
        sld_id_list.remove(item)

    for item in kept_items:
        sld_id_list.append(item)


def import_slide(dest_prs: Presentation, src_slide) -> Any:
    blank_layout = find_blank_layout(dest_prs)
    dest_slide = dest_prs.slides.add_slide(blank_layout)

    for shape in list(dest_slide.shapes):
        remove_shape(shape)

    for shape in src_slide.shapes:
        new_element = deepcopy(shape.element)
        dest_slide.shapes._spTree.insert_element_before(new_element, "p:extLst")  # pyright: ignore[reportAttributeAccessIssue]

    source_rels = {rel.rId: rel for rel in src_slide.part.rels.values()}
    referenced_rids = {
        attr_value
        for element in dest_slide._element.iter()
        for attr_value in element.attrib.values()
        if attr_value in source_rels
    }
    excluded_reltypes = {
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout",
    }
    rid_mapping: dict[str, str] = {}
    for old_rid in referenced_rids:
        rel = source_rels[old_rid]
        if rel.reltype in excluded_reltypes:
            continue
        if rel.is_external:
            new_rid = dest_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel._target)
        elif rel.reltype == RT.IMAGE:
            _, new_rid = dest_slide.part.get_or_add_image_part(BytesIO(rel.target_part.blob))
        else:
            new_rid = dest_slide.part.rels.get_or_add(rel.reltype, rel._target)
        rid_mapping[old_rid] = new_rid

    if rid_mapping:
        for element in dest_slide._element.iter():
            for attr_name, attr_value in list(element.attrib.items()):
                if attr_value in rid_mapping:
                    element.set(attr_name, rid_mapping[attr_value])
    return dest_slide


def iter_text_shapes(slide):
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            yield shape


def get_shape_by_index(slide, shape_index: int):
    shapes = list(slide.shapes)
    if shape_index < 1 or shape_index > len(shapes):
        return None
    return shapes[shape_index - 1]


def replace_placeholder_text(prs: Presentation, mapping: dict[str, str]) -> None:
    for slide in prs.slides:
        for shape in iter_text_shapes(slide):
            text = getattr(shape, "text", "")
            normalized = " ".join(text.split())
            for source, target in mapping.items():
                source_normalized = " ".join(source.split())
                if normalized == source_normalized or source_normalized in normalized:
                    shape.text = target
                    break


def set_all_fonts(prs: Presentation, font_family: str) -> None:
    for slide in prs.slides:
        for shape in iter_text_shapes(slide):
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = font_family


def clear_slide_except(slide, keep_indices: set[int]) -> None:
    original_shapes = list(slide.shapes)
    for index, shape in enumerate(original_shapes, start=1):
        if index not in keep_indices:
            remove_shape(shape)


def protect_tokens(text: str, tokens: Iterable[str] | None = None) -> str:
    tokens = list(tokens or [])
    protected = text
    for token in tokens:
        if " " in token:
            protected = protected.replace(token, token.replace(" ", "\u00a0"))
    return protected


def wrap_text_for_layout(
    text: str,
    *,
    max_chars_per_line: int | None = None,
    protect_token_list: Iterable[str] | None = None,
) -> str:
    if not text:
        return ""

    protected = protect_tokens(text, protect_token_list)
    if max_chars_per_line is None or max_chars_per_line <= 0:
        return protected

    wrapped_paragraphs: list[str] = []
    for paragraph in protected.splitlines():
        words = paragraph.split(" ")
        if len(words) <= 1:
            wrapped_paragraphs.append(paragraph)
            continue

        lines: list[str] = []
        current = ""
        for word in words:
            trial = word if not current else f"{current} {word}"
            visible_len = len(trial.replace("\u00a0", " "))
            if current and visible_len > max_chars_per_line:
                lines.append(current)
                current = word
            else:
                current = trial
        if current:
            lines.append(current)
        wrapped_paragraphs.append("\n".join(lines))
    return "\n".join(wrapped_paragraphs)


def set_text(
    shape,
    text: str,
    *,
    font_name: str,
    font_size: float,
    font_color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN | None = None,
    vertical: MSO_ANCHOR = MSO_ANCHOR.TOP,
    max_chars_per_line: int | None = None,
    protect_token_list: Iterable[str] | None = None,
    is_placeholder: bool = False,
    fit_strategy: str = "preserve_template",
) -> None:
    if not hasattr(shape, "text_frame"):
        return

    text_frame = shape.text_frame
    template_run_style: dict[str, Any] = {}
    template_alignment = None
    if is_placeholder:
        for paragraph in text_frame.paragraphs:
            if template_alignment is None:
                template_alignment = paragraph.alignment
            for run in paragraph.runs:
                template_run_style = {
                    "name": run.font.name,
                    "size": run.font.size,
                    "bold": run.font.bold,
                    "italic": run.font.italic,
                    "underline": run.font.underline,
                }
                try:
                    template_run_style["rgb"] = run.font.color.rgb
                except AttributeError:
                    template_run_style["rgb"] = None
                break
            if template_run_style:
                break
    text_frame.clear()
    
    # 템플릿 퍼스트(Template-First) 지원: 
    # 슬라이드 마스터의 Placeholder인 경우, 강제 리사이징 끄지 않음
    if is_placeholder and fit_strategy == "shrink":
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    if not is_placeholder:
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        text_frame.vertical_anchor = vertical
        text_frame.margin_left = Pt(4)
        text_frame.margin_right = Pt(4)
        text_frame.margin_top = Pt(3)
        text_frame.margin_bottom = Pt(3)

    prepared = wrap_text_for_layout(
        text,
        max_chars_per_line=max_chars_per_line if (not is_placeholder or fit_strategy == "manual_wrap") else None,
        protect_token_list=protect_token_list,
    )
    lines = prepared.split("\n") if prepared else [""]
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        if not is_placeholder or align is not None:
            paragraph.alignment = align if align is not None else PP_ALIGN.LEFT
        elif template_alignment is not None:
            paragraph.alignment = template_alignment
        run = paragraph.add_run()
        run.text = line
        
        # Placeholder가 아닐 때만 하드코딩 폰트/스타일 강제 주입
        if is_placeholder:
            if template_run_style.get("name") is not None:
                run.font.name = template_run_style["name"]
            if template_run_style.get("size") is not None:
                run.font.size = template_run_style["size"]
            if template_run_style.get("bold") is not None:
                run.font.bold = template_run_style["bold"]
            if template_run_style.get("italic") is not None:
                run.font.italic = template_run_style["italic"]
            if template_run_style.get("underline") is not None:
                run.font.underline = template_run_style["underline"]
            if template_run_style.get("rgb") is not None:
                run.font.color.rgb = template_run_style["rgb"]
        else:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = font_color
            run.font.bold = bold


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_name: str,
    font_size: float,
    font_color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN | None = None,
    vertical: MSO_ANCHOR = MSO_ANCHOR.TOP,
    max_chars_per_line: int | None = None,
    protect_token_list: Iterable[str] | None = None,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    set_text(
        box,
        text,
        font_name=font_name,
        font_size=font_size,
        font_color=font_color,
        bold=bold,
        align=align,
        vertical=vertical,
        max_chars_per_line=max_chars_per_line,
        protect_token_list=protect_token_list,
    )


@dataclass
class ThemeConfig:
    name: str
    font_family: str
    colors: dict[str, RGBColor]
    sizes: dict[str, float]

    @classmethod
    def from_json(cls, path: str | Path) -> "ThemeConfig":
        data = json.loads(Path(path).read_text(encoding="utf-8"))
        colors = {key: hex_to_rgb(value) for key, value in data["colors"].items()}
        return cls(
            name=data["name"],
            font_family=data["font_family"],
            colors=colors,
            sizes=data["sizes"],
        )


def add_simple_bar_chart(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    categories: list[str],
    values: list[float],
    theme: ThemeConfig,
) -> None:
    if len(categories) != len(values):
        raise ValueError("categories and values length mismatch")
    if not values:
        return

    max_value = max(values)
    bar_width = width / max(len(values), 1)

    for idx, (category, value) in enumerate(zip(categories, values, strict=False)):
        bar_height = 0 if max_value == 0 else (value / max_value) * (height - 0.5)
        x = left + idx * bar_width + 0.1
        y = top + height - bar_height - 0.2
        w = bar_width - 0.2

        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(bar_height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = theme.colors["accent_1"]
        shape.line.fill.background()

        add_textbox(
            slide,
            x,
            top + height - 0.18,
            w,
            0.18,
            category,
            font_name=theme.font_family,
            font_size=max(theme.sizes["body"] - 1, 10),
            font_color=theme.colors["gray"],
            align=PP_ALIGN.CENTER,
            protect_token_list=DEFAULT_PROTECTED_TOKENS,
        )


def set_table_cell_text(
    cell,
    text: str,
    *,
    font_name: str,
    font_size: float,
    font_color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    text_frame = cell.text_frame
    text_frame.clear()
    text_frame.margin_left = Pt(6)
    text_frame.margin_right = Pt(5)
    text_frame.margin_top = Pt(2)
    text_frame.margin_bottom = Pt(2)
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold


def add_simple_table(
    slide,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    headers: list[str],
    rows: list[list[str | int | float]],
    theme: ThemeConfig,
) -> None:
    if not rows:
        raise ValueError("table rows must not be empty")
    column_count = len(headers) if headers else len(rows[0])
    if column_count <= 0:
        raise ValueError("table must have at least one column")
    for row in rows:
        if len(row) != column_count:
            raise ValueError("table rows must be rectangular")

    row_count = len(rows) + (1 if headers else 0)
    table_shape = slide.shapes.add_table(
        row_count,
        column_count,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    table = table_shape.table
    for column in table.columns:
        column.width = Inches(width / column_count)
    for row in table.rows:
        row.height = Inches(height / row_count)

    body_font_size = max(theme.sizes.get("body", 11) - 1, 8.5)
    header_font_size = max(theme.sizes.get("body_strong", body_font_size), body_font_size)
    start_row = 0
    if headers:
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme.colors["primary"]
            set_table_cell_text(
                cell,
                str(header),
                font_name=theme.font_family,
                font_size=header_font_size,
                font_color=theme.colors["white"],
                bold=True,
                align=PP_ALIGN.CENTER,
            )
        start_row = 1

    for row_idx, row_values in enumerate(rows, start=start_row):
        for col_idx, value in enumerate(row_values):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme.colors["white"] if (row_idx - start_row) % 2 == 0 else theme.colors["light"]
            align = PP_ALIGN.RIGHT if isinstance(value, (int, float)) else PP_ALIGN.LEFT
            set_table_cell_text(
                cell,
                str(value),
                font_name=theme.font_family,
                font_size=body_font_size,
                font_color=theme.colors["primary"],
                align=align,
            )


def add_image(slide, image_path: str | Path, *, left: float, top: float, width: float, height: float):
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def build_curated_template(
    *,
    source_template: str | Path,
    output_path: str | Path,
    keep_slides: list[int],
    placeholder_map: dict[str, str],
    theme: ThemeConfig,
) -> Path:
    source_template = Path(source_template)
    output_path = Path(output_path)
    shutil.copyfile(source_template, output_path)
    prs = Presentation(str(output_path))

    reorder_and_trim_slides(prs, keep_slides)
    replace_placeholder_text(prs, placeholder_map)
    set_all_fonts(prs, theme.font_family)
    prs.save(str(output_path))
    return output_path


def load_catalog(path: str | Path) -> dict[str, Any]:
    return json.loads(Path(path).read_text(encoding="utf-8"))


def is_picture_shape(shape) -> bool:
    return shape.shape_type == MSO_SHAPE_TYPE.PICTURE

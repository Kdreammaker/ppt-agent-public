from __future__ import annotations

import datetime as dt
import hashlib
import html
import json
import re
import shutil
import sys
import time
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path, PurePosixPath
from typing import Any, Literal

import jsonschema
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pydantic import BaseModel, ConfigDict, Field, model_validator

BASE_DIR = Path(__file__).resolve().parents[1]
GUIDE_SCHEMA_PATH = BASE_DIR / "config" / "ppt-maker-design-guide-packet.schema.json"
STRATEGY_REGISTRY_PATH = BASE_DIR / "config" / "variant_strategy_registry.json"
PALETTE_PRESET_REGISTRY_PATH = BASE_DIR / "config" / "palette_preset_registry.json"
DEFAULT_PROJECTS_DIR = BASE_DIR / "outputs" / "projects"

APPROVED_ASSET_BLOCKING_REASONS = {
    "approved_asset_file_missing",
    "approved_asset_checksum_mismatch",
    "approved_asset_size_mismatch",
    "approved_asset_type_unsupported_for_slot",
    "approved_asset_ref_not_found_in_manifest",
    "approved_asset_path_unsafe",
    "approved_asset_insert_failed",
    "policy_blocked",
    "license_blocked",
}

PRIVATE_KEY_PATTERNS = [
    re.compile(r"\b[a-zA-Z]:\\"),
    re.compile(r"[A-Za-z]:\\(?:[^\\/:*?\"<>|\r\n]+\\)*[^\\/:*?\"<>|\r\n]*"),
    re.compile(r"/Users/"),
    re.compile(r"/(?:Users|home)/[^\s\"']+"),
    re.compile(r"/home/"),
    re.compile(r"\bdrive[_-]?id\b", re.IGNORECASE),
    re.compile(r"\btoken\b", re.IGNORECASE),
    re.compile(r"\basset_uid\b", re.IGNORECASE),
    re.compile(r"\bsource_attachment\b", re.IGNORECASE),
]

PROMPT_LITERAL_COMMAND_PATTERNS = [
    re.compile(pattern, re.IGNORECASE)
    for pattern in [
        r"처음부터",
        r"자동으로",
        r"만들어\s*줘",
        r"작성해\s*줘",
        r"발표\s*자료",
        r"슬라이드\s*만들어",
        r"PPT\s*만들어",
        r"보고서\s*만들어",
        r"create\s+(?:an?\s+)?\d{1,2}\s*[- ]?\s*(?:slide|slides|page|pages)(?:\s+(?:deck|presentation|proposal))?",
        r"make\s+(?:an?\s+)?\d{1,2}\s*[- ]?\s*(?:(?:slide|slides|page|pages)\s+)?(?:deck|presentation|proposal)",
        r"generate\s+\d{1,2}\s*(?:slides|pages)",
        r"build\s+(?:an?\s+)?\d{1,2}\s*[- ]?\s*(?:(?:slide|slides|page|pages)\s+)?(?:deck|presentation|proposal)",
        r"\d{1,2}\s*(?:장|쪽|페이지|슬라이드).*(?:만들어|작성|제작)",
        r"produce\s+two\s+visually\s+distinct\s+variants",
        r"please\s+create",
        r"make\s+a\s+deck",
        r"generate\s+slides",
    ]
]

PROMPT_LITERAL_COMMAND_LABELS = [
    "처음부터",
    "자동으로",
    "만들어줘",
    "작성해줘",
    "발표자료",
    "슬라이드 만들어",
    "PPT 만들어",
    "보고서 만들어",
    "create numbered slide request",
    "make numbered deck request",
    "generate numbered slides request",
    "build numbered presentation request",
    "Korean numbered deck request",
    "request for two visual variants",
    "please create",
    "make a deck",
    "generate slides",
]

COMPOSITION_FAMILY_ORDER = [
    "luxury_editorial",
    "food_product_launch",
    "internal_strategy_report",
    "product_launch",
    "marketing_campaign",
    "executive_brief",
    "retail_go_to_market",
    "screen_product_tour",
    "public_institution_report",
    "investor_pitch",
]

COMPOSITION_FAMILY_PROFILES: dict[str, dict[str, Any]] = {
    "luxury_editorial": {
        "preferred_archetypes": ["editorial_masthead", "product_hero_stage", "material_detail_spread"],
        "recipe_families": ["editorial_masthead", "product_hero_stage", "material_detail_spread", "curated_detail_cards"],
        "major_visual_asset_treatment": "editorial negative space with material/product detail surfaces",
        "avoid": ["generic_phone_mockup", "generic_three_card_summary"],
    },
    "food_product_launch": {
        "preferred_archetypes": ["product_hero_stage", "ingredient_detail_spread", "retail_shelf_story", "campaign_calendar"],
        "recipe_families": ["product_hero_stage", "ingredient_detail_spread", "retail_shelf_story", "campaign_calendar"],
        "major_visual_asset_treatment": "product stage, ingredient detail, retail shelf, and campaign timing",
        "avoid": ["unrelated_phone_mockup", "abstract_generic_bars_only"],
    },
    "internal_strategy_report": {
        "preferred_archetypes": ["executive_summary_memo", "decision_matrix", "risk_action_map", "roadmap_timeline"],
        "recipe_families": ["executive_summary_memo", "decision_matrix", "risk_action_map", "roadmap_timeline"],
        "major_visual_asset_treatment": "memo, decision, risk, and action-map primitives",
        "avoid": ["product_hero_stage", "phone_mockup", "campaign_visual"],
    },
    "product_launch": {
        "preferred_archetypes": ["product_hero_stage", "feature_story", "roadmap_timeline"],
        "recipe_families": [],
        "major_visual_asset_treatment": "product narrative with staged proof moments",
        "avoid": [],
    },
    "marketing_campaign": {
        "preferred_archetypes": ["campaign_calendar", "channel_grid", "audience_hook"],
        "recipe_families": ["campaign_calendar", "channel_grid", "audience_hook"],
        "major_visual_asset_treatment": "campaign beats and channel calendar",
        "avoid": [],
    },
    "executive_brief": {
        "preferred_archetypes": ["executive_summary_memo", "decision_matrix", "metric_brief"],
        "recipe_families": ["executive_summary_memo", "decision_matrix", "metric_brief"],
        "major_visual_asset_treatment": "executive decision surfaces",
        "avoid": [],
    },
    "retail_go_to_market": {
        "preferred_archetypes": ["retail_shelf_story", "channel_grid", "campaign_calendar"],
        "recipe_families": ["retail_shelf_story", "channel_grid", "campaign_calendar"],
        "major_visual_asset_treatment": "shelf and channel execution story",
        "avoid": [],
    },
    "screen_product_tour": {
        "preferred_archetypes": ["screen_product_tour", "journey_flow", "feature_story"],
        "recipe_families": ["screen_product_tour", "journey_flow", "feature_story"],
        "major_visual_asset_treatment": "screen walkthrough and user journey",
        "avoid": [],
    },
    "public_institution_report": {
        "preferred_archetypes": ["formal_memo", "compliance_grid", "roadmap_timeline"],
        "recipe_families": ["formal_memo", "compliance_grid", "roadmap_timeline"],
        "major_visual_asset_treatment": "formal memo and accountable evidence tables",
        "avoid": [],
    },
    "investor_pitch": {
        "preferred_archetypes": ["editorial_masthead", "market_proof", "roadmap_timeline"],
        "recipe_families": ["investor_hero", "market_proof", "roadmap_timeline"],
        "major_visual_asset_treatment": "pitch hero and selective proof",
        "avoid": [],
    },
}


def utc_now() -> str:
    return dt.datetime.now(dt.UTC).replace(microsecond=0).isoformat()


def json_dump(path: Path, payload: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(payload, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")


def safe_rel(path: Path, base_dir: Path = BASE_DIR) -> str:
    path = path.resolve()
    try:
        return path.relative_to(base_dir.resolve()).as_posix()
    except ValueError:
        return "[external-path-redacted]"


def safe_string(value: Any) -> Any:
    if isinstance(value, dict):
        return {
            str(key): safe_string(item)
            for key, item in value.items()
            if not any(pattern.search(str(key)) for pattern in PRIVATE_KEY_PATTERNS)
        }
    if isinstance(value, list):
        return [safe_string(item) for item in value]
    if isinstance(value, str):
        cleaned = value
        for pattern in PRIVATE_KEY_PATTERNS:
            cleaned = pattern.sub("[redacted]", cleaned)
        return cleaned
    return value


def prompt_literal_matches(text: Any) -> list[str]:
    raw = str(text or "")
    return [
        label
        for label, pattern in zip(PROMPT_LITERAL_COMMAND_LABELS, PROMPT_LITERAL_COMMAND_PATTERNS, strict=True)
        if pattern.search(raw)
    ]


def pptx_prompt_literal_matches(pptx_path: Path) -> list[dict[str, Any]]:
    matches: list[dict[str, Any]] = []
    if not pptx_path.exists():
        return matches
    try:
        with zipfile.ZipFile(pptx_path) as archive:
            slide_names = sorted(
                name
                for name in archive.namelist()
                if re.match(r"ppt/slides/slide\d+\.xml$", name)
            )
            for slide_name in slide_names:
                slide_no_match = re.search(r"slide(\d+)\.xml$", slide_name)
                slide_no = int(slide_no_match.group(1)) if slide_no_match else None
                raw = archive.read(slide_name).decode("utf-8", errors="ignore")
                text_items: list[str] = []
                try:
                    xml_root = ET.fromstring(raw)
                    for node in xml_root.iter():
                        if node.tag.endswith("}t") and node.text:
                            text_items.append(node.text)
                except ET.ParseError:
                    text_items.append(raw)
                for text in text_items:
                    for match in prompt_literal_matches(text):
                        matches.append({"slide_no": slide_no, "field": "pptx_slide_xml", "phrase": match})
    except (OSError, zipfile.BadZipFile):
        return matches
    return matches


def sanitize_prompt_literal_text(value: Any, *, fallback: str = "Untitled Deck") -> str:
    text = str(safe_string(value) or "")
    for pattern in PROMPT_LITERAL_COMMAND_PATTERNS:
        text = pattern.sub(" ", text)
    text = re.sub(r"\b(PPT|deck|slides?)\b\s*$", " ", text, flags=re.IGNORECASE)
    text = re.sub(r"\s+", " ", text)
    text = re.sub(r"\s+([.,;:!?])", r"\1", text)
    text = text.strip(" \t\r\n-_/.,;:!?。،")
    text = re.sub(r"^(를|을|은|는|이|가|및|and)\s+", "", text).strip()
    text = re.sub(r"\s+(를|을)(?=[.。]|$)", "", text).strip()
    text = re.sub(r"\s+(를|을)\s*[.。]\s*", ". ", text).strip()
    text = re.split(r"(?<=[.。])\s+", text, maxsplit=1)[0].strip(" .。")
    lowered = text.lower()
    if all(token in text for token in ["식품", "신상품"]) and "런칭" in text:
        return "신제품 런칭 캠페인 전략"
    if "분기" in text and "내부" in text and "전략" in text and "보고" in text:
        return "분기 내부 전략 보고"
    if ("프리미엄" in text or "럭셔리" in text) and "향수" in text and ("런칭" in text or "launch" in lowered):
        return "프리미엄 향수 브랜드 런칭 전략"
    text = re.sub(r"\s+(를|을)$", "", text).strip()
    return text or fallback


def sanitize_visible_value(value: Any, *, fallback: str = "Untitled Deck") -> str:
    return sanitize_prompt_literal_text(value, fallback=fallback)


def sanitize_report_text(value: Any) -> Any:
    if isinstance(value, dict):
        return {key: sanitize_report_text(item) for key, item in value.items()}
    if isinstance(value, list):
        return [sanitize_report_text(item) for item in value]
    if not isinstance(value, str):
        return value
    text = value.replace(str(BASE_DIR), "[workspace]")
    text = re.sub(r"[A-Za-z]:\\(?:[^\\/:*?\"<>|\r\n]+\\)*[^\\/:*?\"<>|\r\n]*", "[local-path-redacted]", text)
    text = re.sub(r"/(?:Users|home)/[^\s\"']+", "[local-path-redacted]", text)
    return text


def hex_to_rgb(value: str) -> RGBColor:
    value = value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


class StrictModel(BaseModel):
    model_config = ConfigDict(extra="forbid")


class FlexibleModel(BaseModel):
    model_config = ConfigDict(extra="allow")


class Contract(StrictModel):
    name: str = "b44.design_guide_packet"
    version: str = "1.0"
    compatibility: str = "additive_to_b44_asset_handoff"


class GuideIdentity(StrictModel):
    guide_id: str
    guide_version: str
    status: str
    language: str
    project_name: str
    topic: str
    slide_count: int = Field(ge=1)
    created_at: str


class ProjectBrief(StrictModel):
    audience: list[str] = Field(min_length=1)
    tone: list[str] = Field(min_length=1)
    objective: str
    constraints: list[str] = Field(default_factory=list)


class ReferenceFile(StrictModel):
    file_name: str
    extension: str
    purpose: str
    public_safe_note: str | None = None


class PaletteSource(StrictModel):
    kind: str
    asset_ref: str | None = None


class PaletteColor(StrictModel):
    role: str
    name: str
    hex: str = Field(pattern=r"^#[0-9A-Fa-f]{6}$")
    rgb: list[int] = Field(min_length=3, max_length=3)
    usage: str
    accessibility_note: str | None = None
    contrast: dict[str, Any] = Field(default_factory=dict)


class Palette(StrictModel):
    palette_id: str
    palette_name: str
    source: PaletteSource
    colors: list[PaletteColor] = Field(min_length=1)


class FontPackageRef(StrictModel):
    family: str
    asset_ref: str | None = None
    license_action: str
    package_status: str


class TypographyRole(StrictModel):
    role: str
    font_stack: list[str] = Field(min_length=1)
    default_pt: float = Field(gt=0)
    min_pt: float = Field(gt=0)
    max_pt: float = Field(gt=0)
    weight: str
    line_limit: int = Field(ge=1)
    fallback: str

    @model_validator(mode="after")
    def validate_size_order(self) -> "TypographyRole":
        if self.min_pt > self.default_pt or self.default_pt > self.max_pt:
            raise ValueError(f"typography role {self.role} must satisfy min <= default <= max")
        return self


class Typography(StrictModel):
    font_package_refs: list[FontPackageRef] = Field(default_factory=list)
    roles: list[TypographyRole] = Field(min_length=1)


class SlideSize(StrictModel):
    size: str
    width_in: float = Field(gt=0)
    height_in: float = Field(gt=0)


class SafeArea(StrictModel):
    left_in: float
    top_in: float
    right_in: float
    bottom_in: float

    @model_validator(mode="after")
    def validate_bounds(self) -> "SafeArea":
        if self.left_in >= self.right_in or self.top_in >= self.bottom_in:
            raise ValueError("safe_area must have left < right and top < bottom")
        return self


class BackgroundPolicy(StrictModel):
    default_mode: str
    max_modes_per_deck: int = Field(ge=1)
    allowed_types: list[str] = Field(min_length=1)
    forbidden_types: list[str] = Field(default_factory=list)


class HeaderFooterSlot(StrictModel):
    mode: str = "optional"
    allowed_content: list[str] = Field(default_factory=list)


class HeaderFooterPolicy(StrictModel):
    header: HeaderFooterSlot
    footer: HeaderFooterSlot


class AssetSlot(StrictModel):
    slot_id: str
    slot_type: str
    required: bool = False
    allowed_sources: list[str] = Field(default_factory=list)
    preferred_asset_ref: str | None = None
    fallback: str = "Use native editable fallback."
    crop_or_mask_policy: str = "preserve_aspect_ratio"
    visible_text_allowed: bool = False


class SlidePlanItem(StrictModel):
    slide_no: int = Field(ge=1)
    layout_archetype: str
    content_brief: str
    visible_content_candidates: list[str] = Field(default_factory=list)
    asset_slots: list[AssetSlot] = Field(default_factory=list)
    qa_checks: list[str] = Field(default_factory=list)
    content_guidance: str | None = None
    implementation_notes: str | None = None
    renderer_metadata: dict[str, Any] = Field(default_factory=dict)


class SlidePlan(StrictModel):
    slides: list[SlidePlanItem] = Field(min_length=1)


class LayoutArchetype(FlexibleModel):
    id: str
    purpose: str | None = None
    description: str | None = None
    native_renderer: str | None = None
    native_assembly_hint: str | None = None
    required_asset_slots: list[str] = Field(default_factory=list)
    notes: str | None = None


class AssetSlotPolicy(FlexibleModel):
    resolution_priority: list[str] = Field(default_factory=list)
    checksum_required_for_package_assets: bool = True
    fallback_required_when_unavailable: bool = True

    def priority(self) -> list[str]:
        if self.resolution_priority:
            return self.resolution_priority
        extra = self.model_extra or {}
        value = extra.get("priority_order")
        return list(value) if isinstance(value, list) else []


class QARule(StrictModel):
    id: str
    severity: str = "error"
    check: str


class FallbackPolicy(FlexibleModel):
    native_shape_fallback_allowed: bool = True
    generated_image_requires_approval: bool = True
    report_every_fallback: bool = True


class ApprovedAssetReference(FlexibleModel):
    asset_ref: str
    slot_id: str | None = None
    asset_type: str | None = None
    asset_role: str | None = None
    manifest_checksum: str | None = None
    approved: bool = False
    local_path: str | None = None
    sha256: str | None = None
    file_size_bytes: int | None = None
    package_status: str = "metadata_only"


class PublicSafety(FlexibleModel):
    public_safe: bool = True
    forbidden_visible_terms: list[str] = Field(default_factory=list)
    forbidden_report_fields: list[str] = Field(default_factory=list)
    html_screenshot_used_in_pptx: bool = False


DEFAULT_FORBIDDEN_VISIBLE_TERMS = [
    "content_guidance",
    "implementation_notes",
    "policy_notes",
    "slot_id",
    "layout_recipe_id",
    "asset_uid",
    "handoff_signal",
]


VARIANT_STRATEGY_PROFILES = {
    "investor_open": {
        "id": "investor_open",
        "description": "Open investor-pitch composition with spacious title hierarchy and classic slide rhythm.",
        "palette_emphasis": ["background", "neutral", "main"],
        "typography_emphasis": ["hero_title", "slide_title", "body"],
    },
    "operator_dense": {
        "id": "operator_dense",
        "description": "Denser operator-style composition with compact hierarchy, stronger token surfaces, and alternate layout recipes.",
        "palette_emphasis": ["main", "support", "accent", "neutral"],
        "typography_emphasis": ["body", "caption", "metric"],
    },
}


def load_strategy_registry() -> dict[str, Any]:
    if not STRATEGY_REGISTRY_PATH.exists():
        return {
            "strategies": list(VARIANT_STRATEGY_PROFILES.values()),
            "aliases": {},
            "fallbacks": {"general_unknown_intent": {"variant_a": "investor_open", "variant_b": "operator_dense"}},
        }
    return json.loads(STRATEGY_REGISTRY_PATH.read_text(encoding="utf-8"))


def load_routed_strategy_pair(
    guide_path: str | Path,
    *,
    output_root: str | Path = DEFAULT_PROJECTS_DIR,
    project_id: str | None = None,
    routing_report_path: str | Path | None = None,
) -> tuple[str | None, str | None, dict[str, Any] | None]:
    candidates: list[Path] = []
    if routing_report_path:
        candidates.append(Path(routing_report_path))
    guide = Path(guide_path).resolve()
    candidates.extend(
        [
            guide.parent / "routing-report.json",
            guide.parent.parent / "routing-report.json",
        ]
    )
    if project_id:
        candidates.append(Path(output_root).resolve() / project_id / "routing-report.json")
    seen: set[Path] = set()
    for candidate in candidates:
        candidate = candidate.resolve()
        if candidate in seen:
            continue
        seen.add(candidate)
        if not candidate.exists():
            continue
        payload = json.loads(candidate.read_text(encoding="utf-8"))
        selected = payload.get("selected", {})
        strategy_a = selected.get("variant_a", {}).get("strategy_id")
        strategy_b = selected.get("variant_b", {}).get("strategy_id")
        if strategy_a and strategy_b:
            return resolve_strategy_id(strategy_a), resolve_strategy_id(strategy_b), safe_string(
                {
                    "routing_report_path": safe_rel(candidate),
                    "mapping_id": selected.get("mapping_id"),
                    "fallback_used": payload.get("fallback_used", False),
                    "confidence": payload.get("confidence"),
                    "strategy_pair": {"variant_a": strategy_a, "variant_b": strategy_b},
                }
            )
    return None, None, None


def resolve_strategy_id(strategy_id: str) -> str:
    registry = load_strategy_registry()
    return registry.get("aliases", {}).get(strategy_id, strategy_id)


def strategy_profile(strategy_id: str) -> dict[str, Any]:
    registry = load_strategy_registry()
    canonical = registry.get("aliases", {}).get(strategy_id, strategy_id)
    for item in registry.get("strategies", []):
        if item.get("strategy_id") == canonical:
            profile = dict(item)
            profile["id"] = canonical
            profile["description"] = profile.get("description") or profile.get("evidence_style") or canonical
            profile["typography_emphasis"] = profile.get("typography_role_bias", [])
            return profile
    fallback = VARIANT_STRATEGY_PROFILES.get(canonical) or VARIANT_STRATEGY_PROFILES["investor_open"]
    return dict(fallback)


def load_palette_registry() -> dict[str, Any]:
    if PALETTE_PRESET_REGISTRY_PATH.exists():
        return json.loads(PALETTE_PRESET_REGISTRY_PATH.read_text(encoding="utf-8"))
    return {
        "presets": [
            {
                "palette_id": "fallback_default",
                "palette_name": "Fallback Default",
                "policy": "restrained",
                "colors": {
                    "main": ["Calm Blue", "#2F5F98", "Safe default primary accent."],
                    "support": ["Gray Blue", "#6B7A90", "Safe default secondary accent."],
                    "accent": ["Muted Green", "#4D8B73", "Safe default positive accent."],
                    "background": ["Off White", "#F8F9FB", "Safe default background."],
                    "neutral": ["Ink", "#18202B", "Safe default text."],
                },
            }
        ]
    }


def hex_to_rgb_list(value: str) -> list[int]:
    value = value.strip().lstrip("#")
    return [int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)]


def palette_preset(palette_id: str) -> dict[str, Any]:
    registry = load_palette_registry()
    presets = {item.get("palette_id"): item for item in registry.get("presets", [])}
    return presets.get(palette_id) or presets.get("fallback_default") or next(iter(presets.values()))


def palette_from_preset(preset: dict[str, Any], *, source_type: str, evidence: list[str]) -> dict[str, Any]:
    colors = []
    for role in ["main", "support", "accent", "background", "neutral"]:
        name, hex_value, usage = preset["colors"][role]
        colors.append(
            {
                "role": role,
                "name": name,
                "hex": hex_value,
                "rgb": hex_to_rgb_list(hex_value),
                "usage": usage,
            }
        )
    return {
        "palette_id": preset["palette_id"],
        "palette_name": preset["palette_name"],
        "source": {
            "kind": "registry_palette" if source_type == "intent_matched_palette" else "project_defined",
            "asset_ref": f"palette-resolution:{source_type}:{preset['palette_id']}:{'|'.join(evidence[:4])}",
        },
        "colors": colors,
    }


def palette_color_from_candidate(role: str, candidate: Any, fallback_usage: str) -> dict[str, Any] | None:
    if isinstance(candidate, str):
        hex_value = candidate
        name = role.replace("_", " ").title()
        usage = fallback_usage
        contrast = {}
    elif isinstance(candidate, list) and candidate:
        name = str(candidate[0] if len(candidate) > 1 else role.replace("_", " ").title())
        hex_value = str(candidate[1] if len(candidate) > 1 else candidate[0])
        usage = str(candidate[2] if len(candidate) > 2 else fallback_usage)
        contrast = {}
    elif isinstance(candidate, dict):
        hex_value = str(candidate.get("hex") or candidate.get("hex_value") or candidate.get("value") or "")
        name = str(candidate.get("name") or candidate.get("color_name") or role.replace("_", " ").title())
        usage = str(candidate.get("usage") or candidate.get("recommended_usage") or fallback_usage)
        contrast = candidate.get("contrast") if isinstance(candidate.get("contrast"), dict) else {}
    else:
        return None
    if not re.match(r"^#[0-9A-Fa-f]{6}$", hex_value):
        return None
    return {
        "role": role,
        "name": name,
        "hex": hex_value,
        "rgb": hex_to_rgb_list(hex_value),
        "usage": usage,
        "contrast": contrast,
    }


def package_palette_selection(package_response: dict[str, Any], package_manifest_id: str | None) -> dict[str, Any] | None:
    selection = package_response.get("palette_selection")
    if not isinstance(selection, dict):
        selection = package_response.get("approved_palette")
    if not isinstance(selection, dict):
        manifest = package_response.get("package_manifest", {})
        if isinstance(manifest, dict):
            selection = manifest.get("palette_selection") or manifest.get("approved_palette")
    if not isinstance(selection, dict):
        return None
    if selection.get("palette_metadata_available") is False:
        return None
    if selection.get("source_type") and selection.get("source_type") != "asset_system_palette":
        return None

    approved_palette = selection.get("approved_palette") if isinstance(selection.get("approved_palette"), dict) else selection
    if not isinstance(approved_palette, dict):
        return None
    colors_source = approved_palette.get("colors") or approved_palette.get("role_colors") or approved_palette.get("palette_colors")
    if not isinstance(colors_source, dict):
        return None

    fallback = palette_preset("fallback_default")
    colors: list[dict[str, Any]] = []
    for role in ["main", "support", "accent", "background", "neutral"]:
        fallback_usage = fallback["colors"][role][2]
        candidate = colors_source.get(role)
        color = palette_color_from_candidate(role, candidate, fallback_usage)
        if color is None:
            return None
        colors.append(color)

    palette_id = str(approved_palette.get("palette_id") or selection.get("palette_id") or "asset-system-approved-palette")
    palette_name = str(approved_palette.get("palette_name") or approved_palette.get("name") or selection.get("palette_name") or "Asset-System Approved Palette")
    confidence = selection.get("selection_confidence")
    evidence = [f"package_manifest_id={package_manifest_id or 'unknown'}"]
    if confidence is not None:
        evidence.append(f"selection_confidence={confidence}")
    if selection.get("palette_policy"):
        evidence.append(f"palette_policy={selection.get('palette_policy')}")
    if selection.get("source_type"):
        evidence.append(f"declared_source_type={selection.get('source_type')}")
    return {
        "palette_id": palette_id,
        "palette_name": palette_name,
        "source": {
            "kind": "asset_system_palette",
            "asset_ref": f"palette-resolution:asset_system_palette:{palette_id}:{'|'.join(evidence)}",
        },
        "colors": colors,
    }


def should_apply_package_palette(raw: dict[str, Any]) -> bool:
    palette = raw.get("palette")
    if not isinstance(palette, dict):
        return True
    source = palette.get("source") if isinstance(palette.get("source"), dict) else {}
    asset_ref = str(source.get("asset_ref") or "")
    source_kind = str(source.get("kind") or "")
    return asset_ref.startswith("palette-resolution:") or source_kind in {"project_defined", "registry_palette"}


def palette_text_has_any(text: str, tokens: list[str]) -> bool:
    for token in tokens:
        token = token.lower()
        if re.search(r"[a-z0-9]", token):
            pattern = r"(?<![a-z0-9])" + re.escape(token).replace(r"\ ", r"\s+") + r"(?![a-z0-9])"
            if re.search(pattern, text, flags=re.IGNORECASE):
                return True
        elif token in text:
            return True
    return False


def resolve_intent_palette(
    intent_profile: dict[str, Any],
    routing_report: dict[str, Any],
    *,
    source_summary: dict[str, Any] | None = None,
) -> tuple[dict[str, Any], dict[str, Any]]:
    family = str(intent_profile.get("deck_family") or "unknown")
    subtype = str(intent_profile.get("sector_subtype") or "general")
    audience = str(intent_profile.get("audience") or "").lower()
    objective = str(intent_profile.get("objective") or "").lower()
    tone_values = intent_profile.get("tone", [])
    tone = " ".join(str(item).lower() for item in (tone_values if isinstance(tone_values, list) else [tone_values]))
    selected = routing_report.get("selected", {})
    strategies = " ".join(
        str(selected.get(variant, {}).get("strategy_id", "")).lower()
        for variant in ["variant_a", "variant_b"]
    )
    text = " ".join([family, subtype, audience, objective, tone, strategies]).lower()
    evidence: list[str] = [f"family={family}", f"subtype={subtype}"]
    palette_id = "fallback_default"
    source_type = "intent_matched_palette"

    if family == "public_institution_report" or palette_text_has_any(text, ["formal_memo", "public", "policy", "procurement", "grant"]):
        palette_id = "public_institution"
        evidence.append("public/formal signal")
    elif family in {"executive_report", "status_update", "sales_proposal"} or palette_text_has_any(text, ["executive", "finance", "b2b", "saas", "enterprise"]):
        palette_id = "neutral_business"
        evidence.append("business/executive signal")
    elif family in {"research_analysis"} or palette_text_has_any(text, ["market", "analysis", "benchmark", "ecosystem"]):
        palette_id = "market_analysis"
        evidence.append("analysis/market signal")
    elif family == "education_training" or palette_text_has_any(text, ["education", "training", "learning", "teach"]):
        palette_id = "education_clear"
        evidence.append("education signal")
    elif family in {"portfolio"} or palette_text_has_any(text, ["portfolio", "luxury", "editorial", "premium", "gallery"]):
        palette_id = "editorial_minimal"
        evidence.append("editorial/portfolio signal")
    elif palette_text_has_any(text, ["youth", "gen z", "teen", "social", "kpop", "campaign", "entertainment", "creator", "mobile"]):
        palette_id = "consumer_bright"
        evidence.append("youth/social/campaign signal")
    elif family in {"product_introduction", "event_travel", "media_entertainment", "marketing_campaign"}:
        palette_id = "product_neutral"
        evidence.append("product or experience signal")
    else:
        source_type = "fallback_default_palette"
        evidence.append("no higher-priority palette signal")

    preset = palette_preset(palette_id)
    resolution = {
        "source_type": source_type,
        "selected_palette_id": preset["palette_id"],
        "selected_palette_name": preset["palette_name"],
        "policy": preset.get("policy", "balanced"),
        "evidence": evidence,
        "fallback_reason": "No guide, asset-system, brand, or intent palette match was available." if source_type == "fallback_default_palette" else None,
    }
    return palette_from_preset(preset, source_type=source_type, evidence=evidence), resolution


def infer_palette_resolution(packet: GuidePacket) -> dict[str, Any]:
    source = packet.palette.source
    asset_ref = source.asset_ref or ""
    source_type = "guide_packet_palette"
    fallback_reason = None
    if asset_ref.startswith("palette-resolution:"):
        parts = asset_ref.split(":")
        if len(parts) >= 3:
            source_type = parts[1]
        if source_type == "fallback_default_palette":
            fallback_reason = "No guide, asset-system, brand, or intent palette match was available."
    elif packet.palette.palette_id in {"assistant-default", "fallback_default"}:
        source_type = "fallback_default_palette"
        fallback_reason = "No higher-priority palette source was recorded."
    elif source.kind == "registry_palette":
        source_type = "intent_matched_palette"
    return {
        "source_type": source_type,
        "selected_palette_id": packet.palette.palette_id,
        "selected_palette_name": packet.palette.palette_name,
        "source_kind": source.kind,
        "asset_ref": asset_ref,
        "fallback_reason": fallback_reason,
    }


OPERATOR_LAYOUT_RECIPES = {
    "cover": "split_color_panel_left_mockup",
    "big_thesis": "sidebar_thesis_with_operator_note",
    "three_cards": "stacked_decision_rows",
    "behavior_shift": "vertical_before_after_lanes",
    "product_mockup": "left_mockup_right_feature_bands",
    "process_flow": "vertical_process_ladder",
    "three_use_cases": "stacked_use_case_rows",
    "market_sizing": "overlapping_market_circles_with_sidebar",
    "traction_metrics": "horizontal_metric_bars",
    "business_model": "vertical_revenue_lanes",
    "positioning_matrix": "quadrant_tile_matrix",
    "flywheel": "hub_and_spoke_loop",
    "roadmap": "stair_step_roadmap",
    "financial_plan": "horizontal_financial_bars",
    "ask": "top_band_ask_with_use_of_funds_cards",
}


INVESTOR_LAYOUT_RECIPES = {
    "cover": "open_brand_hero_with_right_mockup",
    "big_thesis": "centered_thesis_statement",
    "three_cards": "equal_three_card_grid",
    "behavior_shift": "side_by_side_shift_arrow",
    "product_mockup": "right_mockup_left_narrative",
    "process_flow": "horizontal_step_flow",
    "three_use_cases": "numbered_three_card_grid",
    "market_sizing": "layered_horizontal_market_bars",
    "traction_metrics": "three_kpi_cards",
    "business_model": "two_by_two_business_grid",
    "positioning_matrix": "axis_positioning_matrix",
    "flywheel": "four_node_loop",
    "roadmap": "classic_horizontal_timeline",
    "financial_plan": "vertical_financial_bars",
    "ask": "centered_fundraising_ask",
}


def recipe_slug(value: Any) -> str:
    text = re.sub(r"[^a-zA-Z0-9_.-]+", "_", str(value or "")).strip("_").lower()
    return text or "recipe"


def strategy_layout_recipe(
    slide: SlidePlanItem,
    profile: dict[str, Any],
    *,
    use_dense_recipe: bool,
    composition_family: str | None = None,
) -> str:
    base_recipes = OPERATOR_LAYOUT_RECIPES if use_dense_recipe else INVESTOR_LAYOUT_RECIPES
    base_recipe = base_recipes.get(slide.layout_archetype, f"{'operator' if use_dense_recipe else 'investor'}_{slide.layout_archetype}")
    preferences = [
        recipe_slug(item)
        for item in profile.get("layout_recipe_preferences", [])
        if str(item or "").strip()
    ]
    family_preference = preferred_recipe_family(str(composition_family or ""), slide.slide_no)
    if family_preference:
        preferences = [recipe_slug(family_preference), *[item for item in preferences if item != recipe_slug(family_preference)]]
    if not preferences:
        return base_recipe
    strategy_id = recipe_slug(profile.get("id") or profile.get("strategy_id") or "strategy")
    preference = preferences[(slide.slide_no - 1) % len(preferences)]
    return f"{strategy_id}__{preference}__{base_recipe}"


def recipe_family_from_recipe(recipe: str) -> str:
    parts = [part for part in str(recipe or "").split("__") if part]
    if len(parts) >= 3:
        return parts[1]
    return parts[0] if parts else "default"


def recipe_tokens_from_recipe(recipe: str) -> set[str]:
    family = recipe_family_from_recipe(recipe)
    return {token for token in re.split(r"[^a-zA-Z0-9]+", family.lower()) if token}


class GuidePacket(StrictModel):
    contract: Contract
    guide_identity: GuideIdentity
    project_brief: ProjectBrief
    reference_files: list[ReferenceFile] = Field(min_length=1)
    palette: Palette
    typography: Typography
    slide_size: SlideSize
    safe_area: SafeArea
    background_policy: BackgroundPolicy
    header_footer_policy: HeaderFooterPolicy
    slide_plan: SlidePlan
    layout_archetypes: list[LayoutArchetype] = Field(min_length=1)
    asset_slot_policy: AssetSlotPolicy
    qa_rules: list[QARule] = Field(min_length=1)
    fallback_policy: FallbackPolicy
    approved_asset_references: list[ApprovedAssetReference] = Field(default_factory=list)
    public_safety: PublicSafety

    @model_validator(mode="after")
    def validate_slide_contract(self) -> "GuidePacket":
        slides = self.slide_plan.slides
        if self.guide_identity.slide_count != len(slides):
            raise ValueError(
                "guide_identity.slide_count must match slide_plan.slides length "
                f"({self.guide_identity.slide_count} != {len(slides)})"
            )
        known = {item.id for item in self.layout_archetypes}
        unknown = sorted({slide.layout_archetype for slide in slides} - known)
        if unknown:
            raise ValueError(f"unknown slide_plan layout_archetype values: {unknown}")
        return self


class BuildContext:
    def __init__(
        self,
        packet: GuidePacket,
        project_dir: Path,
        mode: str,
        deck_plan: dict[str, Any],
        variant: str | None = None,
    ) -> None:
        self.packet = packet
        self.project_dir = project_dir
        self.mode = mode
        self.variant = variant
        self.deck_plan = deck_plan
        self.variant_strategy = deck_plan.get("variant_strategy", {}).get("id", "investor_open")
        self.composition_profile = deck_plan.get("composition_profile", {})
        self.primary_family = str(deck_plan.get("primary_family") or self.composition_profile.get("primary_family") or "product_launch")
        self.composition_family = str(deck_plan.get("composition_family") or self.composition_profile.get("composition_family") or self.primary_family)
        self.layout_strategy_by_slide = {
            int(slide["slide_no"]): slide.get("renderer_metadata", {}).get("layout_strategy", {})
            for slide in deck_plan.get("slides", [])
        }
        self.palette = {item.role: item.hex for item in packet.palette.colors}
        self.typography = {item.role: item for item in packet.typography.roles}
        self.used_palette_roles: set[str] = set()
        self.used_typography_roles: set[str] = set()
        self.omitted_header_footer: list[dict[str, Any]] = []
        self.asset_events: list[dict[str, Any]] = []
        self.asset_validation_failures: list[dict[str, Any]] = []
        self.native_renderer_events: list[dict[str, Any]] = []
        self.text_boxes: list[dict[str, Any]] = []
        self.visual_zones: list[dict[str, Any]] = []
        self.color_surfaces: list[dict[str, Any]] = []
        self.current_slide_no: int | None = None

    def color(self, role: str, fallback: str = "neutral") -> RGBColor:
        key = role if role in self.palette else fallback
        self.used_palette_roles.add(key)
        return hex_to_rgb(self.palette[key])

    def font(self, role: str = "body") -> tuple[str, float, bool]:
        item = self.typography.get(role) or self.typography.get("body") or next(iter(self.typography.values()))
        self.used_typography_roles.add(item.role)
        return item.fallback, item.default_pt, item.weight.lower() in {"bold", "700", "semibold"}

    def layout_strategy(self, item: SlidePlanItem) -> dict[str, Any]:
        return self.layout_strategy_by_slide.get(item.slide_no, {})

    def layout_recipe(self, item: SlidePlanItem) -> str:
        return str(self.layout_strategy(item).get("layout_recipe", ""))

    def recipe_family(self, item: SlidePlanItem) -> str:
        return recipe_family_from_recipe(self.layout_recipe(item))

    def recipe_tokens(self, item: SlidePlanItem) -> set[str]:
        return recipe_tokens_from_recipe(self.layout_recipe(item))

    def is_operator_layout(self, item: SlidePlanItem) -> bool:
        strategy = self.layout_strategy(item)
        return strategy.get("renderer_variant") in {"dense", "operator", "dashboard", "workbook"}

    def uses_family(self, *families: str) -> bool:
        return self.primary_family in set(families) or self.composition_family in set(families)

    def add_visual_zone(
        self,
        slide_no: int | None,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        zone_type: str,
        source: str,
    ) -> None:
        if slide_no is None:
            return
        self.visual_zones.append(
            {
                "slide_no": slide_no,
                "zone_type": zone_type,
                "source": source,
                "bounds_in": {"x": x, "y": y, "w": w, "h": h},
            }
        )

    def add_color_surface(self, role: str, x: float, y: float, w: float, h: float, *, kind: str) -> None:
        if self.current_slide_no is None:
            return
        self.color_surfaces.append(
            {
                "slide_no": self.current_slide_no,
                "role": role,
                "kind": kind,
                "area_in2": round(max(0.0, w) * max(0.0, h), 4),
                "bounds_in": {"x": x, "y": y, "w": w, "h": h},
            }
        )


def locate_guide_packet(path_value: str | Path) -> Path:
    path = Path(path_value).resolve()
    if path.is_dir():
        candidate = path / "guide-data.public.json"
    else:
        candidate = path
    if not candidate.exists():
        raise FileNotFoundError(f"guide-data.public.json not found at {safe_rel(candidate)}")
    return candidate


def load_guide_packet(path_value: str | Path) -> tuple[GuidePacket, Path, dict[str, Any]]:
    if not GUIDE_SCHEMA_PATH.exists():
        raise FileNotFoundError(f"Guide schema is missing: {safe_rel(GUIDE_SCHEMA_PATH)}")
    path = locate_guide_packet(path_value)
    raw = json.loads(path.read_text(encoding="utf-8"))
    schema = json.loads(GUIDE_SCHEMA_PATH.read_text(encoding="utf-8"))
    jsonschema.Draft202012Validator(schema).validate(raw)
    packet = GuidePacket.model_validate(raw)
    return packet, path, raw


def export_guide_packet_schema(path: Path = GUIDE_SCHEMA_PATH) -> Path:
    if GUIDE_SCHEMA_PATH.exists():
        path.parent.mkdir(parents=True, exist_ok=True)
        if path.resolve() != GUIDE_SCHEMA_PATH.resolve():
            shutil.copy2(GUIDE_SCHEMA_PATH, path)
        return path
    schema = GuidePacket.model_json_schema()
    schema["$id"] = "https://local.ppt-test/ppt-maker-design-guide-packet.schema.json"
    schema["title"] = "PPT Maker Design Guide Packet"
    json_dump(path, schema)
    return path


def normalized_packet_objects(packet: GuidePacket) -> dict[str, Any]:
    data = packet.model_dump(mode="json")
    keys = [
        "guide_identity",
        "project_brief",
        "palette",
        "typography",
        "slide_size",
        "safe_area",
        "background_policy",
        "header_footer_policy",
        "slide_plan",
        "layout_archetypes",
        "asset_slot_policy",
        "qa_rules",
        "fallback_policy",
        "approved_asset_references",
        "public_safety",
    ]
    return {key: data[key] for key in keys}


def visible_content(slide: SlidePlanItem) -> dict[str, Any]:
    candidates = [sanitize_visible_value(item, fallback=f"Slide {slide.slide_no}") for item in slide.visible_content_candidates]
    title = candidates[0] if candidates else sanitize_visible_value(slide.content_brief, fallback=f"Slide {slide.slide_no}")
    subtitle = candidates[1] if len(candidates) > 1 else sanitize_visible_value(slide.content_brief, fallback=title)
    return {
        "title": title,
        "subtitle": subtitle,
        "candidates": candidates,
        "brief": sanitize_visible_value(slide.content_brief, fallback=subtitle),
    }


def composition_profile_for_family(primary_family: str, variant_strategy: str | None = None) -> dict[str, Any]:
    profile = dict(COMPOSITION_FAMILY_PROFILES.get(primary_family, COMPOSITION_FAMILY_PROFILES["product_launch"]))
    profile["primary_family"] = primary_family
    profile["composition_family"] = primary_family
    profile["public_family"] = primary_family
    profile["variant_strategy_id"] = variant_strategy
    return profile


def infer_public_composition_family(packet: GuidePacket, variant_strategy: str | None = None) -> str:
    strategy = str(variant_strategy or "").lower()
    text = " ".join(
        [
            str(packet.guide_identity.project_name),
            str(packet.guide_identity.topic),
            " ".join(packet.project_brief.audience),
            " ".join(packet.project_brief.tone),
            str(packet.project_brief.objective),
            " ".join(packet.project_brief.constraints),
            strategy,
        ]
    ).lower()
    if "lumaloop" in text:
        return "product_launch"
    if "luxury_editorial" in strategy or any(token in text for token in ["luxury", "럭셔리", "premium", "프리미엄", "향수", "watch"]):
        return "luxury_editorial"
    if any(token in text for token in ["food", "beverage", "food_beverage", "식품", "신상품", "유통", "프로모션", "retail"]) and any(token in text for token in ["launch", "런칭", "홍보", "promotion", "campaign", "product_introduction"]):
        return "food_product_launch"
    if any(token in text for token in ["internal", "내부", "분기", "strategy report", "전략 보고", "경영진", "executive_report", "status_update"]):
        return "internal_strategy_report"
    if "screen_product_tour" in strategy:
        return "screen_product_tour"
    if "launch_campaign" in strategy or "marketing_campaign" in text:
        return "marketing_campaign"
    if "public_institution_report" in text or "formal_memo" in strategy:
        return "public_institution_report"
    if "ir_pitch" in text or "investor" in strategy:
        return "investor_pitch"
    if any(token in text for token in ["sales", "retail", "go to market", "gtm"]):
        return "retail_go_to_market"
    if any(token in text for token in ["executive", "board", "brief", "report"]):
        return "executive_brief"
    return "product_launch"


def preferred_recipe_family(primary_family: str, slide_no: int) -> str | None:
    families = COMPOSITION_FAMILY_PROFILES.get(primary_family, {}).get("recipe_families", [])
    if not families:
        return None
    return str(families[(slide_no - 1) % len(families)])


def layout_strategy_for_slide(slide: SlidePlanItem, variant_strategy: str, composition_profile: dict[str, Any] | None = None) -> dict[str, Any]:
    profile = strategy_profile(variant_strategy)
    strategy_id = profile.get("id") or profile.get("strategy_id") or variant_strategy
    renderer_variant = profile.get("renderer_variant", "open")
    use_dense_recipe = renderer_variant in {"dense", "operator", "dashboard", "workbook"}
    primary_family = str((composition_profile or {}).get("primary_family") or "product_launch")
    recipe = strategy_layout_recipe(slide, profile, use_dense_recipe=use_dense_recipe, composition_family=primary_family)
    preferred_public_archetypes = list((composition_profile or {}).get("preferred_archetypes", []))
    public_archetype = preferred_public_archetypes[(slide.slide_no - 1) % len(preferred_public_archetypes)] if preferred_public_archetypes else slide.layout_archetype
    family_payload = {
        "primary_family": primary_family,
        "composition_family": primary_family,
        "public_archetype": public_archetype,
        "preferred_public_archetypes": preferred_public_archetypes,
        "major_visual_asset_treatment": (composition_profile or {}).get("major_visual_asset_treatment"),
        "avoid": (composition_profile or {}).get("avoid", []),
    }
    if use_dense_recipe:
        return {
            "strategy_id": strategy_id,
            "layout_recipe": recipe,
            "renderer_variant": renderer_variant,
            "content_emphasis": profile.get("slide_rhythm", "evidence_forward"),
            "evidence_treatment": profile.get("evidence_style", "dense_evidence"),
            "visual_asset_role": profile.get("visual_placement", "alternate_operator_layout"),
            "palette_emphasis": profile.get("palette_emphasis", ["main", "support", "accent"]),
            "typography_role_bias": profile.get("typography_role_bias", ["body", "caption", "metric"]),
            "chart_or_table_style": profile.get("chart_or_table_style", "lane_charts_and_bars"),
            "density_budget": profile.get("density", "dense"),
            "title_hierarchy": profile.get("title_hierarchy", "compact_label_plus_title"),
            "visual_placement": profile.get("visual_placement", "alternate_operator_layout"),
            **family_payload,
        }
    return {
        "strategy_id": strategy_id,
        "layout_recipe": recipe,
        "renderer_variant": renderer_variant,
        "content_emphasis": profile.get("slide_rhythm", "story_forward"),
        "evidence_treatment": profile.get("evidence_style", "selective_proof_points"),
        "visual_asset_role": profile.get("visual_placement", "classic_pitch_layout"),
        "palette_emphasis": profile.get("palette_emphasis", ["background", "neutral", "main"]),
        "typography_role_bias": profile.get("typography_role_bias", ["hero_title", "slide_title", "body"]),
        "chart_or_table_style": profile.get("chart_or_table_style", "simple_growth_shapes"),
        "density_budget": profile.get("density", "open"),
        "title_hierarchy": profile.get("title_hierarchy", "large_investor_title"),
        "visual_placement": profile.get("visual_placement", "classic_pitch_layout"),
        **family_payload,
    }


def generate_deck_plan(packet: GuidePacket, source_mode: str, variant_strategy: str = "investor_open") -> dict[str, Any]:
    profile = strategy_profile(resolve_strategy_id(variant_strategy))
    strategy_id = profile.get("id") or profile.get("strategy_id") or variant_strategy
    composition_profile = composition_profile_for_family(infer_public_composition_family(packet, strategy_id), strategy_id)
    return {
        "contract": "ppt-maker.deck-plan.v1",
        "source_mode": source_mode,
        "variant_strategy": profile,
        "composition_profile": composition_profile,
        "primary_family": composition_profile["primary_family"],
        "composition_family": composition_profile["composition_family"],
        "guide_id": packet.guide_identity.guide_id,
        "project_name": sanitize_visible_value(packet.guide_identity.project_name, fallback="Untitled Deck"),
        "slide_count": packet.guide_identity.slide_count,
        "slides": [
            slide_deck_plan_item(slide, strategy_id, composition_profile)
            for slide in packet.slide_plan.slides
        ],
    }


def slide_deck_plan_item(slide: SlidePlanItem, strategy_id: str, composition_profile: dict[str, Any] | None = None) -> dict[str, Any]:
    layout_strategy = layout_strategy_for_slide(slide, strategy_id, composition_profile)
    return {
        "slide_no": slide.slide_no,
        "layout_archetype": slide.layout_archetype,
        "public_archetype": layout_strategy.get("public_archetype"),
        "primary_family": layout_strategy.get("primary_family"),
        "composition_family": layout_strategy.get("composition_family"),
        "strategy_id": layout_strategy["strategy_id"],
        "layout_recipe": layout_strategy["layout_recipe"],
        "content_emphasis": layout_strategy["content_emphasis"],
        "evidence_treatment": layout_strategy["evidence_treatment"],
        "visual_asset_role": layout_strategy["visual_asset_role"],
        "palette_emphasis": layout_strategy["palette_emphasis"],
        "typography_role_bias": layout_strategy["typography_role_bias"],
        "chart_or_table_style": layout_strategy["chart_or_table_style"],
        "density_budget": layout_strategy["density_budget"],
        "visible_content": visible_content(slide),
        "asset_slot_ids": [slot.slot_id for slot in slide.asset_slots],
        "qa_checks": slide.qa_checks,
        "renderer_metadata": {
            "layout_strategy": layout_strategy,
            "source_renderer_metadata": safe_string(slide.renderer_metadata),
        },
        "guidance": {
            "content_brief": safe_string(slide.content_brief),
            "content_guidance": safe_string(slide.content_guidance),
            "implementation_notes": safe_string(slide.implementation_notes),
        },
    }


def generate_renderer_contract(packet: GuidePacket, source_mode: str, deck_plan: dict[str, Any]) -> dict[str, Any]:
    archetype_map = {item.id: item.model_dump(mode="json") for item in packet.layout_archetypes}
    asset_requirements: list[dict[str, Any]] = []
    for slide in packet.slide_plan.slides:
        for slot in slide.asset_slots:
            asset_requirements.append(
                {
                    "slide_no": slide.slide_no,
                    "layout_archetype": slide.layout_archetype,
                    "slot_id": slot.slot_id,
                    "slot_type": slot.slot_type,
                    "required": slot.required,
                    "crop_or_mask_policy": slot.crop_or_mask_policy,
                    "fallback": slot.fallback,
                    "allowed_sources": slot.allowed_sources,
                }
            )
    return {
        "contract": "ppt-maker.renderer-contract.v1",
        "source_mode": source_mode,
        "variant_strategy": deck_plan.get("variant_strategy", {}),
        "composition_profile": deck_plan.get("composition_profile", {}),
        "primary_family": deck_plan.get("primary_family"),
        "composition_family": deck_plan.get("composition_family"),
        "strategy_contract": {
            "strategy_id": deck_plan.get("variant_strategy", {}).get("id")
            or deck_plan.get("variant_strategy", {}).get("strategy_id"),
            "density": deck_plan.get("variant_strategy", {}).get("density"),
            "title_hierarchy": deck_plan.get("variant_strategy", {}).get("title_hierarchy"),
            "slide_rhythm": deck_plan.get("variant_strategy", {}).get("slide_rhythm"),
            "visual_placement": deck_plan.get("variant_strategy", {}).get("visual_placement"),
            "evidence_style": deck_plan.get("variant_strategy", {}).get("evidence_style"),
            "palette_emphasis": deck_plan.get("variant_strategy", {}).get("palette_emphasis", []),
            "typography_role_bias": deck_plan.get("variant_strategy", {}).get("typography_role_bias", []),
            "chart_or_table_style": deck_plan.get("variant_strategy", {}).get("chart_or_table_style"),
            "primary_family": deck_plan.get("primary_family"),
            "composition_family": deck_plan.get("composition_family"),
        },
        "slide_size": packet.slide_size.model_dump(mode="json"),
        "safe_area": packet.safe_area.model_dump(mode="json"),
        "style_tokens": {
            "palette": {
                color.role: {"hex": color.hex, "name": color.name, "usage": color.usage}
                for color in packet.palette.colors
            },
            "typography": {
                role.role: {
                    "font_stack": role.font_stack,
                    "fallback": role.fallback,
                    "default_pt": role.default_pt,
                    "min_pt": role.min_pt,
                    "max_pt": role.max_pt,
                    "weight": role.weight,
                    "line_limit": role.line_limit,
                }
                for role in packet.typography.roles
            },
        },
        "layout_archetype_map": archetype_map,
        "layout_strategy_map": {
            str(slide["slide_no"]): slide.get("renderer_metadata", {}).get("layout_strategy", {})
            for slide in deck_plan.get("slides", [])
        },
        "asset_slot_requirements": asset_requirements,
        "policies": {
            "background_policy": packet.background_policy.model_dump(mode="json"),
            "header_footer_policy": packet.header_footer_policy.model_dump(mode="json"),
            "asset_slot_policy": packet.asset_slot_policy.model_dump(mode="json"),
            "fallback_policy": packet.fallback_policy.model_dump(mode="json"),
            "html_screenshot_used_in_pptx": False,
        },
    }


def text_units(text: str) -> float:
    units = 0.0
    for char in str(text or ""):
        if char in "\r\n":
            continue
        if "\uac00" <= char <= "\ud7af" or "\u3040" <= char <= "\u30ff" or "\u4e00" <= char <= "\u9fff":
            units += 1.05
        elif char.isspace():
            units += 0.38
        else:
            units += 0.55
    return units


def estimated_text_bounds(text: str, x: float, y: float, w: float, h: float, font_pt: float) -> dict[str, float]:
    usable_w = max(0.3, w - 0.08)
    chars_per_line = max(4.0, usable_w * 72 / max(font_pt, 1) * 1.72)
    explicit_lines = str(text or "").split("\n") or [""]
    line_count = 0
    for line in explicit_lines:
        line_count += max(1, int((text_units(line) + chars_per_line - 0.01) // chars_per_line))
    estimated_h = max(h, (font_pt / 72.0) * 1.18 * line_count + 0.08)
    return {"x": x, "y": y, "w": w, "h": estimated_h}


def add_text(
    slide,
    ctx: BuildContext,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    role: str = "body",
    color_role: str = "neutral",
    align: PP_ALIGN | None = None,
    vertical: MSO_ANCHOR = MSO_ANCHOR.TOP,
    qa_text_role: str = "content",
) -> None:
    font_name, size, bold = ctx.font(role)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(3)
    tf.margin_right = Pt(3)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = vertical
    lines = str(text or "").split("\n") or [""]
    for idx, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.alignment = align if align is not None else PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = ctx.color(color_role)
    if ctx.current_slide_no is not None:
        actual_bounds = {"x": x, "y": y, "w": w, "h": h}
        ctx.text_boxes.append(
            {
                "slide_no": ctx.current_slide_no,
                "text_role": qa_text_role,
                "typography_role": role,
                "text_excerpt": safe_string(str(text or "")[:80]),
                "bounds_in": actual_bounds,
                "estimated_text_bounds_in": estimated_text_bounds(str(text or ""), x, y, w, h, size),
            }
        )


def add_rect(slide, ctx: BuildContext, x: float, y: float, w: float, h: float, fill: str, line: str | None = None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ctx.color(fill)
    ctx.add_color_surface(fill, x, y, w, h, kind="rect")
    if line:
        shape.line.color.rgb = ctx.color(line)
        shape.line.width = Pt(1.0)
    else:
        shape.line.fill.background()
    return shape


def add_oval(slide, ctx: BuildContext, x: float, y: float, w: float, h: float, fill: str, line: str | None = None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ctx.color(fill)
    ctx.add_color_surface(fill, x, y, w, h, kind="oval")
    if line:
        shape.line.color.rgb = ctx.color(line)
        shape.line.width = Pt(1.0)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, ctx: BuildContext, x1: float, y1: float, x2: float, y2: float, color: str = "main") -> None:
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = ctx.color(color)
    line.line.width = Pt(2)


def recipe_composition_profile(ctx: BuildContext, item: SlidePlanItem) -> dict[str, Any]:
    family = ctx.recipe_family(item).lower()
    tokens = ctx.recipe_tokens(item)
    role_sets = [
        ("main", "support", "accent"),
        ("support", "accent", "main"),
        ("accent", "main", "support"),
        ("neutral", "main", "accent"),
    ]
    role_index = int(hashlib.sha1(family.encode("utf-8")).hexdigest()[:2], 16) % len(role_sets)
    primary, secondary, tertiary = role_sets[role_index]
    style = "hero_right"
    if {"editorial", "premium", "material"} & tokens:
        style = "minimal_editorial"
    elif {"sensory", "detail", "showcase"} & tokens:
        style = "sensory_showcase"
    elif {"portfolio", "gallery", "feature"} & tokens:
        style = "gallery_grid"
    elif {"journey", "itinerary", "route", "timeline", "tour", "learning", "service"} & tokens:
        style = "journey_path"
    elif {"matrix", "grid", "table", "comparison", "checklist", "compliance", "benchmark", "spec"} & tokens:
        style = "matrix_grid"
    elif {"blueprint", "ownership", "ecosystem", "chain", "map"} & tokens:
        style = "system_map"
    elif {"campaign", "calendar", "channel"} & tokens:
        style = "campaign_grid"
    elif {"memo", "sections", "workbook", "steps"} & tokens:
        style = "memo_sections"
    elif {"proof", "roi", "case", "benefit", "option", "decision", "risk", "action"} & tokens:
        style = "proof_cards"
    elif {"story", "character", "world"} & tokens:
        style = "story_stage"
    return {
        "family": family,
        "tokens": sorted(tokens),
        "style": style,
        "roles": (primary, secondary, tertiary),
    }


RIGHT_VISUAL_COMPATIBLE_ARCHETYPES = {"cover", "big_thesis", "behavior_shift", "product_mockup", "ask"}


def right_visual_reserved_x(ctx: BuildContext, item: SlidePlanItem) -> float:
    profile = recipe_composition_profile(ctx, item)
    w = ctx.packet.slide_size.width_in
    style = profile["style"]
    if style == "sensory_showcase":
        return w - 4.25
    if style == "gallery_grid" and item.layout_archetype in RIGHT_VISUAL_COMPATIBLE_ARCHETYPES:
        return 6.35
    if style == "matrix_grid" and item.layout_archetype in RIGHT_VISUAL_COMPATIBLE_ARCHETYPES:
        return 6.1
    if style == "system_map" and item.layout_archetype in RIGHT_VISUAL_COMPATIBLE_ARCHETYPES:
        return 5.85
    if style == "campaign_grid" and item.layout_archetype in RIGHT_VISUAL_COMPATIBLE_ARCHETYPES:
        return 5.05
    if style in {"hero_right", "proof_cards"}:
        return w - 3.8
    return ctx.packet.safe_area.right_in


def text_width_before_right_visual(ctx: BuildContext, item: SlidePlanItem, x: float, padding: float = 0.25) -> float:
    limit = right_visual_reserved_x(ctx, item) - padding
    return max(1.0, min(ctx.packet.safe_area.right_in, limit) - x)


def apply_recipe_canvas(slide, ctx: BuildContext, item: SlidePlanItem) -> dict[str, Any]:
    profile = recipe_composition_profile(ctx, item)
    primary, secondary, tertiary = profile["roles"]
    style = profile["style"]
    w = ctx.packet.slide_size.width_in
    h = ctx.packet.slide_size.height_in
    area = ctx.packet.safe_area

    if ctx.uses_family("internal_strategy_report"):
        if style in {"memo_sections", "matrix_grid", "proof_cards", "system_map", "journey_path"}:
            add_rect(slide, ctx, 0.0, 0.0, 0.18, h, primary)
            add_line(slide, ctx, area.left_in, h - 0.82, area.right_in, h - 0.82, secondary)
            return profile
    luxury_full_width_text_layouts = {"big_thesis", "three_cards", "three_use_cases", "traction_metrics", "business_model"}
    if (
        ctx.uses_family("luxury_editorial")
        and style in {"sensory_showcase", "minimal_editorial", "hero_right"}
        and item.layout_archetype != "cover"
        and item.layout_archetype not in luxury_full_width_text_layouts
    ):
        add_line(slide, ctx, area.left_in, 1.05, w - 1.2, 1.05, tertiary)
        add_rect(slide, ctx, w - 2.65, 1.55, 1.15, 1.15, "accent")
        add_rect(slide, ctx, w - 2.25, 3.15, 0.65, 2.25, "background", "neutral")
        add_line(slide, ctx, w - 3.05, 1.25, w - 3.05, h - 1.4, secondary)
        ctx.add_visual_zone(item.slide_no, w - 2.35, 1.45, 1.45, 4.1, zone_type="product_or_material_stage", source="luxury_editorial_detail")
        return profile

    if ctx.is_operator_layout(item) and style in {"gallery_grid", "matrix_grid", "system_map", "campaign_grid"}:
        if style == "matrix_grid":
            add_rect(slide, ctx, 0.0, 0.0, w, 0.24, primary)
            add_rect(slide, ctx, 0.0, 0.0, 0.18, h, secondary)
        elif style == "campaign_grid":
            add_rect(slide, ctx, 0.0, 0.0, w, 0.28, tertiary)
            add_rect(slide, ctx, w - 0.38, 0.0, 0.38, h, primary)
            add_rect(slide, ctx, 0.0, h - 0.2, w, 0.2, secondary)
        elif style == "gallery_grid":
            add_rect(slide, ctx, 0.0, h - 0.34, w, 0.34, secondary)
        else:
            add_rect(slide, ctx, 0.0, 0.0, 0.32, h, primary)
            add_rect(slide, ctx, 0.0, h - 0.3, w, 0.3, secondary)
        return profile

    if ctx.uses_family("food_product_launch") and style in {"hero_right", "proof_cards", "minimal_editorial"}:
        add_rect(slide, ctx, 0.0, h - 0.28, w, 0.28, secondary)
        add_oval(slide, ctx, w - 2.85, 1.05, 1.15, 1.15, tertiary)
        add_oval(slide, ctx, w - 1.72, 1.76, 0.62, 0.62, primary)
        add_rect(slide, ctx, w - 3.35, 5.45, 2.35, 0.3, primary)
        add_rect(slide, ctx, w - 2.85, 5.92, 1.8, 0.18, tertiary)
        return profile

    if style == "minimal_editorial":
        add_rect(slide, ctx, 0.0, 0.0, w, 0.18, primary)
        add_rect(slide, ctx, w - 1.15, 0.7, 0.18, h - 1.35, secondary)
        add_line(slide, ctx, area.left_in, h - 1.05, w - 1.25, h - 1.05, tertiary)
    elif style == "sensory_showcase":
        if ctx.uses_family("luxury_editorial"):
            add_line(slide, ctx, w - 4.0, 1.1, w - 4.0, h - 1.2, secondary)
            add_oval(slide, ctx, w - 2.95, 1.15, 1.05, 1.05, tertiary)
            add_oval(slide, ctx, w - 1.75, h - 1.72, 0.82, 0.82, secondary)
            add_rect(slide, ctx, w - 3.35, 3.05, 1.2, 0.22, primary)
            add_rect(slide, ctx, w - 2.75, 3.58, 1.2, 0.22, tertiary)
            add_line(slide, ctx, area.left_in, h - 1.05, w - 1.2, h - 1.05, tertiary)
            if not ctx.is_operator_layout(item) and item.layout_archetype in {"cover", "product_mockup", "ask"}:
                ctx.add_visual_zone(item.slide_no, w - 4.05, 1.0, 3.45, h - 2.0, zone_type="product_or_material_stage", source=style)
            return profile
        if ctx.uses_family("food_product_launch"):
            add_rect(slide, ctx, 0.0, h - 0.28, w, 0.28, secondary)
            add_oval(slide, ctx, w - 2.85, 1.05, 1.15, 1.15, tertiary)
            add_oval(slide, ctx, w - 1.75, 1.8, 0.62, 0.62, primary)
            add_rect(slide, ctx, w - 3.35, 5.45, 2.35, 0.3, primary)
            add_rect(slide, ctx, w - 2.85, 5.92, 1.8, 0.18, tertiary)
            return profile
        add_rect(slide, ctx, w - 3.55, 0.0, 3.55, h, primary)
        add_rect(slide, ctx, w - 4.25, 1.0, 0.48, h - 2.0, secondary)
        add_oval(slide, ctx, w - 2.8, 1.05, 1.2, 1.2, tertiary)
        add_oval(slide, ctx, w - 1.65, h - 1.75, 0.95, 0.95, secondary)
        if not ctx.is_operator_layout(item) and item.layout_archetype in {"big_thesis", "behavior_shift", "product_mockup", "cover"}:
            ctx.add_visual_zone(item.slide_no, w - 4.25, 0.0, 4.25, h, zone_type="right_edge_decorative", source=style)
    elif style == "gallery_grid":
        add_rect(slide, ctx, 0.0, h - 0.34, w, 0.34, secondary)
        if item.layout_archetype in RIGHT_VISUAL_COMPATIBLE_ARCHETYPES:
            for row in range(2):
                for col in range(3):
                    add_rect(slide, ctx, 6.65 + col * 1.75, 1.15 + row * 1.65, 1.3, 1.05, [primary, secondary, tertiary][(row + col) % 3])
            ctx.add_visual_zone(item.slide_no, 6.35, 1.0, w - 6.35, 3.6, zone_type="right_mockup_or_gallery", source=style)
    elif style == "journey_path":
        y = h - 1.55
        add_line(slide, ctx, area.left_in + 0.25, y, area.right_in - 0.3, y, primary)
        for idx in range(4):
            x = area.left_in + 0.5 + idx * 2.7
            add_oval(slide, ctx, x, y - 0.22, 0.44, 0.44, [primary, secondary, tertiary, "neutral"][idx])
        add_rect(slide, ctx, 0.0, 0.0, 0.32, h, secondary)
    elif style == "matrix_grid":
        add_rect(slide, ctx, 0.0, 0.0, w, 0.24, primary)
        if item.layout_archetype in RIGHT_VISUAL_COMPATIBLE_ARCHETYPES:
            x0, y0 = 6.35, 1.15
            for row in range(3):
                for col in range(3):
                    add_rect(slide, ctx, x0 + col * 1.75, y0 + row * 1.05, 1.46, 0.72, [primary, secondary, tertiary][(row * 3 + col) % 3])
            ctx.add_visual_zone(item.slide_no, x0 - 0.25, y0 - 0.15, w - x0, 3.35, zone_type="right_mockup_or_grid", source=style)
    elif style == "system_map":
        add_rect(slide, ctx, 0.0, h - 0.3, w, 0.3, secondary)
        if item.layout_archetype in RIGHT_VISUAL_COMPATIBLE_ARCHETYPES:
            nodes = [(6.1, 1.55), (8.55, 2.35), (6.9, 4.35), (10.1, 4.75)]
            for idx, (x, y) in enumerate(nodes):
                add_oval(slide, ctx, x, y, 0.64, 0.64, [primary, secondary, tertiary, "neutral"][idx])
            for start, end in [(nodes[0], nodes[1]), (nodes[1], nodes[2]), (nodes[2], nodes[3])]:
                add_line(slide, ctx, start[0] + 0.32, start[1] + 0.32, end[0] + 0.32, end[1] + 0.32, tertiary)
            ctx.add_visual_zone(item.slide_no, 5.85, 1.3, w - 5.85, 4.25, zone_type="right_mockup_or_map", source=style)
    elif style == "campaign_grid":
        add_rect(slide, ctx, 0.0, 0.0, w, 0.28, tertiary)
        if item.layout_archetype in RIGHT_VISUAL_COMPATIBLE_ARCHETYPES:
            for col in range(5):
                add_rect(slide, ctx, 5.2 + col * 1.25, 1.2, 0.92, 4.45, [primary, secondary, tertiary, "neutral", primary][col])
            ctx.add_visual_zone(item.slide_no, 5.05, 1.05, w - 5.05, 4.85, zone_type="right_mockup_or_grid", source=style)
    elif style == "memo_sections":
        add_rect(slide, ctx, area.left_in - 0.2, 1.25, 2.15, h - 2.4, "background", primary)
        for idx in range(4):
            add_line(slide, ctx, 4.2, 1.55 + idx * 0.9, w - 1.3, 1.55 + idx * 0.9, [primary, secondary, tertiary, "neutral"][idx])
        add_rect(slide, ctx, 0.0, 0.0, 0.2, h, primary)
    elif style == "proof_cards":
        for idx in range(3):
            add_rect(slide, ctx, 1.1 + idx * 3.85, h - 1.55, 3.05, 0.78, [primary, secondary, tertiary][idx])
        add_rect(slide, ctx, w - 0.32, 0.0, 0.32, h, secondary)
    elif style == "story_stage":
        add_rect(slide, ctx, 0.0, 0.0, w, 0.42, primary)
        add_rect(slide, ctx, 0.0, h - 0.7, w, 0.7, secondary)
        add_oval(slide, ctx, w - 2.35, 1.15, 1.45, 1.45, tertiary)
    else:
        add_rect(slide, ctx, w - 3.25, 0.0, 3.25, h, primary)
        add_rect(slide, ctx, w - 3.55, 0.95, 0.24, h - 1.9, secondary)
        if not ctx.is_operator_layout(item) and item.layout_archetype in {"big_thesis", "behavior_shift", "product_mockup", "cover"}:
            ctx.add_visual_zone(item.slide_no, w - 3.55, 0.0, 3.55, h, zone_type="right_edge_decorative", source=style)

    return profile


def split_items(slide_item: SlidePlanItem, count: int) -> list[str]:
    candidates = [str(item) for item in slide_item.visible_content_candidates]
    if not candidates:
        candidates = [slide_item.content_brief]
    while len(candidates) < count:
        candidates.append(f"Point {len(candidates) + 1}")
    return candidates[:count]


def slide_title(slide, ctx: BuildContext, item: SlidePlanItem) -> None:
    area = ctx.packet.safe_area
    title = visible_content(item)["title"]
    if ctx.is_operator_layout(item):
        add_rect(slide, ctx, area.left_in + 0.35, area.top_in + 0.18, 0.16, 0.82, "support")
        add_text(slide, ctx, f"{item.slide_no:02d}", area.left_in + 0.6, area.top_in + 0.28, 0.55, 0.24, role="caption", color_role="main")
        add_text(slide, ctx, title, area.left_in + 1.1, area.top_in + 0.16, 7.0, 0.65, role="body")
        add_line(slide, ctx, area.left_in + 1.1, area.top_in + 0.95, area.right_in - 0.35, area.top_in + 0.95, "main")
        return
    add_text(slide, ctx, title, area.left_in, area.top_in + 0.2, text_width_before_right_visual(ctx, item, area.left_in, padding=0.7), 0.7, role="slide_title")


def add_footer(slide, ctx: BuildContext, slide_no: int) -> None:
    policy = ctx.packet.header_footer_policy.footer.mode
    if policy == "disabled":
        ctx.omitted_header_footer.append({"slide_no": slide_no, "kind": "footer", "reason": "policy_disabled"})
        return
    area = ctx.packet.safe_area
    add_text(
        slide,
        ctx,
        f"{sanitize_visible_value(ctx.packet.guide_identity.project_name, fallback='Untitled Deck')}  |  {slide_no}/{ctx.packet.guide_identity.slide_count}",
        area.left_in,
        area.bottom_in - 0.22,
        area.right_in - area.left_in,
        0.18,
        role="caption",
        color_role="neutral",
        align=PP_ALIGN.RIGHT,
        qa_text_role="footer",
    )


def normalized_sha256(value: str | None) -> str | None:
    if not value:
        return None
    text = str(value).strip().lower()
    return text.split(":", 1)[1] if text.startswith("sha256:") else text


def sha256_uri(value: str | None) -> str:
    normalized = normalized_sha256(value) or ""
    return f"sha256:{normalized}" if normalized and not normalized.startswith("sha256:") else normalized


def package_asset_manifest_info(asset: ApprovedAssetReference) -> dict[str, Any]:
    extra = asset.model_extra or {}
    insertion_role = extra.get("allowed_insertion_role") or extra.get("insert_as") or asset.slot_id or asset.asset_role
    return {
        "package_manifest_id": extra.get("package_manifest_id"),
        "relative_package_path": public_package_path_or_none(extra.get("relative_package_path")),
        "sha256": sha256_uri(asset.sha256),
        "insertion_role": insertion_role,
        "media_type": extra.get("media_type"),
    }


def add_asset_validation_failure(ctx: BuildContext, failure: dict[str, Any]) -> None:
    reason = canonical_approved_asset_reason(str(failure.get("reason") or "approved_asset_insert_failed"))
    failure["reason"] = reason
    key = (
        failure.get("slot_id"),
        failure.get("asset_ref"),
        failure.get("relative_package_path"),
        failure.get("reason"),
    )
    for item in ctx.asset_validation_failures:
        if (
            item.get("slot_id"),
            item.get("asset_ref"),
            item.get("relative_package_path"),
            item.get("reason"),
        ) == key:
            return
    ctx.asset_validation_failures.append(failure)


def canonical_approved_asset_reason(reason: str) -> str:
    generic_reasons = {
        "checksum_or_file_size_validation_failed": "approved_asset_checksum_mismatch",
        "no valid approved package asset": "approved_asset_insert_failed",
        "no approved package asset with valid checksum and size": "approved_asset_insert_failed",
    }
    normalized = generic_reasons.get(reason, reason)
    return normalized if normalized in APPROVED_ASSET_BLOCKING_REASONS else "approved_asset_insert_failed"


def is_safe_relative_package_path(value: str | None) -> bool:
    return normalize_package_relative_path(value) is not None


def normalize_package_relative_path(value: str | None) -> str | None:
    if not value:
        return None
    normalized = str(value).strip().replace("\\", "/")
    while normalized.startswith("./"):
        normalized = normalized[2:]
    if normalized.startswith("/") or re.match(r"^[A-Za-z]:", normalized) or "://" in normalized:
        return None
    parts = PurePosixPath(normalized).parts
    if not parts or any(part in {"", ".", ".."} for part in parts):
        return None
    return "/".join(parts)


def public_package_path_or_none(value: Any) -> str | None:
    return normalize_package_relative_path(str(value)) if value else None


def public_asset_event(event: dict[str, Any]) -> dict[str, Any]:
    item = dict(event)
    for field in ["relative_package_path", "source_asset_ref"]:
        if field not in item:
            continue
        normalized = public_package_path_or_none(item.get(field))
        if normalized:
            item[field] = normalized
        elif field == "relative_package_path":
            item.pop(field, None)
    return item


def public_asset_failure(failure: dict[str, Any]) -> dict[str, Any]:
    item = dict(failure)
    for field in ["asset_ref", "relative_package_path", "source_asset_ref"]:
        if field not in item or not item.get(field):
            continue
        normalized = public_package_path_or_none(item.get(field))
        item[field] = normalized if normalized else "[unsafe-package-path-redacted]"
    return item


def approved_asset_type_usable_for_slot(asset: ApprovedAssetReference, slot: AssetSlot) -> bool:
    extra = asset.model_extra or {}
    media_type = str(extra.get("media_type") or asset.asset_type or "").lower()
    insert_as = str(extra.get("insert_as") or asset.asset_type or "").lower()
    role = str(extra.get("allowed_insertion_role") or asset.asset_role or "").lower()
    slot_text = f"{slot.slot_id} {slot.slot_type} {slot.crop_or_mask_policy}".lower()
    if "policy_blocked" in str(extra.get("license_action") or "").lower():
        return False
    if media_type and not media_type.startswith("image/") and insert_as not in {"icon", "mockup", "approved_image"}:
        return False
    if any(token in slot_text for token in ["image", "visual", "mockup", "screen", "background", "texture", "logo", "icon", "packshot", "ingredient", "shelf", "campaign", "hero", "chart", "roadmap", "risk"]):
        return True
    return bool(media_type.startswith("image/") or insert_as in {"icon", "mockup", "approved_image"} or "image" in role)


def resolve_asset(ctx: BuildContext, slot: AssetSlot) -> dict[str, Any]:
    matches = [
        asset
        for asset in ctx.packet.approved_asset_references
        if (asset.approved or asset.manifest_checksum)
        and (asset.slot_id == slot.slot_id or asset.asset_ref == slot.preferred_asset_ref or asset.asset_role == slot.slot_type)
    ]
    for asset in matches:
        extra = asset.model_extra or {}
        explicit_failure = extra.get("package_validation_failure")
        if explicit_failure:
            reason = canonical_approved_asset_reason(str(explicit_failure))
            add_asset_validation_failure(
                ctx,
                {
                    "slot_id": slot.slot_id,
                    "asset_ref": asset.asset_ref,
                    "relative_package_path": extra.get("relative_package_path") or asset.asset_ref,
                    "reason": reason,
                    "status": "fail",
                },
            )
            return {"status": "native_fallback", "asset_ref": asset.asset_ref, "path": None, "fallback_reason": reason}
        raw_relative_package_path = str(extra.get("relative_package_path") or asset.asset_ref or "")
        relative_package_path = public_package_path_or_none(raw_relative_package_path)
        if not relative_package_path:
            reason = "approved_asset_path_unsafe"
            add_asset_validation_failure(
                ctx,
                {
                    "slot_id": slot.slot_id,
                    "asset_ref": asset.asset_ref,
                    "relative_package_path": raw_relative_package_path,
                    "reason": reason,
                    "status": "fail",
                },
            )
            return {"status": "native_fallback", "asset_ref": asset.asset_ref, "path": None, "fallback_reason": reason}
        if not approved_asset_type_usable_for_slot(asset, slot):
            reason = "approved_asset_type_unsupported_for_slot"
            add_asset_validation_failure(
                ctx,
                {
                    "slot_id": slot.slot_id,
                    "asset_ref": asset.asset_ref,
                    "relative_package_path": relative_package_path,
                    "reason": reason,
                    "status": "fail",
                    "media_type": extra.get("media_type"),
                    "slot_type": slot.slot_type,
                },
            )
            return {"status": "native_fallback", "asset_ref": asset.asset_ref, "path": None, "fallback_reason": reason}
        if not asset.local_path:
            reason = "approved_asset_file_missing"
            add_asset_validation_failure(
                ctx,
                {
                    "slot_id": slot.slot_id,
                    "asset_ref": asset.asset_ref,
                    "relative_package_path": relative_package_path,
                    "reason": reason,
                    "status": "fail",
                },
            )
            return {"status": "native_fallback", "asset_ref": asset.asset_ref, "path": None, "fallback_reason": reason}
        path = Path(asset.local_path)
        if not path.is_absolute():
            path = (BASE_DIR / path).resolve()
        checksum_valid = None
        size_valid = None
        if path.exists():
            if asset.sha256:
                checksum_valid = hashlib.sha256(path.read_bytes()).hexdigest().lower() == normalized_sha256(asset.sha256)
            if asset.file_size_bytes is not None:
                size_valid = path.stat().st_size == asset.file_size_bytes
            if checksum_valid is not False and size_valid is not False:
                manifest_info = package_asset_manifest_info(asset)
                return {
                    "status": "approved_asset",
                    "asset_ref": asset.asset_ref,
                    "path": path,
                    "checksum_valid": checksum_valid,
                    "file_size_valid": size_valid,
                    "file_size_bytes": asset.file_size_bytes,
                    **manifest_info,
                }
            reason = "approved_asset_checksum_mismatch" if checksum_valid is False else "approved_asset_size_mismatch"
            add_asset_validation_failure(
                ctx,
                {
                    "slot_id": slot.slot_id,
                    "asset_ref": asset.asset_ref,
                    "relative_package_path": public_package_path_or_none((asset.model_extra or {}).get("relative_package_path")),
                    "reason": reason,
                    "checksum_valid": checksum_valid,
                    "file_size_valid": size_valid,
                    "status": "fail",
                },
            )
        else:
            add_asset_validation_failure(
                ctx,
                {
                    "slot_id": slot.slot_id,
                    "asset_ref": asset.asset_ref,
                    "relative_package_path": public_package_path_or_none((asset.model_extra or {}).get("relative_package_path")),
                    "reason": "approved_package_file_missing",
                    "status": "fail",
                },
            )
            return {"status": "native_fallback", "asset_ref": asset.asset_ref, "path": None, "fallback_reason": "approved_asset_file_missing"}
    return {
        "status": "native_fallback",
        "asset_ref": None,
        "path": None,
        "fallback_reason": "no_approved_package_asset_declared_for_slot",
    }


def asset_event_exists(ctx: BuildContext, slot_id: str) -> bool:
    return any(event.get("slot_id") == slot_id for event in ctx.asset_events)


def record_asset_slot_fallback(ctx: BuildContext, slide_no: int, slot: AssetSlot, handling: str) -> None:
    if asset_event_exists(ctx, slot.slot_id):
        return
    resolved = resolve_asset(ctx, slot)
    if resolved["status"] == "approved_asset":
        reason = "approved_asset_insert_failed"
        add_asset_validation_failure(
            ctx,
            {
                "slot_id": slot.slot_id,
                "asset_ref": resolved["asset_ref"],
                "relative_package_path": resolved.get("relative_package_path"),
                "reason": reason,
                "status": "fail",
                "details": "A valid file-backed approved package asset was available for this slot but the renderer did not insert it.",
            },
        )
        ctx.asset_events.append(
            {
                "slide_no": slide_no,
                "slot_id": slot.slot_id,
                "slot_type": slot.slot_type,
                "required": slot.required,
                "resolution": "approved_package_asset_not_inserted",
                "source_asset_ref": resolved["asset_ref"],
                "path": safe_rel(resolved["path"]),
                "checksum_valid": resolved["checksum_valid"],
                "file_size_valid": resolved["file_size_valid"],
                "fallback_used": True,
                "fallback_reason": reason,
                "crop_or_mask_policy": slot.crop_or_mask_policy,
                "native_handling": handling,
                "package_manifest_id": resolved.get("package_manifest_id"),
                "relative_package_path": resolved.get("relative_package_path"),
                "sha256": resolved.get("sha256"),
                "source_type": "approved_package_file",
            }
        )
        return
    ctx.asset_events.append(
        {
            "slide_no": slide_no,
            "slot_id": slot.slot_id,
            "slot_type": slot.slot_type,
            "required": slot.required,
            "resolution": "native_shape_or_mockup_fallback",
            "source_asset_ref": None,
            "path": None,
            "fallback_used": True,
            "fallback_reason": resolved["fallback_reason"],
            "crop_or_mask_policy": slot.crop_or_mask_policy,
            "native_handling": handling,
        }
    )


def record_unhandled_asset_slots(ctx: BuildContext, item: SlidePlanItem) -> None:
    for slot in item.asset_slots:
        if slot.slot_id == "edge_bleed_visual":
            handling = "native right-edge bleed shape cluster; all text remains inside safe area"
        elif slot.slot_id == "background_texture":
            handling = "solid background plus native shape texture; text contrast preserved"
        elif slot.slot_type == "icon":
            handling = "native icon-like shapes in card headers"
        elif slot.slot_type == "background":
            handling = "native background fallback"
        else:
            handling = "native editable fallback or explicit skip recorded"
        record_asset_slot_fallback(ctx, item.slide_no, slot, handling)


def svg_colors(path: Path) -> list[str]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    values = re.findall(r"#[0-9A-Fa-f]{6}", text)
    return values or ["#F7F8FC", "#315CFF", "#FF4DA6", "#22C7A9", "#101828"]


def fixture_svg_derivative(path: Path, output: Path, role: str, size: tuple[int, int] = (900, 600)) -> Path:
    output.parent.mkdir(parents=True, exist_ok=True)
    colors = svg_colors(path)
    bg = colors[0]
    image = Image.new("RGBA", size, bg)
    draw = ImageDraw.Draw(image, "RGBA")
    role_text = role.lower()
    if "phone" in role_text or "mockup" in role_text:
        image = Image.new("RGBA", (900, 1600), "#101828")
        draw = ImageDraw.Draw(image, "RGBA")
        gradient = colors[1:4] or ["#315CFF", "#FF4DA6", "#22C7A9"]
        for y in range(70, 1530):
            ratio = (y - 70) / max(1, 1460)
            first = Image.new("RGB", (1, 1), gradient[0]).getpixel((0, 0))
            last = Image.new("RGB", (1, 1), gradient[min(2, len(gradient) - 1)]).getpixel((0, 0))
            rgb = tuple(int(first[i] * (1 - ratio) + last[i] * ratio) for i in range(3))
            draw.line([(58, y), (842, y)], fill=rgb + (255,), width=1)
        draw.rounded_rectangle([58, 70, 842, 1530], radius=70, outline="#101828", width=0)
        draw.rounded_rectangle([120, 160, 780, 252], radius=46, fill=(255, 255, 255, 238))
        draw.rounded_rectangle([120, 330, 780, 660], radius=52, fill=(255, 255, 255, 230))
        draw.rounded_rectangle([120, 730, 420, 980], radius=44, fill=(255, 255, 255, 220))
        draw.rounded_rectangle([480, 730, 780, 980], radius=44, fill=(255, 255, 255, 205))
        draw.rounded_rectangle([120, 1050, 780, 1310], radius=52, fill=(255, 255, 255, 220))
    elif "edge" in role_text or "bleed" in role_text:
        image = Image.new("RGBA", (1200, 900), colors[0])
        draw = ImageDraw.Draw(image, "RGBA")
        for idx, color in enumerate(colors[1:4]):
            rgb = Image.new("RGB", (1, 1), color).getpixel((0, 0))
            draw.ellipse([760 + idx * 80, 40 + idx * 230, 1320 + idx * 20, 600 + idx * 230], fill=rgb + (210,))
    elif "background" in role_text or "texture" in role_text:
        image = Image.new("RGBA", (1600, 900), colors[0])
        draw = ImageDraw.Draw(image, "RGBA")
        line_rgb = Image.new("RGB", (1, 1), colors[1] if len(colors) > 1 else "#315CFF").getpixel((0, 0))
        for x in range(160, 1600, 160):
            draw.line([(x, 0), (x, 900)], fill=line_rgb + (28,), width=2)
        for y in range(120, 900, 120):
            draw.line([(0, y), (1600, y)], fill=line_rgb + (28,), width=2)
    else:
        image = Image.new("RGBA", (256, 256), (255, 255, 255, 0))
        draw = ImageDraw.Draw(image, "RGBA")
        for idx, color in enumerate(colors[:3]):
            rgb = Image.new("RGB", (1, 1), color).getpixel((0, 0))
            draw.arc([38 + idx * 12, 55 + idx * 18, 218 - idx * 12, 201 - idx * 18], 20, 330, fill=rgb + (255,), width=18)
    image.save(output)
    return output


def pptx_ready_asset_path(ctx: BuildContext, source_path: Path, role: str) -> Path:
    if source_path.suffix.lower() not in {".svg", ".svgz"}:
        return source_path
    digest = hashlib.sha256(source_path.read_bytes()).hexdigest()[:12]
    output = ctx.project_dir / "_asset_derivatives" / f"{source_path.stem}-{digest}.png"
    if not output.exists():
        output.parent.mkdir(parents=True, exist_ok=True)
        try:
            import cairosvg  # type: ignore[import-not-found]

            cairosvg.svg2png(url=str(source_path), write_to=str(output))
        except Exception:
            fixture_svg_derivative(source_path, output, role)
    return output


def approved_asset_event(resolved: dict[str, Any], slot: AssetSlot, slide_no: int | None, inserted_path: Path | None = None) -> dict[str, Any]:
    event = {
        "slot_id": slot.slot_id,
        "slide_no": slide_no,
        "resolution": "approved_package_asset",
        "source_asset_ref": resolved["asset_ref"],
        "path": safe_rel(resolved["path"]),
        "checksum_valid": resolved["checksum_valid"],
        "file_size_valid": resolved["file_size_valid"],
        "fallback_used": False,
        "slot_type": slot.slot_type,
        "required": slot.required,
        "crop_or_mask_policy": slot.crop_or_mask_policy,
        "package_manifest_id": resolved.get("package_manifest_id"),
        "relative_package_path": resolved.get("relative_package_path"),
        "sha256": resolved.get("sha256"),
        "source_type": "approved_package_file",
        "file_backed": True,
        "insertion_role": resolved.get("insertion_role") or slot.slot_id,
    }
    if inserted_path and inserted_path.resolve() != Path(resolved["path"]).resolve():
        event["pptx_inserted_derivative"] = safe_rel(inserted_path)
        event["derivative_reason"] = "svg_converted_to_png_for_powerpoint_insertion"
    return event


def add_approved_asset_picture(
    slide,
    ctx: BuildContext,
    slot: AssetSlot | None,
    slide_no: int,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    zone_type: str,
) -> bool:
    if not slot:
        return False
    resolved = resolve_asset(ctx, slot)
    if resolved["status"] != "approved_asset":
        return False
    inserted_path = pptx_ready_asset_path(ctx, resolved["path"], resolved.get("insertion_role") or slot.slot_id)
    try:
        slide.shapes.add_picture(str(inserted_path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    except Exception as exc:
        add_asset_validation_failure(
            ctx,
            {
                "slot_id": slot.slot_id,
                "asset_ref": resolved["asset_ref"],
                "relative_package_path": resolved.get("relative_package_path"),
                "reason": "approved_asset_insert_failed",
                "status": "fail",
                "details": str(exc),
            },
        )
        return False
    ctx.asset_events.append(approved_asset_event(resolved, slot, slide_no, inserted_path))
    ctx.add_visual_zone(slide_no, x, y, w, h, zone_type=zone_type, source=slot.slot_id)
    return True


def add_phone_mockup(
    slide,
    ctx: BuildContext,
    x: float,
    y: float,
    w: float,
    h: float,
    slot: AssetSlot | None = None,
    slide_no: int | None = None,
) -> None:
    if slot:
        resolved = resolve_asset(ctx, slot)
        if resolved["status"] == "approved_asset":
            inserted_path = pptx_ready_asset_path(ctx, resolved["path"], resolved.get("insertion_role") or slot.slot_id)
            try:
                slide.shapes.add_picture(str(inserted_path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
                ctx.asset_events.append(approved_asset_event(resolved, slot, slide_no, inserted_path))
                ctx.add_visual_zone(slide_no, x, y, w, h, zone_type="phone_or_mockup", source=slot.slot_id)
                return
            except Exception as exc:
                add_asset_validation_failure(
                    ctx,
                    {
                        "slot_id": slot.slot_id,
                        "asset_ref": resolved["asset_ref"],
                        "relative_package_path": resolved.get("relative_package_path"),
                        "reason": "approved_asset_insert_failed",
                        "status": "fail",
                        "details": str(exc),
                    },
                )
                resolved = {**resolved, "fallback_reason": "approved_asset_insert_failed"}
        ctx.asset_events.append(
            {
                "slot_id": slot.slot_id,
                "slide_no": slide_no,
                "resolution": "native_shape_or_mockup_fallback",
                "source_asset_ref": None,
                "path": None,
                "fallback_used": True,
                "fallback_reason": resolved["fallback_reason"],
                "slot_type": slot.slot_type,
                "required": slot.required,
                "crop_or_mask_policy": slot.crop_or_mask_policy,
                "native_handling": "native editable phone/mockup fallback",
            }
        )
    frame = add_rect(slide, ctx, x, y, w, h, "neutral")
    frame.adjustments[0] = 0.12
    screen = add_rect(slide, ctx, x + 0.18, y + 0.35, w - 0.36, h - 0.72, "background")
    screen.adjustments[0] = 0.08
    add_rect(slide, ctx, x + 0.38, y + 0.75, w - 0.76, 0.46, "main")
    add_rect(slide, ctx, x + 0.38, y + 1.4, w - 0.76, 0.34, "support")
    add_rect(slide, ctx, x + 0.38, y + 1.92, w - 0.76, 0.34, "accent")
    ctx.add_visual_zone(slide_no, x, y, w, h, zone_type="phone_or_mockup", source=slot.slot_id if slot else "native_phone_mockup")
    add_text(slide, ctx, native_mockup_label(ctx), x + 0.45, y + 0.81, w - 0.9, 0.22, role="caption", color_role="background", align=PP_ALIGN.CENTER, qa_text_role="decorative")


def first_slot(item: SlidePlanItem, kinds: set[str]) -> AssetSlot | None:
    matches: list[AssetSlot] = []
    for slot in item.asset_slots:
        slot_text = f"{slot.slot_id} {slot.slot_type} {slot.crop_or_mask_policy}".lower()
        if slot.slot_type in kinds or any(kind in slot_text for kind in kinds):
            matches.append(slot)
    if matches:
        return next((slot for slot in matches if slot.preferred_asset_ref), matches[0])
    return next((slot for slot in item.asset_slots if slot.preferred_asset_ref), item.asset_slots[0]) if item.asset_slots else None


def native_mockup_label(ctx: BuildContext) -> str:
    text = " ".join(
        [
            ctx.packet.guide_identity.project_name,
            ctx.packet.guide_identity.topic,
            ctx.packet.project_brief.objective,
        ]
    ).lower()
    if "triage" in text or "emergency department" in text:
        return "Flow signal"
    if "light" in text or "exhibition" in text or "sponsor" in text:
        return "Journey cue"
    if "robot" in text or "safety" in text:
        return "Safety loop"
    if "food" in text or "retail" in text:
        return "Launch cue"
    return "Proof loop"


def cover_chip_labels(ctx: BuildContext) -> list[str]:
    text = " ".join(
        [
            ctx.packet.guide_identity.project_name,
            ctx.packet.guide_identity.topic,
            ctx.packet.project_brief.objective,
        ]
    ).lower()
    if "triage" in text or "emergency department" in text:
        return ["Triage", "Staffing", "Waits"]
    if "light" in text or "exhibition" in text or "sponsor" in text:
        return ["Journey", "Brand", "Impact"]
    if "robot" in text or "safety" in text:
        return ["Zones", "Handoffs", "Checklist"]
    return ["Audience", "Evidence", "Action"]


def render_cover(slide, ctx: BuildContext, item: SlidePlanItem) -> None:
    area = ctx.packet.safe_area
    if ctx.uses_family("internal_strategy_report"):
        add_rect(slide, ctx, 0.0, 0.0, 3.15, ctx.packet.slide_size.height_in, "background", "main")
        add_text(slide, ctx, "CONFIDENTIAL BRIEF", 0.78, 1.0, 1.9, 0.22, role="caption", color_role="main", align=PP_ALIGN.CENTER)
        add_line(slide, ctx, 3.55, 1.0, 3.55, 6.4, "main")
        title = visible_content(item)["title"]
        add_text(slide, ctx, title, 4.05, 1.05, 7.75, 1.9, role="slide_title", color_role="neutral")
        add_text(slide, ctx, visible_content(item)["subtitle"], 4.08, 3.45, 6.6, 0.78, role="body", color_role="neutral")
        for idx, label in enumerate(["Decision", "Risk", "Action"]):
            add_rect(slide, ctx, 4.1 + idx * 2.05, 4.95, 1.58, 0.52, ["main", "support", "accent"][idx])
            add_text(slide, ctx, label, 4.22 + idx * 2.05, 5.12, 1.3, 0.16, role="caption", color_role="background", align=PP_ALIGN.CENTER)
        return
    if ctx.uses_family("luxury_editorial"):
        add_rect(slide, ctx, 8.6, 0.85, 2.75, 5.7, "background", "neutral")
        add_rect(slide, ctx, 9.05, 1.3, 1.85, 4.8, "accent")
        add_line(slide, ctx, 1.0, 1.05, 7.2, 1.05, "accent")
        add_text(slide, ctx, visible_content(item)["title"], 1.0, 1.55, 6.65, 1.55, role="hero_title", color_role="neutral")
        add_text(slide, ctx, visible_content(item)["subtitle"], 1.05, 4.35, 5.8, 0.65, role="body", color_role="neutral")
        add_text(slide, ctx, "Craft / Detail / Launch", 1.05, 5.55, 3.1, 0.2, role="caption", color_role="accent")
        ctx.add_visual_zone(item.slide_no, 8.4, 0.7, 3.2, 6.0, zone_type="product_or_material_stage", source="luxury_editorial")
        return
    if ctx.uses_family("food_product_launch"):
        hero_slot = first_slot(item, {"packshot", "product", "hero", "image"})
        if not add_approved_asset_picture(slide, ctx, hero_slot, item.slide_no, 7.85, 1.0, 3.55, 3.95, zone_type="product_stage"):
            add_oval(slide, ctx, 8.15, 1.05, 3.0, 3.0, "accent")
            add_rect(slide, ctx, 8.5, 3.9, 2.35, 1.25, "background", "main")
        add_rect(slide, ctx, 7.85, 5.25, 3.55, 0.32, "support")
        title = visible_content(item)["title"]
        long_title = len(title) > 18 or bool(re.search(r"[\uac00-\ud7af]", title))
        add_text(slide, ctx, title, 0.95, 1.1, 6.55, 1.9 if long_title else 1.35, role="slide_title" if long_title else "hero_title", color_role="neutral")
        add_text(slide, ctx, visible_content(item)["subtitle"], 1.0, 3.65 if long_title else 3.25, 5.8, 0.82, role="body", color_role="neutral")
        for idx, label in enumerate(["Taste", "Shelf", "Promo"]):
            add_rect(slide, ctx, 1.0 + idx * 1.55, 5.15, 1.08, 0.38, ["main", "support", "accent"][idx])
            add_text(slide, ctx, label, 1.08 + idx * 1.55, 5.28, 0.92, 0.12, role="caption", color_role="background", align=PP_ALIGN.CENTER)
        ctx.add_visual_zone(item.slide_no, 7.65, 0.95, 3.95, 4.75, zone_type="product_stage", source="food_product_launch")
        return
    if ctx.is_operator_layout(item):
        add_rect(slide, ctx, 0.0, 0.0, 4.8, ctx.packet.slide_size.height_in, "main")
        add_phone_mockup(slide, ctx, 1.25, 1.05, 2.25, 5.45, first_slot(item, {"mockup", "phone"}), item.slide_no)
        title = visible_content(item)["title"]
        long_title = len(title) > 52 or bool(re.search(r"[\uac00-\ud7af]", title))
        title_y = 1.0 if long_title else 1.2
        title_h = 2.05 if long_title else 1.15
        subtitle_y = 3.85 if long_title else 2.85
        chip_y = 5.15 if long_title else 4.55
        add_text(slide, ctx, title, 5.45, title_y, 6.7, title_h, role="slide_title" if long_title else "hero_title", color_role="neutral")
        add_text(slide, ctx, visible_content(item)["subtitle"], 5.5, subtitle_y, 5.9, 0.95, role="body", color_role="neutral")
        for idx, label in enumerate(cover_chip_labels(ctx)):
            add_rect(slide, ctx, 5.55 + idx * 1.65, chip_y, 1.28, 0.34, ["support", "accent", "neutral"][idx])
            add_text(slide, ctx, label, 5.62 + idx * 1.65, chip_y + 0.09, 1.12, 0.1, role="caption", color_role="background", align=PP_ALIGN.CENTER)
        return
    cover_phone_x = 9.25
    title_x = area.left_in + 0.25
    title_w = min(text_width_before_right_visual(ctx, item, title_x, padding=0.35), cover_phone_x - title_x - 0.35)
    title_text = visible_content(item)["title"]
    narrow_right_visual = title_w < 6.25 or len(title_text) > 48
    title_role = "slide_title" if narrow_right_visual else "hero_title"
    title_y = 1.35 if narrow_right_visual else 1.55
    title_h = 2.55 if narrow_right_visual else 1.3
    add_text(slide, ctx, title_text, title_x, title_y, title_w, title_h, role=title_role, color_role="neutral")
    subtitle_x = area.left_in + 0.3
    subtitle_w = min(text_width_before_right_visual(ctx, item, subtitle_x, padding=0.45), cover_phone_x - subtitle_x - 0.45)
    subtitle_y = 4.25 if narrow_right_visual else 3.15
    accent_y = 5.28 if narrow_right_visual else 4.35
    add_text(slide, ctx, visible_content(item)["subtitle"], subtitle_x, subtitle_y, subtitle_w, 0.75, role="body", color_role="neutral")
    add_rect(slide, ctx, area.left_in + 0.3, accent_y, 2.3, 0.18, "main")
    add_rect(slide, ctx, area.left_in + 2.75, accent_y, 1.2, 0.18, "support")
    add_rect(slide, ctx, area.left_in + 4.1, accent_y, 1.7, 0.18, "accent")
    add_phone_mockup(slide, ctx, cover_phone_x, 0.95, 2.2, 5.55, first_slot(item, {"mockup", "phone"}), item.slide_no)


def render_big_thesis(slide, ctx: BuildContext, item: SlidePlanItem) -> None:
    if ctx.is_operator_layout(item):
        if not (ctx.uses_family("food_product_launch") and add_approved_asset_picture(slide, ctx, first_slot(item, {"ingredient", "material", "product", "image"}), item.slide_no, 0.8, 1.1, 3.1, 4.9, zone_type="ingredient_detail_asset")):
            add_rect(slide, ctx, 0.8, 1.1, 3.1, 4.9, "neutral")
        add_text(slide, ctx, "THESIS", 1.15, 1.55, 1.6, 0.25, role="caption", color_role="background")
        add_text(slide, ctx, visible_content(item)["title"], 4.55, 1.55, 7.4, 1.55, role="slide_title", color_role="neutral")
        add_text(slide, ctx, visible_content(item)["brief"], 4.6, 3.65, 6.2, 0.9, role="body", color_role="neutral")
        add_rect(slide, ctx, 4.6, 5.25, 5.9, 0.22, "support")
        return
    if ctx.uses_family("food_product_launch") and item.asset_slots:
        add_approved_asset_picture(slide, ctx, first_slot(item, {"ingredient", "material", "product", "image"}), item.slide_no, 8.45, 1.35, 3.25, 4.35, zone_type="ingredient_detail_asset")
    x = 1.05
    safe_w = text_width_before_right_visual(ctx, item, x, padding=0.35)
    align = PP_ALIGN.LEFT if safe_w < 9.5 else PP_ALIGN.CENTER
    add_text(slide, ctx, visible_content(item)["title"], x, 1.55, safe_w, 2.1, role="hero_title", color_role="neutral", align=align)
    body_x = x + (0.7 if safe_w >= 9.5 else 0.0)
    body_w = max(1.0, safe_w - (1.4 if safe_w >= 9.5 else 0.0))
    add_text(slide, ctx, visible_content(item)["brief"], body_x, 4.1, body_w, 0.9, role="body", color_role="neutral", align=align)
    add_rect(slide, ctx, x, 5.55, min(3.9, safe_w), 0.16, "main")


def render_three_cards(slide, ctx: BuildContext, item: SlidePlanItem) -> None:
    slide_title(slide, ctx, item)
    icon_slot = first_slot(item, {"icon"})
    if ctx.is_operator_layout(item):
        add_text(slide, ctx, visible_content(item)["brief"], 0.9, 1.65, 3.25, 1.25, role="body", color_role="neutral")
        for idx, text in enumerate(split_items(item, 3)):
            y = 2.05 + idx * 1.35
            add_rect(slide, ctx, 4.8, y, 6.85, 0.92, ["main", "support", "accent"][idx])
            if idx == 0:
                add_approved_asset_picture(slide, ctx, icon_slot, item.slide_no, 5.02, y + 0.18, 0.42, 0.42, zone_type="icon")
            add_text(slide, ctx, f"{idx + 1}", 5.05, y + 0.22, 0.35, 0.18, role="caption", color_role="background", align=PP_ALIGN.CENTER)
            add_text(slide, ctx, text, 5.7, y + 0.22, 5.3, 0.24, role="body", color_role="background")
        return
    for idx, text in enumerate(split_items(item, 3)):
        x = 0.75 + idx * 4.15
        add_rect(slide, ctx, x, 2.1, 3.55, 3.25, "background", "main")
        if idx == 0 and add_approved_asset_picture(slide, ctx, icon_slot, item.slide_no, x + 0.25, 2.4, 0.55, 0.55, zone_type="icon"):
            pass
        else:
            add_rect(slide, ctx, x + 0.25, 2.4, 0.55, 0.55, ["main", "support", "accent"][idx])
        add_text(slide, ctx, text, x + 0.25, 3.2, 3.05, 1.15, role="body", color_role="neutral")


def render_behavior_shift(slide, ctx: BuildContext, item: SlidePlanItem) -> None:
    edge_slot = first_slot(item, {"image", "bleed"})
    if not ctx.is_operator_layout(item):
        add_approved_asset_picture(slide, ctx, edge_slot, item.slide_no, 9.05, 0.0, 4.28, ctx.packet.slide_size.height_in, zone_type="right_edge_asset")
    slide_title(slide, ctx, item)
    left, right = split_items(item, 2)
    if ctx.is_operator_layout(item):
        add_rect(slide, ctx, 1.0, 2.1, 10.8, 1.25, "background", "neutral")
        add_rect(slide, ctx, 1.0, 4.25, 10.8, 1.25, "main")
        add_text(slide, ctx, left, 1.35, 2.48, 9.9, 0.36, role="body", color_role="neutral", align=PP_ALIGN.CENTER)
        add_text(slide, ctx, right, 1.35, 4.63, 9.9, 0.36, role="body", color_role="background", align=PP_ALIGN.CENTER)
        add_line(slide, ctx, 6.4, 3.48, 6.4, 4.15, "support")
        return
    content_w = text_width_before_right_visual(ctx, item, 1.0, padding=0.4)
    add_rect(slide, ctx, 1.0, 2.25, content_w, 1.18, "background", "neutral")
    add_rect(slide, ctx, 1.0, 4.05, content_w, 1.18, "background", "main")
    add_text(slide, ctx, left, 1.35, 2.58, max(1.0, content_w - 0.7), 0.32, role="body", color_role="neutral", align=PP_ALIGN.CENTER)
    add_text(slide, ctx, right, 1.35, 4.38, max(1.0, content_w - 0.7), 0.32, role="body", color_role="neutral", align=PP_ALIGN.CENTER)
    add_line(slide, ctx, 1.4, 3.65, min(1.0 + content_w - 0.4, 8.4), 3.65, "main")


def render_product_mockup(slide, ctx: BuildContext, item: SlidePlanItem) -> None:
    slide_title(slide, ctx, item)
    if ctx.uses_family("luxury_editorial"):
        add_rect(slide, ctx, 1.05, 2.15, 4.0, 2.85, "background", "neutral")
        add_rect(slide, ctx, 5.45, 2.15, 2.1, 2.85, "accent")
        add_line(slide, ctx, 8.0, 2.2, 8.0, 5.15, "main")
        add_text(slide, ctx, visible_content(item)["brief"], 1.08, 5.35, 6.25, 0.62, role="body", color_role="neutral")
        ctx.add_visual_zone(item.slide_no, 1.0, 2.0, 6.7, 3.25, zone_type="material_detail_spread", source="luxury_editorial")
        return
    if ctx.uses_family("food_product_launch"):
        ingredient_slot = first_slot(item, {"ingredient", "packshot", "product", "image", "material", "context"})
        if not add_approved_asset_picture(slide, ctx, ingredient_slot, item.slide_no, 1.05, 2.0, 4.2, 3.2, zone_type="ingredient_product_stage"):
            add_oval(slide, ctx, 1.05, 2.0, 2.35, 2.35, "accent")
            add_oval(slide, ctx, 3.65, 2.55, 1.35, 1.35, "support")
        add_rect(slide, ctx, 5.55, 2.0, 2.25, 3.2, "background", "main")
        for idx, text in enumerate(split_items(item, 3)):
            add_rect(slide, ctx, 8.25, 2.0 + idx * 1.05, 3.05, 0.66, ["main", "support", "accent"][idx])
            add_text(slide, ctx, text, 8.45, 2.2 + idx * 1.05, 2.65, 0.18, role="caption", color_role="background", align=PP_ALIGN.CENTER)
        ctx.add_visual_zone(item.slide_no, 0.95, 1.9, 6.95, 3.45, zone_type="ingredient_product_stage", source="food_product_launch")
        return
    if ctx.uses_family("internal_strategy_report"):
        for idx, text in enumerate(split_items(item, 3)):
            y = 1.8 + idx * 1.25
            add_rect(slide, ctx, 1.05, y, 10.55, 0.86, "background", ["main", "support", "accent"][idx])
            add_text(slide, ctx, f"Option {idx + 1}", 1.35, y + 0.26, 1.25, 0.16, role="caption", color_role=["main", "support", "accent"][idx])
            add_text(slide, ctx, text, 3.0, y + 0.24, 7.4, 0.2, role="body", color_role="neutral")
        return
    if ctx.is_operator_layout(item):
        add_phone_mockup(slide, ctx, 1.2, 1.45, 2.15, 5.35, first_slot(item, {"phone", "mockup", "screen", "image"}), item.slide_no)
        for idx, text in enumerate(split_items(item, 3)):
            add_rect(slide, ctx, 4.35, 2.0 + idx * 1.25, 6.7, 0.82, ["main", "support", "accent"][idx])
            add_text(slide, ctx, text, 4.65, 2.22 + idx * 1.25, 5.95, 0.2, role="body", color_role="background")
        return
    content_x = 0.9
    content_w = min(text_width_before_right_visual(ctx, item, content_x, padding=0.45), 8.4 - content_x - 0.45)
    add_text(slide, ctx, visible_content(item)["brief"], content_x, 2.15, content_w, 1.0, role="body", color_role="neutral")
    add_phone_mockup(slide, ctx, 8.4, 1.35, 2.45, 5.45, first_slot(item, {"phone", "mockup", "screen", "image"}), item.slide_no)


def render_process_flow(slide, ctx: BuildContext, item: SlidePlanItem) -> None:
    slide_title(slide, ctx, item)
    items = split_items(item, 4)
    if ctx.is_operator_layout(item):
        for idx, text in enumerate(items):
            y = 1.75 + idx * 1.1
            add_oval(slide, ctx, 2.0, y, 0.78, 0.78, ["main", "support", "accent", "neutral"][idx])
            add_text(slide, ctx, text, 3.15, y + 0.18, 6.8, 0.25, role="body", color_role="neutral")
            if idx < len(items) - 1:
                add_line(slide, ctx, 2.39, y + 0.8, 2.39, y + 1.08, "neutral")
        return
    for idx, text in enumerate(items):
        x = 0.8 + idx * 3.1
        add_rect(slide, ctx, x, 2.8, 2.25, 1.35, ["main", "support", "accent", "neutral"][idx % 4])
        add_text(slide, ctx, text, x + 0.18, 3.15, 1.9, 0.45, role="caption", color_role="background" if idx != 3 else "background", align=PP_ALIGN.CENTER)
        if idx < len(items) - 1:
            add_line(slide, ctx, x + 2.3, 3.48, x + 2.95, 3.48, "neutral")


def render_three_use_cases(slide, ctx: BuildContext, item: SlidePlanItem) -> None:
    slide_title(slide, ctx, item)
    if ctx.is_operator_layout(item):
        for idx, text in enumerate(split_items(item, 3)):
            y = 1.85 + idx * 1.45
            add_rect(slide, ctx, 0.95, y, 10.95, 0.95, "background", ["main", "support", "accent"][idx])
            add_text(slide, ctx, f"USE {idx + 1}", 1.25, y + 0.3, 1.25, 0.18, role="caption", color_role=["main", "support", "accent"][idx])
            add_text(slide, ctx, text, 2.85, y + 0.28, 7.6, 0.24, role="body", color_role="neutral")
        return
    for idx, text in enumerate(split_items(item, 3)):
        x = 1.0 + idx * 4.0
        add_rect(slide, ctx, x, 2.0, 3.1, 3.55, "background", ["main", "support", "accent"][idx])
        add_text(slide, ctx, f"0{idx + 1}", x + 0.25, 2.25, 0.8, 0.35, role="metric", color_role=["main", "support", "accent"][idx])
        add_text(slide, ctx, text, x + 0.3, 3.15, 2.45, 1.1, role="body", color_role="neutral")


def render_market_sizing(slide, ctx: BuildContext, item: SlidePlanItem) -> None:
    bg_slot = first_slot(item, {"background", "texture"})
    add_approved_asset_picture(slide, ctx, bg_slot, item.slide_no, 0.0, 0.0, ctx.packet.slide_size.width_in, ctx.packet.slide_size.height_in, zone_type="background_asset")
    slide_title(slide, ctx, item)
    items = split_items(item, 3)
    if ctx.is_operator_layout(item):
        circles = [(1.4, 2.0, 3.9, "main"), (4.55, 2.35, 3.2, "support"), (7.0, 2.8, 2.45, "accent")]
        for idx, (x, y, size, color) in enumerate(circles):
            add_oval(slide, ctx, x, y, size, size, color)
            add_text(slide, ctx, items[idx], x + 0.35, y + size / 2 - 0.2, size - 0.7, 0.28, role="caption", color_role="background", align=PP_ALIGN.CENTER)
        add_text(slide, ctx, visible_content(item)["brief"], 9.7, 2.65, 2.35, 1.1, role="body", color_role="neutral")
        return
    widths = [8.5, 6.1, 3.7]
    colors = ["main", "support", "accent"]
    for idx, text in enumerate(items):
        add_rect(slide, ctx, 2.4 + idx * 0.6, 2.25 + idx * 0.75, widths[idx], 0.8, colors[idx])
        add_text(slide, ctx, text, 2.65 + idx * 0.6, 2.43 + idx * 0.75, widths[idx] - 0.5, 0.25, role="body", color_role="background", align=PP_ALIGN.CENTER)


def render_traction_metrics(slide, ctx: BuildContext, item: SlidePlanItem) -> None:
    slide_title(slide, ctx, item)
    if ctx.uses_family("luxury_editorial"):
        for idx, text in enumerate(split_items(item, 3)):
            y = 2.15 + idx * 1.0
            add_text(slide, ctx, f"0{idx + 1}", 1.15, y + 0.08, 0.55, 0.18, role="caption", color_role=["neutral", "accent", "main"][idx])
            add_text(slide, ctx, text, 1.95, y, 5.2, 0.38, role="body", color_role="neutral")
            add_line(slide, ctx, 1.15, y + 0.72, 7.25, y + 0.72, ["main", "support", "accent"][idx])
        return
    if ctx.is_operator_layout(item):
        for idx, text in enumerate(split_items(item, 3)):
            y = 1.85 + idx * 1.35
            add_text(slide, ctx, text.split(" ")[0], 1.1, y, 2.2, 0.55, role="metric", color_role=["main", "support", "accent"][idx])
            add_rect(slide, ctx, 3.7, y + 0.18, 5.2 + idx * 0.65, 0.34, ["main", "support", "accent"][idx])
            add_text(slide, ctx, " ".join(text.split(" ")[1:]) or text, 9.8, y + 0.18, 1.8, 0.22, role="caption", color_role="neutral")
        return
    for idx, text in enumerate(split_items(item, 3)):
        x = 1.0 + idx * 4.0
        add_text(slide, ctx, text.split(" ")[0], x, 2.25, 3.1, 0.8, role="metric", color_role=["main", "support", "accent"][idx], align=PP_ALIGN.CENTER)
        add_text(slide, ctx, " ".join(text.split(" ")[1:]) or text, x, 3.2, 3.1, 0.7, role="body", color_role="neutral", align=PP_ALIGN.CENTER)


def render_business_model(slide, ctx: BuildContext, item: SlidePlanItem) -> None:
    slide_title(slide, ctx, item)
    items = split_items(item, 4)
    if ctx.is_operator_layout(item):
        for idx, text in enumerate(items):
            x = 1.0 + idx * 2.9
            add_rect(slide, ctx, x, 2.0, 2.15, 3.75, ["main", "support", "accent", "neutral"][idx])
            add_text(slide, ctx, text, x + 0.22, 3.25, 1.72, 0.72, role="caption", color_role="background", align=PP_ALIGN.CENTER)
        return
    for idx, text in enumerate(items):
        row, col = divmod(idx, 2)
        x = 1.05 + col * 5.85
        y = 2.05 + row * 1.85
        add_rect(slide, ctx, x, y, 5.2, 1.3, "background", "main" if idx % 2 == 0 else "support")
        add_text(slide, ctx, text, x + 0.25, y + 0.38, 4.7, 0.45, role="body", color_role="neutral", align=PP_ALIGN.CENTER)


def render_positioning_matrix(slide, ctx: BuildContext, item: SlidePlanItem) -> None:
    slide_title(slide, ctx, item)
    if ctx.is_operator_layout(item):
        add_rect(slide, ctx, 1.35, 2.0, 4.85, 1.55, "background", "main")
        add_rect(slide, ctx, 6.35, 2.0, 4.85, 1.55, "background", "support")
        add_rect(slide, ctx, 1.35, 3.8, 4.85, 1.55, "background", "accent")
        add_rect(slide, ctx, 6.35, 3.8, 4.85, 1.55, "main")
        add_text(slide, ctx, split_items(item, 1)[0], 6.65, 4.35, 4.2, 0.32, role="body", color_role="background", align=PP_ALIGN.CENTER)
        return
    add_line(slide, ctx, 2.1, 5.55, 11.2, 5.55, "neutral")
    add_line(slide, ctx, 6.65, 1.8, 6.65, 6.25, "neutral")
    add_rect(slide, ctx, 7.25, 2.2, 2.65, 1.05, "main")
    add_text(slide, ctx, split_items(item, 1)[0], 7.42, 2.48, 2.3, 0.35, role="caption", color_role="background", align=PP_ALIGN.CENTER)
    add_text(slide, ctx, "Social", 10.4, 5.75, 1.2, 0.25, role="caption", color_role="neutral")
    add_text(slide, ctx, "AI native", 6.8, 1.55, 1.2, 0.25, role="caption", color_role="neutral")


def render_flywheel(slide, ctx: BuildContext, item: SlidePlanItem) -> None:
    slide_title(slide, ctx, item)
    items = split_items(item, 4)
    if ctx.is_operator_layout(item):
        add_oval(slide, ctx, 5.25, 2.85, 2.35, 2.35, "background", "main")
        add_text(slide, ctx, "Loop", 5.75, 3.75, 1.3, 0.22, role="body", color_role="main", align=PP_ALIGN.CENTER)
        coords = [(5.45, 1.65), (8.55, 3.2), (5.45, 5.45), (2.3, 3.2)]
        for idx, (x, y) in enumerate(coords):
            add_rect(slide, ctx, x, y, 1.95, 0.66, ["main", "support", "accent", "neutral"][idx])
            add_text(slide, ctx, items[idx], x + 0.12, y + 0.18, 1.7, 0.16, role="caption", color_role="background", align=PP_ALIGN.CENTER)
            add_line(slide, ctx, 6.42, 4.0, x + 0.95, y + 0.33, ["main", "support", "accent", "neutral"][idx])
        return
    coords = [(5.6, 2.0), (8.0, 3.55), (5.6, 5.1), (3.2, 3.55)]
    for idx, (x, y) in enumerate(coords):
        add_rect(slide, ctx, x, y, 2.0, 0.85, ["main", "support", "accent", "neutral"][idx])
        add_text(slide, ctx, items[idx], x + 0.12, y + 0.24, 1.75, 0.25, role="caption", color_role="background", align=PP_ALIGN.CENTER)
    add_line(slide, ctx, 7.6, 2.55, 8.0, 3.55, "main")
    add_line(slide, ctx, 8.0, 4.2, 7.6, 5.1, "support")
    add_line(slide, ctx, 5.6, 5.55, 5.2, 4.2, "accent")
    add_line(slide, ctx, 5.2, 2.45, 5.6, 2.0, "neutral")


def render_roadmap(slide, ctx: BuildContext, item: SlidePlanItem) -> None:
    slide_title(slide, ctx, item)
    if ctx.is_operator_layout(item):
        for idx, text in enumerate(split_items(item, 4)):
            y = 1.75 + idx * 1.1
            add_rect(slide, ctx, 1.3 + idx * 0.28, y, 8.8 - idx * 0.55, 0.72, ["main", "support", "accent", "neutral"][idx])
            add_text(slide, ctx, text, 1.65 + idx * 0.28, y + 0.21, 5.2, 0.18, role="caption", color_role="background")
        return
    add_line(slide, ctx, 1.2, 3.75, 12.0, 3.75, "neutral")
    for idx, text in enumerate(split_items(item, 4)):
        x = 1.4 + idx * 2.8
        add_rect(slide, ctx, x, 3.25, 0.52, 0.52, ["main", "support", "accent", "neutral"][idx])
        add_text(slide, ctx, text, x - 0.2, 4.15, 2.2, 0.8, role="caption", color_role="neutral", align=PP_ALIGN.CENTER)


def render_financial_plan(slide, ctx: BuildContext, item: SlidePlanItem) -> None:
    slide_title(slide, ctx, item)
    items = split_items(item, 4)
    if ctx.is_operator_layout(item):
        for idx, text in enumerate(items):
            y = 2.0 + idx * 0.86
            add_text(slide, ctx, text, 1.25, y, 2.6, 0.24, role="caption", color_role="neutral")
            add_rect(slide, ctx, 4.15, y, 3.6 + idx * 1.05, 0.38, ["main", "support", "accent", "neutral"][idx])
        return
    max_h = 2.6
    for idx, text in enumerate(items):
        h = max_h * (idx + 1) / len(items)
        x = 2.0 + idx * 2.25
        add_rect(slide, ctx, x, 5.55 - h, 1.3, h, ["main", "support", "accent", "neutral"][idx])
        add_text(slide, ctx, text, x - 0.35, 5.85, 2.0, 0.35, role="caption", color_role="neutral", align=PP_ALIGN.CENTER)


def render_ask(slide, ctx: BuildContext, item: SlidePlanItem) -> None:
    if ctx.uses_family("luxury_editorial"):
        add_line(slide, ctx, 1.05, 1.05, 10.95, 1.05, "accent")
        add_text(slide, ctx, visible_content(item)["title"], 1.05, 1.55, 6.8, 0.95, role="slide_title", color_role="neutral")
        add_text(slide, ctx, visible_content(item)["subtitle"], 1.08, 2.9, 6.05, 0.65, role="body", color_role="neutral")
        for idx, text in enumerate(split_items(item, 3)):
            y = 4.2 + idx * 0.55
            add_line(slide, ctx, 1.08, y, 6.8, y, ["main", "support", "accent"][idx])
            add_text(slide, ctx, text, 1.12, y + 0.08, 5.4, 0.22, role="caption", color_role="neutral")
        add_rect(slide, ctx, 8.55, 1.65, 1.85, 3.85, "background", "accent")
        add_rect(slide, ctx, 10.62, 2.2, 0.58, 2.75, "accent")
        ctx.add_visual_zone(item.slide_no, 8.35, 1.45, 3.1, 4.3, zone_type="product_or_material_stage", source="luxury_editorial_detail")
        return
    if ctx.uses_family("internal_strategy_report"):
        add_rect(slide, ctx, 0.95, 1.35, 10.9, 0.95, "background", "main")
        add_text(slide, ctx, visible_content(item)["title"], 1.25, 1.68, 10.3, 0.24, role="slide_title", color_role="neutral", align=PP_ALIGN.CENTER)
        for idx, text in enumerate(split_items(item, 3)):
            y = 3.0 + idx * 0.9
            add_rect(slide, ctx, 1.35, y, 9.9, 0.55, ["main", "support", "accent"][idx])
            add_text(slide, ctx, text, 1.65, y + 0.17, 9.25, 0.14, role="caption", color_role="background", align=PP_ALIGN.CENTER)
        add_text(slide, ctx, visible_content(item)["subtitle"], 2.2, 6.0, 7.9, 0.32, role="body", color_role="neutral", align=PP_ALIGN.CENTER)
        return
    if ctx.is_operator_layout(item):
        add_rect(slide, ctx, 0.0, 0.0, ctx.packet.slide_size.width_in, 2.35, "main")
        add_text(slide, ctx, visible_content(item)["title"], 1.15, 0.92, 10.7, 0.72, role="slide_title", color_role="background", align=PP_ALIGN.CENTER)
        for idx, text in enumerate(split_items(item, 3)):
            add_rect(slide, ctx, 1.35 + idx * 3.6, 3.45, 2.85, 1.25, ["support", "accent", "neutral"][idx])
            add_text(slide, ctx, text, 1.55 + idx * 3.6, 3.88, 2.45, 0.22, role="caption", color_role="background", align=PP_ALIGN.CENTER)
        add_text(slide, ctx, visible_content(item)["subtitle"], 3.3, 5.55, 6.6, 0.45, role="body", color_role="neutral", align=PP_ALIGN.CENTER)
        return
    title_x = 1.2
    title_w = text_width_before_right_visual(ctx, item, title_x, padding=0.35)
    align = PP_ALIGN.LEFT if title_w < 8.0 else PP_ALIGN.CENTER
    add_text(slide, ctx, visible_content(item)["title"], title_x, 1.35, title_w, 1.1, role="hero_title", color_role="neutral", align=align)
    card_w = min(4.55, max(2.8, title_w - 1.0))
    card_x = title_x + max(0.0, (title_w - card_w) / 2)
    add_rect(slide, ctx, card_x, 3.15, card_w, 1.3, "main")
    add_text(slide, ctx, visible_content(item)["subtitle"], card_x + 0.25, 3.52, max(1.0, card_w - 0.5), 0.45, role="body", color_role="background", align=PP_ALIGN.CENTER)
    for idx, text in enumerate(split_items(item, 3)):
        add_text(slide, ctx, text, 1.4 + idx * 3.75, 5.35, 3.1, 0.55, role="caption", color_role="neutral", align=PP_ALIGN.CENTER)


RENDERERS = {
    "cover": render_cover,
    "big_thesis": render_big_thesis,
    "three_cards": render_three_cards,
    "behavior_shift": render_behavior_shift,
    "product_mockup": render_product_mockup,
    "process_flow": render_process_flow,
    "three_use_cases": render_three_use_cases,
    "market_sizing": render_market_sizing,
    "traction_metrics": render_traction_metrics,
    "business_model": render_business_model,
    "positioning_matrix": render_positioning_matrix,
    "flywheel": render_flywheel,
    "roadmap": render_roadmap,
    "financial_plan": render_financial_plan,
    "ask": render_ask,
}


def apply_background(slide, ctx: BuildContext) -> None:
    bg_role = "background" if "background" in ctx.palette else "neutral"
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = ctx.color(bg_role)
    if any(strategy.get("renderer_variant") in {"dense", "operator", "dashboard", "workbook"} for strategy in ctx.layout_strategy_by_slide.values()):
        add_rect(slide, ctx, 0, 0, 0.24, ctx.packet.slide_size.height_in, "main")


def render_packet_pptx(
    packet: GuidePacket,
    project_dir: Path,
    mode: str,
    deck_plan: dict[str, Any],
    variant: str | None = None,
) -> tuple[Path, BuildContext]:
    ctx = BuildContext(packet, project_dir, mode, deck_plan, variant)
    prs = Presentation()
    prs.slide_width = Inches(packet.slide_size.width_in)
    prs.slide_height = Inches(packet.slide_size.height_in)
    blank = prs.slide_layouts[6]

    for item in packet.slide_plan.slides:
        ctx.current_slide_no = item.slide_no
        slide = prs.slides.add_slide(blank)
        apply_background(slide, ctx)
        composition_profile = apply_recipe_canvas(slide, ctx, item)
        renderer = RENDERERS.get(item.layout_archetype)
        if renderer is None:
            raise ValueError(f"Native renderer is missing for layout archetype: {item.layout_archetype}")
        renderer(slide, ctx, item)
        record_unhandled_asset_slots(ctx, item)
        add_footer(slide, ctx, item.slide_no)
        ctx.native_renderer_events.append(
            {
                "slide_no": item.slide_no,
                "layout_archetype": item.layout_archetype,
                "renderer": renderer.__name__,
                "layout_recipe": ctx.layout_recipe(item),
                "recipe_family": composition_profile["family"],
                "recipe_composition_style": composition_profile["style"],
                "native_powerpoint_objects": True,
            }
        )
    ctx.current_slide_no = None

    pptx_path = project_dir / "generated.pptx"
    project_dir.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_path))
    return pptx_path, ctx


def synthetic_previews(packet: GuidePacket, output_dir: Path) -> list[Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    paths: list[Path] = []
    width, height = 1280, 720
    try:
        font_big = ImageFont.truetype("arial.ttf", 44)
        font_small = ImageFont.truetype("arial.ttf", 24)
    except OSError:
        font_big = ImageFont.load_default()
        font_small = ImageFont.load_default()
    bg = next((c.hex for c in packet.palette.colors if c.role == "background"), "#F7F8FC")
    ink = next((c.hex for c in packet.palette.colors if c.role == "neutral"), "#101828")
    accent = next((c.hex for c in packet.palette.colors if c.role == "main"), "#315CFF")
    for slide in packet.slide_plan.slides:
        image = Image.new("RGB", (width, height), bg)
        draw = ImageDraw.Draw(image)
        draw.rectangle([0, 0, 18, height], fill=accent)
        draw.text((70, 90), visible_content(slide)["title"][:70], fill=ink, font=font_big)
        draw.text((72, 180), slide.layout_archetype, fill=accent, font=font_small)
        for idx, item in enumerate(split_items(slide, 4)):
            draw.rounded_rectangle([90 + idx * 285, 320, 320 + idx * 285, 470], radius=14, outline=accent, width=3)
            draw.text((110 + idx * 285, 350), item[:24], fill=ink, font=font_small)
        path = output_dir / f"slide_{slide.slide_no:02d}.png"
        image.save(path)
        paths.append(path)
    return paths


def build_preview_contact_sheet(paths: list[Path], output_path: Path) -> Path | None:
    if not paths:
        return None
    thumbs: list[Image.Image] = []
    for path in paths:
        with Image.open(path).convert("RGB") as image:
            image.thumbnail((320, 180))
            canvas = Image.new("RGB", (320, 180), "white")
            canvas.paste(image, ((320 - image.width) // 2, (180 - image.height) // 2))
            thumbs.append(canvas)
    cols = min(4, len(thumbs))
    rows = (len(thumbs) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * 320, rows * 180), "white")
    for idx, thumb in enumerate(thumbs):
        sheet.paste(thumb, ((idx % cols) * 320, (idx // cols) * 180))
    output_path.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(output_path)
    return output_path


def render_previews(pptx_path: Path, packet: GuidePacket, project_dir: Path) -> dict[str, Any]:
    preview_dir = project_dir / "previews"
    cleanup_warnings: list[dict[str, Any]] = []

    def cleanup_preview_dir(reason: str) -> None:
        try:
            if preview_dir.exists():
                shutil.rmtree(preview_dir)
        except OSError as exc:
            cleanup_warnings.append(
                {
                    "reason": reason,
                    "path": safe_rel(preview_dir),
                    "error_type": exc.__class__.__name__,
                    "message": sanitize_report_text(str(exc)),
                }
            )
        preview_dir.mkdir(parents=True, exist_ok=True)
        stale_pdf = preview_dir / "_pdf"
        try:
            if stale_pdf.exists():
                shutil.rmtree(stale_pdf)
        except OSError as exc:
            cleanup_warnings.append(
                {
                    "reason": f"{reason}_stale_pdf_cleanup",
                    "path": safe_rel(stale_pdf),
                    "error_type": exc.__class__.__name__,
                    "message": sanitize_report_text(str(exc)),
                }
            )

    cleanup_preview_dir("initial")
    last_error: Exception | None = None
    attempt_details: list[dict[str, Any]] = []
    for attempt in range(3):
        try:
            from scripts.validate_visual_smoke import render_pptx

            paths = render_pptx(pptx_path, preview_dir)
            expected = packet.guide_identity.slide_count
            if len(paths) != expected:
                raise ValueError(f"preview slide count mismatch: rendered={len(paths)}, expected={expected}")
            attempt_details.append(
                {
                    "attempt": attempt + 1,
                    "status": "pass",
                    "rendered_slide_count": len(paths),
                    "expected_slide_count": expected,
                }
            )
            contact_sheet = build_preview_contact_sheet(paths, project_dir / "preview-contact-sheet.png")
            return {
                "status": "rendered",
                "method": "libreoffice_pdf_png",
                "paths": [safe_rel(path) for path in paths],
                "contact_sheet": safe_rel(contact_sheet) if contact_sheet else None,
                "rendered_slide_count": len(paths),
                "fallback_used": False,
                "render_attempts": attempt + 1,
                "render_attempt_details": attempt_details,
                "cleanup_warnings": cleanup_warnings,
            }
        except Exception as exc:
            last_error = exc
            details = getattr(exc, "details", None)
            attempt_details.append(
                {
                    "attempt": attempt + 1,
                    "status": "fail",
                    "error_type": exc.__class__.__name__,
                    "message": sanitize_report_text(str(exc)),
                    "details": sanitize_report_text(details) if isinstance(details, dict) else None,
                }
            )
            if attempt < 2:
                time.sleep(0.75 * (attempt + 1))
                cleanup_preview_dir(f"retry_{attempt + 2}")
    cleanup_preview_dir("synthetic_fallback")
    paths = synthetic_previews(packet, preview_dir)
    contact_sheet = build_preview_contact_sheet(paths, project_dir / "preview-contact-sheet.png")
    reason = f"PPTX preview rendering failed ({last_error.__class__.__name__ if last_error else 'unknown'}); local command details redacted."
    return {
        "status": "blocked_render_failed",
        "method": "synthetic_layout_preview",
        "paths": [safe_rel(path) for path in paths],
        "contact_sheet": safe_rel(contact_sheet) if contact_sheet else None,
        "rendered_slide_count": len(paths),
        "fallback_used": True,
        "validation_blocker": True,
        "fallback_reason": reason,
        "render_attempts": 3,
        "render_attempt_details": attempt_details,
        "cleanup_warnings": cleanup_warnings,
    }


def generate_html_guide(packet: GuidePacket, project_dir: Path, mode: str, requested: bool) -> dict[str, Any]:
    if not requested:
        return {
            "generated": False,
            "reason": "HTML guide skipped until explicitly requested or final build is approved.",
            "file_path": None,
            "rendered_url": None,
            "viewport_size": None,
            "screenshot_path": None,
            "javascript_runtime_errors": [],
            "button_control_count": 0,
            "html_screenshot_used_in_pptx": False,
        }
    html_dir = project_dir / "html"
    html_dir.mkdir(parents=True, exist_ok=True)
    path = html_dir / "guide.html"
    palette = "".join(
        f"<li><span style='background:{html.escape(color.hex)}'></span>{html.escape(color.role)} {html.escape(color.hex)}</li>"
        for color in packet.palette.colors
    )
    slides = "".join(
        f"<section><h2>{slide.slide_no}. {html.escape(slide.layout_archetype)}</h2>"
        f"<p>{html.escape(slide.content_brief)}</p></section>"
        for slide in packet.slide_plan.slides
    )
    path.write_text(
        "<!doctype html><html><head><meta charset='utf-8'><title>Guide</title>"
        "<style>body{font-family:Arial,sans-serif;margin:32px;background:#f7f8fc;color:#101828}"
        "section{border:1px solid #d0d5dd;padding:16px;margin:12px 0;background:white}"
        "span{display:inline-block;width:16px;height:16px;margin-right:8px;vertical-align:middle}</style></head><body>"
        f"<h1>{html.escape(packet.guide_identity.project_name)}</h1><ul>{palette}</ul>{slides}</body></html>",
        encoding="utf-8",
    )
    button_count = len(re.findall(r"<button\b", path.read_text(encoding="utf-8"), flags=re.IGNORECASE))
    return {
        "generated": True,
        "reason": "Assistant mode guide review evidence.",
        "file_path": safe_rel(path),
        "rendered_url": None,
        "viewport_size": {"width": 1440, "height": 960},
        "screenshot_path": None,
        "javascript_runtime_errors": [],
        "button_control_count": button_count,
        "html_screenshot_used_in_pptx": False,
    }


def write_planning_artifacts(
    packet: GuidePacket,
    project_dir: Path,
    source_mode: str,
    variant_strategy: str = "investor_open",
) -> tuple[dict[str, Any], dict[str, Any]]:
    deck_plan = generate_deck_plan(packet, source_mode, variant_strategy)
    renderer_contract = generate_renderer_contract(packet, source_mode, deck_plan)
    json_dump(project_dir / "deck-plan.json", deck_plan)
    json_dump(project_dir / "renderer-contract.json", renderer_contract)
    json_dump(project_dir / "asset-slot-plan.json", {"asset_slot_requirements": renderer_contract["asset_slot_requirements"]})
    json_dump(
        project_dir / "qa-plan.json",
        {
            "rules": [rule.model_dump(mode="json") for rule in packet.qa_rules],
            "required_reports": [
                "guide-compliance-report.json",
                "final-qa.json",
                "used-assets-report.json",
                "html-guide-render-report.json",
            ],
        },
    )
    brief = [
        f"# {packet.guide_identity.project_name} Draft Design Brief",
        "",
        f"- Objective: {packet.project_brief.objective}",
        f"- Audience: {', '.join(packet.project_brief.audience)}",
        f"- Tone: {', '.join(packet.project_brief.tone)}",
        f"- Slides: {packet.guide_identity.slide_count}",
        "",
        "This brief is machine-facing planning evidence. HTML guide content is not used as PPTX slide content.",
    ]
    (project_dir / "draft_design_brief.md").write_text("\n".join(brief) + "\n", encoding="utf-8")
    return deck_plan, renderer_contract


def guide_compliance_report(packet: GuidePacket, source_path: Path, source_mode: str) -> dict[str, Any]:
    return {
        "status": "pass",
        "validated_at": utc_now(),
        "schema_path": safe_rel(GUIDE_SCHEMA_PATH),
        "source_path": safe_rel(source_path),
        "source_mode": source_mode,
        "normalized_objects": sorted(normalized_packet_objects(packet)),
        "slide_count_matches": packet.guide_identity.slide_count == len(packet.slide_plan.slides),
        "layout_archetypes": sorted({item.id for item in packet.layout_archetypes}),
        "unknown_layout_archetypes": [],
        "blocked": False,
        "public_safety": packet.public_safety.model_dump(mode="json"),
    }


def bounds_intersect(a: dict[str, float], b: dict[str, float], padding: float = 0.0) -> bool:
    ax1 = a["x"] + padding
    ay1 = a["y"] + padding
    ax2 = a["x"] + a["w"] - padding
    ay2 = a["y"] + a["h"] - padding
    bx1 = b["x"]
    by1 = b["y"]
    bx2 = b["x"] + b["w"]
    by2 = b["y"] + b["h"]
    return ax1 < bx2 and ax2 > bx1 and ay1 < by2 and ay2 > by1


def intersection_area(a: dict[str, float], b: dict[str, float]) -> float:
    x1 = max(a["x"], b["x"])
    y1 = max(a["y"], b["y"])
    x2 = min(a["x"] + a["w"], b["x"] + b["w"])
    y2 = min(a["y"] + a["h"], b["y"] + b["h"])
    return max(0.0, x2 - x1) * max(0.0, y2 - y1)


def qa_bounds(box: dict[str, Any]) -> dict[str, float]:
    return box.get("estimated_text_bounds_in") or box["bounds_in"]


def visual_collision_scan(packet: GuidePacket, ctx: BuildContext) -> dict[str, Any]:
    collision_zone_types = {
        "right_edge_decorative",
        "right_edge_asset",
        "right_mockup_or_gallery",
        "right_mockup_or_grid",
        "right_mockup_or_map",
        "phone_or_mockup",
        "product_or_material_stage",
    }
    issues: list[dict[str, Any]] = []
    content_boxes = [box for box in ctx.text_boxes if box.get("text_role") == "content"]
    zones = [zone for zone in ctx.visual_zones if zone.get("zone_type") in collision_zone_types]
    for box in content_boxes:
        for zone in zones:
            if box["slide_no"] != zone["slide_no"]:
                continue
            box_bounds = qa_bounds(box)
            if bounds_intersect(box_bounds, zone["bounds_in"], padding=0.02):
                issues.append(
                    {
                        "type": "text_visual_collision",
                        "slide_no": box["slide_no"],
                        "text_excerpt": box.get("text_excerpt"),
                        "visual_zone_type": zone.get("zone_type"),
                        "visual_zone_source": zone.get("source"),
                        "text_bounds_in": box_bounds,
                        "visual_bounds_in": zone["bounds_in"],
                        "severity": "fail",
                    }
                )
    major_roles = {"hero_title", "slide_title", "body", "metric"}
    major_content_boxes = [
        box
        for box in content_boxes
        if box.get("typography_role") in major_roles and str(box.get("text_excerpt") or "").strip()
    ]
    checked_text_pairs = 0
    for idx, box in enumerate(major_content_boxes):
        for other in major_content_boxes[idx + 1:]:
            if box["slide_no"] != other["slide_no"]:
                continue
            checked_text_pairs += 1
            a = qa_bounds(box)
            b = qa_bounds(other)
            area = intersection_area(a, b)
            if area <= 0.015:
                continue
            smaller = min(max(0.001, a["w"] * a["h"]), max(0.001, b["w"] * b["h"]))
            if area / smaller < 0.06:
                continue
            issues.append(
                {
                    "type": "text_text_collision",
                    "slide_no": box["slide_no"],
                    "text_excerpt": box.get("text_excerpt"),
                    "other_text_excerpt": other.get("text_excerpt"),
                    "text_bounds_in": a,
                    "other_text_bounds_in": b,
                    "overlap_area_in2": round(area, 4),
                    "severity": "fail",
                }
            )
    area = packet.safe_area
    for box in content_boxes:
        b = qa_bounds(box)
        if b["x"] < area.left_in - 0.05 or b["y"] < area.top_in - 0.05 or b["x"] + b["w"] > area.right_in + 0.05 or b["y"] + b["h"] > area.bottom_in + 0.05:
            issues.append(
                {
                    "type": "text_outside_safe_area",
                    "slide_no": box["slide_no"],
                    "text_excerpt": box.get("text_excerpt"),
                    "text_bounds_in": b,
                    "safe_area": area.model_dump(mode="json"),
                    "severity": "warn",
                }
            )
    status = "pass" if not any(issue["severity"] == "fail" for issue in issues) else "fail"
    if status == "pass" and issues:
        status = "warning"
    return {
        "status": status,
        "method": "native text bounds against reserved decorative/mockup/image zones and estimated text-to-text overlap",
        "checked_text_boxes": len(content_boxes),
        "checked_visual_zones": len(zones),
        "checked_text_pairs": checked_text_pairs,
        "issues": issues,
    }


def color_spread_scan(packet: GuidePacket, ctx: BuildContext) -> dict[str, Any]:
    resolution = infer_palette_resolution(packet)
    preset = palette_preset(packet.palette.palette_id)
    policy = preset.get("policy", "guide_defined")
    expressive_allowed = policy == "expressive" or resolution["source_type"] == "guide_packet_palette"
    accent_roles = {"main", "support", "accent"}
    slide_area = packet.slide_size.width_in * packet.slide_size.height_in
    issues: list[dict[str, Any]] = []
    by_slide: dict[int, dict[str, float]] = {}
    for surface in ctx.color_surfaces:
        role = surface.get("role")
        if role not in accent_roles:
            continue
        by_slide.setdefault(int(surface["slide_no"]), {}).setdefault(role, 0.0)
        by_slide[int(surface["slide_no"])][role] += float(surface.get("area_in2") or 0.0)
    for slide_no, role_areas in sorted(by_slide.items()):
        large_roles = [
            role
            for role, area in role_areas.items()
            if area / max(slide_area, 1.0) >= 0.08
        ]
        total_accent_ratio = sum(role_areas.values()) / max(slide_area, 1.0)
        if not expressive_allowed and len(large_roles) >= 3 and total_accent_ratio >= 0.22:
            issues.append(
                {
                    "type": "excessive_accent_spread",
                    "slide_no": slide_no,
                    "large_accent_roles": sorted(large_roles),
                    "accent_surface_ratio": round(total_accent_ratio, 3),
                    "palette_policy": policy,
                    "severity": "warn",
                }
            )
    return {
        "status": "warning" if issues else "pass",
        "method": "native shape surface role spread scan",
        "palette_policy": policy,
        "expressive_palette_allowed": expressive_allowed,
        "issues": issues,
    }


def relative_luminance(hex_value: str) -> float:
    channels = [value / 255 for value in hex_to_rgb_list(hex_value)]
    linear = []
    for channel in channels:
        if channel <= 0.03928:
            linear.append(channel / 12.92)
        else:
            linear.append(((channel + 0.055) / 1.055) ** 2.4)
    return 0.2126 * linear[0] + 0.7152 * linear[1] + 0.0722 * linear[2]


def contrast_ratio(foreground_hex: str, background_hex: str) -> float:
    foreground = relative_luminance(foreground_hex)
    background = relative_luminance(background_hex)
    lighter = max(foreground, background)
    darker = min(foreground, background)
    return (lighter + 0.05) / (darker + 0.05)


def palette_contrast_scan(packet: GuidePacket) -> dict[str, Any]:
    colors = {color.role: color.hex for color in packet.palette.colors}
    required_roles = ["main", "support", "accent", "background", "neutral"]
    issues: list[dict[str, Any]] = []
    for role in required_roles:
        if role not in colors:
            issues.append({"type": "missing_role_color", "role": role, "severity": "fail"})
    if issues:
        return {
            "status": "fail",
            "method": "WCAG role-pair palette contrast scan",
            "checked_pairs": [],
            "issues": issues,
        }

    pair_specs = [
        ("neutral", "background", 4.5, "fail", "default body text on deck background"),
        ("main", "background", 4.5, "fail", "primary accent text or metric on deck background"),
        ("background", "main", 3.0, "fail", "short label text on primary accent fills"),
        ("background", "support", 3.0, "warn", "short label text on support accent fills"),
        ("background", "accent", 3.0, "warn", "short label text on accent fills"),
    ]
    checked_pairs = []
    for foreground_role, background_role, minimum, severity, reason in pair_specs:
        ratio = round(contrast_ratio(colors[foreground_role], colors[background_role]), 2)
        checked_pairs.append(
            {
                "key": f"{foreground_role}_on_{background_role}",
                "foreground_role": foreground_role,
                "background_role": background_role,
                "ratio": ratio,
                "minimum_ratio": minimum,
            }
        )
        if ratio < minimum:
            issues.append(
                {
                    "type": "low_contrast",
                    "key": f"{foreground_role}_on_{background_role}",
                    "foreground_role": foreground_role,
                    "background_role": background_role,
                    "ratio": ratio,
                    "minimum_ratio": minimum,
                    "reason": reason,
                    "severity": severity,
                }
            )
    status = "pass"
    if any(issue["severity"] == "fail" for issue in issues):
        status = "fail"
    elif issues:
        status = "warning"
    return {
        "status": status,
        "method": "WCAG role-pair palette contrast scan",
        "checked_pairs": checked_pairs,
        "issues": issues,
    }


def asset_system_status(internal_status: str) -> str:
    return {"warning": "warn"}.get(internal_status, internal_status if internal_status in {"pass", "fail", "blocked"} else "fail")


def palette_asset_ref_metadata(asset_ref: str | None) -> dict[str, str]:
    if not asset_ref or not asset_ref.startswith("palette-resolution:"):
        return {}
    parts = asset_ref.split(":", 3)
    if len(parts) < 4:
        return {}
    metadata: dict[str, str] = {}
    for item in parts[3].split("|"):
        key, separator, value = item.partition("=")
        if separator and key:
            metadata[key] = value
    return metadata


def approved_palette_evidence(packet: GuidePacket) -> dict[str, Any] | None:
    role_colors: dict[str, dict[str, Any]] = {}
    for color in packet.palette.colors:
        if color.role not in {"main", "support", "accent", "background", "neutral"}:
            continue
        role_colors[color.role] = {
            "name": color.name,
            "hex": color.hex,
            "rgb": color.rgb,
            "usage": color.usage,
        }
    if set(role_colors) != {"main", "support", "accent", "background", "neutral"}:
        return None
    return {
        "palette_id": packet.palette.palette_id,
        "palette_name": packet.palette.palette_name,
        "colors": {role: role_colors[role] for role in ["main", "support", "accent", "background", "neutral"]},
    }


def contrast_level(ratio: Any) -> str:
    try:
        value = float(ratio)
    except (TypeError, ValueError):
        return "fail"
    if value >= 7:
        return "AAA"
    if value >= 4.5:
        return "AA"
    return "fail"


def palette_resolution_reason(palette_resolution: dict[str, Any]) -> str:
    source_type = palette_resolution.get("source_type")
    if palette_resolution.get("fallback_reason"):
        return str(palette_resolution["fallback_reason"])
    if source_type == "asset_system_palette":
        return "Asset-system package palette was used as the deck-level palette."
    if source_type == "guide_packet_palette":
        return "Guide packet palette was used as the deck-level palette."
    if source_type == "intent_matched_palette":
        return "PPT maker selected a local palette from the request intent and routing profile."
    if source_type == "fallback_default_palette":
        return "No higher-priority palette source was available; fallback default palette was used."
    return "Palette source was resolved by the PPT maker."


def palette_resolution_evidence(
    packet: GuidePacket,
    palette_resolution: dict[str, Any],
    contrast_scan: dict[str, Any],
) -> dict[str, Any]:
    source_type = str(palette_resolution.get("source_type") or "fallback_default_palette")
    metadata = palette_asset_ref_metadata(palette_resolution.get("asset_ref"))
    evidence: dict[str, Any] = {
        "source_type": source_type,
        "palette_metadata_available": source_type != "fallback_default_palette",
        "selected_palette_id": palette_resolution.get("selected_palette_id"),
        "palette_policy": str(palette_resolution.get("policy") or metadata.get("palette_policy") or palette_preset(packet.palette.palette_id).get("policy", "guide_defined")),
        "reason": palette_resolution_reason(palette_resolution),
        "notes": [
            "dominant_colors are treated only as asset-level visual traits, not deck-level palette metadata.",
            "Structured palette evidence belongs in validation_results[].evidence.palette_resolution.",
        ],
    }
    if metadata.get("selection_confidence"):
        try:
            evidence["selection_confidence"] = float(metadata["selection_confidence"])
        except ValueError:
            pass
    approved_palette = approved_palette_evidence(packet)
    if approved_palette:
        evidence["approved_palette"] = approved_palette
    evidence["text_color_recommendation"] = {
        "default_text_role": "neutral",
        "on_main": "background",
        "on_support": "background",
        "on_accent": "neutral",
        "on_background": "neutral",
    }
    wcag_contrast: dict[str, dict[str, Any]] = {}
    for item in contrast_scan.get("checked_pairs", []):
        key = str(item.get("key") or f"{item.get('foreground_role')}_on_{item.get('background_role')}")
        wcag_contrast[key] = {
            "foreground_role": str(item.get("foreground_role")),
            "background_role": str(item.get("background_role")),
            "ratio": float(item.get("ratio") or 1),
            "level": contrast_level(item.get("ratio")),
        }
    if wcag_contrast:
        evidence["wcag_contrast"] = wcag_contrast
    return evidence


def layout_recipes_used(ctx: BuildContext) -> list[dict[str, Any]]:
    return [
        {
            "layout_recipe_id": event.get("layout_recipe") or "native_layout",
            "module_id": event.get("layout_archetype") or "native_renderer",
            "composition_id": event.get("recipe_composition_style") or event.get("renderer") or "native",
            "rendered_slide": int(event["slide_no"]),
        }
        for event in ctx.native_renderer_events
    ]


def approved_package_assets_for_asset_system(ctx: BuildContext) -> list[dict[str, Any]]:
    grouped: dict[tuple[str, str, str], set[int]] = {}
    for event in ctx.asset_events:
        if event.get("fallback_used") or event.get("resolution") != "approved_package_asset":
            continue
        relative_package_path = public_package_path_or_none(event.get("relative_package_path"))
        if not relative_package_path:
            continue
        key = (
            str(event.get("package_manifest_id") or "package-manifest:unknown"),
            relative_package_path,
            str(event.get("sha256") or ""),
        )
        if not key[2]:
            continue
        grouped.setdefault(key, set()).add(int(event.get("slide_no") or 1))
    return [
        {
            "package_manifest_id": key[0],
            "relative_package_path": key[1],
            "sha256": key[2],
            "used_on_slides": sorted(slides),
        }
        for key, slides in sorted(grouped.items())
    ]


def approved_ref_to_declared_slot_ids(packet: GuidePacket) -> dict[str, set[str]]:
    slots_by_ref: dict[str, set[str]] = {}
    for reference in packet.approved_asset_references:
        extra = reference.model_extra or {}
        if extra.get("approved_asset_ref_declared_by_package_slot") is not True:
            continue
        slot_id = str(extra.get("package_slot_id") or reference.slot_id or "").strip()
        approved_ref = normalize_package_relative_path(str(extra.get("approved_asset_ref") or reference.asset_ref or ""))
        if not slot_id or not approved_ref:
            continue
        slots_by_ref.setdefault(approved_ref, set()).add(slot_id)
    return slots_by_ref


def approved_events_for_asset(ctx: BuildContext, asset: dict[str, Any]) -> list[dict[str, Any]]:
    asset_path = normalize_package_relative_path(str(asset.get("relative_package_path") or ""))
    if not asset_path:
        return []
    return [
        event
        for event in ctx.asset_events
        if not event.get("fallback_used")
        and event.get("resolution") == "approved_package_asset"
        and normalize_package_relative_path(str(event.get("relative_package_path") or "")) == asset_path
        and str(event.get("package_manifest_id") or "package-manifest:unknown") == str(asset.get("package_manifest_id") or "package-manifest:unknown")
        and str(event.get("sha256") or "") == str(asset.get("sha256") or "")
    ]


def inserted_asset_summary(asset: dict[str, Any], events: list[dict[str, Any]], slot_id: str | None = None) -> dict[str, Any]:
    role = next((event.get("insertion_role") for event in events if event.get("insertion_role")), "approved_package_file")
    used_on_slides = sorted({int(event.get("slide_no") or 1) for event in events}) or list(asset.get("used_on_slides") or [])
    summary = {
        **asset,
        "used_on_slides": used_on_slides,
        "source_type": "approved_package_file",
        "file_backed": True,
        "insertion_role": role,
    }
    if slot_id:
        summary["slot_id"] = slot_id
    return summary


def fallback_summary(event: dict[str, Any]) -> dict[str, Any]:
    allowed_fields = [
        "slide_no",
        "slot_id",
        "slot_type",
        "required",
        "resolution",
        "source_asset_ref",
        "path",
        "needed_asset_intent",
        "fallback_used",
        "fallback_type",
        "fallback_reason",
        "reason",
        "crop_or_mask_policy",
        "native_handling",
        "impact",
    ]
    return {
        field: event[field]
        for field in allowed_fields
        if field in event and event[field] is not None
    }


def missing_asset_expectations(packet: GuidePacket, ctx: BuildContext) -> list[dict[str, Any]]:
    expectations: list[dict[str, Any]] = []
    events_by_slot: dict[str, list[dict[str, Any]]] = {}
    for event in ctx.asset_events:
        events_by_slot.setdefault(str(event.get("slot_id")), []).append(event)

    family_intents = {
        "luxury_editorial": {
            "needed_asset_intent": "Premium product or material detail imagery suitable for editorial launch slides.",
            "needed_asset_type": "product_or_material_image",
            "recommended_package_addition": "Add approved hero product cutout and macro material/detail images.",
        },
        "food_product_launch": {
            "needed_asset_intent": "Food product packshot, ingredient macro, shelf context, or campaign visual.",
            "needed_asset_type": "food_product_packshot_or_ingredient_image",
            "recommended_package_addition": "Add approved packshot, ingredient detail, retail shelf, and campaign key visual assets.",
        },
        "internal_strategy_report": {
            "needed_asset_intent": "Internal evidence chart, risk map, decision table, or roadmap visual.",
            "needed_asset_type": "strategy_evidence_visual",
            "recommended_package_addition": "Add approved internal metric chart, decision matrix, risk map, or roadmap image where native evidence is insufficient.",
        },
    }
    default_intent = {
        "needed_asset_intent": "Topic-specific visual asset for the requested composition family.",
        "needed_asset_type": "topic_visual_asset",
        "recommended_package_addition": "Add an approved package asset matching the slide slot intent.",
    }
    family_payload = family_intents.get(ctx.primary_family, default_intent)
    for slide in packet.slide_plan.slides:
        for slot in slide.asset_slots:
            slot_events = events_by_slot.get(slot.slot_id, [])
            fallback_event = next((event for event in slot_events if event.get("fallback_used")), None)
            missing_event = not slot_events
            if not fallback_event and not missing_event:
                continue
            importance = "required" if slot.required else ("recommended" if ctx.primary_family in family_intents else "optional")
            fallback_used = bool(fallback_event)
            expectations.append(
                {
                    "slot_id": slot.slot_id,
                    "needed_asset_intent": family_payload["needed_asset_intent"],
                    "needed_asset_type": family_payload["needed_asset_type"],
                    "importance": importance,
                    "affected_slides": [slide.slide_no],
                    "fallback_used": fallback_used,
                    "impact": (
                        "Safe native fallback used; topic specificity may be lower than an approved asset package."
                        if fallback_used
                        else "No asset event was recorded for this slot; verify renderer coverage."
                    ),
                    "recommended_package_addition": family_payload["recommended_package_addition"],
                }
            )
    return expectations


def used_assets_report_payload(packet: GuidePacket, ctx: BuildContext) -> dict[str, Any]:
    inserted_assets = []
    declared_slots_by_ref = approved_ref_to_declared_slot_ids(packet)
    for asset in approved_package_assets_for_asset_system(ctx):
        events = approved_events_for_asset(ctx, asset)
        normalized_path = normalize_package_relative_path(str(asset.get("relative_package_path") or ""))
        declared_slot_ids = declared_slots_by_ref.get(normalized_path or "", set())
        if len(declared_slot_ids) == 1:
            inserted_assets.append(inserted_asset_summary(asset, events, next(iter(declared_slot_ids))))
            continue
        if len(declared_slot_ids) > 1:
            event_slots = {str(event.get("slot_id") or "") for event in events}
            event_slots.discard("")
            split_slot_ids = sorted(declared_slot_ids & event_slots)
            if split_slot_ids:
                for slot_id in split_slot_ids:
                    slot_events = [event for event in events if event.get("slot_id") == slot_id]
                    inserted_assets.append(inserted_asset_summary(asset, slot_events, slot_id))
                continue
            summary = inserted_asset_summary(asset, events)
            summary["asset_intent"] = "ambiguous approved package ref shared by multiple declared asset slots; slot_id omitted"
            inserted_assets.append(summary)
            continue
        inserted_assets.append(inserted_asset_summary(asset, events))
    fallback_events = [
        {
            "slide_no": event.get("slide_no"),
            "slot_id": event.get("slot_id"),
            "slot_type": event.get("slot_type"),
            "reason": event.get("fallback_reason") or "native_shape_or_mockup_fallback",
        }
        for event in ctx.asset_events
        if event.get("fallback_used")
    ]
    expectations = missing_asset_expectations(packet, ctx)
    public_events = [public_asset_event(event) for event in ctx.asset_events]
    return {
        "contract": {
            "name": "b44.ppt_maker_used_assets_report",
            "version": "1.0",
            "compatibility": "additive_to_b44_asset_handoff",
        },
        "report_id": f"used-assets:{re.sub(r'[^a-zA-Z0-9_.:-]+', '-', packet.guide_identity.guide_id)}",
        "deck_id": f"deck:{re.sub(r'[^a-zA-Z0-9_.:-]+', '-', packet.guide_identity.guide_id)}",
        "status": "fail" if ctx.asset_validation_failures else ("warning" if any(event.get("fallback_used") for event in ctx.asset_events) else "pass"),
        "generated_at": utc_now(),
        "asset_resolution_priority": packet.asset_slot_policy.priority(),
        "total_asset_slots": sum(len(slide.asset_slots) for slide in packet.slide_plan.slides),
        "recorded_asset_slots": len({event.get("slot_id") for event in ctx.asset_events}),
        "events": public_events,
        "fallbacks": [fallback_summary(event) for event in public_events if event.get("fallback_used")],
        "inserted_assets": inserted_assets,
        "fallback_events": fallback_events,
        "missing_asset_expectations": expectations,
        "public_safety": {
            "private_refs_redacted": True,
            "contains_drive_ids": False,
            "contains_private_paths": False,
            "contains_tokens": False,
        },
    }


def asset_system_final_qa_report(
    packet: GuidePacket,
    ctx: BuildContext,
    qa: dict[str, Any],
) -> dict[str, Any]:
    fonts = sorted({role.fallback for role in packet.typography.roles if role.fallback})
    approved_assets = approved_package_assets_for_asset_system(ctx)
    palette_resolution = qa.get("palette_resolution", infer_palette_resolution(packet))
    color_scan = qa.get("color_spread_scan", {})
    contrast_scan = qa.get("contrast_scan", {})
    palette_detail_parts = [
        f"source_type={palette_resolution.get('source_type')}",
        f"palette_id={palette_resolution.get('selected_palette_id')}",
        f"palette_name={palette_resolution.get('selected_palette_name')}",
        f"source_kind={palette_resolution.get('source_kind')}",
        f"fallback_used={str(palette_resolution.get('source_type') == 'fallback_default_palette').lower()}",
    ]
    if palette_resolution.get("fallback_reason"):
        palette_detail_parts.append(f"fallback_reason={palette_resolution.get('fallback_reason')}")
    if palette_resolution.get("asset_ref"):
        palette_detail_parts.append(f"asset_ref={palette_resolution.get('asset_ref')}")
    validation_results = [
        {
            "check": "approved_asset_checksums",
            "status": "fail" if ctx.asset_validation_failures else ("pass" if approved_assets else "not_run"),
            "details": "Approved package assets were validated before insertion." if approved_assets else "No approved package assets were inserted.",
        },
        {
            "check": "visual_collision",
            "status": asset_system_status(qa.get("overflow_overlap_scan", {}).get("status", "fail")),
            "details": qa.get("overflow_overlap_scan", {}).get("details", "Text bounds were checked against reserved visual zones."),
        },
        {
            "check": "public_boundary",
            "status": "pass",
            "details": "Returned report uses public-safe relative paths and no Drive identifiers or tokens.",
        },
        {
            "check": "preview_render",
            "status": "pass" if qa.get("preview_report", {}).get("status") == "rendered" else "fail",
            "details": f"Rendered previews: {qa.get('preview_report', {}).get('rendered_slide_count', 0)}.",
        },
        {
            "check": "palette_resolution",
            "status": "pass" if palette_resolution.get("source_type") != "fallback_default_palette" else "warn",
            "details": "; ".join(str(part) for part in palette_detail_parts if part),
            "evidence": {
                "palette_resolution": palette_resolution_evidence(packet, palette_resolution, contrast_scan),
            },
        },
        {
            "check": "color_spread",
            "status": asset_system_status(color_scan.get("status", "not_run")),
            "details": (
                f"status={color_scan.get('status', 'not_run')}; "
                f"palette_policy={color_scan.get('palette_policy', 'unknown')}; "
                f"issues={len(color_scan.get('issues', []))}"
            ),
        },
        {
            "check": "palette_contrast",
            "status": asset_system_status(contrast_scan.get("status", "not_run")),
            "details": (
                f"status={contrast_scan.get('status', 'not_run')}; "
                f"method={contrast_scan.get('method', 'not_run')}; "
                f"issues={len(contrast_scan.get('issues', []))}"
            ),
        },
        {
            "check": "prompt_literal_leakage",
            "status": asset_system_status(qa.get("prompt_literal_leakage_scan", {}).get("status", "fail")),
            "details": (
                f"status={qa.get('prompt_literal_leakage_scan', {}).get('status', 'not_run')}; "
                f"matches={len(qa.get('prompt_literal_leakage_scan', {}).get('matches', []))}; "
                "checked visible title/body/footer text; structured evidence is retained in internal final-qa.json until schema expansion."
            ),
        },
        {
            "check": "composition_family",
            "status": "pass",
            "details": (
                f"primary_family={qa.get('composition_family_summary', {}).get('primary_family')}; "
                f"composition_family={qa.get('composition_family_summary', {}).get('composition_family')}; "
                f"archetypes={','.join(qa.get('composition_family_summary', {}).get('preferred_public_archetypes', []))}"
            ),
        },
    ]
    remaining_risks = []
    if not approved_assets:
        remaining_risks.append("No approved package assets were inserted; native fallback or no asset slot was used.")
    if qa.get("asset_slots_fallbacks", {}).get("count", 0):
        remaining_risks.append("Some asset slots used native fallback; review fallback events before final use.")
    return {
        "contract": {
            "name": "b44.ppt_maker_asset_handoff",
            "version": "1.0",
            "compatibility": "additive_to_b43_1",
        },
        "qa_report_id": f"qa-report:{re.sub(r'[^a-zA-Z0-9_.:-]+', '-', packet.guide_identity.guide_id)}",
        "created_at": utc_now(),
        "deck": {
            "deck_id": f"deck:{re.sub(r'[^a-zA-Z0-9_.:-]+', '-', packet.guide_identity.guide_id)}",
            "slide_count": packet.guide_identity.slide_count,
            "locale": packet.guide_identity.language,
        },
        "status": asset_system_status(qa["status"]),
        "fonts_used": [
            {"family_name": font, "source": "system_font", "verified_in_pptx": True}
            for font in fonts
        ],
        "font_fallbacks": [],
        "embedded_media_count": len(approved_assets),
        "approved_package_assets": approved_assets,
        "layout_recipes_used": layout_recipes_used(ctx),
        "unresolved_metadata_only_items": [],
        "manual_actions": [
            {
                "action_type": "none",
                "description": "No manual actions required by the PPT maker for this returned report.",
                "required_before_final_use": False,
            }
        ],
        "remaining_risks": remaining_risks,
        "validation_results": validation_results,
    }


def resolve_package_file(package_response_path: Path, relative_package_path: str) -> Path:
    candidates = [
        package_response_path.parent / relative_package_path,
        package_response_path.parent.parent / relative_package_path,
        BASE_DIR / relative_package_path,
    ]
    for candidate in candidates:
        if candidate.exists():
            return candidate.resolve()
    return candidates[0].resolve()


def package_slot_type(slot: dict[str, Any], manifest_asset: dict[str, Any] | None = None) -> str:
    text = " ".join(
        str(value or "")
        for value in [
            slot.get("slot_id"),
            slot.get("slot_role"),
            slot.get("asset_intent"),
            (manifest_asset or {}).get("insert_as"),
            (manifest_asset or {}).get("allowed_insertion_role"),
        ]
    ).lower()
    if "icon" in text:
        return "icon"
    if "logo" in text or "brand" in text:
        return "logo"
    if "background" in text or "texture" in text:
        return "background"
    if "roadmap" in text or "risk" in text or "strategy" in text or "evidence" in text:
        return "chart"
    return "image"


def package_slot_aliases(slot_id: str) -> set[str]:
    normalized = slot_id.replace("-", "_").lower()
    aliases = {normalized}
    families = {
        "product_hero": {"product_hero", "food_product_packshot", "hero_mockup"},
        "food_product_packshot": {"food_product_packshot", "product_hero", "packshot", "hero_mockup"},
        "ingredient_detail": {"ingredient_detail", "material_texture", "ingredient_visual"},
        "retail_shelf_context": {"retail_shelf_context", "context_visual", "market_context"},
        "campaign_calendar_visual": {"campaign_calendar_visual", "roadmap_or_risk_visual"},
        "luxury_product_detail": {"luxury_product_detail", "material_texture", "product_hero"},
        "material_texture": {"material_texture", "background_texture", "ingredient_detail"},
        "strategy_evidence_visual": {"strategy_evidence_visual", "roadmap_or_risk_visual"},
        "roadmap_or_risk_visual": {"roadmap_or_risk_visual", "strategy_evidence_visual"},
        "phone_screen": {"phone_screen", "hero_mockup"},
        "edge_bleed": {"edge_bleed", "edge_bleed_visual", "right_edge_bleed"},
        "edge_bleed_visual": {"edge_bleed", "edge_bleed_visual", "right_edge_bleed"},
        "icon": {"icon", "problem_icons"},
        "logo": {"logo", "brand_logo"},
        "background": {"background", "background_texture"},
        "background_texture": {"background", "background_texture", "material_texture"},
    }
    for key, values in families.items():
        if normalized == key or normalized in values or any(value in normalized for value in values):
            aliases.update(values)
    return aliases


def package_slot_target_slide(slides: list[dict[str, Any]], slot_id: str) -> dict[str, Any] | None:
    normalized = slot_id.replace("-", "_").lower()
    preferred_layouts: list[str]
    if "detail" in normalized:
        preferred_layouts = ["product_mockup", "big_thesis", "cover"]
    elif any(term in normalized for term in ["packshot", "product_hero", "luxury_product", "logo"]):
        preferred_layouts = ["cover", "product_mockup"]
    elif any(term in normalized for term in ["ingredient", "material", "phone_screen"]):
        preferred_layouts = ["product_mockup", "three_cards", "cover"]
    elif any(term in normalized for term in ["retail", "strategy", "evidence", "icon"]):
        preferred_layouts = ["three_cards", "market_sizing", "business_model"]
    elif any(term in normalized for term in ["campaign", "roadmap", "risk"]):
        preferred_layouts = ["roadmap", "process_flow", "financial_plan"]
    elif any(term in normalized for term in ["edge_bleed", "background", "texture"]):
        preferred_layouts = ["behavior_shift", "market_sizing", "cover"]
    else:
        preferred_layouts = ["cover", "product_mockup", "three_cards"]
    for layout in preferred_layouts:
        for slide in slides:
            if slide.get("layout_archetype") == layout:
                return slide
    return slides[0] if slides else None


def apply_package_asset_slots(raw: dict[str, Any], package_slots: list[dict[str, Any]], manifest_assets: dict[str, dict[str, Any]]) -> None:
    slides = raw.get("slide_plan", {}).get("slides", [])
    if not isinstance(slides, list):
        return
    for package_slot in package_slots:
        slot_id = str(package_slot.get("slot_id") or "").strip()
        approved_ref = normalize_package_relative_path(str(package_slot.get("approved_asset_ref") or "").strip())
        if not slot_id or not approved_ref:
            continue
        aliases = package_slot_aliases(slot_id)
        matched = False
        for slide in slides:
            slots = slide.get("asset_slots") or []
            if not isinstance(slots, list):
                continue
            for slot in slots:
                existing_id = str(slot.get("slot_id") or "").replace("-", "_").lower()
                if existing_id in aliases:
                    slot["slot_id"] = slot_id
                    slot["preferred_asset_ref"] = approved_ref
                    slot["crop_or_mask_policy"] = package_slot.get("crop_or_mask_policy") or slot.get("crop_or_mask_policy") or "preserve_aspect_ratio"
                    sources = list(slot.get("allowed_sources") or [])
                    if "approved_package_asset" not in sources:
                        sources.insert(0, "approved_package_asset")
                    slot["allowed_sources"] = sources
                    matched = True
        if matched:
            continue
        target = package_slot_target_slide(slides, slot_id)
        if not target:
            continue
        target_slots = target.setdefault("asset_slots", [])
        manifest_asset = manifest_assets.get(approved_ref)
        target_slots.append(
            {
                "slot_id": slot_id,
                "slot_type": package_slot_type(package_slot, manifest_asset),
                "required": bool(package_slot.get("visual_priority") == "high" or package_slot.get("required")),
                "allowed_sources": ["approved_package_asset", "user_supplied_approved_asset", "generated_image_approved_for_deck"],
                "preferred_asset_ref": approved_ref,
                "fallback": "Use native editable fallback and report the package consumption blocker.",
                "crop_or_mask_policy": package_slot.get("crop_or_mask_policy") or (manifest_asset or {}).get("recommended_fit_mode") or "preserve_aspect_ratio",
            }
        )


def asset_slot_ids_for_package_asset(asset: dict[str, Any]) -> list[str]:
    allowed_role = str(asset.get("allowed_insertion_role") or asset.get("asset_role") or asset.get("insert_as") or "").lower()
    insert_as = str(asset.get("insert_as") or "").lower()
    if "phone_screen_or_product_mockup" in allowed_role:
        return ["hero_mockup", "phone_screen"]
    if "edge_bleed_visual" in allowed_role or "right-edge" in allowed_role or "right edge" in allowed_role:
        return ["edge_bleed_visual"]
    if "background_texture" in allowed_role or "background" in allowed_role or "texture" in allowed_role:
        return ["background_texture"]
    if insert_as == "icon" or "icon" in allowed_role:
        return ["problem_icons"]
    return [str(asset.get("asset_role") or insert_as or "approved_asset")]


def inject_approved_package_references(raw: dict[str, Any], approved_package_path: str | Path | None) -> dict[str, Any]:
    if not approved_package_path:
        return raw
    package_path = Path(approved_package_path).resolve()
    package_response = json.loads(package_path.read_text(encoding="utf-8-sig"))
    manifest = package_response.get("package_manifest", {})
    package_manifest_id = manifest.get("package_manifest_id")
    manifest_assets = {
        normalized_path: asset
        for asset in manifest.get("assets", [])
        if (normalized_path := normalize_package_relative_path(str(asset.get("relative_package_path") or "")))
    }
    package_slots = [slot for slot in package_response.get("asset_slots", []) if isinstance(slot, dict)]
    apply_package_asset_slots(raw, package_slots, manifest_assets)
    approved_palette = package_palette_selection(package_response, package_manifest_id)
    if approved_palette and should_apply_package_palette(raw):
        raw["palette"] = approved_palette
        constraints = raw.setdefault("project_brief", {}).setdefault("constraints", [])
        if isinstance(constraints, list):
            constraints.append(f"Palette resolution: asset_system_palette / {approved_palette['palette_id']}.")
    references = list(raw.get("approved_asset_references") or [])

    def append_reference(asset: dict[str, Any] | None, slot_id: str, asset_ref: str, package_slot: dict[str, Any] | None = None, failure: str | None = None) -> None:
        relative_path = public_package_path_or_none((asset or {}).get("relative_package_path") or asset_ref) or str(asset_ref or "")
        local_path = resolve_package_file(package_path, relative_path) if asset and not failure else None
        references.append(
            {
                "asset_ref": public_package_path_or_none(asset_ref) or relative_path,
                "slot_id": slot_id,
                "asset_type": (asset or {}).get("insert_as") or (asset or {}).get("media_type") or package_slot_type(package_slot or {}, asset),
                "asset_role": (asset or {}).get("insert_as") or (asset or {}).get("asset_role") or (package_slot or {}).get("slot_role"),
                "approved": True,
                "local_path": str(local_path) if local_path else None,
                "sha256": (asset or {}).get("sha256"),
                "file_size_bytes": (asset or {}).get("size_bytes"),
                "package_status": "approved_package_file",
                "package_manifest_id": package_manifest_id,
                "relative_package_path": relative_path,
                "approved_asset_ref": public_package_path_or_none(asset_ref) or relative_path,
                "package_slot_id": (package_slot or {}).get("slot_id") or slot_id,
                "approved_asset_ref_declared_by_package_slot": bool(package_slot),
                "allowed_insertion_role": (asset or {}).get("allowed_insertion_role") or (asset or {}).get("insert_as") or (package_slot or {}).get("slot_role"),
                "media_type": (asset or {}).get("media_type"),
                "insert_as": (asset or {}).get("insert_as"),
                "recommended_fit_mode": (asset or {}).get("recommended_fit_mode"),
                "allowed_crop_mode": (asset or {}).get("allowed_crop_mode"),
                "mask_shape": (asset or {}).get("mask_shape"),
                "bleed_allowance": (asset or {}).get("bleed_allowance"),
                "package_crop_or_mask_policy": (package_slot or {}).get("crop_or_mask_policy"),
                "license_action": (asset or {}).get("license_action"),
                "package_validation_failure": failure,
            }
        )

    handled_manifest_paths: set[str] = set()
    for package_slot in package_slots:
        slot_id = str(package_slot.get("slot_id") or "").strip()
        raw_approved_ref = str(package_slot.get("approved_asset_ref") or "").strip()
        approved_ref = normalize_package_relative_path(raw_approved_ref)
        if not slot_id or not raw_approved_ref:
            continue
        if not approved_ref:
            append_reference(None, slot_id, raw_approved_ref, package_slot, "approved_asset_path_unsafe")
            continue
        asset = manifest_assets.get(approved_ref)
        if not asset:
            append_reference(None, slot_id, approved_ref, package_slot, "approved_asset_ref_not_found_in_manifest")
            continue
        handled_manifest_paths.add(approved_ref)
        append_reference(asset, slot_id, approved_ref, package_slot)

    for asset in manifest.get("assets", []):
        relative_path = normalize_package_relative_path(str(asset.get("relative_package_path") or ""))
        if not relative_path:
            continue
        if package_slots and relative_path in handled_manifest_paths:
            continue
        local_path = resolve_package_file(package_path, relative_path)
        for slot_id in asset_slot_ids_for_package_asset(asset):
            append_reference(asset, slot_id, relative_path)
    raw["approved_asset_references"] = references
    return raw


def deck_plan_summary(packet: GuidePacket, ctx: BuildContext) -> dict[str, Any]:
    slot_results = []
    for slide in packet.slide_plan.slides:
        for slot in slide.asset_slots:
            event = next((item for item in ctx.asset_events if item.get("slide_no") == slide.slide_no and item.get("slot_id") == slot.slot_id), None)
            status = "missing"
            if event and event.get("resolution") == "approved_package_asset":
                status = "approved_package_inserted"
            elif event and event.get("fallback_used"):
                status = "fallback_used"
            slot_results.append(
                {
                    "slide_no": slide.slide_no,
                    "slot_id": slot.slot_id,
                    "slot_type": slot.slot_type,
                    "status": status,
                    "relative_package_path": event.get("relative_package_path") if event else None,
                }
            )
    return {
        "contract": {
            "name": "b44.ppt_maker_deck_plan_summary",
            "version": "1.0",
            "compatibility": "additive_to_b44_design_guide_packet",
        },
        "deck_id": f"deck:{re.sub(r'[^a-zA-Z0-9_.:-]+', '-', packet.guide_identity.guide_id)}",
        "slide_count": packet.guide_identity.slide_count,
        "primary_family": ctx.primary_family,
        "composition_family": ctx.composition_family,
        "preferred_public_archetypes": ctx.composition_profile.get("preferred_archetypes", []),
        "required_slot_results": slot_results,
    }


def renderer_contract_summary(ctx: BuildContext, packet: GuidePacket) -> dict[str, Any]:
    return {
        "contract": {
            "name": "b44.ppt_maker_renderer_contract_summary",
            "version": "1.0",
            "compatibility": "additive_to_b44_asset_handoff",
        },
        "deck_id": f"deck:{re.sub(r'[^a-zA-Z0-9_.:-]+', '-', packet.guide_identity.guide_id)}",
        "primary_family": ctx.primary_family,
        "composition_family": ctx.composition_family,
        "archetype_recipe_summary": [
            {
                "slide_no": event.get("slide_no"),
                "public_archetype": ctx.layout_strategy_by_slide.get(int(event.get("slide_no") or 0), {}).get("public_archetype"),
                "layout_archetype": event.get("layout_archetype"),
                "layout_recipe": event.get("layout_recipe"),
                "recipe_family": event.get("recipe_family"),
                "recipe_composition_style": event.get("recipe_composition_style"),
            }
            for event in ctx.native_renderer_events
        ],
        "inserted_assets": [
            {"relative_package_path": item["relative_package_path"]}
            for item in approved_package_assets_for_asset_system(ctx)
        ],
    }


def final_qa_report(
    packet: GuidePacket,
    ctx: BuildContext,
    pptx_path: Path,
    preview_report: dict[str, Any],
    html_report: dict[str, Any],
    project_dir: Path,
    deck_plan: dict[str, Any],
) -> dict[str, Any]:
    required_archetypes = {slide.layout_archetype for slide in packet.slide_plan.slides}
    rendered_archetypes = {event["layout_archetype"] for event in ctx.native_renderer_events}
    fallback_count = sum(1 for event in ctx.asset_events if event.get("fallback_used"))
    approved_events = [event for event in ctx.asset_events if not event.get("fallback_used")]
    fallback_events = [event for event in ctx.asset_events if event.get("fallback_used")]
    required_slots = [slot for slide in packet.slide_plan.slides for slot in slide.asset_slots if slot.required]
    unresolved_blockers = []
    recorded_slot_ids = {event.get("slot_id") for event in ctx.asset_events}
    for slot in required_slots:
        if slot.slot_id not in recorded_slot_ids:
            unresolved_blockers.append(
                {
                    "type": "required_asset_slot_missing",
                    "slot_id": slot.slot_id,
                    "status": "fail",
                }
            )
    if preview_report.get("validation_blocker"):
        unresolved_blockers.append(
            {
                "type": "pptx_render_failed",
                "status": "fail",
                "reason": preview_report.get("fallback_reason"),
            }
        )
    visual_scan = visual_collision_scan(packet, ctx)
    color_scan = color_spread_scan(packet, ctx)
    contrast_scan = palette_contrast_scan(packet)
    palette_resolution = infer_palette_resolution(packet)
    if visual_scan["status"] == "fail":
        unresolved_blockers.append(
            {
                "type": "visual_collision",
                "status": "fail",
                "issue_count": len(visual_scan["issues"]),
            }
        )
    if contrast_scan["status"] == "fail":
        unresolved_blockers.append(
            {
                "type": "palette_contrast",
                "status": "fail",
                "issue_count": len(contrast_scan["issues"]),
            }
        )
    for failure in ctx.asset_validation_failures:
        unresolved_blockers.append(public_asset_failure(failure))
    missing_visible_blockers = []
    for slide in deck_plan["slides"]:
        raw = json.dumps(slide["visible_content"], ensure_ascii=False)
        for term in packet.public_safety.forbidden_visible_terms or DEFAULT_FORBIDDEN_VISIBLE_TERMS:
            if term and term in raw:
                missing_visible_blockers.append({"slide_no": slide["slide_no"], "term": term})
    prompt_literal_matches_found = []
    visible_scan_items = []
    for slide in deck_plan["slides"]:
        payload = slide.get("visible_content", {})
        for field in ["title", "subtitle", "brief"]:
            text = str(payload.get(field) or "")
            visible_scan_items.append({"slide_no": slide["slide_no"], "field": field, "text": text})
            for match in prompt_literal_matches(text):
                prompt_literal_matches_found.append({"slide_no": slide["slide_no"], "field": field, "phrase": match})
    project_name_matches = prompt_literal_matches(sanitize_visible_value(packet.guide_identity.project_name, fallback="Untitled Deck"))
    for match in project_name_matches:
        prompt_literal_matches_found.append({"slide_no": None, "field": "footer_project_name", "phrase": match})
    prompt_literal_matches_found.extend(pptx_prompt_literal_matches(pptx_path))
    if prompt_literal_matches_found:
        unresolved_blockers.append(
            {
                "type": "prompt_literal_leakage",
                "status": "fail",
                "issue_count": len(prompt_literal_matches_found),
            }
        )
    status = "pass" if not missing_visible_blockers and not unresolved_blockers else "fail"
    checksum_checked = sum(
        1
        for event in ctx.asset_events
        if event.get("checksum_valid") is not None or event.get("file_size_valid") is not None
    )
    return {
        "status": status,
        "generated_at": utc_now(),
        "pptx_path": safe_rel(pptx_path),
        "slide_count": packet.guide_identity.slide_count,
        "native_powerpoint_rendering": True,
        "html_screenshot_used_in_pptx": False,
        "variant_strategy_used": ctx.deck_plan.get("variant_strategy", {}),
        "composition_family_summary": {
            "primary_family": ctx.primary_family,
            "composition_family": ctx.composition_family,
            "preferred_public_archetypes": ctx.composition_profile.get("preferred_archetypes", []),
            "major_visual_asset_treatment": ctx.composition_profile.get("major_visual_asset_treatment"),
            "avoid": ctx.composition_profile.get("avoid", []),
        },
        "layout_strategy_map_used": {
            str(slide_no): strategy
            for slide_no, strategy in sorted(ctx.layout_strategy_by_slide.items())
        },
        "recipe_composition_map_used": {
            str(event["slide_no"]): {
                "layout_recipe": event.get("layout_recipe"),
                "recipe_family": event.get("recipe_family"),
                "recipe_composition_style": event.get("recipe_composition_style"),
            }
            for event in ctx.native_renderer_events
        },
        "all_fixture_archetypes_rendered": required_archetypes.issubset(rendered_archetypes),
        "rendered_archetypes": sorted(rendered_archetypes),
        "layout_archetypes_used": {
            "status": "pass" if required_archetypes.issubset(rendered_archetypes) else "fail",
            "items": sorted(rendered_archetypes),
        },
        "palette_roles_used": sorted(ctx.used_palette_roles),
        "palette_resolution": palette_resolution,
        "typography_roles_used": sorted(ctx.used_typography_roles),
        "safe_area_applied": packet.safe_area.model_dump(mode="json"),
        "background_policy_applied": packet.background_policy.model_dump(mode="json"),
        "header_footer_omissions": ctx.omitted_header_footer,
        "asset_fallback_count": fallback_count,
        "approved_assets_used": {
            "status": "pass" if approved_events else "not_applicable",
            "count": len(approved_events),
            "items": approved_events,
        },
        "asset_slots_fallbacks": {
            "status": "warning" if fallback_events else "pass",
            "count": len(fallback_events),
            "items": fallback_events,
        },
        "checksum_validation_result": {
            "status": "pass" if checksum_checked else "not_applicable",
            "checked_assets": checksum_checked,
            "failed_assets": [
                event
                for event in ctx.asset_events
                if event.get("checksum_valid") is False or event.get("file_size_valid") is False
            ],
        },
        "overflow_overlap_scan": {
            **visual_scan,
            "details": "Text boxes were checked against right-edge decorative blocks, phone/mockup frames, edge-bleed assets, and safe-area bounds.",
        },
        "contrast_scan": contrast_scan,
        "color_spread_scan": color_scan,
        "unresolved_blockers": unresolved_blockers,
        "preview_report": preview_report,
        "html_report": html_report,
        "visible_content_blockers": missing_visible_blockers,
        "prompt_literal_leakage_scan": {
            "status": "fail" if prompt_literal_matches_found else "pass",
            "blocked_as": "P1",
            "command_phrases": PROMPT_LITERAL_COMMAND_LABELS,
            "matches": prompt_literal_matches_found,
            "checked_visible_fields": visible_scan_items,
        },
        "output_manifest": {
            "generated.pptx": safe_rel(project_dir / "generated.pptx"),
            "deck-plan.json": safe_rel(project_dir / "deck-plan.json"),
            "guide-compliance-report.json": safe_rel(project_dir / "guide-compliance-report.json"),
            "final-qa.json": safe_rel(project_dir / "final-qa.json"),
            "asset-system-final-qa.json": safe_rel(project_dir / "asset-system-final-qa.json"),
            "used-assets-report.json": safe_rel(project_dir / "used-assets-report.json"),
            "deck-plan-summary.json": safe_rel(project_dir / "deck-plan-summary.json"),
            "renderer-contract-summary.json": safe_rel(project_dir / "renderer-contract-summary.json"),
            "html-guide-render-report.json": safe_rel(project_dir / "html-guide-render-report.json"),
            "approved-package-response.json": safe_rel(project_dir / "approved-package-response.json") if (project_dir / "approved-package-response.json").exists() else None,
            "previews": [safe_rel(path) for path in sorted((project_dir / "previews").glob("*.png"))],
        },
    }


def build_from_guide_packet(
    guide_path: str | Path,
    *,
    mode: Literal["assistant", "auto"] = "assistant",
    output_root: str | Path = DEFAULT_PROJECTS_DIR,
    project_id: str | None = None,
    html_guide_requested: bool = False,
    variant: str | None = None,
    variant_strategy: str = "investor_open",
    approved_package_path: str | Path | None = None,
) -> dict[str, Any]:
    packet, source_path, raw = load_guide_packet(guide_path)
    if approved_package_path:
        raw = inject_approved_package_references(raw, approved_package_path)
        packet = GuidePacket.model_validate(raw)
    source_mode = mode if source_path.name == "guide-data.public.json" else "explicit_guide_packet"
    project_id = project_id or re.sub(r"[^a-zA-Z0-9_.-]+", "-", packet.guide_identity.project_name.lower()).strip("-")
    project_dir = Path(output_root).resolve() / project_id
    project_dir.mkdir(parents=True, exist_ok=True)
    if approved_package_path:
        shutil.copyfile(Path(approved_package_path).resolve(), project_dir / "approved-package-response.json")
    json_dump(project_dir / "guide-data.public.json", safe_string(raw))
    deck_plan, renderer_contract = write_planning_artifacts(packet, project_dir, source_mode, variant_strategy)
    json_dump(project_dir / "guide-compliance-report.json", guide_compliance_report(packet, source_path, source_mode))
    html_report = generate_html_guide(packet, project_dir, mode, html_guide_requested)
    json_dump(project_dir / "html-guide-render-report.json", html_report)
    pptx_path, ctx = render_packet_pptx(packet, project_dir, mode, deck_plan, variant)
    preview_report = render_previews(pptx_path, packet, project_dir)
    used_assets = used_assets_report_payload(packet, ctx)
    json_dump(project_dir / "used-assets-report.json", used_assets)
    qa = final_qa_report(packet, ctx, pptx_path, preview_report, html_report, project_dir, deck_plan)
    json_dump(project_dir / "final-qa.json", qa)
    json_dump(project_dir / "asset-system-final-qa.json", asset_system_final_qa_report(packet, ctx, qa))
    json_dump(project_dir / "deck-plan-summary.json", deck_plan_summary(packet, ctx))
    json_dump(project_dir / "renderer-contract-summary.json", renderer_contract_summary(ctx, packet))
    return {
        "status": "built" if qa["status"] == "pass" else "blocked",
        "mode": mode,
        "project_dir": safe_rel(project_dir),
        "pptx_path": safe_rel(pptx_path),
        "reports": {
            "deck_plan": safe_rel(project_dir / "deck-plan.json"),
            "renderer_contract": safe_rel(project_dir / "renderer-contract.json"),
            "guide_compliance": safe_rel(project_dir / "guide-compliance-report.json"),
            "final_qa": safe_rel(project_dir / "final-qa.json"),
            "asset_system_final_qa": safe_rel(project_dir / "asset-system-final-qa.json"),
            "used_assets": safe_rel(project_dir / "used-assets-report.json"),
            "deck_plan_summary": safe_rel(project_dir / "deck-plan-summary.json"),
            "renderer_contract_summary": safe_rel(project_dir / "renderer-contract-summary.json"),
            "html_guide": safe_rel(project_dir / "html-guide-render-report.json"),
        },
    }


def plan_from_guide_packet(
    guide_path: str | Path,
    *,
    mode: Literal["assistant"] = "assistant",
    output_root: str | Path = DEFAULT_PROJECTS_DIR,
    project_id: str | None = None,
    html_guide_requested: bool = True,
    variant_strategy: str = "investor_open",
    approved_package_path: str | Path | None = None,
) -> dict[str, Any]:
    packet, source_path, raw = load_guide_packet(guide_path)
    if approved_package_path:
        raw = inject_approved_package_references(raw, approved_package_path)
        packet = GuidePacket.model_validate(raw)
    source_mode = mode if source_path.name == "guide-data.public.json" else "explicit_guide_packet"
    project_id = project_id or re.sub(r"[^a-zA-Z0-9_.-]+", "-", packet.guide_identity.project_name.lower()).strip("-")
    project_dir = Path(output_root).resolve() / project_id
    project_dir.mkdir(parents=True, exist_ok=True)
    if approved_package_path:
        shutil.copyfile(Path(approved_package_path).resolve(), project_dir / "approved-package-response.json")
    json_dump(project_dir / "guide-data.public.json", safe_string(raw))
    write_planning_artifacts(packet, project_dir, source_mode, variant_strategy)
    json_dump(project_dir / "guide-compliance-report.json", guide_compliance_report(packet, source_path, source_mode))
    html_report = generate_html_guide(packet, project_dir, mode, html_guide_requested)
    json_dump(project_dir / "html-guide-render-report.json", html_report)
    checkpoint = {
        "status": "waiting_for_approval",
        "mode": mode,
        "approval_required": True,
        "next_action": "Re-run with --build-approved after reviewing planning artifacts.",
        "artifacts": {
            "guide_packet": safe_rel(project_dir / "guide-data.public.json"),
            "deck_plan": safe_rel(project_dir / "deck-plan.json"),
            "renderer_contract": safe_rel(project_dir / "renderer-contract.json"),
            "asset_slot_plan": safe_rel(project_dir / "asset-slot-plan.json"),
            "qa_plan": safe_rel(project_dir / "qa-plan.json"),
            "guide_compliance": safe_rel(project_dir / "guide-compliance-report.json"),
            "html_guide": safe_rel(project_dir / "html" / "guide.html") if html_report.get("generated") else None,
            "html_guide_report": safe_rel(project_dir / "html-guide-render-report.json"),
            "draft_design_brief": safe_rel(project_dir / "draft_design_brief.md"),
        },
    }
    json_dump(project_dir / "assistant-build-checkpoint.json", checkpoint)
    return {
        "status": "waiting_for_approval",
        "mode": mode,
        "project_dir": safe_rel(project_dir),
        "approval_required": True,
        "next_action": "review_planning_artifacts_then_rerun_with_build_approved",
        "reports": checkpoint["artifacts"],
    }


def load_json_file(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def variant_distinctiveness_report(root: Path, profile_a: dict[str, Any], profile_b: dict[str, Any]) -> dict[str, Any]:
    variant_dirs = {"variant_a": root / "variant-a", "variant_b": root / "variant-b"}
    plans = {key: load_json_file(path / "deck-plan.json") for key, path in variant_dirs.items()}
    contracts = {key: load_json_file(path / "renderer-contract.json") for key, path in variant_dirs.items()}

    def layout_signature(plan: dict[str, Any]) -> list[str]:
        return [
            f"{slide.get('layout_archetype')}|{slide.get('public_archetype')}|{slide.get('layout_recipe')}"
            for slide in plan.get("slides", [])
        ]

    def recipe_families(plan: dict[str, Any]) -> list[str]:
        return [
            recipe_family_from_recipe(str(slide.get("layout_recipe") or ""))
            for slide in plan.get("slides", [])
        ]

    def density(plan: dict[str, Any]) -> str:
        values = [str(slide.get("density_budget") or "") for slide in plan.get("slides", [])]
        return max(set(values), key=values.count) if values else "unknown"

    signatures = {key: layout_signature(plan) for key, plan in plans.items()}
    recipes = {key: recipe_families(plan) for key, plan in plans.items()}
    changed_layouts = sum(1 for a, b in zip(signatures["variant_a"], signatures["variant_b"], strict=False) if a != b)
    changed_recipes = sum(1 for a, b in zip(recipes["variant_a"], recipes["variant_b"], strict=False) if a != b)
    comparisons = {
        "strategy_id": {
            "variant_a": profile_a.get("id") or profile_a.get("strategy_id"),
            "variant_b": profile_b.get("id") or profile_b.get("strategy_id"),
            "different": (profile_a.get("id") or profile_a.get("strategy_id")) != (profile_b.get("id") or profile_b.get("strategy_id")),
        },
        "composition_family": {
            "variant_a": plans["variant_a"].get("composition_family"),
            "variant_b": plans["variant_b"].get("composition_family"),
            "different": plans["variant_a"].get("composition_family") != plans["variant_b"].get("composition_family"),
        },
        "layout_signature": {
            "changed_slides": changed_layouts,
            "slide_count": max(len(signatures["variant_a"]), len(signatures["variant_b"])),
            "different": changed_layouts > 0,
        },
        "recipe_family": {
            "variant_a": recipes["variant_a"],
            "variant_b": recipes["variant_b"],
            "changed_slides": changed_recipes,
            "different": changed_recipes > 0,
        },
        "major_visual_asset_treatment": {
            "variant_a": plans["variant_a"].get("composition_profile", {}).get("major_visual_asset_treatment"),
            "variant_b": plans["variant_b"].get("composition_profile", {}).get("major_visual_asset_treatment"),
            "different": plans["variant_a"].get("composition_profile", {}).get("major_visual_asset_treatment") != plans["variant_b"].get("composition_profile", {}).get("major_visual_asset_treatment"),
        },
        "information_density": {
            "variant_a": density(plans["variant_a"]),
            "variant_b": density(plans["variant_b"]),
            "different": density(plans["variant_a"]) != density(plans["variant_b"]),
        },
    }
    different_dimensions = [key for key, value in comparisons.items() if value.get("different")]
    status = "pass" if len(different_dimensions) >= 3 and changed_layouts >= 1 else "warn"
    warnings = []
    if not comparisons["composition_family"]["different"]:
        warnings.append("Variants share the same public composition family; human review should confirm visual diversity.")
    if len(different_dimensions) <= 1:
        warnings.append("Variants differ only by strategy id or one dimension.")
    if changed_layouts == 0:
        warnings.append("Layout signatures are identical.")
    human_review_recommended = bool(warnings) or status != "pass"
    return {
        "status": status,
        "checked_dimensions": comparisons,
        "different_dimensions": different_dimensions,
        "warnings": warnings,
        "human_review_recommended": human_review_recommended,
        "contracts_compared": {
            "variant_a_strategy_contract": contracts["variant_a"].get("strategy_contract", {}),
            "variant_b_strategy_contract": contracts["variant_b"].get("strategy_contract", {}),
        },
    }


def append_auto_distinctiveness_to_reports(root: Path, distinctiveness: dict[str, Any]) -> None:
    detail = (
        f"status={distinctiveness.get('status')}; "
        f"different_dimensions={','.join(distinctiveness.get('different_dimensions', []))}; "
        f"warnings={len(distinctiveness.get('warnings', []))}"
    )
    for variant_dir in [root / "variant-a", root / "variant-b"]:
        qa_path = variant_dir / "final-qa.json"
        if qa_path.exists():
            qa = load_json_file(qa_path)
            qa["auto_visual_distinctiveness"] = distinctiveness
            json_dump(qa_path, qa)
        asset_qa_path = variant_dir / "asset-system-final-qa.json"
        if asset_qa_path.exists():
            asset_qa = load_json_file(asset_qa_path)
            asset_qa.setdefault("validation_results", []).append(
                {
                    "check": "auto_visual_distinctiveness",
                    "status": "pass" if distinctiveness.get("status") == "pass" else "warn",
                    "details": detail + "; structured evidence is retained in variant-comparison-report.json until schema expansion.",
                }
            )
            json_dump(asset_qa_path, asset_qa)


def build_auto_variants(
    guide_path: str | Path,
    *,
    output_root: str | Path = DEFAULT_PROJECTS_DIR,
    project_id: str | None = None,
    html_guide_requested: bool = False,
    variant_strategy_a: str = "investor_open",
    variant_strategy_b: str = "operator_dense",
    routing_report_path: str | Path | None = None,
    approved_package_path: str | Path | None = None,
) -> dict[str, Any]:
    packet, _, raw = load_guide_packet(guide_path)
    project_id = project_id or re.sub(r"[^a-zA-Z0-9_.-]+", "-", packet.guide_identity.project_name.lower()).strip("-")
    root = Path(output_root).resolve() / project_id
    root.mkdir(parents=True, exist_ok=True)
    if approved_package_path:
        shutil.copyfile(Path(approved_package_path).resolve(), root / "approved-package-response.json")
    routing_source = None
    if routing_report_path or (variant_strategy_a == "investor_open" and variant_strategy_b == "operator_dense"):
        routed_a, routed_b, routing_source = load_routed_strategy_pair(
            guide_path,
            output_root=output_root,
            project_id=project_id,
            routing_report_path=routing_report_path,
        )
        if routed_a and routed_b:
            variant_strategy_a = routed_a
            variant_strategy_b = routed_b
    variant_a = build_from_guide_packet(
        guide_path,
        mode="auto",
        output_root=root,
        project_id="variant-a",
        html_guide_requested=html_guide_requested,
        variant="a",
        variant_strategy=resolve_strategy_id(variant_strategy_a),
        approved_package_path=approved_package_path,
    )
    variant_b = build_from_guide_packet(
        guide_path,
        mode="auto",
        output_root=root,
        project_id="variant-b",
        html_guide_requested=html_guide_requested,
        variant="b",
        variant_strategy=resolve_strategy_id(variant_strategy_b),
        approved_package_path=approved_package_path,
    )
    profile_a = strategy_profile(variant_strategy_a)
    profile_b = strategy_profile(variant_strategy_b)
    distinctiveness = variant_distinctiveness_report(root, profile_a, profile_b)
    comparison = {
        "status": "pass" if variant_a["status"] == "built" and variant_b["status"] == "built" and distinctiveness["status"] == "pass" else ("warn" if variant_a["status"] == "built" and variant_b["status"] == "built" else "blocked"),
        "generated_at": utc_now(),
        "variant_a": variant_a,
        "variant_b": variant_b,
        "strategy_pair": {
            "variant_a": profile_a.get("id") or profile_a.get("strategy_id"),
            "variant_b": profile_b.get("id") or profile_b.get("strategy_id"),
        },
        "routing_source": routing_source,
        "visual_distinctiveness": distinctiveness,
        "differences": [
            f"Variant A uses `{profile_a.get('id') or profile_a.get('strategy_id')}` with {profile_a.get('slide_rhythm', 'an open slide rhythm')} and {profile_a.get('evidence_style', 'selective evidence')}.",
            f"Variant B uses `{profile_b.get('id') or profile_b.get('strategy_id')}` with {profile_b.get('slide_rhythm', 'an alternate slide rhythm')} and {profile_b.get('evidence_style', 'alternate evidence treatment')}.",
            "The two variants differ first in deck-plan.json and renderer-contract.json strategy, layout recipe, content emphasis, evidence treatment, palette emphasis, typography bias, and chart/table style.",
            "The renderer consumes those machine-facing plan fields and does not branch on the variant letter.",
        ],
        "recommendation": "variant-a",
        "recommendation_reason": "Variant A is treated as the primary interpretation; use Variant B when its alternate evidence or density treatment better fits the audience.",
    }
    json_dump(root / "variant-comparison-report.json", comparison)
    append_auto_distinctiveness_to_reports(root, distinctiveness)
    (root / "auto-mode-recommendation.md").write_text(
        "# Auto Mode Recommendation\n\n"
        "Recommended variant: `variant-a`.\n\n"
        f"Variant A uses `{profile_a.get('id') or profile_a.get('strategy_id')}`. "
        f"Variant B uses `{profile_b.get('id') or profile_b.get('strategy_id')}`. "
        "Choose Variant B when its alternate density, evidence treatment, or visual rhythm better matches the audience.\n",
        encoding="utf-8",
    )
    return {
        "status": "pass" if comparison["status"] in {"pass", "warn"} else comparison["status"],
        "mode": "auto",
        "project_dir": safe_rel(root),
        "variant_a_pptx": variant_a["pptx_path"],
        "variant_b_pptx": variant_b["pptx_path"],
        "comparison_report": safe_rel(root / "variant-comparison-report.json"),
        "recommendation": safe_rel(root / "auto-mode-recommendation.md"),
        "routing_source": routing_source,
    }


def compose_default_packet_from_prompt(prompt: str, *, slide_count: int = 5, project_name: str | None = None) -> dict[str, Any]:
    raw_project_name = project_name or (prompt.split(".")[0].strip()[:72] if prompt.strip() else "Untitled Deck")
    project_name = sanitize_prompt_literal_text(raw_project_name, fallback="Untitled Deck")[:72]
    archetypes = ["cover", "big_thesis", "three_cards", "process_flow", "ask"]
    fallback_palette = palette_from_preset(
        palette_preset("fallback_default"),
        source_type="fallback_default_palette",
        evidence=["compose_default_packet_from_prompt"],
    )
    slides = []
    for index, archetype in enumerate(archetypes[:slide_count], start=1):
        slides.append(
            {
                "slide_no": index,
                "layout_archetype": archetype,
                "content_brief": sanitize_prompt_literal_text(prompt, fallback=project_name) if index == 1 else f"{project_name} planning point {index}",
                "visible_content_candidates": [
                    project_name if index == 1 else f"{project_name} point {index}",
                    "A concise machine-planned slide generated from sparse intake.",
                ],
                "asset_slots": [],
                "qa_checks": ["No HTML screenshot", "Native PPT objects only"],
            }
        )
    return {
        "contract": {"name": "b44.design_guide_packet", "version": "1.0", "compatibility": "additive_to_b44_asset_handoff"},
        "guide_identity": {
            "guide_id": f"assistant:{re.sub(r'[^a-z0-9]+', '-', project_name.lower()).strip('-') or 'deck'}",
            "guide_version": "1.0.0",
            "status": "draft",
            "language": "en-US",
            "project_name": project_name,
            "topic": sanitize_prompt_literal_text(prompt, fallback=project_name),
            "slide_count": len(slides),
            "created_at": utc_now(),
        },
        "project_brief": {
            "audience": ["general business audience"],
            "tone": ["clear", "credible", "modern"],
            "objective": sanitize_prompt_literal_text(prompt, fallback=project_name),
            "constraints": ["Use public-safe local artifacts only.", "Do not use HTML screenshots in PPTX."],
        },
        "reference_files": [{"file_name": "sparse-intake", "extension": ".txt", "purpose": "User natural-language request"}],
        "palette": fallback_palette,
        "typography": {
            "font_package_refs": [],
            "roles": [
                {"role": "hero_title", "font_stack": ["Aptos", "Arial"], "default_pt": 46, "min_pt": 36, "max_pt": 54, "weight": "bold", "line_limit": 2, "fallback": "Aptos Display"},
                {"role": "slide_title", "font_stack": ["Aptos", "Arial"], "default_pt": 30, "min_pt": 24, "max_pt": 36, "weight": "bold", "line_limit": 2, "fallback": "Aptos"},
                {"role": "metric", "font_stack": ["Aptos", "Arial"], "default_pt": 38, "min_pt": 30, "max_pt": 46, "weight": "bold", "line_limit": 1, "fallback": "Aptos Display"},
                {"role": "body", "font_stack": ["Aptos", "Arial"], "default_pt": 17, "min_pt": 12, "max_pt": 21, "weight": "regular", "line_limit": 4, "fallback": "Aptos"},
                {"role": "caption", "font_stack": ["Aptos", "Arial"], "default_pt": 11, "min_pt": 9, "max_pt": 13, "weight": "regular", "line_limit": 2, "fallback": "Aptos"},
            ],
        },
        "slide_size": {"size": "widescreen_16_9", "width_in": 13.333, "height_in": 7.5},
        "safe_area": {"left_in": 0.45, "top_in": 0.35, "right_in": 12.88, "bottom_in": 7.1},
        "background_policy": {"default_mode": "single_solid_cloud", "max_modes_per_deck": 3, "allowed_types": ["solid", "approved_image", "generated_image"], "forbidden_types": ["html_preview_image", "unknown_stock_image"]},
        "header_footer_policy": {"header": {"mode": "optional", "allowed_content": ["section label"]}, "footer": {"mode": "optional", "allowed_content": ["page number"]}},
        "slide_plan": {"slides": slides},
        "layout_archetypes": [
            {
                "id": item,
                "description": item.replace("_", " "),
                "native_assembly_hint": "Use native editable PowerPoint text boxes, shapes, connectors, and chart-like primitives.",
            }
            for item in archetypes[:slide_count]
        ],
        "asset_slot_policy": {
            "priority_order": [
                "approved_package_asset",
                "user_supplied_approved_asset",
                "generated_image_approved_for_deck",
                "native_shape_or_mockup_fallback",
            ],
            "notes": ["Use native fallback when no approved asset package is available."],
        },
        "qa_rules": [{"id": "native-ppt", "severity": "blocking", "check": "Slides must use native PowerPoint objects."}],
        "fallback_policy": {
            "missing_asset": "Use native shape or mockup fallback and report it.",
            "font_unavailable": "Use the role fallback font and report the fallback.",
            "overflow": "Summarize text before shrinking; report unresolved overflow.",
        },
        "approved_asset_references": [],
        "public_safety": {
            "private_refs_redacted": True,
            "contains_raw_assets": False,
            "contains_drive_ids": False,
            "contains_private_paths": False,
            "contains_generated_private_reports": False,
            "contains_tokens": False,
        },
    }


INTENT_ARCHETYPE_SEQUENCES = {
    "ir_pitch": ["cover", "big_thesis", "market_sizing", "product_mockup", "traction_metrics", "business_model", "roadmap", "ask"],
    "executive_report": ["cover", "big_thesis", "traction_metrics", "three_cards", "process_flow", "financial_plan", "ask"],
    "public_institution_report": ["cover", "big_thesis", "three_cards", "process_flow", "financial_plan", "roadmap", "ask"],
    "portfolio": ["cover", "three_cards", "product_mockup", "process_flow", "traction_metrics", "ask"],
    "product_introduction": ["cover", "big_thesis", "product_mockup", "three_cards", "three_use_cases", "traction_metrics", "ask"],
    "sales_proposal": ["cover", "big_thesis", "three_cards", "traction_metrics", "positioning_matrix", "roadmap", "ask"],
    "education_training": ["cover", "big_thesis", "process_flow", "three_cards", "three_use_cases", "roadmap", "ask"],
    "research_analysis": ["cover", "big_thesis", "market_sizing", "positioning_matrix", "traction_metrics", "three_cards", "ask"],
    "marketing_campaign": ["cover", "big_thesis", "three_use_cases", "process_flow", "roadmap", "traction_metrics", "ask"],
    "event_travel": ["cover", "big_thesis", "process_flow", "three_cards", "roadmap", "ask"],
    "media_entertainment": ["cover", "big_thesis", "three_cards", "three_use_cases", "roadmap", "ask"],
    "status_update": ["cover", "big_thesis", "traction_metrics", "process_flow", "roadmap", "ask"],
    "unknown": ["cover", "big_thesis", "three_cards", "process_flow", "ask"],
}


def clean_visible_phrase(value: Any, *, fallback: str = "Audience-ready point", limit: int = 96) -> str:
    text = sanitize_visible_value(value, fallback=fallback)
    text = re.sub(r"\s+", " ", text).strip(" .,:;-")
    return (text or fallback)[:limit]


def display_title(value: str) -> str:
    small_words = {"a", "an", "and", "as", "at", "by", "for", "in", "of", "on", "or", "the", "to", "with"}
    words = re.findall(r"[A-Za-z0-9]+(?:[-'][A-Za-z0-9]+)?|[가-힣]+", str(value or ""))
    output = []
    for index, word in enumerate(words):
        lowered = word.lower()
        if index > 0 and lowered in small_words:
            output.append(lowered)
        elif word.isupper() and len(word) <= 4:
            output.append(word)
        else:
            output.append(word[:1].upper() + word[1:].lower())
    return " ".join(output)


def split_focus_for_pairing(items: list[str], target_count: int) -> list[str]:
    cleaned = [clean_visible_phrase(item, fallback="Operating focus", limit=64) for item in items if clean_visible_phrase(item, fallback="", limit=64)]
    if target_count <= 0:
        return []
    if not cleaned:
        return []
    if len(cleaned) <= target_count:
        return cleaned
    grouped: list[str] = []
    remaining = list(cleaned)
    while remaining and len(grouped) < target_count:
        slots_left = target_count - len(grouped)
        take = max(1, len(remaining) // slots_left)
        chunk = remaining[:take]
        remaining = remaining[take:]
        grouped.append(" and ".join(chunk[:2]) if len(chunk) <= 2 else ", ".join(chunk[:-1]) + f", and {chunk[-1]}")
    return grouped


def compact_focus_title(value: str, *, limit: int = 48) -> str:
    title = display_title(value)
    if len(title) <= limit:
        return title
    parts = [
        re.sub(r"\b(?:measurable|patient|visitor|daily|regional|operational|corporate)\b", "", part, flags=re.IGNORECASE).strip()
        for part in re.split(r"\s+and\s+|,\s*", title)
        if part.strip()
    ]
    compact_parts = []
    for part in parts:
        words = [word for word in part.split() if word.lower() not in {"and", "the", "for"}]
        compact_parts.append(words[-1] if len(words) > 2 else " ".join(words))
    compact_parts = [part for part in compact_parts if part]
    if len(compact_parts) >= 3:
        compact = ", ".join(compact_parts[:-1]) + f", and {compact_parts[-1]}"
    else:
        compact = " and ".join(compact_parts)
    return (compact or title)[:limit].strip(" ,")


def requested_slide_count_from_intake(request_intake: dict[str, Any] | None, default: int) -> int:
    payload = (request_intake or {}).get("requested_slide_count")
    value = payload.get("value") if isinstance(payload, dict) else payload
    try:
        count = int(value)
    except (TypeError, ValueError):
        return default
    return count if 1 <= count <= 40 else default


def archetypes_for_count(family: str, count: int) -> list[str]:
    base = list(INTENT_ARCHETYPE_SEQUENCES.get(family, INTENT_ARCHETYPE_SEQUENCES["unknown"]))
    if count <= len(base):
        if count == 1:
            return ["cover"]
        return [base[0], *base[1:count - 1], "ask"]
    extension_pool = ["big_thesis", "process_flow", "three_cards", "three_use_cases", "roadmap", "traction_metrics"]
    archetypes = list(base)
    while len(archetypes) < count:
        insert_at = max(1, len(archetypes) - 1)
        archetypes.insert(insert_at, extension_pool[(len(archetypes) - len(base)) % len(extension_pool)])
    return archetypes[:count]


def sparse_slide_labels(
    *,
    topic: str,
    family: str,
    subtype: str,
    transformed: dict[str, Any],
    count: int,
) -> list[tuple[str, str, str]]:
    objective = clean_visible_phrase(transformed.get("objective") or f"Clarify the plan for {topic}", fallback=f"Clarify the plan for {topic}", limit=120)
    audience = clean_visible_phrase(transformed.get("audience") or "the intended audience", fallback="the intended audience", limit=100)
    focus_areas = [display_title(item) for item in transformed.get("focus_areas", []) if str(item).strip()]
    family_openers = {
        "sales_proposal": "A partner-ready story for the decision team.",
        "education_training": "A practical training path for the people who run the work.",
        "executive_report": "A concise operating story for leadership review.",
        "public_institution_report": "A clear civic story with accountable next steps.",
        "marketing_campaign": "A campaign story that connects audience, channel, and outcome.",
        "product_introduction": "A product story that connects value, proof, and activation.",
    }
    family_closers = {
        "sales_proposal": "Partnership Ask and Success Measures",
        "education_training": "Practice Commitments and Adoption Metrics",
        "executive_report": "Decision and Operating Metrics",
        "public_institution_report": "Implementation Commitments",
        "marketing_campaign": "Launch Plan and Metrics",
        "product_introduction": "Activation Plan and Metrics",
    }

    def subtitle_for_focus(focus: str, index: int) -> str:
        lowered = focus.lower()
        if any(token in lowered for token in ["metric", "outcome", "success", "adoption"]):
            return f"Define how {audience} will measure progress and know what changed."
        if any(token in lowered for token in ["journey", "communication", "handoff", "flow"]):
            return f"Show the sequence clearly so {audience} can act without ambiguity."
        if any(token in lowered for token in ["risk", "incident", "escalation", "safety"]):
            return f"Make ownership, thresholds, and response behavior easy to follow."
        if any(token in lowered for token in ["staffing", "operations", "lane", "protocol", "checklist"]):
            return f"Turn {focus.lower()} into repeatable operating behavior."
        return f"Make {focus.lower()} concrete for {audience}."

    grouped_focus = split_focus_for_pairing(focus_areas, max(0, count - 2))
    base = [("Opening", topic, family_openers.get(family, objective))]
    for index, focus in enumerate(grouped_focus, start=1):
        title = compact_focus_title(focus)
        base.append((title, title, subtitle_for_focus(focus, index)))
    base.append(("Close", family_closers.get(family, "Next Action and Success Measures"), f"Align {audience} on ownership, timing, and the evidence of progress."))

    while len(base) < count:
        next_no = len(base)
        if focus_areas:
            focus = focus_areas[(next_no - 1) % len(focus_areas)]
            title = compact_focus_title(focus)
            base.insert(-1 if len(base) > 1 else len(base), (title, title, subtitle_for_focus(focus, next_no)))
        else:
            title = f"{subtype.replace('_', ' ').title()} Detail"
            base.insert(-1 if len(base) > 1 else len(base), (title, title, "Turn the request into audience-ready evidence."))
    return base[:count]


def compose_guide_packet_from_intent(
    intent_profile: dict[str, Any],
    source_summary: dict[str, Any],
    routing_report: dict[str, Any],
    *,
    request_intake: dict[str, Any] | None = None,
    project_name: str | None = None,
) -> dict[str, Any]:
    transformed = (request_intake or {}).get("transformed_request") or {}
    topic = clean_visible_phrase(project_name or transformed.get("deck_title") or intent_profile.get("topic") or "Untitled Deck", fallback="Untitled Deck", limit=80)
    family = str(intent_profile.get("deck_family") or "unknown")
    subtype = str(intent_profile.get("sector_subtype") or "general")
    requested_count = requested_slide_count_from_intake(request_intake, default=5)
    archetypes = archetypes_for_count(family, requested_count)
    strategy_a = routing_report.get("selected", {}).get("variant_a", {}).get("strategy_id", "board_brief")
    strategy_b = routing_report.get("selected", {}).get("variant_b", {}).get("strategy_id", "demo_story")
    packet = compose_default_packet_from_prompt(str(topic), slide_count=requested_count, project_name=str(topic))
    resolved_palette, palette_resolution = resolve_intent_palette(intent_profile, routing_report, source_summary=source_summary)
    packet["palette"] = resolved_palette
    packet["guide_identity"]["guide_id"] = f"intent:{re.sub(r'[^a-z0-9]+', '-', str(topic).lower()).strip('-') or 'deck'}"
    packet["guide_identity"]["language"] = str((request_intake or {}).get("detected_language") or "en-US")
    packet["guide_identity"]["topic"] = str(topic)
    packet["project_brief"] = {
        "audience": [str(intent_profile.get("audience") or "general audience")],
        "tone": [str(item) for item in intent_profile.get("tone", ["clear", "credible"])],
        "objective": str(intent_profile.get("objective") or "create"),
        "constraints": [
            "Use public-safe local artifacts only.",
            "Do not use HTML screenshots in PPTX.",
            "Keep visible content separate from router metadata and source evidence.",
            "Do not place raw user command text on slides, footers, or public visible-content reports.",
            f"Requested slide count: {requested_count}.",
            f"Intent family: {family}; subtype: {subtype}.",
            f"Palette resolution: {palette_resolution['source_type']} / {palette_resolution['selected_palette_id']}.",
        ],
    }
    source_refs = []
    for source in source_summary.get("sources", []):
        source_refs.append(
            {
                "file_name": str(source.get("public_label") or source.get("source_id")),
                "extension": f".{source.get('kind', 'source')}",
                "purpose": "Bounded sparse-intake source summary",
                "public_safe_note": "Private source location and raw payload are not included.",
            }
        )
    packet["reference_files"] = source_refs or [
        {"file_name": "sparse-intake", "extension": ".txt", "purpose": "User natural-language request"}
    ]
    slides = []
    labels = sparse_slide_labels(topic=topic, family=family, subtype=subtype, transformed=transformed, count=requested_count)
    public_family_hint = infer_public_composition_family(GuidePacket.model_validate(packet), strategy_a)

    def beta_asset_slots_for_slide(slide_no: int) -> list[dict[str, Any]]:
        if public_family_hint == "luxury_editorial" and slide_no in {1, 2, 3}:
            return [
                {
                    "slot_id": f"luxury_product_detail_{slide_no}",
                    "slot_type": "image",
                    "required": slide_no == 1,
                    "allowed_sources": ["approved_package_asset", "user_supplied_approved_asset", "generated_image_approved_for_deck"],
                    "fallback": "Use native editorial material/product-stage fallback and report the missing expectation.",
                    "crop_or_mask_policy": "product_or_material_detail_crop",
                }
            ]
        if public_family_hint == "food_product_launch" and slide_no in {1, 2, 3, 4}:
            slot_ids = {
                1: "food_product_packshot",
                2: "ingredient_detail",
                3: "retail_shelf_context",
                4: "campaign_calendar_visual",
            }
            return [
                {
                    "slot_id": slot_ids[slide_no],
                    "slot_type": "image",
                    "required": slide_no == 1,
                    "allowed_sources": ["approved_package_asset", "user_supplied_approved_asset", "generated_image_approved_for_deck"],
                    "fallback": "Use native food/product launch fallback and report the missing expectation.",
                    "crop_or_mask_policy": "product_launch_visual_crop",
                }
            ]
        if public_family_hint == "internal_strategy_report" and slide_no in {2, 3, 4}:
            return [
                {
                    "slot_id": f"strategy_evidence_visual_{slide_no}",
                    "slot_type": "chart",
                    "required": False,
                    "allowed_sources": ["approved_package_asset", "user_supplied_approved_asset"],
                    "fallback": "Use native memo, matrix, risk, or roadmap primitives and report the expectation.",
                    "crop_or_mask_policy": "native_evidence_surface",
                }
            ]
        return []

    for index, archetype in enumerate(archetypes, start=1):
        label, title, subtitle = labels[min(index - 1, len(labels) - 1)]
        audience_text = clean_visible_phrase(transformed.get("audience") or intent_profile.get("audience") or "intended audience", fallback="intended audience", limit=90)
        support_candidate = f"For {audience_text}" if audience_text != "general audience" else f"{topic} focus"
        slides.append(
            {
                "slide_no": index,
                "layout_archetype": archetype,
                "content_brief": sanitize_prompt_literal_text(f"{label}: {subtitle}", fallback=label),
                "visible_content_candidates": [
                    sanitize_prompt_literal_text(title, fallback=label),
                    sanitize_prompt_literal_text(subtitle, fallback="Audience-ready review"),
                    support_candidate,
                    f"{label}: {sanitize_prompt_literal_text(title, fallback=label)}",
                ],
                "asset_slots": beta_asset_slots_for_slide(index),
                "qa_checks": ["No HTML screenshot", "Native PPT objects only", "No private source paths on visible slides"],
            }
        )
    packet["guide_identity"]["slide_count"] = len(slides)
    packet["slide_plan"]["slides"] = slides
    packet["layout_archetypes"] = [
        {
            "id": item,
            "description": item.replace("_", " "),
            "native_assembly_hint": "Use native editable PowerPoint text boxes, shapes, connectors, tables, and chart-like primitives.",
        }
        for item in dict.fromkeys(archetypes)
    ]
    packet["public_safety"].update(
        {
            "contains_raw_assets": False,
            "contains_drive_ids": False,
            "contains_private_paths": False,
            "contains_tokens": False,
        }
    )
    GuidePacket.model_validate(packet)
    return safe_string(packet)

from __future__ import annotations

import math
import unicodedata
from pathlib import Path
from typing import Any

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from system.pptx_system import (
    DEFAULT_PROTECTED_TOKENS,
    ThemeConfig,
    add_image,
    add_simple_bar_chart,
    add_simple_table,
    get_shape_by_index,
    remove_shape,
    resolve_color,
    set_text,
)
from system.text_summarizer import llm_summary_enabled, summarize_to_budget

ALIGNMENTS = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}
VALID_FIT_STRATEGIES = {"preserve_template", "shrink", "manual_wrap"}
CJK_UNIT_WIDTH = 1.7


class SlotResolutionError(ValueError):
    pass


def fallback_text_slot_bounds(slot_def: dict[str, Any], slide) -> dict[str, float]:
    slot_name = str(slot_def.get("slot", "")).lower()
    existing_text_shapes = [shape for shape in slide.shapes if hasattr(shape, "text_frame")]
    if slot_name == "title":
        return {"left": 0.72, "top": 0.55, "width": 11.9, "height": 0.72}
    if slot_name in {"subtitle", "eyebrow", "section"}:
        return {"left": 0.82, "top": 1.28, "width": 11.3, "height": 0.55}
    if slot_name in {"footer", "source", "note"}:
        return {"left": 0.82, "top": 6.88, "width": 11.3, "height": 0.28}
    top = min(6.25, 1.95 + max(0, len(existing_text_shapes) - 2) * 0.7)
    return {"left": 0.82, "top": top, "width": 11.3, "height": 0.62}


def font_size_for_role(theme: ThemeConfig, role: str | None, fallback: str = "body") -> float:
    if role is None:
        return theme.sizes[fallback]
    return theme.sizes.get(role, theme.sizes[fallback])


def slot_fit_strategy(slot_def: dict[str, Any], override: dict[str, Any]) -> str:
    fit_strategy = override.get("fit_strategy") or slot_def.get("fit_strategy") or "preserve_template"
    if fit_strategy not in VALID_FIT_STRATEGIES:
        raise ValueError(f"Unsupported fit_strategy for slot '{slot_def['slot']}': {fit_strategy}")
    return fit_strategy


def slot_text_budget(slot_def: dict[str, Any], override: dict[str, Any] | None = None) -> int | None:
    override = override or {}
    max_chars_per_line = override.get("max_chars_per_line") or slot_def.get("max_chars_per_line")
    max_lines = override.get("max_lines") or slot_def.get("max_lines")
    if not max_chars_per_line or not max_lines:
        return None
    return max(1, int(max_chars_per_line) * int(max_lines))


def char_budget_units(char: str) -> float:
    return CJK_UNIT_WIDTH if unicodedata.east_asian_width(char) in {"F", "W"} else 1.0


def text_budget_units(text: str) -> int:
    return int(math.ceil(sum(char_budget_units(char) for char in text)))


def truncate_text_to_budget_units(text: str, budget: int, suffix: str = "...") -> str:
    suffix_units = text_budget_units(suffix)
    if budget <= suffix_units:
        return suffix[: max(1, budget)]

    allowed_units = budget - suffix_units
    used_units = 0.0
    kept: list[str] = []
    for char in text:
        next_units = used_units + char_budget_units(char)
        if next_units > allowed_units:
            break
        kept.append(char)
        used_units = next_units
    return "".join(kept).rstrip() + suffix


def coerce_text_to_slot_budget(
    value: str,
    slot_def: dict[str, Any],
    override: dict[str, Any] | None = None,
) -> tuple[str, dict[str, Any] | None]:
    text = str(value)
    budget = slot_text_budget(slot_def, override)
    original_units = text_budget_units(text)
    if budget is None or original_units <= budget:
        return text, None

    deterministic = truncate_text_to_budget_units(text, budget)

    final_text = deterministic
    resolution = "deterministic_cutoff"
    llm_metadata: dict[str, Any] = {
        "llm_attempted": False,
        "llm_used": False,
    }
    if llm_summary_enabled():
        protected_tokens = slot_def.get("protect_tokens") or DEFAULT_PROTECTED_TOKENS
        fit_strategy = slot_fit_strategy(slot_def, override or {})
        summary, llm_metadata = summarize_to_budget(
            text=text,
            budget=budget,
            slot_name=str(slot_def.get("slot")),
            slide_context={
                "slot": slot_def.get("slot"),
                "fit_strategy": fit_strategy,
            },
            protected_tokens=protected_tokens,
            allow_line_breaks=fit_strategy == "manual_wrap",
        )
        if summary is not None:
            final_text = summary
            resolution = "llm_summary"
            if text_budget_units(final_text) > budget:
                final_text = deterministic
                resolution = "deterministic_cutoff"
                llm_metadata["fallback_reason"] = "llm_summary_exceeded_cjk_weighted_budget"
    event = {
        "slot": slot_def.get("slot"),
        "budget": budget,
        "original_chars": len(text),
        "original_units": original_units,
        "truncated_chars": len(final_text),
        "summary_units": text_budget_units(final_text),
        "language_weight": "cjk_weighted" if original_units > len(text) else "default",
        "original_preview": text[:120],
        "truncated_text": final_text,
        "resolution": resolution,
        "summary_chars": len(final_text),
        **llm_metadata,
    }
    if resolution == "deterministic_cutoff" and llm_metadata.get("llm_attempted") and not llm_metadata.get("fallback_reason"):
        event["fallback_reason"] = "deterministic_cutoff"
    return final_text, event


def shape_matches_kind(shape, expected_kind: str) -> bool:
    if expected_kind == "text":
        return hasattr(shape, "text_frame")
    if expected_kind == "image":
        return shape.shape_type == MSO_SHAPE_TYPE.PICTURE or bool(getattr(shape, "is_placeholder", False))
    raise ValueError(f"Unsupported slot shape kind: {expected_kind}")


def expected_slot_shape_name(slot_name: str) -> str:
    return f"slot:{slot_name}"


def unique_candidate(
    candidates: list[Any],
    *,
    slide_id: str,
    slot_name: str,
    strategy: str,
) -> Any | None:
    if not candidates:
        return None
    if len(candidates) == 1:
        return candidates[0]
    names = [getattr(shape, "name", "<unnamed>") for shape in candidates[:5]]
    raise SlotResolutionError(
        f"Ambiguous {strategy} match for template slot {slide_id}:{slot_name}; "
        f"matched {len(candidates)} shapes: {names}"
    )


def resolve_slot_shape(
    slide,
    slot_def: dict[str, Any],
    expected_kind: str,
    *,
    slide_id: str = "<unknown>",
    required: bool = False,
) -> Any | None:
    slot_name = slot_def["slot"]
    expected_name = expected_slot_shape_name(slot_name)
    blueprint_shape_name = slot_def.get("shape_name")
    if blueprint_shape_name != expected_name:
        raise SlotResolutionError(
            f"Template slot {slide_id}:{slot_name} has invalid strict shape identity; "
            f"expected shape_name={expected_name!r}, blueprint shape_name={blueprint_shape_name!r}, "
            f"kind={expected_kind!r}"
        )

    shapes = [shape for shape in slide.shapes if shape_matches_kind(shape, expected_kind)]
    match = unique_candidate(
        [shape for shape in shapes if getattr(shape, "name", None) == expected_name],
        slide_id=slide_id,
        slot_name=slot_name,
        strategy="strict shape_name",
    )
    if match is not None:
        return match

    available_names = [getattr(shape, "name", "<unnamed>") for shape in shapes[:12]]
    raise SlotResolutionError(
        f"Unable to resolve template slot {slide_id}:{slot_name} as {expected_kind} "
        f"with strict shape_name={expected_name!r}; available {expected_kind} shapes={available_names}"
    )


def apply_text_slot(
    slide,
    slot_def: dict[str, Any],
    value: str,
    theme: ThemeConfig,
    override: dict[str, Any] | None = None,
    slide_id: str = "<unknown>",
) -> None:
    override = override or {}
    try:
        shape = resolve_slot_shape(slide, slot_def, "text", slide_id=slide_id)
    except SlotResolutionError:
        bounds = slot_def.get("bounds") or fallback_text_slot_bounds(slot_def, slide)
        shape = slide.shapes.add_textbox(
            Inches(bounds["left"]),
            Inches(bounds["top"]),
            Inches(bounds["width"]),
            Inches(bounds["height"]),
        )
        shape.name = expected_slot_shape_name(slot_def["slot"])
    font_role = override.get("font_role") or slot_def.get("font_role") or "body"
    font_size = override.get("font_size") or font_size_for_role(theme, font_role)
    font_color = resolve_color(theme, override.get("color") or slot_def.get("color_token"), "primary")
    bold = override.get("bold", slot_def.get("bold", False))
    align_value = override.get("align") or slot_def.get("align")
    align = ALIGNMENTS.get(align_value) if align_value else None
    fit_strategy = slot_fit_strategy(slot_def, override)
    max_chars_per_line = override.get("max_chars_per_line") or slot_def.get("max_chars_per_line")
    protected_tokens = slot_def.get("protect_tokens") or DEFAULT_PROTECTED_TOKENS

    if shape is not None and hasattr(shape, "text_frame"):
        set_text(
            shape,
            value,
            font_name=theme.font_family,
            font_size=font_size,
            font_color=font_color,
            bold=bold,
            align=align,
            max_chars_per_line=max_chars_per_line,
            protect_token_list=protected_tokens,
            is_placeholder=True,
            fit_strategy=fit_strategy,
        )
        return


def apply_image_slot(slide, slot_def: dict[str, Any], image_path: str | Path, *, slide_id: str = "<unknown>") -> None:
    try:
        shape = resolve_slot_shape(slide, slot_def, "image", slide_id=slide_id)
    except SlotResolutionError:
        shape = None
    bounds = slot_def.get("bounds")
    if bounds is None and shape is not None:
        bounds = {
            "left": shape.left / 914400,
            "top": shape.top / 914400,
            "width": shape.width / 914400,
            "height": shape.height / 914400,
        }
    if bounds is None:
        raise ValueError(f"Image slot '{slot_def['slot']}' has no bounds")
    if shape is not None:
        remove_shape(shape)
    add_image(
        slide,
        image_path,
        left=bounds["left"],
        top=bounds["top"],
        width=bounds["width"],
        height=bounds["height"],
    )


def resolve_image_slot_shape_for_clear(slide, slot_def: dict[str, Any], *, slide_id: str) -> Any | None:
    try:
        return resolve_slot_shape(slide, slot_def, "image", slide_id=slide_id)
    except SlotResolutionError:
        pass

    bounds = slot_def.get("bounds")
    if not bounds:
        return None
    candidates = [shape for shape in slide.shapes if shape_matches_kind(shape, "image")]
    for shape in candidates:
        actual = {
            "left": shape.left / 914400,
            "top": shape.top / 914400,
            "width": shape.width / 914400,
            "height": shape.height / 914400,
        }
        if all(abs(float(actual[key]) - float(bounds[key])) <= 0.04 for key in ("left", "top", "width", "height")):
            return shape
    return None


def apply_chart_slot(slide, slot_def: dict[str, Any], chart_data: dict[str, Any], theme: ThemeConfig) -> None:
    bounds = slot_def.get("bounds")
    if bounds is None:
        raise ValueError(f"Chart slot '{slot_def['slot']}' has no bounds")
    chart_type = chart_data.get("chart_type", "bar")
    if chart_type != "bar":
        raise ValueError(f"Unsupported chart_type for slot '{slot_def['slot']}': {chart_type}")
    add_simple_bar_chart(
        slide,
        left=bounds["left"],
        top=bounds["top"],
        width=bounds["width"],
        height=bounds["height"],
        categories=chart_data["categories"],
        values=chart_data["values"],
        theme=theme,
    )


def apply_table_slot(slide, slot_def: dict[str, Any], table_data: dict[str, Any], theme: ThemeConfig) -> None:
    bounds = slot_def.get("bounds")
    if bounds is None:
        raise ValueError(f"Table slot '{slot_def['slot']}' has no bounds")
    add_simple_table(
        slide,
        left=bounds["left"],
        top=bounds["top"],
        width=bounds["width"],
        height=bounds["height"],
        headers=table_data.get("headers", []),
        rows=table_data["rows"],
        theme=theme,
    )


def clear_residual_text(slide_or_group, patterns: list[str]) -> None:
    lowered_patterns = [pattern.lower() for pattern in patterns if pattern]
    if not lowered_patterns:
        return
    for shape in list(slide_or_group.shapes):
        if hasattr(shape, "shapes"):
            clear_residual_text(shape, patterns)
        text = getattr(shape, "text", "")
        normalized = " ".join(text.split()).lower()
        if normalized and any(pattern in normalized for pattern in lowered_patterns):
            remove_shape(shape)


def render_template_slide(
    slide,
    slide_spec: dict[str, Any],
    theme: ThemeConfig,
    blueprint: dict[str, Any],
    overflow_events: list[dict[str, Any]] | None = None,
) -> None:
    for shape_index in blueprint.get("remove_shapes", []):
        shape = get_shape_by_index(slide, shape_index)
        if shape is not None:
            remove_shape(shape)

    text_slots = {item["slot"]: item for item in blueprint.get("editable_text_slots", [])}
    text_values = dict(slide_spec.get("text_slots", {}))
    sequential_values = slide_spec.get("text_values", [])
    sequential_slots = [item["slot"] for item in blueprint.get("editable_text_slots", []) if item["slot"] not in text_values]
    for slot_name, value in zip(sequential_slots, sequential_values, strict=False):
        text_values[slot_name] = value

    clear_unfilled_slots = slide_spec.get("clear_unfilled_slots", True)
    for slot_name, slot_def in text_slots.items():
        value = text_values.get(slot_name, "" if clear_unfilled_slots else None)
        if value is None:
            continue
        slot_override = slide_spec.get("slot_overrides", {}).get(slot_name, {})
        coerced_value, event = coerce_text_to_slot_budget(value, slot_def, slot_override)
        if event is not None and overflow_events is not None:
            event.update(
                {
                    "slide_id": blueprint["slide_id"],
                    "template_key": blueprint["template_key"],
                    "slot_kind": "text",
                }
            )
            overflow_events.append(event)
        apply_text_slot(slide, slot_def, coerced_value, theme, slot_override, slide_id=blueprint["slide_id"])

    image_slots = {item["slot"]: item for item in blueprint.get("editable_image_slots", [])}
    provided_image_slots = set(slide_spec.get("image_slots", {}))
    for slot_name, image_path in slide_spec.get("image_slots", {}).items():
        if slot_name not in image_slots:
            raise KeyError(f"Unknown image slot '{slot_name}' for slide {blueprint['slide_id']}")
        apply_image_slot(slide, image_slots[slot_name], image_path, slide_id=blueprint["slide_id"])
    if slide_spec.get("clear_unfilled_image_slots", False):
        for slot_name, slot_def in image_slots.items():
            if slot_name in provided_image_slots:
                continue
            shape = resolve_image_slot_shape_for_clear(slide, slot_def, slide_id=blueprint["slide_id"])
            if shape is not None:
                remove_shape(shape)

    chart_slots = {item["slot"]: item for item in blueprint.get("editable_chart_slots", [])}
    for slot_name, chart_data in slide_spec.get("chart_slots", {}).items():
        if slot_name not in chart_slots:
            raise KeyError(f"Unknown chart slot '{slot_name}' for slide {blueprint['slide_id']}")
        apply_chart_slot(slide, chart_slots[slot_name], chart_data, theme)

    table_slots = {item["slot"]: item for item in blueprint.get("editable_table_slots", [])}
    for slot_name, table_data in slide_spec.get("table_slots", {}).items():
        if slot_name not in table_slots:
            raise KeyError(f"Unknown table slot '{slot_name}' for slide {blueprint['slide_id']}")
        apply_table_slot(slide, table_slots[slot_name], table_data, theme)

    residual_patterns = []
    if slide_spec.get("clear_residual_placeholders", True):
        residual_patterns.extend(
            [
                "please fill",
                "please write",
                "please enter",
                "input the",
                "provide a concise",
                "provide a brief",
                "write a detailed",
                "write the detailed",
                "copyrights",
                "office suite",
            ]
        )
    residual_patterns.extend(slide_spec.get("clear_residual_text_patterns", []))
    clear_residual_text(slide, residual_patterns)
